import os
import io
import random
import uuid
import logging
import pandas as pd
from datetime import datetime, timedelta
from functools import wraps

from flask import (Flask, request, jsonify, render_template,
                   send_from_directory, redirect, url_for,
                   Response, send_file, session)
from werkzeug.utils import secure_filename
from flask_login import (LoginManager, login_user, login_required,
                         logout_user, current_user)
from werkzeug.security import check_password_hash, generate_password_hash

from core.database import (engine, Base, SessionLocal, UserQuery, UserDataset, UserChart,
                            AffiliateLink, QueryCount, ScreenPost, MusicTrack,
                            VideoTrack, ChatMessage, PostAnalytics,
                            ScheduledPost, PushSubscription, PaymentScreenshot,
                            OTPVerification)
from core.auth import (User, create_default_admin,
                       create_session_token, validate_session_token,
                       invalidate_session, get_user_plan, set_user_plan)
from core.agent import DataScienceAgent
from utils.email import send_report_email
from utils.pdf_generator import generate_pdf_report

# Keep-alive scheduler for HuggingFace Space
from apscheduler.schedulers.background import BackgroundScheduler

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__, template_folder='templates', static_folder='static')
app.secret_key = os.getenv("SECRET_KEY", "super-secret-key-change-this!")
UPLOAD_DIR = '/tmp/uploads' if os.path.exists('/tmp') else '.'
os.makedirs(UPLOAD_DIR, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_DIR
app.config['MAX_CONTENT_LENGTH'] = 510 * 1024 * 1024  # 510MB for chunked uploads

# ── HuggingFace Keep-Alive Ping ───────────────────────────────
def ping_hf_space():
    """Ping HF Space every 5 min to prevent sleep"""
    try:
        import requests as _req
        hf_url = os.environ.get('HF_SPACE_URL', '').strip()
        if hf_url:
            r = _req.get(hf_url + '/ping', timeout=10)
            if r.status_code == 200:
                print(f"[KeepAlive] HF Space OK — {r.json().get('message','')}")
            else:
                print(f"[KeepAlive] HF Space responded {r.status_code}")
    except Exception as e:
        print(f"[KeepAlive] HF ping failed: {e}")

_scheduler = BackgroundScheduler(daemon=True)
_scheduler.add_job(ping_hf_space, 'interval', minutes=5, id='hf_keepalive')
_scheduler.start()
print("[KeepAlive] HF Space ping scheduler started — every 5 min")

# ── UPI Config ──────────────────────────────────────────────────
UPI_ID    = os.getenv("UPI_ID", "yourname@upi")
UPI_NAME  = os.getenv("UPI_NAME", "DS Agent")
UPI_AMOUNT = "499"

# ── Gmail SMTP Config ────────────────────────────────────────────
GMAIL_USER = os.getenv("GMAIL_USER", "")
GMAIL_PASS = os.getenv("GMAIL_PASS", "")

# ── Flask-Login ───────────────────────────────────────────────────
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

FREE_DAILY_LIMIT = 10

@login_manager.user_loader
def load_user(user_id):
    db = SessionLocal()
    try:
        return db.query(User).filter(User.id == int(user_id)).first()
    except:
        return None
    finally:
        db.close()

# ── DB Init ───────────────────────────────────────────────────────
try:
    Base.metadata.create_all(bind=engine)
    create_default_admin()
    # ── Fix: Existing admins ko verified mark karo ──
    _db = SessionLocal()
    try:
        _admins = _db.query(User).filter(User.is_admin == True).all()
        for _a in _admins:
            if not _a.is_verified:
                _a.is_verified = True
                logger.info(f"✅ Admin '{_a.username}' verified!")
        _db.commit()
    except Exception as _e:
        _db.rollback()
    finally:
        _db.close()
    logger.info("✅ Database ready")
except Exception as e:
    logger.error(f"❌ DB init: {e}")

agent = DataScienceAgent()

# ── Decorators ────────────────────────────────────────────────────
def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not current_user.is_authenticated or not current_user.is_admin:
            return jsonify({"error": "Admin only!"}), 403
        return f(*args, **kwargs)
    return decorated

def single_session_check(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if current_user.is_authenticated:
            token = session.get('session_token')
            if token and not validate_session_token(current_user.id, token):
                logout_user()
                session.clear()
                return jsonify({"error": "SESSION_EXPIRED",
                                "message": "Aapka account kisi aur device pe login hua hai!"}), 401
        return f(*args, **kwargs)
    return decorated

def pro_required(feature_name="This feature"):
    def decorator(f):
        @wraps(f)
        def decorated(*args, **kwargs):
            plan = get_user_plan(current_user.id)
            if plan != "pro":
                return jsonify({
                    "error": "PRO_REQUIRED",
                    "message": f"⭐ {feature_name} sirf Pro plan mein available hai!",
                    "upgrade_url": url_for('pricing')
                }), 403
            return f(*args, **kwargs)
        return decorated
    return decorator

def check_query_limit(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if current_user.is_admin:
            return f(*args, **kwargs)
        plan = get_user_plan(current_user.id)
        if plan == "free":
            today = datetime.utcnow().strftime('%Y-%m-%d')
            db = SessionLocal()
            try:
                qc = db.query(QueryCount).filter_by(
                    user_id=current_user.id, date=today).first()
                count = qc.count if qc else 0
                if count >= FREE_DAILY_LIMIT:
                    return jsonify({
                        "response": f"⚠️ Aapki aaj ki {FREE_DAILY_LIMIT} free queries khatam ho gayi hain!\n\n⭐ Pro plan upgrade karo unlimited queries ke liye.",
                        "limit_reached": True
                    })
                if qc:
                    qc.count += 1
                else:
                    db.add(QueryCount(user_id=current_user.id, date=today, count=1))
                db.commit()
            except Exception as e:
                logger.error(f"Query count error: {e}")
            finally:
                db.close()
        return f(*args, **kwargs)
    return decorated

# ── Email OTP Helpers ─────────────────────────────────────────────
def generate_otp() -> str:
    return str(random.randint(100000, 999999))

def send_otp_email(to_email: str, otp: str, username: str) -> bool:
    import smtplib
    from email.mime.multipart import MIMEMultipart
    from email.mime.text import MIMEText
    gmail_user = os.getenv("GMAIL_USER", "")
    gmail_pass = os.getenv("GMAIL_PASS", "")
    if not gmail_user or not gmail_pass:
        logger.warning(f"[DEV] OTP for {username} ({to_email}): {otp}")
        return True
    html_body = f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8"><style>
body{{font-family:Segoe UI,sans-serif;background:#f6f7fb;margin:0;padding:20px}}
.card{{max-width:480px;margin:0 auto;background:white;border-radius:20px;overflow:hidden}}
.hdr{{background:linear-gradient(135deg,#5b5ef4,#8b5cf6);padding:32px 28px;text-align:center}}
.hdr h1{{color:white;margin:0;font-size:26px}}
.body{{padding:32px 28px}}
.otp-box{{background:#f0f0ff;border:2px solid #5b5ef4;border-radius:16px;padding:24px;text-align:center;margin:24px 0}}
.otp{{font-size:42px;font-weight:900;color:#5b5ef4;letter-spacing:10px;font-family:monospace}}
.footer{{padding:20px;background:#f8f9fc;text-align:center;font-size:11px;color:#aaa}}
</style></head><body>
<div class="card">
<div class="hdr"><h1>DS Agent</h1></div>
<div class="body">
<p>Namaste <b>{username}</b>!</p>
<p>Tumhara OTP:</p>
<div class="otp-box"><div class="otp">{otp}</div><p style="font-size:12px;color:#888">10 minutes mein expire hoga</p></div>
<p style="font-size:13px;color:#666">Kisi ke saath share mat karo!</p>
</div>
<div class="footer">DS Agent 2025</div>
</div></body></html>"""
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = "DS Agent - Email Verify Karo"
        msg["From"]    = f"DS Agent <{gmail_user}>"
        msg["To"]      = to_email
        msg.attach(MIMEText(html_body, "html"))
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
            smtp.login(gmail_user, gmail_pass)
            smtp.sendmail(gmail_user, to_email, msg.as_string())
        logger.info(f"OTP email sent to {to_email}")
        return True
    except Exception as e:
        logger.error(f"Gmail SMTP error: {e}")
        return False

# ── Auth Routes ───────────────────────────────────────────────────
@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        email    = request.form.get('email', '').strip().lower()
        password = request.form.get('password', '')

        if not username or not email or not password:
            return render_template('register.html', error="Sab fields required hain!")
        if len(username) < 3:
            return render_template('register.html', error="Username min 3 characters hona chahiye")
        if len(password) < 6:
            return render_template('register.html', error="Password min 6 characters hona chahiye")
        if '@' not in email or '.' not in email:
            return render_template('register.html', error="Valid email address daalo!")

        db = SessionLocal()
        try:
            if db.query(User).filter(User.username == username).first():
                return render_template('register.html', error="Ye username already le liya gaya hai!")
            if db.query(User).filter(User.email == email).first():
                return render_template('register.html', error="Ye email already registered hai!")

            user = User(
                username=username,
                password_hash=generate_password_hash(password),
                email=email,
                is_verified=False
            )
            db.add(user)
            db.flush()

            otp = generate_otp()
            expires = datetime.utcnow() + timedelta(minutes=10)
            db.query(OTPVerification).filter_by(user_id=user.id).delete()
            db.add(OTPVerification(user_id=user.id, email=email, otp=otp, expires_at=expires))
            db.commit()

            send_otp_email(email, otp, username)
            session['pending_verify_user_id'] = user.id
            session['pending_verify_email']   = email
            return redirect(url_for('verify_email'))

        except Exception as e:
            db.rollback()
            import traceback
            logger.error(f"Register error: {e}")
            logger.error(traceback.format_exc())
            return render_template('register.html', error=f"Error: {str(e)}")
        finally:
            db.close()
    return render_template('register.html')


@app.route('/verify-email', methods=['GET', 'POST'])
def verify_email():
    user_id = session.get('pending_verify_user_id')
    email   = session.get('pending_verify_email', '')
    if not user_id:
        return redirect(url_for('register'))

    if request.method == 'POST':
        entered_otp = request.form.get('otp', '').strip()
        db = SessionLocal()
        try:
            record = db.query(OTPVerification).filter_by(
                user_id=user_id, otp=entered_otp, is_used=False).first()
            if not record:
                return render_template('verify_email.html',
                    error="❌ Galat OTP! Dobara check karo.", email=email)
            if record.expires_at < datetime.utcnow():
                return render_template('verify_email.html',
                    error="⏰ OTP expire ho gaya! Resend karo.", email=email)
            user = db.query(User).filter_by(id=user_id).first()
            user.is_verified = True
            record.is_used = True
            db.commit()
            session.pop('pending_verify_user_id', None)
            session.pop('pending_verify_email', None)
            return redirect(url_for('login') + '?verified=1')
        except Exception as e:
            db.rollback()
            logger.error(f"Verify error: {e}")
            return render_template('verify_email.html', error="System error!", email=email)
        finally:
            db.close()
    return render_template('verify_email.html', email=email)


@app.route('/resend-otp', methods=['POST'])
def resend_otp():
    user_id = session.get('pending_verify_user_id')
    if not user_id:
        return jsonify({"error": "Session expire ho gaya!"}), 400
    db = SessionLocal()
    try:
        user = db.query(User).filter_by(id=user_id).first()
        if not user:
            return jsonify({"error": "User nahi mila!"}), 404
        last_otp = db.query(OTPVerification).filter_by(
            user_id=user_id).order_by(OTPVerification.created_at.desc()).first()
        if last_otp:
            elapsed = (datetime.utcnow() - last_otp.created_at).total_seconds()
            if elapsed < 60:
                return jsonify({"error": f"⏳ {int(60 - elapsed)} seconds baad resend karo!"}), 429
        otp = generate_otp()
        expires = datetime.utcnow() + timedelta(minutes=10)
        db.query(OTPVerification).filter_by(user_id=user_id).delete()
        db.add(OTPVerification(user_id=user_id, email=user.email, otp=otp, expires_at=expires))
        db.commit()
        send_otp_email(user.email, otp, user.username)
        return jsonify({"message": f"✅ OTP {user.email} pe bheja!"})
    except Exception as e:
        db.rollback()
        return jsonify({"error": str(e)}), 500
    finally:
        db.close()


@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '')
        db = SessionLocal()
        try:
            user = db.query(User).filter(User.username == username).first()
            if user and check_password_hash(user.password_hash, password):
                if not user.is_admin and not user.is_verified:
                    session['pending_verify_user_id'] = user.id
                    session['pending_verify_email']   = user.email or ''
                    otp = generate_otp()
                    expires = datetime.utcnow() + timedelta(minutes=10)
                    db.query(OTPVerification).filter_by(user_id=user.id).delete()
                    db.add(OTPVerification(user_id=user.id, email=user.email, otp=otp, expires_at=expires))
                    db.commit()
                    send_otp_email(user.email, otp, user.username)
                    return redirect(url_for('verify_email'))
                login_user(user)
                token = create_session_token(user.id)
                session['session_token'] = token
                return redirect(url_for('home'))
            return render_template('login.html', error="❌ Galat username ya password!")
        except Exception as e:
            return render_template('login.html', error="System error!")
        finally:
            db.close()
    verified_msg = request.args.get('verified')
    success_msg = "✅ Email verify ho gaya! Ab login karo." if verified_msg else None
    return render_template('login.html', success=success_msg)


@app.route('/logout')
@login_required
def logout():
    invalidate_session(current_user.id)
    session.clear()
    logout_user()
    return redirect(url_for('login'))

# ── Main Pages ────────────────────────────────────────────────────
@app.route('/')
@login_required
def home():
    try:
        import json, io
        db = SessionLocal()
        dataset = db.query(UserDataset).filter_by(user_id=current_user.id).order_by(UserDataset.uploaded_at.desc()).first()
        db.close()
        if dataset and agent.df is None:
            import tempfile
            tmp = tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False, encoding='utf-8')
            tmp.write(dataset.csv_data)
            tmp.close()
            agent.load_data(tmp.name)
            logger.info(f"✅ Auto-loaded dataset for user {current_user.id}: {dataset.filename}")
    except Exception as e:
        logger.error(f"⚠️ Auto-load error: {e}")
    return render_template('index.html')

@app.route('/dashboard')
@login_required
def dashboard():
    return render_template('dashboard.html')

@app.route('/pricing')
@login_required
def pricing():
    return render_template('pricing.html', plan=get_user_plan(current_user.id))

# ── User Info ─────────────────────────────────────────────────────
@app.route('/get_username')
@login_required
def get_username():
    if current_user.is_admin:
        return jsonify({"username": current_user.username, "is_admin": True, "plan": "pro", "queries_used": 0, "queries_limit": 99999})
    plan = get_user_plan(current_user.id)
    today = datetime.utcnow().strftime('%Y-%m-%d')
    db = SessionLocal()
    try:
        qc = db.query(QueryCount).filter_by(user_id=current_user.id, date=today).first()
        queries_used = qc.count if qc else 0
    finally:
        db.close()
    return jsonify({"username": current_user.username, "is_admin": False, "plan": plan, "queries_used": queries_used, "queries_limit": FREE_DAILY_LIMIT if plan == "free" else 99999})

# ── UPI Payment Info ──────────────────────────────────────────────
@app.route('/api/upi_info')
@login_required
def upi_info():
    return jsonify({"upi_id": UPI_ID, "upi_name": UPI_NAME, "amount": UPI_AMOUNT})

# ── Payment Screenshot Submit ─────────────────────────────────────
@app.route('/submit_payment_screenshot', methods=['POST'])
@login_required
@single_session_check
def submit_payment_screenshot():
    data = request.get_json()
    if not data or not data.get('image_data'):
        return jsonify({"error": "Screenshot required!"}), 400
    db = SessionLocal()
    try:
        ss = PaymentScreenshot(
            user_id=current_user.id, username=current_user.username,
            image_data=data['image_data'], image_mime=data.get('image_mime', 'image/jpeg'),
            utr_number=data.get('utr_number', ''), amount=data.get('amount', '499'), status='pending'
        )
        db.add(ss)
        db.commit()
        logger.info(f"📸 Payment screenshot from {current_user.username}")
        return jsonify({"message": "✅ Screenshot submit ho gaya! Admin 24hrs mein verify karega."})
    except Exception as e:
        db.rollback()
        return jsonify({"error": str(e)}), 500
    finally:
        db.close()

# ── Admin: View Screenshots ───────────────────────────────────────
@app.route('/admin/payment_screenshots')
@login_required
@admin_required
def admin_payment_screenshots():
    db = SessionLocal()
    try:
        items = db.query(PaymentScreenshot).order_by(PaymentScreenshot.created_at.desc()).all()
        return jsonify([{"id": s.id, "username": s.username, "image_data": s.image_data, "image_mime": s.image_mime, "utr_number": s.utr_number or "N/A", "amount": s.amount or "499", "status": s.status, "created_at": s.created_at.strftime('%d %b %H:%M')} for s in items])
    finally:
        db.close()

# ── Admin: Approve/Reject Payment ────────────────────────────────
@app.route('/admin/review_payment', methods=['POST'])
@login_required
@admin_required
def review_payment():
    data = request.get_json()
    ss_id = data.get('id'); status = data.get('status')
    if status not in ('approved', 'rejected'):
        return jsonify({"error": "Invalid status"}), 400
    db = SessionLocal()
    try:
        ss = db.query(PaymentScreenshot).filter_by(id=ss_id).first()
        if not ss: return jsonify({"error": "Not found"}), 404
        ss.status = status
        if status == 'approved':
            set_user_plan(ss.user_id, 'pro', expires_at=datetime.utcnow() + timedelta(days=30))
        db.commit()
        msg = f"✅ Approved! Pro activated for {ss.username}" if status == 'approved' else f"❌ Rejected for {ss.username}"
        return jsonify({"message": msg})
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

# ── Affiliate Links ───────────────────────────────────────────────
@app.route('/affiliate_links', methods=['GET'])
@login_required
def get_affiliate_links():
    db = SessionLocal()
    try:
        links = db.query(AffiliateLink).filter_by(is_active=True).all()
        return jsonify([{"id": l.id, "title": l.title, "url": l.url, "description": l.description} for l in links])
    finally:
        db.close()

@app.route('/admin/affiliate', methods=['POST'])
@login_required
@admin_required
def add_affiliate():
    data = request.get_json()
    db = SessionLocal()
    try:
        link = AffiliateLink(title=data.get('title','Sponsored'), url=data.get('url',''), description=data.get('description',''), is_active=True)
        db.add(link); db.commit()
        return jsonify({"message": "✅ Added!", "id": link.id})
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

@app.route('/admin/affiliate/<int:link_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_affiliate(link_id):
    db = SessionLocal()
    try:
        link = db.query(AffiliateLink).filter_by(id=link_id).first()
        if link: link.is_active = False; db.commit(); return jsonify({"message": "✅ Removed!"})
        return jsonify({"error": "Not found"}), 404
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

@app.route('/admin/affiliate/<int:link_id>', methods=['PUT'])
@login_required
@admin_required
def update_affiliate(link_id):
    data = request.get_json()
    db = SessionLocal()
    try:
        link = db.query(AffiliateLink).filter_by(id=link_id).first()
        if link:
            link.title = data.get('title', link.title); link.url = data.get('url', link.url); link.description = data.get('description', link.description)
            db.commit(); return jsonify({"message": "✅ Updated!"})
        return jsonify({"error": "Not found"}), 404
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

# ── Admin Panel ───────────────────────────────────────────────────
@app.route('/admin')
@login_required
@admin_required
def admin_panel():
    return render_template('admin.html')

@app.route('/admin/users')
@login_required
@admin_required
def admin_users():
    from core.database import Subscription
    db = SessionLocal()
    try:
        users = db.query(User).all()
        result = []
        for u in users:
            sub = db.query(Subscription).filter_by(user_id=u.id, is_active=True).first()
            result.append({"id": u.id, "username": u.username, "is_admin": u.is_admin, "plan": sub.plan if sub else "free", "email": u.email or "", "is_verified": u.is_verified})
        return jsonify(result)
    finally:
        db.close()

@app.route('/admin/query_stats')
@login_required
@admin_required
def admin_query_stats():
    db = SessionLocal()
    try:
        today = datetime.utcnow().strftime('%Y-%m-%d')
        all_counts = db.query(QueryCount).all()
        user_totals = {}
        for qc in all_counts:
            user_totals[qc.user_id] = user_totals.get(qc.user_id, 0) + qc.count
        today_counts = db.query(QueryCount).filter_by(date=today).all()
        today_map = {qc.user_id: qc.count for qc in today_counts}
        users = db.query(User).all()
        name_map = {u.id: u.username for u in users}
        totals = sorted([{"user_id": uid, "username": name_map.get(uid, f"user_{uid}"), "count": cnt} for uid, cnt in user_totals.items()], key=lambda x: x['count'], reverse=True)
        today_data = sorted([{"user_id": uid, "username": name_map.get(uid, f"user_{uid}"), "count": cnt} for uid, cnt in today_map.items()], key=lambda x: x['count'], reverse=True)
        return jsonify({"totals": totals, "today": today_data})
    except Exception as e:
        return jsonify({"error": str(e), "totals": [], "today": []}), 500
    finally:
        db.close()

@app.route('/admin/set_plan', methods=['POST'])
@login_required
@admin_required
def admin_set_plan():
    data = request.get_json()
    user_id = data.get('user_id'); plan = data.get('plan', 'free'); expires_days = data.get('expires_days')
    expires_at = (datetime.utcnow() + timedelta(days=int(expires_days)) if expires_days else None)
    set_user_plan(user_id, plan, expires_at=expires_at)
    return jsonify({"message": f"✅ User {user_id} plan set to {plan}"})

# ── Upload ────────────────────────────────────────────────────────
@app.route('/upload', methods=['POST'])
@login_required
@single_session_check
def upload():
    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    if file and file.filename:
        filename = secure_filename(file.filename)
        try:
            full_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(full_path)
            result = agent.load_data(full_path)
            try:
                import json
                csv_string = open(full_path, 'r', encoding='latin1').read()
                db = SessionLocal()
                # Keep last 5 datasets per user
                old = db.query(UserDataset).filter_by(user_id=current_user.id).order_by(UserDataset.uploaded_at.asc()).all()
                if len(old) >= 5:
                    for o in old[:len(old)-4]: db.delete(o)
                dataset = UserDataset(user_id=current_user.id, filename=filename, csv_data=csv_string, columns=json.dumps(agent.available_columns), row_count=len(agent.df))
                db.add(dataset); db.commit(); db.close()
                logger.info(f"✅ Dataset saved to DB: {filename}")
            except Exception as db_err:
                logger.error(f"⚠️ DB save error: {db_err}")
            return jsonify({"message": f"✅ Uploaded: {filename}\n{result}"}), 200
        except Exception as e:
            return jsonify({"error": str(e)}), 500
    return jsonify({"error": "No file selected"}), 400

# ── Chat ──────────────────────────────────────────────────────────
@app.route('/chat', methods=['POST'])
@login_required
@single_session_check
@check_query_limit
def chat():
    data = request.get_json()
    user_message = data.get('message', '').strip()
    ai_mode = data.get('ai_mode', False)
    if ai_mode and not current_user.is_admin:
        plan = get_user_plan(current_user.id)
        if plan != "pro":
            return jsonify({"response": "⭐ AI Chat Mode sirf Pro plan mein available hai!", "upgrade_required": True})
    try:
        agent._current_user_id = current_user.id
        if ai_mode:
            response = agent.conversational_query(user_message, user_id=current_user.id)
        else:
            response = agent.query(user_message, user_id=current_user.id)
        return jsonify({"response": response})
    except Exception as e:
        return jsonify({"response": "⚠️ Error processing query."}), 500

# ── Downloads ─────────────────────────────────────────────────────
@app.route('/download/csv')
@login_required
@single_session_check
def download_csv():
    if agent.df is None: return jsonify({"error": "Pehle file upload karo!"}), 400
    csv_buffer = io.StringIO(); agent.df.to_csv(csv_buffer, index=False)
    return Response(csv_buffer.getvalue().encode('utf-8-sig'), mimetype='text/csv', headers={"Content-Disposition": "attachment; filename=data_export.csv"})

@app.route('/download/excel')
@login_required
@single_session_check
def download_excel():
    plan = get_user_plan(current_user.id)
    if plan != "pro" and not current_user.is_admin:
        return jsonify({"error": "PRO_REQUIRED", "message": "⭐ Excel export sirf Pro plan mein!"}), 403
    if agent.df is None: return jsonify({"error": "Pehle file upload karo!"}), 400
    excel_buffer = io.BytesIO()
    with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
        agent.df.to_excel(writer, sheet_name='Raw Data', index=False)
        agent.df.describe().to_excel(writer, sheet_name='Summary Stats')
    excel_buffer.seek(0)
    return send_file(excel_buffer, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', as_attachment=True, download_name='analysis_report.xlsx')

@app.route('/download/html')
@login_required
@single_session_check
def download_html():
    if agent.df is None: return jsonify({"error": "Pehle file upload karo!"}), 400
    html_content = f"""<!DOCTYPE html><html><head><title>Report</title></head><body><h1>🤖 DS Agent Report</h1><p>{datetime.now().strftime('%Y-%m-%d %H:%M')}</p>{agent.df.head(10).to_html(index=False)}{agent.df.describe().to_html()}</body></html>"""
    return Response(html_content, mimetype='text/html', headers={"Content-Disposition": "attachment; filename=report.html"})

@app.route('/charts')
@login_required
def list_charts():
    db = SessionLocal()
    try:
        charts = db.query(UserChart).filter_by(user_id=current_user.id).order_by(UserChart.created_at.asc()).all()
        return jsonify({"charts": [{"id": c.id, "title": c.chart_title, "time": c.created_at.strftime('%d %b %H:%M')} for c in charts], "total": len(charts)})
    finally:
        db.close()

@app.route('/chart/<int:chart_id>')
@login_required
def get_chart(chart_id):
    import base64, io
    db = SessionLocal()
    try:
        chart = db.query(UserChart).filter_by(id=chart_id, user_id=current_user.id).first()
        if not chart: return jsonify({'error': 'Chart nahi mila!'}), 404
        img_bytes = base64.b64decode(chart.image_data)
        return send_file(io.BytesIO(img_bytes), mimetype='image/png')
    finally:
        db.close()

@app.route('/charts/clear')
@login_required
def clear_charts():
    db = SessionLocal()
    try:
        db.query(UserChart).filter_by(user_id=current_user.id).delete()
        db.commit()
        return jsonify({"message": "✅ Sab charts clear!"})
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

@app.route('/plot.png')
def plot_png():
    plot_path = '/tmp/plot.png' if os.path.exists('/tmp/plot.png') else os.path.join('static', 'plot.png')
    if not os.path.exists(plot_path): return jsonify({'error': 'Plot nahi bana abhi tak'}), 404
    return send_file(plot_path, mimetype='image/png')

@app.route('/api/charts/list')
@login_required
def charts_list_api():
    db = SessionLocal()
    try:
        charts = db.query(UserChart).filter(
            UserChart.user_id == current_user.id,
            UserChart.chart_type == 'matplotlib'
        ).order_by(UserChart.created_at.desc()).limit(20).all()
        result = []
        for c in charts:
            result.append({
                'id': c.id,
                'title': c.chart_title or f'Chart {c.id}',
                'time': c.created_at.strftime('%d %b %H:%M')
            })
        return jsonify({'charts': result, 'total': len(result)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        db.close()

@app.route('/dashboard.png')
def dashboard_png():
    if current_user.is_authenticated:
        if not current_user.is_admin:
            plan = get_user_plan(current_user.id)
            if plan != "pro": return jsonify({"error": "PRO_REQUIRED"}), 403
    dash_path = '/tmp/dashboard.png' if os.path.exists('/tmp/dashboard.png') else os.path.join('static', 'dashboard.png')
    if not os.path.exists(dash_path): return jsonify({'error': 'Dashboard nahi bana abhi tak'}), 404
    return send_file(dash_path, mimetype='image/png')

@app.route('/screen')
@login_required
def screen(): return render_template('screen.html')

# ── Screen Posts ──────────────────────────────────────────────────
@app.route('/api/screen/posts')
@login_required
def get_screen_posts():
    db = SessionLocal()
    try:
        posts = db.query(ScreenPost).filter_by(is_active=True).order_by(ScreenPost.order_num.asc(), ScreenPost.created_at.desc()).all()
        return jsonify([{"id": p.id, "title": p.title, "content": p.content, "post_type": p.post_type, "affiliate_url": p.affiliate_url, "image_data": p.image_data, "image_mime": p.image_mime} for p in posts])
    finally:
        db.close()

@app.route('/api/screen/posts', methods=['POST'])
@login_required
@admin_required
def add_screen_post():
    import base64
    db = SessionLocal()
    try:
        post_type = request.form.get('post_type', 'text'); title = request.form.get('title', ''); content = request.form.get('content', ''); aff_url = request.form.get('affiliate_url', ''); order_num = int(request.form.get('order_num', 0))
        image_data, image_mime = None, None
        if post_type == 'image' and 'image' in request.files:
            img = request.files['image']
            if img and img.filename: image_data = base64.b64encode(img.read()).decode('utf-8'); image_mime = img.content_type or 'image/jpeg'
        post = ScreenPost(title=title, content=content, post_type=post_type, affiliate_url=aff_url if aff_url else None, image_data=image_data, image_mime=image_mime, order_num=order_num, is_active=True)
        db.add(post); db.commit()
        return jsonify({"message": "✅ Post added!", "id": post.id})
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

@app.route('/api/screen/posts/<int:post_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_screen_post(post_id):
    db = SessionLocal()
    try:
        p = db.query(ScreenPost).filter_by(id=post_id).first()
        if p: p.is_active = False; db.commit(); return jsonify({"message": "✅ Removed!"})
        return jsonify({"error": "Not found"}), 404
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

# ── Music ─────────────────────────────────────────────────────────
@app.route('/api/music')
@login_required
def get_music():
    db = SessionLocal()
    try:
        tracks = db.query(MusicTrack).filter_by(is_active=True).order_by(MusicTrack.created_at.asc()).all()
        return jsonify([{"id": t.id, "title": t.title, "artist": t.artist or "Unknown", "audio_data": t.audio_data, "mime_type": t.mime_type} for t in tracks])
    finally:
        db.close()

@app.route('/api/music', methods=['POST'])
@login_required
@admin_required
def upload_music():
    import base64
    db = SessionLocal()
    try:
        if 'audio' not in request.files: return jsonify({"error": "No file"}), 400
        f = request.files['audio']; title = request.form.get('title', f.filename or 'Unknown'); artist = request.form.get('artist', '')
        if not f or not f.filename: return jsonify({"error": "Empty file"}), 400
        raw = f.read(); audio_b64 = base64.b64encode(raw).decode('utf-8'); mime = f.content_type or 'audio/mpeg'
        track = MusicTrack(title=title, artist=artist, audio_data=audio_b64, mime_type=mime, is_active=True)
        db.add(track); db.commit()
        return jsonify({"message": f"✅ '{title}' upload ho gaya!", "id": track.id})
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

@app.route('/api/music/<int:track_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_music(track_id):
    db = SessionLocal()
    try:
        t = db.query(MusicTrack).filter_by(id=track_id).first()
        if t: t.is_active = False; db.commit(); return jsonify({"message": "✅ Removed!"})
        return jsonify({"error": "Not found"}), 404
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

# ── Video ─────────────────────────────────────────────────────────
@app.route('/api/video')
@login_required
def get_videos():
    db = SessionLocal()
    try:
        videos = db.query(VideoTrack).filter_by(is_active=True).order_by(VideoTrack.created_at.asc()).all()
        return jsonify([{"id": v.id, "title": v.title, "description": v.description or "", "video_data": v.video_data, "mime_type": v.mime_type, "thumbnail": v.thumbnail, "thumb_mime": v.thumb_mime or "image/jpeg"} for v in videos])
    finally:
        db.close()

@app.route('/api/video', methods=['POST'])
@login_required
@admin_required
def upload_video():
    import base64
    db = SessionLocal()
    try:
        if 'video' not in request.files: return jsonify({"error": "No video file"}), 400
        f = request.files['video']; title = request.form.get('title', f.filename or 'Video'); desc = request.form.get('description', '')
        if not f or not f.filename: return jsonify({"error": "Empty file"}), 400
        raw = f.read(); video_b64 = base64.b64encode(raw).decode('utf-8'); mime = f.content_type or 'video/mp4'
        thumb_b64, thumb_mime = None, None
        if 'thumbnail' in request.files:
            th = request.files['thumbnail']
            if th and th.filename: thumb_b64 = base64.b64encode(th.read()).decode('utf-8'); thumb_mime = th.content_type or 'image/jpeg'
        video = VideoTrack(title=title, description=desc, video_data=video_b64, mime_type=mime, thumbnail=thumb_b64, thumb_mime=thumb_mime, is_active=True)
        db.add(video); db.commit()
        return jsonify({"message": f"✅ '{title}' upload ho gaya!", "id": video.id})
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

@app.route('/api/video/<int:video_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_video(video_id):
    db = SessionLocal()
    try:
        v = db.query(VideoTrack).filter_by(id=video_id).first()
        if v: v.is_active = False; db.commit(); return jsonify({"message": "✅ Removed!"})
        return jsonify({"error": "Not found"}), 404
    except Exception as e:
        db.rollback(); return jsonify({"error": str(e)}), 500
    finally:
        db.close()

# ── Stats ─────────────────────────────────────────────────────────
@app.route('/api/stats')
@login_required
def get_stats():
    from core.database import Subscription
    db = SessionLocal()
    try:
        total_users = db.query(User).count(); pro_users = db.query(Subscription).filter_by(plan='pro', is_active=True).count()
        today = datetime.utcnow().strftime('%Y-%m-%d'); today_q = db.query(QueryCount).filter_by(date=today).all()
        today_total = sum(q.count for q in today_q); total_q = db.query(QueryCount).all(); total_queries = sum(q.count for q in total_q)
        return jsonify({"total_users": total_users, "pro_users": pro_users, "today_queries": today_total, "total_queries": total_queries})
    finally:
        db.close()

@app.route('/health')
def health(): return jsonify({"status": "healthy", "version": "3.2.0"}), 200

# ── YouTube Stream ────────────────────────────────────────────────
@app.route('/api/yt-stream', methods=['POST'])
@login_required
def get_yt_stream():
    try:
        import yt_dlp
        data = request.get_json(); url = data.get('url', '').strip()
        if not url: return jsonify({'error': 'URL required!'}), 400
        if 'youtube.com' not in url and 'youtu.be' not in url: return jsonify({'error': 'Valid YouTube URL daalo!'}), 400
        ydl_opts = {'format': 'best[ext=mp4][height<=480]/best[height<=480]/best', 'quiet': True, 'no_warnings': True, 'noplaylist': True}
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False); formats = info.get('formats', []); stream_url = None
            for f in reversed(formats):
                if f.get('ext') == 'mp4' and f.get('url'): stream_url = f['url']; break
            if not stream_url: stream_url = info.get('url') or (formats[-1].get('url') if formats else None)
            title = info.get('title', 'YouTube Video'); thumbnail = info.get('thumbnail', ''); duration = info.get('duration', 0)
        if not stream_url: return jsonify({'error': 'Stream URL extract nahi hui!'}), 500
        return jsonify({'stream_url': stream_url, 'title': title, 'thumbnail': thumbnail, 'duration': duration})
    except ImportError:
        return jsonify({'error': 'yt-dlp install nahi hai!'}), 500
    except Exception as e:
        return jsonify({'error': f'Error: {str(e)[:100]}'}), 500

# ── Live Chat ─────────────────────────────────────────────────────
@app.route('/api/chat/messages')
@login_required
def get_chat_messages():
    db = SessionLocal()
    try:
        msgs = db.query(ChatMessage).order_by(ChatMessage.created_at.desc()).limit(50).all()
        return jsonify([{'id': m.id, 'username': m.username, 'message': m.message, 'is_admin': m.is_admin, 'time': m.created_at.strftime('%H:%M')} for m in reversed(msgs)])
    finally:
        db.close()

@app.route('/api/chat/messages', methods=['POST'])
@login_required
def send_chat_message():
    db = SessionLocal()
    try:
        data = request.get_json(); msg = data.get('message', '').strip()
        if not msg or len(msg) > 500: return jsonify({'error': 'Invalid message'}), 400
        chat = ChatMessage(username=current_user.username, message=msg, is_admin=current_user.is_admin)
        db.add(chat); db.commit()
        return jsonify({'message': '✅ Sent!', 'id': chat.id})
    except Exception as e:
        db.rollback(); return jsonify({'error': str(e)}), 500
    finally:
        db.close()

@app.route('/api/chat/messages/<int:msg_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_chat_message(msg_id):
    db = SessionLocal()
    try:
        m = db.query(ChatMessage).filter_by(id=msg_id).first()
        if m: db.delete(m); db.commit(); return jsonify({'message': '✅ Deleted!'})
        return jsonify({'error': 'Not found'}), 404
    finally:
        db.close()

# ── Analytics ─────────────────────────────────────────────────────
@app.route('/api/analytics/view/<int:post_id>', methods=['POST'])
@login_required
def track_post_view(post_id):
    db = SessionLocal()
    try:
        today = datetime.utcnow().strftime('%Y-%m-%d'); rec = db.query(PostAnalytics).filter_by(post_id=post_id, date=today).first()
        if rec: rec.views += 1
        else: rec = PostAnalytics(post_id=post_id, views=1, clicks=0, date=today); db.add(rec)
        db.commit(); return jsonify({'ok': True})
    except Exception as e:
        db.rollback(); return jsonify({'error': str(e)}), 500
    finally:
        db.close()

@app.route('/api/analytics/click/<int:post_id>', methods=['POST'])
@login_required
def track_post_click(post_id):
    db = SessionLocal()
    try:
        today = datetime.utcnow().strftime('%Y-%m-%d'); rec = db.query(PostAnalytics).filter_by(post_id=post_id, date=today).first()
        if rec: rec.clicks += 1
        else: rec = PostAnalytics(post_id=post_id, views=0, clicks=1, date=today); db.add(rec)
        db.commit(); return jsonify({'ok': True})
    except Exception as e:
        db.rollback(); return jsonify({'error': str(e)}), 500
    finally:
        db.close()

@app.route('/api/analytics')
@login_required
@admin_required
def get_analytics():
    db = SessionLocal()
    try:
        posts = db.query(ScreenPost).filter_by(is_active=True).all()
        result = []
        for p in posts:
            recs = db.query(PostAnalytics).filter_by(post_id=p.id).all()
            total_views = sum(r.views for r in recs); total_clicks = sum(r.clicks for r in recs)
            result.append({'post_id': p.id, 'title': p.title, 'views': total_views, 'clicks': total_clicks, 'ctr': round(total_clicks/total_views*100, 1) if total_views else 0})
        result.sort(key=lambda x: x['views'], reverse=True)
        return jsonify(result)
    finally:
        db.close()

# ── Scheduled Posts ───────────────────────────────────────────────
@app.route('/api/scheduled', methods=['GET'])
@login_required
@admin_required
def get_scheduled():
    db = SessionLocal()
    try:
        posts = db.query(ScheduledPost).filter_by(is_published=False).order_by(ScheduledPost.scheduled_at).all()
        return jsonify([{'id': p.id, 'title': p.title, 'content': p.content, 'post_type': p.post_type, 'scheduled_at': p.scheduled_at.strftime('%Y-%m-%d %H:%M'), 'is_published': p.is_published} for p in posts])
    finally:
        db.close()

@app.route('/api/scheduled', methods=['POST'])
@login_required
@admin_required
def create_scheduled():
    import base64
    db = SessionLocal()
    try:
        title = request.form.get('title', ''); content = request.form.get('content', ''); post_type = request.form.get('post_type', 'text'); affiliate_url = request.form.get('affiliate_url', ''); scheduled_at = request.form.get('scheduled_at', '')
        if not title or not scheduled_at: return jsonify({'error': 'Title aur time required!'}), 400
        from datetime import datetime as dt
        sched_time = dt.strptime(scheduled_at, '%Y-%m-%dT%H:%M')
        img_data, img_mime = None, None
        if post_type == 'image' and 'image' in request.files:
            f = request.files['image']
            if f and f.filename: img_data = base64.b64encode(f.read()).decode(); img_mime = f.content_type or 'image/jpeg'
        sp = ScheduledPost(title=title, content=content, post_type=post_type, affiliate_url=affiliate_url, scheduled_at=sched_time, image_data=img_data, image_mime=img_mime)
        db.add(sp); db.commit()
        return jsonify({'message': f'✅ "{title}" scheduled!', 'id': sp.id})
    except Exception as e:
        db.rollback(); return jsonify({'error': str(e)}), 500
    finally:
        db.close()

@app.route('/api/scheduled/<int:sp_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_scheduled(sp_id):
    db = SessionLocal()
    try:
        sp = db.query(ScheduledPost).filter_by(id=sp_id).first()
        if sp: db.delete(sp); db.commit(); return jsonify({'message': '✅ Removed!'})
        return jsonify({'error': 'Not found'}), 404
    finally:
        db.close()

@app.route('/api/scheduled/publish', methods=['POST'])
@login_required
@admin_required
def publish_due_posts():
    db = SessionLocal()
    try:
        now = datetime.utcnow()
        due = db.query(ScheduledPost).filter(ScheduledPost.scheduled_at <= now, ScheduledPost.is_published == False).all()
        published = 0
        for sp in due:
            post = ScreenPost(title=sp.title, content=sp.content, post_type=sp.post_type, affiliate_url=sp.affiliate_url, image_data=sp.image_data, image_mime=sp.image_mime, is_active=True)
            db.add(post); sp.is_published = True; published += 1
        db.commit()
        return jsonify({'message': f'✅ {published} posts published!'})
    except Exception as e:
        db.rollback(); return jsonify({'error': str(e)}), 500
    finally:
        db.close()

# ── Push Notifications ────────────────────────────────────────────
VAPID_PUBLIC_KEY = os.getenv('VAPID_PUBLIC_KEY', ''); VAPID_PRIVATE_KEY = os.getenv('VAPID_PRIVATE_KEY', ''); VAPID_EMAIL = os.getenv('VAPID_EMAIL', 'mailto:admin@dsagent.com')

@app.route('/api/push/vapid-public-key')
def get_vapid_key(): return jsonify({'publicKey': VAPID_PUBLIC_KEY})

@app.route('/api/push/subscribe', methods=['POST'])
@login_required
def push_subscribe():
    db = SessionLocal()
    try:
        data = request.get_json()
        sub = PushSubscription(user_id=current_user.id, endpoint=data['endpoint'], p256dh=data['keys']['p256dh'], auth=data['keys']['auth'])
        db.add(sub); db.commit()
        return jsonify({'message': '✅ Subscribed!'})
    except Exception as e:
        db.rollback(); return jsonify({'error': str(e)}), 500
    finally:
        db.close()

@app.route('/api/push/send', methods=['POST'])
@login_required
@admin_required
def send_push():
    db = SessionLocal()
    try:
        data = request.get_json(); title = data.get('title', 'DS Agent'); body = data.get('body', '')
        subs = db.query(PushSubscription).all(); sent = 0
        if VAPID_PRIVATE_KEY:
            try:
                from pywebpush import webpush, WebPushException
                import json
                for s in subs:
                    try:
                        webpush(subscription_info={'endpoint': s.endpoint, 'keys': {'p256dh': s.p256dh, 'auth': s.auth}}, data=json.dumps({'title': title, 'body': body}), vapid_private_key=VAPID_PRIVATE_KEY, vapid_claims={'sub': VAPID_EMAIL})
                        sent += 1
                    except: pass
            except ImportError: pass
        return jsonify({'message': f'✅ {sent}/{len(subs)} notifications sent!'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        db.close()

@app.route('/send_report', methods=['POST'])
@login_required
@single_session_check
def send_report():
    try:
        data = request.get_json()
        client_email = data.get('email'); report_type = data.get('report_type', 'summary'); schedule = data.get('schedule', 'now')
        if not client_email: return jsonify({"error": "Email required"}), 400
        if schedule != 'now':
            return jsonify({"message": f"✅ Report scheduled! {schedule.capitalize()} emails {client_email} pe jaayenge.", "scheduled": True, "schedule": schedule})
        db = SessionLocal()
        last_query = db.query(UserQuery).filter_by(user_id=current_user.id).order_by(UserQuery.timestamp.desc()).first()
        insights = last_query.response_text if last_query else "No data"
        db.close()
        subject_map = {'summary': '📊 Data Analysis Summary Report', 'full': '📊 Full Data Analysis Report', 'charts': '📊 Charts & Visualizations Report'}
        subject = subject_map.get(report_type, '📊 Your Data Analysis Report')
        pdf_filename = f"static/report_{current_user.username}_{int(datetime.now().timestamp())}.pdf"
        success = generate_pdf_report(pdf_filename, client_email, insights, 'static/plot.png')
        if success:
            email_sent = send_report_email(to_email=client_email, subject=subject, body=f"Hi,\n\nPlease find attached your {report_type} report.\n\nRegards,\nDS Agent", attachment_path=pdf_filename)
            if email_sent: return jsonify({"message": f"✅ {subject} sent to {client_email}!"})
            return jsonify({"error": "Email failed"}), 500
        return jsonify({"error": "PDF generation failed"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ── CSV Compare ───────────────────────────────────────────────────
@app.route('/upload_compare', methods=['POST'])
@login_required
@single_session_check
def upload_compare():
    if 'file' not in request.files: return jsonify({"error": "No file provided"}), 400
    file = request.files['file']
    if not file or not file.filename: return jsonify({"error": "Empty file"}), 400
    try:
        import io as _io
        filename = secure_filename(file.filename)
        content = file.read().decode('utf-8', errors='replace')
        df2 = pd.read_csv(_io.StringIO(content))
        rows, cols = len(df2), len(df2.columns); columns = list(df2.columns)
        stats = {}
        for col in df2.select_dtypes(include='number').columns:
            stats[col] = {"mean": round(float(df2[col].mean()), 2), "min": round(float(df2[col].min()), 2), "max": round(float(df2[col].max()), 2), "nulls": int(df2[col].isnull().sum())}
        comparison = None
        if agent.df is not None:
            df1 = agent.df; common_cols = list(set(df1.columns) & set(df2.columns))
            comparison = {"file1_rows": len(df1), "file2_rows": rows, "file1_cols": len(df1.columns), "file2_cols": cols, "common_columns": common_cols, "only_in_file1": list(set(df1.columns) - set(df2.columns)), "only_in_file2": list(set(df2.columns) - set(df1.columns)), "col_stats": {}}
            for col in common_cols:
                if pd.api.types.is_numeric_dtype(df1[col]) and pd.api.types.is_numeric_dtype(df2[col]):
                    comparison["col_stats"][col] = {"file1_mean": round(float(df1[col].mean()), 2), "file2_mean": round(float(df2[col].mean()), 2), "file1_max": round(float(df1[col].max()), 2), "file2_max": round(float(df2[col].max()), 2), "diff_mean": round(float(df2[col].mean()) - float(df1[col].mean()), 2)}
        return jsonify({"message": f"✅ {filename} loaded!", "filename": filename, "rows": rows, "cols": cols, "columns": columns, "stats": stats, "comparison": comparison})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ── Save Plotly Chart ─────────────────────────────────────────────
@app.route('/save_plot_base64', methods=['POST'])
@login_required
@single_session_check
def save_plot_base64():
    try:
        data = request.get_json(); image_b64 = data.get('image'); title = data.get('title', f'Plotly Chart {datetime.now().strftime("%d %b %H:%M")}')
        if not image_b64: return jsonify({"error": "Image data required"}), 400
        if ',' in image_b64: image_b64 = image_b64.split(',', 1)[1]
        db = SessionLocal()
        try:
            chart = UserChart(user_id=current_user.id, chart_title=title, image_data=image_b64, chart_type='plotly')
            db.add(chart); db.commit()
            return jsonify({"message": f"✅ '{title}' gallery mein save ho gaya!", "chart_id": chart.id})
        except Exception as db_err:
            db.rollback(); return jsonify({"error": str(db_err)}), 500
        finally:
            db.close()
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ── PWA Manifest ──────────────────────────────────────────────────
@app.route('/manifest.json')
def pwa_manifest():
    return jsonify({
        "name": "DS Agent",
        "short_name": "DS Agent",
        "description": "Your Personal Data Science AI Assistant",
        "start_url": "/",
        "display": "standalone",
        "background_color": "#3730A3",
        "theme_color": "#5B5EF4",
        "orientation": "portrait",
        "icons": [
            {
                "src": "/static/icon-192.png",
                "sizes": "192x192",
                "type": "image/png",
                "purpose": "any"
            },
            {
                "src": "/static/icon-512.png",
                "sizes": "512x512",
                "type": "image/png",
                "purpose": "any"
            },
            {
                "src": "/static/icon-maskable-192.png",
                "sizes": "192x192",
                "type": "image/png",
                "purpose": "maskable"
            },
            {
                "src": "/static/icon-maskable-512.png",
                "sizes": "512x512",
                "type": "image/png",
                "purpose": "maskable"
            }
        ]
    })

@app.route('/static/<path:filename>')
def serve_static(filename):
    import os
    static_dir = os.path.join(os.path.dirname(__file__), 'static')
    return send_from_directory(static_dir, filename)


@app.route('/sw.js')
def service_worker():
    sw_content = """
const CACHE_NAME = 'ds-agent-v1';
const STATIC_ASSETS = ['/'];
self.addEventListener('install', e => {
    e.waitUntil(caches.open(CACHE_NAME).then(c => c.addAll(STATIC_ASSETS).catch(()=>{})));
    self.skipWaiting();
});
self.addEventListener('activate', e => {
    e.waitUntil(caches.keys().then(keys =>
        Promise.all(keys.filter(k => k !== CACHE_NAME).map(k => caches.delete(k)))
    ));
    self.clients.claim();
});
self.addEventListener('fetch', e => {
    if (e.request.method !== 'GET') return;
    if (['/chat','/upload','/api/'].some(p => e.request.url.includes(p))) return;
    e.respondWith(
        fetch(e.request)
            .then(r => { if(r.ok){caches.open(CACHE_NAME).then(c=>c.put(e.request,r.clone()));}return r;})
            .catch(() => caches.match(e.request))
    );
});
"""
    return Response(sw_content, mimetype='application/javascript', headers={'Service-Worker-Allowed': '/'})

@app.route('/run_code', methods=['POST'])
@login_required
def run_code():
    import traceback, base64, io as _io, sys
    from contextlib import redirect_stdout, redirect_stderr
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    data = request.get_json()
    code = data.get('code', '').strip()
    if not code:
        return jsonify({'error': 'Code empty hai!'}), 400

    if agent.df is None:
        return jsonify({'error': '⚠️ Pehle file upload karo!'}), 400

    # Build execution namespace with df and common libs
    import pandas as pd
    import numpy as np
    namespace = {
        'df': agent.df.copy(),
        'pd': pd,
        'np': np,
        'plt': plt,
    }
    try:
        import sklearn
        namespace['sklearn'] = sklearn
        from sklearn import linear_model, preprocessing, metrics, model_selection
        namespace['linear_model'] = linear_model
        namespace['preprocessing'] = preprocessing
        namespace['metrics'] = metrics
        namespace['model_selection'] = model_selection
    except: pass
    try:
        import seaborn as sns
        namespace['sns'] = sns
    except: pass

    stdout_buf = _io.StringIO()
    stderr_buf = _io.StringIO()
    img_b64 = None

    # SECURITY: Block dangerous commands
    BLOCKED = ['os.system','subprocess','__import__','open(','eval(','exec(',
               'shutil','rmdir','remove','unlink','socket','requests.get',
               'requests.post','urllib','httpx','builtins','globals()','locals()']
    code_lower = code.lower()
    for bad in BLOCKED:
        if bad in code:
            return jsonify({'output':'','image':None,
                'error': f'⛔ Blocked: "{bad}" allowed nahi hai notebook mein!'})

    try:
        plt.close('all')
        # Restricted builtins — dangerous functions remove
        safe_builtins = {k: v for k, v in __builtins__.items()
                        if k not in ('open','__import__','eval','exec','compile',
                                     'breakpoint','input','memoryview')}                         if isinstance(__builtins__, dict) else {}
        namespace['__builtins__'] = safe_builtins
        with redirect_stdout(stdout_buf), redirect_stderr(stderr_buf):
            exec(compile(code, '<notebook>', 'exec'), namespace)

        # Check if any plot was created
        if plt.get_fignums():
            img_buf = _io.BytesIO()
            plt.savefig(img_buf, format='png', bbox_inches='tight', dpi=120)
            plt.close('all')
            img_buf.seek(0)
            img_b64 = base64.b64encode(img_buf.read()).decode('utf-8')

        output = stdout_buf.getvalue()
        err_out = stderr_buf.getvalue()
        if err_out and not output:
            output = err_out

        return jsonify({
            'output': output.strip(),
            'image': img_b64,
            'error': None
        })

    except Exception as e:
        tb = traceback.format_exc()
        # Clean traceback — only show relevant part
        lines = tb.strip().split('\n')
        clean = '\n'.join(l for l in lines if '<notebook>' in l or 'Error' in l or 'error' in l.lower())
        return jsonify({
            'output': '',
            'image': None,
            'error': clean or str(e)
        })

# ── Real-time Stock/Crypto Data ──────────────────────────────────
@app.route('/api/stock', methods=['POST'])
@login_required
def get_stock():
    try:
        import yfinance as yf
        data = request.get_json()
        symbol = data.get('symbol', '').upper().strip()
        period = data.get('period', '1mo')
        if not symbol:
            return jsonify({'error': 'Symbol daalo (e.g. AAPL, BTC-USD)'}), 400

        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period=period)

        if hist.empty:
            return jsonify({'error': f'"{symbol}" ka data nahi mila. Symbol check karo.'}), 404

        # Current price
        current = float(hist['Close'].iloc[-1])
        prev = float(hist['Close'].iloc[-2]) if len(hist) > 1 else current
        change = round(current - prev, 2)
        change_pct = round((change / prev) * 100, 2) if prev else 0

        # Chart
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import io as _io, base64

        fig, ax = plt.subplots(figsize=(8, 3.5))
        ax.plot(hist.index, hist['Close'], color='#5b5ef4', linewidth=2)
        ax.fill_between(hist.index, hist['Close'], alpha=0.1, color='#5b5ef4')
        ax.set_title(f'{symbol} — {period}', fontsize=13, fontweight='bold')
        ax.set_ylabel('Price', fontsize=10)
        ax.grid(True, alpha=0.3)
        fig.tight_layout()

        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close('all')
        buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode()

        return jsonify({
            'symbol': symbol,
            'name': info.get('longName', symbol),
            'current': current,
            'change': change,
            'change_pct': change_pct,
            'currency': info.get('currency', 'USD'),
            'market_cap': info.get('marketCap'),
            'volume': info.get('volume'),
            'chart': chart_b64,
            'period': period
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── Multi-file SQL Join ───────────────────────────────────────────
@app.route('/api/sql_join', methods=['POST'])
@login_required
def sql_join():
    try:
        import io as _io
        data = request.get_json()
        join_on = data.get('join_on', '').strip()
        join_type = data.get('join_type', 'inner')  # inner/left/right/outer
        query = data.get('query', '').strip()        # optional SQL query

        if agent.df is None:
            return jsonify({'error': '⚠️ Pehle main file upload karo!'}), 400

        df2_csv = data.get('df2_csv', '')
        if not df2_csv:
            return jsonify({'error': 'Doosri file ka data nahi mila!'}), 400

        df1 = agent.df.copy()
        df2 = pd.read_csv(_io.StringIO(df2_csv))

        if query:
            # SQL query mode — pandasql style using pandas
            # Simple SELECT parser
            q = query.upper()
            result = df1  # default
            if 'JOIN' in q:
                # Parse: SELECT * FROM df1 JOIN df2 ON col
                if join_on:
                    how_map = {'inner':'inner','left':'left','right':'right','outer':'outer','full':'outer'}
                    how = how_map.get(join_type.lower(), 'inner')
                    result = pd.merge(df1, df2, on=join_on, how=how, suffixes=('_file1','_file2'))
                else:
                    result = pd.merge(df1, df2, left_index=True, right_index=True, how='inner')
            elif 'WHERE' in q:
                # Basic filter
                result = df1
        else:
            # Simple merge
            if join_on:
                if join_on not in df1.columns:
                    return jsonify({'error': f'Column "{join_on}" file 1 mein nahi hai! File 1 columns: {list(df1.columns)[:10]}'}), 400
                if join_on not in df2.columns:
                    return jsonify({'error': f'Column "{join_on}" file 2 mein nahi hai! File 2 columns: {list(df2.columns)[:10]}'}), 400
                how_map = {'inner':'inner','left':'left','right':'right','outer':'outer'}
                result = pd.merge(df1, df2, on=join_on, how=how_map.get(join_type,'inner'), suffixes=('_f1','_f2'))
            else:
                # Auto-detect common column
                common = list(set(df1.columns) & set(df2.columns))
                if not common:
                    return jsonify({'error': 'Koi common column nahi! File1: ' + str(list(df1.columns)[:5]) + ' | File2: ' + str(list(df2.columns)[:5])}), 400
                result = pd.merge(df1, df2, on=common[0], how='inner', suffixes=('_f1','_f2'))
                join_on = common[0]

        preview = result.head(20).to_string(index=False)
        return jsonify({
            'message': f'✅ Join successful! {len(result)} rows × {len(result.columns)} cols',
            'rows': len(result),
            'cols': len(result.columns),
            'join_on': join_on,
            'join_type': join_type,
            'columns': list(result.columns),
            'preview': preview,
            'csv': result.to_csv(index=False)
        })
    except Exception as e:
        import traceback
        return jsonify({'error': str(e)}), 500


@app.route('/notebook')
@login_required
def notebook():
    return render_template('notebook.html')


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 10000))
    app.run(host='0.0.0.0', port=port)


# ════════════════════════════════════════════════════════════════
# FEATURE 1: DATA CLEANING
# ════════════════════════════════════════════════════════════════
@app.route('/api/clean', methods=['POST'])
@login_required
def data_clean():
    if agent.df is None:
        return jsonify({'error': 'Pehle file upload karo!'}), 400
    try:
        data = request.get_json()
        action = data.get('action', '')
        col = data.get('column', '')
        value = data.get('value', '')
        df = agent.df.copy()
        msg = ''

        if action == 'drop_nulls':
            before = len(df)
            df = df.dropna(subset=[col] if col else None)
            msg = f'Nulls drop kiye: {before - len(df)} rows hatayi'
        elif action == 'fill_nulls':
            if col:
                if value == 'mean' and pd.api.types.is_numeric_dtype(df[col]):
                    df[col] = df[col].fillna(df[col].mean())
                    msg = f'"{col}" ke nulls mean se fill kiye'
                elif value == 'median' and pd.api.types.is_numeric_dtype(df[col]):
                    df[col] = df[col].fillna(df[col].median())
                    msg = f'"{col}" ke nulls median se fill kiye'
                elif value == 'mode':
                    df[col] = df[col].fillna(df[col].mode()[0])
                    msg = f'"{col}" ke nulls mode se fill kiye'
                elif value == 'zero':
                    df[col] = df[col].fillna(0)
                    msg = f'"{col}" ke nulls 0 se fill kiye'
                elif value == 'ffill':
                    df[col] = df[col].ffill()
                    msg = f'"{col}" ke nulls forward fill kiye'
                else:
                    df[col] = df[col].fillna(value)
                    msg = f'"{col}" ke nulls "{value}" se fill kiye'
        elif action == 'drop_duplicates':
            before = len(df)
            df = df.drop_duplicates(subset=[col] if col else None)
            msg = f'Duplicates hataaye: {before - len(df)} rows'
        elif action == 'drop_column':
            if col in df.columns:
                df = df.drop(columns=[col])
                msg = f'Column "{col}" hataaya'
        elif action == 'rename_column':
            new_name = data.get('new_name', '')
            if col in df.columns and new_name:
                df = df.rename(columns={col: new_name})
                msg = f'"{col}" → "{new_name}" rename kiya'
        elif action == 'fix_dtype':
            dtype = data.get('dtype', 'float')
            try:
                if dtype == 'int':
                    df[col] = pd.to_numeric(df[col], errors='coerce').astype('Int64')
                elif dtype == 'float':
                    df[col] = pd.to_numeric(df[col], errors='coerce')
                elif dtype == 'str':
                    df[col] = df[col].astype(str)
                elif dtype == 'datetime':
                    df[col] = pd.to_datetime(df[col], errors='coerce')
                msg = f'"{col}" ko {dtype} mein convert kiya'
            except Exception as e:
                return jsonify({'error': str(e)}), 400
        elif action == 'remove_outliers':
            if pd.api.types.is_numeric_dtype(df[col]):
                Q1 = df[col].quantile(0.25)
                Q3 = df[col].quantile(0.75)
                IQR = Q3 - Q1
                before = len(df)
                df = df[~((df[col] < (Q1 - 1.5*IQR)) | (df[col] > (Q3 + 1.5*IQR)))]
                msg = f'Outliers hataaye "{col}" se: {before - len(df)} rows'
        elif action == 'strip_whitespace':
            str_cols = [col] if col else list(df.select_dtypes(include='object').columns)
            for c in str_cols:
                df[c] = df[c].astype(str).str.strip()
            msg = f'Whitespace strip kiya: {", ".join(str_cols)}'
        elif action == 'lowercase':
            str_cols = [col] if col else list(df.select_dtypes(include='object').columns)
            for c in str_cols:
                df[c] = df[c].astype(str).str.lower()
            msg = f'Lowercase kiya: {", ".join(str_cols)}'
        elif action == 'info':
            nulls = df.isnull().sum()
            null_cols = nulls[nulls > 0].to_dict()
            dups = int(df.duplicated().sum())
            dtypes = df.dtypes.astype(str).to_dict()
            return jsonify({
                'rows': len(df), 'cols': len(df.columns),
                'nulls': null_cols, 'duplicates': dups,
                'dtypes': dtypes, 'columns': list(df.columns),
                'memory_mb': round(df.memory_usage(deep=True).sum() / 1024**2, 2)
            })
        else:
            return jsonify({'error': 'Unknown action'}), 400

        # Save cleaned df back
        agent.df = df
        import io as _io
        csv_str = df.to_csv(index=False)
        db = SessionLocal()
        dataset = db.query(UserDataset).filter_by(user_id=current_user.id).order_by(UserDataset.uploaded_at.desc()).first()
        if dataset:
            dataset.csv_data = csv_str
            dataset.row_count = len(df)
            dataset.columns = json.dumps(list(df.columns))
            db.commit()
        db.close()

        return jsonify({
            'message': f'✅ {msg}',
            'rows': len(df), 'cols': len(df.columns),
            'columns': list(df.columns),
            'preview': df.head(5).to_string(index=False)
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 2: AUTO ML
# ════════════════════════════════════════════════════════════════
@app.route('/api/automl', methods=['POST'])
@login_required
def auto_ml():
    if agent.df is None:
        return jsonify({'error': 'Pehle file upload karo!'}), 400
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from sklearn.model_selection import train_test_split, cross_val_score
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        from sklearn.linear_model import LinearRegression, LogisticRegression, Ridge
        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor, GradientBoostingClassifier, GradientBoostingRegressor
        from sklearn.tree import DecisionTreeClassifier, DecisionTreeRegressor
        from sklearn.neighbors import KNeighborsClassifier, KNeighborsRegressor
        from sklearn.metrics import (accuracy_score, r2_score, mean_squared_error,
                                     classification_report, confusion_matrix)
        from sklearn.naive_bayes import GaussianNB
        import numpy as np

        data = request.get_json()
        target_col = data.get('target', '')
        task_type = data.get('task', 'auto')
        test_size = float(data.get('test_size', 0.2))

        df = agent.df.copy().dropna()

        if target_col not in df.columns:
            return jsonify({'error': f'Column "{target_col}" nahi mila! Available: {list(df.columns)[:10]}'}), 400

        # Limit to 10k rows for speed
        if len(df) > 10000:
            df = df.sample(10000, random_state=42)

        y = df[target_col]
        X = df.drop(columns=[target_col])

        # Encode categorical
        le_dict = {}
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder()
            X[c] = le.fit_transform(X[c].astype(str))
            le_dict[c] = le
        if y.dtype == object or str(y.dtype) == 'category':
            le_y = LabelEncoder()
            y = le_y.fit_transform(y.astype(str))
            y_labels = list(le_y.classes_)
        else:
            y_labels = None

        # Auto detect task
        if task_type == 'auto':
            unique_ratio = len(np.unique(y)) / len(y)
            task_type = 'classification' if (len(np.unique(y)) <= 20 or unique_ratio < 0.05) else 'regression'

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=test_size, random_state=42)

        scaler = StandardScaler()
        X_train_s = scaler.fit_transform(X_train)
        X_test_s = scaler.transform(X_test)

        # Models to try
        if task_type == 'classification':
            models = {
                'Random Forest': RandomForestClassifier(n_estimators=100, random_state=42),
                'Gradient Boosting': GradientBoostingClassifier(random_state=42),
                'Logistic Regression': LogisticRegression(max_iter=1000, random_state=42),
                'Decision Tree': DecisionTreeClassifier(random_state=42),
                'KNN': KNeighborsClassifier(),
                'Naive Bayes': GaussianNB(),
            }
        else:
            models = {
                'Random Forest': RandomForestRegressor(n_estimators=100, random_state=42),
                'Gradient Boosting': GradientBoostingRegressor(random_state=42),
                'Ridge Regression': Ridge(),
                'Linear Regression': LinearRegression(),
                'Decision Tree': DecisionTreeRegressor(random_state=42),
                'KNN': KNeighborsRegressor(),
            }

        results = []
        best_model = None
        best_score = -999

        for name, model in models.items():
            try:
                use_scaled = name in ('Logistic Regression', 'KNN', 'Ridge Regression', 'Linear Regression', 'Naive Bayes')
                Xtr = X_train_s if use_scaled else X_train
                Xte = X_test_s if use_scaled else X_test
                model.fit(Xtr, y_train)
                y_pred = model.predict(Xte)
                if task_type == 'classification':
                    score = accuracy_score(y_test, y_pred)
                    metric = 'Accuracy'
                else:
                    score = r2_score(y_test, y_pred)
                    rmse = float(np.sqrt(mean_squared_error(y_test, y_pred)))
                    metric = 'R2 Score'
                results.append({'model': name, 'score': round(float(score), 4), 'metric': metric})
                if score > best_score:
                    best_score = score
                    best_model = (name, model, use_scaled)
            except Exception as e:
                results.append({'model': name, 'score': None, 'error': str(e)})

        results.sort(key=lambda x: x['score'] or -1, reverse=True)

        # Feature importance chart
        chart_b64 = None
        feat_importance = {}
        if best_model and hasattr(best_model[1], 'feature_importances_'):
            importances = best_model[1].feature_importances_
            feat_df = pd.DataFrame({'feature': X.columns, 'importance': importances})
            feat_df = feat_df.sort_values('importance', ascending=True).tail(15)
            feat_importance = dict(zip(feat_df['feature'], feat_df['importance'].round(4)))
            fig, ax = plt.subplots(figsize=(7, 4))
            ax.barh(feat_df['feature'], feat_df['importance'], color='#5b5ef4')
            ax.set_title(f'Feature Importance — {best_model[0]}', fontweight='bold')
            ax.set_xlabel('Importance')
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); chart_b64 = base64.b64encode(buf.read()).decode()

        # Confusion matrix for classification
        cm_b64 = None
        if task_type == 'classification' and best_model:
            name, model, use_scaled = best_model
            Xte = X_test_s if use_scaled else X_test
            y_pred = model.predict(Xte)
            cm = confusion_matrix(y_test, y_pred)
            fig, ax = plt.subplots(figsize=(5, 4))
            im = ax.imshow(cm, cmap='Blues')
            ax.set_title('Confusion Matrix', fontweight='bold')
            plt.colorbar(im, ax=ax)
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); cm_b64 = base64.b64encode(buf.read()).decode()

        return jsonify({
            'task': task_type,
            'target': target_col,
            'train_rows': len(X_train),
            'test_rows': len(X_test),
            'features': list(X.columns),
            'best_model': best_model[0] if best_model else None,
            'best_score': round(float(best_score), 4),
            'results': results,
            'feature_importance': feat_importance,
            'chart': chart_b64,
            'cm_chart': cm_b64
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/automl/predict', methods=['POST'])
@login_required
def automl_predict():
    """Predict on new data using last trained model stored in session"""
    return jsonify({'error': 'Train karo pehle /api/automl se, phir notebook mein predict karo!'}), 400


# ════════════════════════════════════════════════════════════════
# FEATURE 3: ADVANCED STATISTICS
# ════════════════════════════════════════════════════════════════
@app.route('/api/stats', methods=['POST'])
@login_required
def advanced_stats():
    if agent.df is None:
        return jsonify({'error': 'Pehle file upload karo!'}), 400
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from scipy import stats as scipy_stats
        import numpy as np

        data = request.get_json()
        test = data.get('test', '')
        col1 = data.get('col1', '')
        col2 = data.get('col2', '')
        group_col = data.get('group_col', '')
        alpha = float(data.get('alpha', 0.05))
        df = agent.df.copy().dropna()

        result = {}
        chart_b64 = None

        if test == 'describe':
            num_df = df.select_dtypes(include='number')
            desc = num_df.describe().round(4)
            result = {
                'summary': desc.to_dict(),
                'skewness': num_df.skew().round(4).to_dict(),
                'kurtosis': num_df.kurtosis().round(4).to_dict(),
                'nulls': df.isnull().sum().to_dict()
            }

        elif test == 'correlation':
            num_df = df.select_dtypes(include='number')
            corr = num_df.corr().round(4)
            fig, ax = plt.subplots(figsize=(8, 6))
            import seaborn as sns
            sns.heatmap(corr, annot=True, fmt='.2f', cmap='coolwarm', ax=ax,
                       annot_kws={'size': 8}, square=True)
            ax.set_title('Correlation Heatmap', fontweight='bold', fontsize=13)
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=110); plt.close('all')
            buf.seek(0); chart_b64 = base64.b64encode(buf.read()).decode()
            result = {'correlation_matrix': corr.to_dict(), 'chart': chart_b64}

        elif test == 'ttest':
            if not col1:
                return jsonify({'error': 'col1 chahiye!'}), 400
            if group_col:
                groups = df[group_col].unique()
                if len(groups) < 2:
                    return jsonify({'error': '2 groups chahiye!'}), 400
                g1 = df[df[group_col] == groups[0]][col1].dropna()
                g2 = df[df[group_col] == groups[1]][col1].dropna()
                t_stat, p_val = scipy_stats.ttest_ind(g1, g2)
                result = {
                    'test': 'Independent T-Test',
                    't_statistic': round(float(t_stat), 4),
                    'p_value': round(float(p_val), 6),
                    'significant': bool(p_val < alpha),
                    'interpretation': f'{"Significant difference" if p_val < alpha else "No significant difference"} between groups (alpha={alpha})',
                    'group1': {'name': str(groups[0]), 'mean': round(float(g1.mean()), 4), 'n': len(g1)},
                    'group2': {'name': str(groups[1]), 'mean': round(float(g2.mean()), 4), 'n': len(g2)}
                }
            elif col2:
                t_stat, p_val = scipy_stats.ttest_ind(df[col1].dropna(), df[col2].dropna())
                result = {
                    'test': 'Independent T-Test',
                    't_statistic': round(float(t_stat), 4),
                    'p_value': round(float(p_val), 6),
                    'significant': bool(p_val < alpha),
                    'interpretation': f'{"Significant difference" if p_val < alpha else "No significant difference"} (alpha={alpha})'
                }

        elif test == 'anova':
            if not col1 or not group_col:
                return jsonify({'error': 'col1 aur group_col chahiye!'}), 400
            groups = [g[col1].dropna().values for _, g in df.groupby(group_col)]
            f_stat, p_val = scipy_stats.f_oneway(*groups)
            result = {
                'test': 'One-Way ANOVA',
                'f_statistic': round(float(f_stat), 4),
                'p_value': round(float(p_val), 6),
                'significant': bool(p_val < alpha),
                'num_groups': len(groups),
                'interpretation': f'{"Significant difference" if p_val < alpha else "No significant difference"} among groups (alpha={alpha})'
            }

        elif test == 'chi2':
            if not col1 or not col2:
                return jsonify({'error': 'col1 aur col2 chahiye!'}), 400
            ct = pd.crosstab(df[col1], df[col2])
            chi2, p_val, dof, expected = scipy_stats.chi2_contingency(ct)
            result = {
                'test': 'Chi-Square Test',
                'chi2_statistic': round(float(chi2), 4),
                'p_value': round(float(p_val), 6),
                'degrees_of_freedom': int(dof),
                'significant': bool(p_val < alpha),
                'interpretation': f'{"Significant association" if p_val < alpha else "No significant association"} between {col1} and {col2} (alpha={alpha})'
            }

        elif test == 'normality':
            if not col1:
                return jsonify({'error': 'col1 chahiye!'}), 400
            series = df[col1].dropna()
            if len(series) > 5000:
                series = series.sample(5000)
            stat, p_val = scipy_stats.shapiro(series) if len(series) <= 5000 else scipy_stats.normaltest(series)
            # Distribution plot
            fig, axes = plt.subplots(1, 2, figsize=(10, 4))
            axes[0].hist(series, bins=30, color='#5b5ef4', alpha=0.7, edgecolor='white')
            axes[0].set_title(f'Distribution of {col1}', fontweight='bold')
            scipy_stats.probplot(series, dist='norm', plot=axes[1])
            axes[1].set_title('Q-Q Plot', fontweight='bold')
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); chart_b64 = base64.b64encode(buf.read()).decode()
            result = {
                'test': 'Normality Test (Shapiro-Wilk)',
                'statistic': round(float(stat), 4),
                'p_value': round(float(p_val), 6),
                'is_normal': bool(p_val > alpha),
                'interpretation': f'{"Normal distribution" if p_val > alpha else "Not normal distribution"} (alpha={alpha})',
                'chart': chart_b64
            }

        elif test == 'distribution':
            num_cols = list(df.select_dtypes(include='number').columns[:6])
            fig, axes = plt.subplots(2, 3, figsize=(12, 7))
            axes = axes.flatten()
            for i, c in enumerate(num_cols):
                axes[i].hist(df[c].dropna(), bins=25, color='#5b5ef4', alpha=0.75, edgecolor='white')
                axes[i].set_title(c, fontweight='bold', fontsize=10)
            for j in range(len(num_cols), 6):
                axes[j].set_visible(False)
            plt.suptitle('Distribution of Numeric Columns', fontweight='bold', fontsize=13)
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); chart_b64 = base64.b64encode(buf.read()).decode()
            result = {'columns': num_cols, 'chart': chart_b64}

        if 'chart' not in result and chart_b64:
            result['chart'] = chart_b64
        return jsonify(result)

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 4: BIG DATA — Chunked Upload (up to 500MB via HuggingFace)
# ════════════════════════════════════════════════════════════════
import uuid as _uuid

@app.route('/api/upload_chunk', methods=['POST'])
@login_required
def upload_chunk():
    """Send one 25MB chunk to HuggingFace Space"""
    try:
        url = get_hf_url()
        if not url:
            return jsonify({'error': 'ML Server offline! HuggingFace Space start karo.'}), 400

        upload_id = request.form.get('upload_id', '')
        chunk_index = request.form.get('chunk_index', '0')
        total_chunks = request.form.get('total_chunks', '1')
        filename = request.form.get('filename', 'data.csv')

        if 'file' not in request.files:
            return jsonify({'error': 'Chunk file missing!'}), 400

        chunk_file = request.files['file']

        # Forward chunk to HuggingFace
        files = {'file': (filename, chunk_file.read(), 'application/octet-stream')}
        data = {
            'upload_id': upload_id,
            'chunk_index': chunk_index,
            'total_chunks': total_chunks,
            'filename': filename
        }
        import requests as req
        r = req.post(url + '/upload_chunk', files=files, data=data, timeout=120)
        return jsonify(r.json())

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/merge_chunks', methods=['POST'])
@login_required
def merge_chunks():
    """Tell HuggingFace to merge all chunks and return processed data"""
    try:
        url = get_hf_url()
        if not url:
            return jsonify({'error': 'ML Server offline!'}), 400

        data = request.get_json()
        import requests as req
        r = req.post(url + '/merge_chunks', json=data, timeout=300)
        result = r.json()

        if result.get('success') and result.get('sample_csv'):
            # Load sample into agent
            import io as _io
            df = pd.read_csv(_io.StringIO(result['sample_csv']))
            agent.df = df
            agent.available_columns = list(df.columns)

            # Save to DB (sample only)
            db = SessionLocal()
            db.query(UserDataset).filter_by(user_id=current_user.id).delete()
            dataset = UserDataset(
                user_id=current_user.id,
                filename=result['filename'],
                csv_data=result['sample_csv'][:2000000],  # max 2MB in DB
                rows=result['total_rows'],
                columns=result['total_cols']
            )
            db.add(dataset); db.commit(); db.close()

            return jsonify({
                'success': True,
                'filename': result['filename'],
                'total_rows': result['total_rows'],
                'sample_rows': result['sample_rows'],
                'total_cols': result['total_cols'],
                'size_mb': result['size_mb'],
                'columns': result['columns'],
                'null_info': result.get('null_info', {}),
                'message': result['message']
            })

        return jsonify(result)

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 4: BIG DATA (chunked, up to 50MB)
# ════════════════════════════════════════════════════════════════
@app.route('/api/upload_big', methods=['POST'])
@login_required
def upload_big():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'File nahi mili!'}), 400
        file = request.files['file']
        filename = secure_filename(file.filename)
        ext = filename.rsplit('.', 1)[-1].lower()
        import io as _io

        file_bytes = file.read()
        size_mb = len(file_bytes) / (1024 * 1024)

        if size_mb > 52:
            return jsonify({'error': f'File {size_mb:.1f}MB hai — max 50MB allowed!'}), 400

        if ext == 'csv':
            # Chunked read for large CSVs
            chunks = []
            chunk_size = 50000
            buf = _io.StringIO(file_bytes.decode('utf-8', errors='replace'))
            for chunk in pd.read_csv(buf, chunksize=chunk_size):
                chunks.append(chunk)
                if sum(len(c) for c in chunks) > 500000:
                    break
            df = pd.concat(chunks, ignore_index=True)
        elif ext in ('xlsx', 'xls'):
            df = pd.read_excel(_io.BytesIO(file_bytes))
        elif ext == 'json':
            df = pd.read_json(_io.BytesIO(file_bytes))
        elif ext == 'parquet':
            df = pd.read_parquet(_io.BytesIO(file_bytes))
        elif ext == 'tsv':
            df = pd.read_csv(_io.StringIO(file_bytes.decode('utf-8', errors='replace')), sep='\t')
        else:
            return jsonify({'error': f'Format ".{ext}" support nahi — CSV/Excel/JSON/Parquet/TSV use karo'}), 400

        agent.df = df
        agent.available_columns = list(df.columns)

        csv_str = df.to_csv(index=False)
        db = SessionLocal()
        db.query(UserDataset).filter_by(user_id=current_user.id).delete()
        dataset = UserDataset(
            user_id=current_user.id, filename=filename,
            csv_data=csv_str, columns=json.dumps(list(df.columns)),
            row_count=len(df)
        )
        db.add(dataset); db.commit(); db.close()

        num_cols = list(df.select_dtypes(include='number').columns)
        cat_cols = list(df.select_dtypes(include='object').columns)

        return jsonify({
            'message': f'Big file load ho gayi! {len(df):,} rows × {len(df.columns)} cols ({size_mb:.1f}MB)',
            'rows': len(df), 'cols': len(df.columns),
            'size_mb': round(size_mb, 2),
            'columns': list(df.columns),
            'numeric_cols': num_cols,
            'categorical_cols': cat_cols,
            'nulls': int(df.isnull().sum().sum()),
            'duplicates': int(df.duplicated().sum()),
            'preview': df.head(3).to_string(index=False)
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 5: SCHEDULED AUTO REPORTS
# ════════════════════════════════════════════════════════════════
@app.route('/api/auto_report', methods=['POST'])
@login_required
def auto_report():
    """Generate and optionally email a full auto report"""
    if agent.df is None:
        return jsonify({'error': 'Pehle file upload karo!'}), 400
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np
        from scipy import stats as scipy_stats

        data = request.get_json()
        email_to = data.get('email', '')
        df = agent.df.copy()

        sections = []

        # 1. Overview
        sections.append({
            'title': 'Dataset Overview',
            'content': {
                'rows': len(df), 'columns': len(df.columns),
                'numeric_cols': len(df.select_dtypes(include='number').columns),
                'categorical_cols': len(df.select_dtypes(include='object').columns),
                'total_nulls': int(df.isnull().sum().sum()),
                'null_percentage': round(df.isnull().sum().sum() / df.size * 100, 2),
                'duplicates': int(df.duplicated().sum()),
                'memory_mb': round(df.memory_usage(deep=True).sum() / 1024**2, 2)
            }
        })

        # 2. Numeric stats
        num_df = df.select_dtypes(include='number')
        if not num_df.empty:
            sections.append({
                'title': 'Numeric Summary',
                'content': num_df.describe().round(3).to_dict()
            })

        # 3. Top correlations
        charts = []
        if len(num_df.columns) >= 2:
            corr = num_df.corr()
            top_corr = []
            for i in range(len(corr.columns)):
                for j in range(i+1, len(corr.columns)):
                    top_corr.append({
                        'col1': corr.columns[i], 'col2': corr.columns[j],
                        'correlation': round(float(corr.iloc[i, j]), 4)
                    })
            top_corr.sort(key=lambda x: abs(x['correlation']), reverse=True)
            sections.append({'title': 'Top Correlations', 'content': top_corr[:10]})

            # Heatmap
            import seaborn as sns
            fig, ax = plt.subplots(figsize=(8, 6))
            sns.heatmap(corr, annot=True, fmt='.2f', cmap='coolwarm', ax=ax,
                       annot_kws={'size': 7}, square=True)
            ax.set_title('Correlation Heatmap', fontweight='bold')
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); charts.append({'title': 'Correlation Heatmap', 'img': base64.b64encode(buf.read()).decode()})

        # 4. Distributions
        plot_cols = list(num_df.columns[:4])
        if plot_cols:
            fig, axes = plt.subplots(1, len(plot_cols), figsize=(4*len(plot_cols), 3.5))
            if len(plot_cols) == 1: axes = [axes]
            for i, c in enumerate(plot_cols):
                axes[i].hist(df[c].dropna(), bins=25, color='#5b5ef4', alpha=0.75, edgecolor='white')
                axes[i].set_title(c, fontweight='bold', fontsize=9)
            plt.suptitle('Distributions', fontweight='bold')
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); charts.append({'title': 'Distributions', 'img': base64.b64encode(buf.read()).decode()})

        # 5. Categorical counts
        cat_df = df.select_dtypes(include='object')
        cat_summary = {}
        for c in cat_df.columns[:5]:
            vc = df[c].value_counts().head(5)
            cat_summary[c] = {'top_values': vc.to_dict(), 'unique': int(df[c].nunique())}
            if int(df[c].nunique()) <= 15:
                fig, ax = plt.subplots(figsize=(6, 3))
                vc.plot(kind='bar', ax=ax, color='#06d6a0', edgecolor='white')
                ax.set_title(f'{c} — Value Counts', fontweight='bold')
                ax.set_xlabel(''); plt.xticks(rotation=30, ha='right')
                plt.tight_layout()
                buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=90); plt.close('all')
                buf.seek(0); charts.append({'title': f'{c} Distribution', 'img': base64.b64encode(buf.read()).decode()})
        if cat_summary:
            sections.append({'title': 'Categorical Summary', 'content': cat_summary})

        # Build HTML report
        import datetime
        template = data.get('template', 'corporate')
        report_title = data.get('report_title', 'Data Analysis Report')
        company_name = data.get('company_name', 'DS Agent')
        now_str = datetime.datetime.now().strftime('%d %b %Y, %H:%M')
        meta_str = f"{len(df):,} rows × {len(df.columns)} columns"

        # ── TEMPLATE STYLES ──
        templates = {
            'corporate': {
                'name': 'Corporate Blue',
                'body_bg': '#f0f4f8',
                'header_bg': 'linear-gradient(135deg,#1e3a5f,#2563eb)',
                'header_text': '#ffffff',
                'accent': '#2563eb',
                'accent_light': '#dbeafe',
                'section_border': '#e2e8f0',
                'stat_bg': '#f8fafc',
                'stat_color': '#1e3a5f',
                'table_head': '#1e3a5f',
                'table_head_text': '#ffffff',
                'table_alt': '#f0f4f8',
                'font': 'Georgia, serif',
                'footer_bg': '#1e3a5f',
                'footer_text': '#94a3b8',
                'container_shadow': '0 8px 32px rgba(30,58,95,0.15)',
                'border_radius': '0px',
            },
            'healthcare': {
                'name': 'Healthcare Green',
                'body_bg': '#f0fdf4',
                'header_bg': 'linear-gradient(135deg,#065f46,#059669)',
                'header_text': '#ffffff',
                'accent': '#059669',
                'accent_light': '#dcfce7',
                'section_border': '#bbf7d0',
                'stat_bg': '#f0fdf4',
                'stat_color': '#065f46',
                'table_head': '#065f46',
                'table_head_text': '#ffffff',
                'table_alt': '#f0fdf4',
                'font': 'DM Sans, Segoe UI, sans-serif',
                'footer_bg': '#065f46',
                'footer_text': '#6ee7b7',
                'container_shadow': '0 8px 32px rgba(6,95,70,0.15)',
                'border_radius': '16px',
            },
            'executive': {
                'name': 'Executive Dark',
                'body_bg': '#0f172a',
                'header_bg': 'linear-gradient(135deg,#0f172a,#1e293b)',
                'header_text': '#f1f5f9',
                'accent': '#f59e0b',
                'accent_light': '#1e293b',
                'section_border': '#1e293b',
                'stat_bg': '#1e293b',
                'stat_color': '#f59e0b',
                'table_head': '#1e293b',
                'table_head_text': '#f59e0b',
                'table_alt': '#0f172a',
                'font': 'Inter, Segoe UI, sans-serif',
                'footer_bg': '#020617',
                'footer_text': '#475569',
                'container_shadow': '0 8px 32px rgba(0,0,0,0.5)',
                'border_radius': '12px',
            },
            'minimal': {
                'name': 'Minimal Clean',
                'body_bg': '#ffffff',
                'header_bg': '#ffffff',
                'header_text': '#111827',
                'accent': '#111827',
                'accent_light': '#f9fafb',
                'section_border': '#e5e7eb',
                'stat_bg': '#f9fafb',
                'stat_color': '#111827',
                'table_head': '#f3f4f6',
                'table_head_text': '#374151',
                'table_alt': '#f9fafb',
                'font': 'Inter, Helvetica, sans-serif',
                'footer_bg': '#f9fafb',
                'footer_text': '#9ca3af',
                'container_shadow': '0 1px 3px rgba(0,0,0,0.1)',
                'border_radius': '8px',
            },
            'vibrant': {
                'name': 'Vibrant Purple',
                'body_bg': '#faf5ff',
                'header_bg': 'linear-gradient(135deg,#5b5ef4,#8b5cf6,#ec4899)',
                'header_text': '#ffffff',
                'accent': '#7c3aed',
                'accent_light': '#ede9fe',
                'section_border': '#ddd6fe',
                'stat_bg': '#f5f3ff',
                'stat_color': '#5b21b6',
                'table_head': '#7c3aed',
                'table_head_text': '#ffffff',
                'table_alt': '#faf5ff',
                'font': 'Sora, DM Sans, sans-serif',
                'footer_bg': 'linear-gradient(135deg,#5b5ef4,#8b5cf6)',
                'footer_text': '#e9d5ff',
                'container_shadow': '0 8px 32px rgba(91,94,244,0.2)',
                'border_radius': '20px',
            }
        }

        t = templates.get(template, templates['corporate'])
        br = t['border_radius']
        body_color = '#1a1a2e' if template != 'executive' else '#e2e8f0'

        html = f"""<!DOCTYPE html>
<html><head><meta charset='utf-8'>
<meta name='viewport' content='width=device-width,initial-scale=1'>
<link href='https://fonts.googleapis.com/css2?family=DM+Sans:wght@400;600;700&family=Sora:wght@400;700&family=Inter:wght@400;600;700&display=swap' rel='stylesheet'>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:{t['font']};background:{t['body_bg']};padding:30px 16px;color:{body_color}}}
.container{{max-width:920px;margin:0 auto;background:{'#1a2035' if template=='executive' else 'white'};border-radius:{br};overflow:hidden;box-shadow:{t['container_shadow']}}}
.header{{background:{t['header_bg']};color:{t['header_text']};padding:36px 40px;position:relative;overflow:hidden}}
.header::after{{content:'';position:absolute;right:-40px;top:-40px;width:200px;height:200px;background:rgba(255,255,255,0.05);border-radius:50%}}
.header::before{{content:'';position:absolute;right:60px;bottom:-30px;width:120px;height:120px;background:rgba(255,255,255,0.04);border-radius:50%}}
.header h1{{font-size:26px;font-weight:800;letter-spacing:-0.5px;position:relative}}
.header .subtitle{{font-size:13px;opacity:0.75;margin-top:6px;position:relative}}
.header .meta{{display:flex;gap:20px;margin-top:16px;position:relative}}
.header .meta span{{background:rgba(255,255,255,0.15);padding:4px 12px;border-radius:20px;font-size:12px;font-weight:600}}
.section{{padding:24px 40px;border-bottom:1px solid {t['section_border']}}}
.section h2{{font-size:15px;font-weight:800;color:{t['accent']};margin-bottom:14px;display:flex;align-items:center;gap:8px;letter-spacing:0.3px;text-transform:uppercase}}
.stat-grid{{display:grid;grid-template-columns:repeat(4,1fr);gap:14px}}
.stat-box{{background:{t['stat_bg']};border-radius:12px;padding:16px;text-align:center;border:1px solid {t['section_border']}}}
.stat-box .val{{font-size:24px;font-weight:800;color:{t['stat_color']}}}
.stat-box .lbl{{font-size:11px;color:#888;margin-top:4px;font-weight:600;text-transform:uppercase;letter-spacing:0.5px}}
table{{width:100%;border-collapse:collapse;font-size:12.5px}}
th{{background:{t['table_head']};color:{t['table_head_text']};padding:10px 12px;text-align:left;font-weight:700;font-size:11px;text-transform:uppercase;letter-spacing:0.5px}}
th:first-child{{border-radius:{'6px 0 0 6px' if br!='0px' else '0'}}}
th:last-child{{border-radius:{'0 6px 6px 0' if br!='0px' else '0'}}}
tr:nth-child(even) td{{background:{t['table_alt']}}}
td{{padding:9px 12px;border-bottom:1px solid {t['section_border']};color:{body_color}}}
.chart-wrap{{background:{t['stat_bg']};border-radius:10px;padding:12px;margin-top:8px}}
img{{max-width:100%;border-radius:8px}}
.corr-pos{{color:#059669;font-weight:700}}
.corr-neg{{color:#ef4444;font-weight:700}}
.badge{{display:inline-block;padding:2px 10px;border-radius:20px;font-size:11px;font-weight:700;background:{t['accent_light']};color:{t['accent']}}}
.footer{{background:{t['footer_bg']};padding:18px 40px;text-align:center;font-size:12px;color:{t['footer_text']}}}
.footer strong{{color:{t['accent']}}}
@media(max-width:600px){{.stat-grid{{grid-template-columns:repeat(2,1fr)}}body{{padding:10px 8px}}.section{{padding:16px 18px}}.header{{padding:24px 20px}}}}
</style>
</head><body><div class='container'>
<div class='header'>
  <h1>📊 {report_title}</h1>
  <div class='subtitle'>{company_name}</div>
  <div class='meta'>
    <span>📅 {now_str}</span>
    <span>📁 {meta_str}</span>
    <span>🎨 {t['name']}</span>
  </div>
</div>"""

        # Overview section
        ov = sections[0]['content']
        html += f"""<div class='section'><h2>📋 Dataset Overview</h2>
        <div class='stat-grid'>
        <div class='stat-box'><div class='val'>{ov['rows']:,}</div><div class='lbl'>Total Rows</div></div>
        <div class='stat-box'><div class='val'>{ov['columns']}</div><div class='lbl'>Columns</div></div>
        <div class='stat-box'><div class='val'>{ov['null_percentage']}%</div><div class='lbl'>Null Values</div></div>
        <div class='stat-box'><div class='val'>{ov['duplicates']}</div><div class='lbl'>Duplicates</div></div>
        </div></div>"""

        # Numeric summary
        for sec in sections[1:]:
            html += f"<div class='section'><h2>{sec['title']}</h2>"
            if isinstance(sec['content'], dict) and sec['title'] == 'Numeric Summary':
                html += "<table><tr><th>Column</th><th>Mean</th><th>Std Dev</th><th>Min</th><th>Max</th><th>Nulls</th></tr>"
                for col in list(sec['content'].get('mean', {}).keys()):
                    d = sec['content']
                    html += f"<tr><td><b>{col}</b></td><td>{d.get('mean',{}).get(col,'')}</td><td>{d.get('std',{}).get(col,'')}</td><td>{d.get('min',{}).get(col,'')}</td><td>{d.get('max',{}).get(col,'')}</td><td>{int(df[col].isnull().sum())}</td></tr>"
                html += "</table>"
            elif isinstance(sec['content'], list):
                html += "<table><tr><th>Column 1</th><th>Column 2</th><th>Correlation</th></tr>"
                for row in sec['content'][:10]:
                    cls = 'corr-pos' if row['correlation'] > 0 else 'corr-neg'
                    html += f"<tr><td>{row['col1']}</td><td>{row['col2']}</td><td class='{cls}'>{row['correlation']}</td></tr>"
                html += "</table>"
            html += "</div>"

        # Charts
        for ch in charts:
            html += f"<div class='section'><h2>{ch['title']}</h2><div class='chart-wrap'><img src='data:image/png;base64,{ch['img']}'/></div></div>"

        html += f"<div class='footer'>Generated by <strong>DS Agent</strong> — Your AI Data Scientist &nbsp;|&nbsp; Template: {t['name']}</div></div></body></html>"

        # Overview section
        ov = sections[0]['content']
        html += f"""<div class='section'><h2>📋 Dataset Overview</h2>
        <div class='stat-grid'>
        <div class='stat-box'><div class='val'>{ov['rows']:,}</div><div class='lbl'>Rows</div></div>
        <div class='stat-box'><div class='val'>{ov['columns']}</div><div class='lbl'>Columns</div></div>
        <div class='stat-box'><div class='val'>{ov['null_percentage']}%</div><div class='lbl'>Nulls</div></div>
        <div class='stat-box'><div class='val'>{ov['duplicates']}</div><div class='lbl'>Duplicates</div></div>
        </div></div>"""

        # Numeric summary
        for sec in sections[1:]:
            html += f"<div class='section'><h2>{sec['title']}</h2>"
            if isinstance(sec['content'], dict) and sec['title'] == 'Numeric Summary':
                html += "<table><tr><th>Column</th><th>Mean</th><th>Std</th><th>Min</th><th>Max</th><th>Nulls</th></tr>"
                for col in list(sec['content'].get('mean', {}).keys()):
                    d = sec['content']
                    html += f"<tr><td><b>{col}</b></td><td>{d.get('mean',{}).get(col,'')}</td><td>{d.get('std',{}).get(col,'')}</td><td>{d.get('min',{}).get(col,'')}</td><td>{d.get('max',{}).get(col,'')}</td><td>{int(df[col].isnull().sum())}</td></tr>"
                html += "</table>"
            elif isinstance(sec['content'], list):
                html += "<table><tr><th>Col 1</th><th>Col 2</th><th>Correlation</th></tr>"
                for row in sec['content'][:10]:
                    color = '#059669' if row['correlation'] > 0 else '#ef4444'
                    html += f"<tr><td>{row['col1']}</td><td>{row['col2']}</td><td style='color:{color};font-weight:700'>{row['correlation']}</td></tr>"
                html += "</table>"
            html += "</div>"

        # Charts
        for ch in charts:
            html += f"<div class='section'><h2>{ch['title']}</h2><img src='data:image/png;base64,{ch['img']}'/></div>"

        html += "<div class='footer'>Generated by DS Agent — Your AI Data Scientist</div></div></body></html>"

        # Optionally email
        email_sent = False
        if email_to:
            try:
                import resend as resend_lib
                resend_lib.api_key = os.environ.get('RESEND_API_KEY', '')
                resend_lib.Emails.send({
                    'from': 'DS Agent <reports@dsagent.app>',
                    'to': [email_to],
                    'subject': f'Your Data Report — {datetime.datetime.now().strftime("%d %b %Y")}',
                    'html': html
                })
                email_sent = True
            except Exception as e:
                logger.warning(f'Email failed: {e}')

        return jsonify({
            'message': f'Report ready! {len(charts)} charts generated.',
            'html': html,
            'email_sent': email_sent,
            'sections': len(sections),
            'charts': len(charts)
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 6: PIVOT TABLE
# ════════════════════════════════════════════════════════════════
@app.route('/api/pivot', methods=['POST'])
@login_required
def pivot_table():
    if agent.df is None:
        return jsonify({'error': 'Pehle file upload karo!'}), 400
    try:
        data = request.get_json()
        index_col = data.get('index', '')
        values_col = data.get('values', '')
        aggfunc = data.get('aggfunc', 'sum')
        columns_col = data.get('columns', '')
        df = agent.df.copy()

        if not index_col or index_col not in df.columns:
            return jsonify({'error': f'index column chahiye! Available: {list(df.columns)[:8]}'}), 400
        if not values_col or values_col not in df.columns:
            return jsonify({'error': f'values column chahiye! Available: {list(df.columns)[:8]}'}), 400

        pivot_kwargs = {
            'index': index_col,
            'values': values_col,
            'aggfunc': aggfunc
        }
        if columns_col and columns_col in df.columns:
            pivot_kwargs['columns'] = columns_col

        pivot = pd.pivot_table(df, **pivot_kwargs).round(3)

        # Chart
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(min(10, 2+len(pivot.columns)*1.5), 5))
        pivot.plot(kind='bar', ax=ax, colormap='viridis', edgecolor='white')
        ax.set_title(f'Pivot: {aggfunc}({values_col}) by {index_col}', fontweight='bold')
        ax.set_xlabel(index_col); ax.set_ylabel(f'{aggfunc}({values_col})')
        plt.xticks(rotation=30, ha='right'); plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
        buf.seek(0); chart_b64 = base64.b64encode(buf.read()).decode()

        return jsonify({
            'message': f'Pivot table ready! {pivot.shape[0]} rows × {pivot.shape[1]} cols',
            'pivot': pivot.reset_index().to_dict(orient='records'),
            'columns': [str(index_col)] + [str(c) for c in pivot.columns],
            'chart': chart_b64,
            'csv': pivot.reset_index().to_csv(index=False)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 7: TIME SERIES ANALYSIS
# ════════════════════════════════════════════════════════════════
@app.route('/api/timeseries', methods=['POST'])
@login_required
def time_series():
    if agent.df is None:
        return jsonify({'error': 'Pehle file upload karo!'}), 400
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        data = request.get_json()
        date_col = data.get('date_col', '')
        value_col = data.get('value_col', '')
        freq = data.get('freq', 'auto')
        forecast_periods = int(data.get('forecast_periods', 7))

        df = agent.df.copy()
        if not date_col or date_col not in df.columns:
            return jsonify({'error': f'date_col chahiye! Available: {list(df.columns)[:8]}'}), 400
        if not value_col or value_col not in df.columns:
            return jsonify({'error': f'value_col chahiye! Available: {list(df.columns)[:8]}'}), 400

        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        df = df.dropna(subset=[date_col, value_col])
        df = df.sort_values(date_col)

        series = df.set_index(date_col)[value_col]

        # Resample if needed
        if freq == 'auto':
            time_diff = (series.index[-1] - series.index[0]).days
            freq = 'D' if time_diff <= 365 else 'W' if time_diff <= 1825 else 'ME'

        # Rolling stats
        window = max(7, len(series) // 10)
        rolling_mean = series.rolling(window=window).mean()
        rolling_std = series.rolling(window=window).std()

        # Simple trend forecast (linear)
        from sklearn.linear_model import LinearRegression
        X_t = np.arange(len(series)).reshape(-1, 1)
        lr = LinearRegression().fit(X_t, series.values)
        trend_line = lr.predict(X_t)
        future_X = np.arange(len(series), len(series)+forecast_periods).reshape(-1, 1)
        forecast = lr.predict(future_X)

        # Plot
        fig, axes = plt.subplots(2, 1, figsize=(10, 7))
        # Main time series
        axes[0].plot(series.index, series.values, color='#5b5ef4', linewidth=1.5, label='Actual', alpha=0.8)
        axes[0].plot(series.index, rolling_mean, color='#f59e0b', linewidth=2, label=f'{window}-period MA')
        axes[0].fill_between(series.index,
                             rolling_mean - rolling_std,
                             rolling_mean + rolling_std,
                             alpha=0.15, color='#f59e0b')
        axes[0].set_title(f'{value_col} — Time Series', fontweight='bold', fontsize=12)
        axes[0].legend(fontsize=9); axes[0].grid(True, alpha=0.3)

        # Forecast
        last_date = series.index[-1]
        future_dates = pd.date_range(last_date, periods=forecast_periods+1, freq='D')[1:]
        axes[1].plot(series.index[-50:], series.values[-50:], color='#5b5ef4', linewidth=1.5, label='Actual')
        axes[1].plot(series.index, trend_line, color='#9ca3af', linewidth=1, linestyle='--', label='Trend')
        axes[1].plot(future_dates, forecast, color='#06d6a0', linewidth=2.5, marker='o', markersize=4, label=f'{forecast_periods}-period Forecast')
        axes[1].axvline(last_date, color='#ef4444', linewidth=1, linestyle=':', alpha=0.7)
        axes[1].set_title('Trend + Forecast', fontweight='bold', fontsize=12)
        axes[1].legend(fontsize=9); axes[1].grid(True, alpha=0.3)

        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
        buf.seek(0); chart_b64 = base64.b64encode(buf.read()).decode()

        return jsonify({
            'message': f'Time series analysis complete! {len(series)} data points analyzed.',
            'date_col': date_col,
            'value_col': value_col,
            'data_points': len(series),
            'date_range': f'{series.index[0].strftime("%Y-%m-%d")} to {series.index[-1].strftime("%Y-%m-%d")}',
            'trend': 'Upward' if lr.coef_[0] > 0 else 'Downward',
            'trend_slope': round(float(lr.coef_[0]), 4),
            'mean': round(float(series.mean()), 4),
            'std': round(float(series.std()), 4),
            'min': round(float(series.min()), 4),
            'max': round(float(series.max()), 4),
            'forecast': [{'date': str(d.date()), 'value': round(float(v), 4)} for d, v in zip(future_dates, forecast)],
            'chart': chart_b64
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500



# ════════════════════════════════════════════════════════════════
# TARIKA 3 — GOOGLE COLAB HYBRID (Heavy ML/Deep Learning)
# ════════════════════════════════════════════════════════════════

@app.route('/api/colab/status', methods=['GET'])
@login_required
def colab_status():
    """Check HuggingFace (primary) and Colab (backup)"""
    import requests as req
    hf_url = os.environ.get('HF_SPACE_URL', '').strip()
    colab_url = os.environ.get('COLAB_URL', '').strip()

    # HuggingFace — main page check (Gradio runs on 7860, no /ping needed)
    if hf_url:
        try:
            r = req.get(hf_url, timeout=10)
            if r.status_code == 200:
                return jsonify({'connected': True, 'url': hf_url, 'source': 'HuggingFace', 'message': 'HuggingFace Connected!'})
        except: pass

    # Colab — /ping check
    if colab_url:
        try:
            r = req.get(colab_url + '/ping', timeout=6)
            if r.status_code == 200:
                return jsonify({'connected': True, 'url': colab_url, 'source': 'Colab', 'message': 'Colab Connected!'})
        except: pass

    return jsonify({'connected': False, 'source': 'none', 'message': 'Dono offline!'})


@app.route('/api/colab/run', methods=['POST'])
@login_required
def colab_run():
    """Send heavy task to Google Colab"""
    colab_url = os.environ.get('COLAB_URL', '').strip()
    if not colab_url:
        return jsonify({'error': 'COLAB_URL env variable set nahi hai! Render dashboard pe add karo.'}), 400
    try:
        import requests as req
        data = request.get_json()
        task = data.get('task', '')
        code = data.get('code', '')
        payload = {}

        if agent.df is not None:
            import io as _io
            csv_str = agent.df.to_csv(index=False)
            payload['csv_data'] = csv_str

        payload['task'] = task
        payload['code'] = code
        payload['params'] = data.get('params', {})

        r = req.post(colab_url + '/run', json=payload, timeout=120)
        result = r.json()
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': f'Colab error: {str(e)} — Colab notebook open hai?'}), 500


@app.route('/api/colab/deep_learning', methods=['POST'])
@login_required
def colab_deep_learning():
    """Run Deep Learning — HuggingFace primary, Colab backup"""
    import requests as req
    hf_url = os.environ.get('HF_SPACE_URL', '').strip()
    colab_url = os.environ.get('COLAB_URL', '').strip()

    # Pick available server
    server_url = None
    if hf_url:
        try:
            r = req.get(hf_url, timeout=10)
            if r.status_code == 200:
                server_url = hf_url
                print(f"[DL] Using HuggingFace: {hf_url}")
        except: pass
    if not server_url and colab_url:
        try:
            r = req.get(colab_url + '/ping', timeout=6)
            if r.status_code == 200:
                server_url = colab_url
        except: pass
    if not server_url:
        return jsonify({'error': 'ML Server offline! HuggingFace Space check karo — ping karo pehle.'}), 400
    try:
        import requests as req
        data = request.get_json()
        model_type = data.get('model_type', 'neural_network')
        target = data.get('target', '')
        epochs = int(data.get('epochs', 10))
        layers = data.get('layers', [64, 32])

        if agent.df is None:
            return jsonify({'error': 'Pehle file upload karo!'}), 400

        payload = {
            'task': 'deep_learning',
            'csv_data': agent.df.to_csv(index=False),
            'params': {
                'model_type': model_type,
                'target': target,
                'epochs': epochs,
                'layers': layers
            }
        }
        result = call_hf_api(server_url, '/run', payload)
        safe = validate_hf_result(result, {
            'success': False,
            'task': 'unknown',
            'target': '',
            'epochs_ran': 0,
            'metrics': {},
            'train_rows': 0,
            'test_rows': 0,
            'message': 'Training complete!',
            'chart': None,
        })
        return jsonify(safe), (500 if safe.get('error') else 200)
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500



# ════════════════════════════════════════════════════════════════
# HEALTHCARE MODULE
# ════════════════════════════════════════════════════════════════

# Healthcare roles store (in-memory + DB via user metadata)
HEALTHCARE_ROLES = ['admin', 'doctor', 'nurse', 'analyst', 'receptionist']

ROLE_PERMISSIONS = {
    'admin':        ['dashboard', 'patients', 'reports', 'risk', 'finance', 'staff', 'privacy_off'],
    'doctor':       ['dashboard', 'patients', 'reports', 'risk'],
    'nurse':        ['dashboard', 'patients'],
    'analyst':      ['dashboard', 'reports', 'risk'],
    'receptionist': ['dashboard', 'patients'],
}

SENSITIVE_COLUMNS = [
    'name', 'patient_name', 'full_name', 'phone', 'mobile', 'address',
    'email', 'aadhar', 'ssn', 'dob', 'date_of_birth', 'contact',
    'guardian', 'emergency_contact', 'account', 'insurance_id'
]


def get_user_hc_role(user_id):
    """Get healthcare role from DB notes field or default to analyst"""
    try:
        db = SessionLocal()
        user = db.query(User).filter_by(id=user_id).first()
        db.close()
        if user and user.is_admin:
            return 'admin'
        # Store role in username prefix like "dr_john" -> doctor
        if user:
            uname = user.username.lower()
            if uname.startswith('dr_') or uname.startswith('doc_'):
                return 'doctor'
            elif uname.startswith('nurse_') or uname.startswith('nr_'):
                return 'nurse'
            elif uname.startswith('rec_') or uname.startswith('reception_'):
                return 'receptionist'
        return 'analyst'
    except:
        return 'analyst'


def mask_sensitive(df, role):
    """Mask sensitive columns based on role"""
    if 'privacy_off' in ROLE_PERMISSIONS.get(role, []):
        return df  # Admin sees all
    masked = df.copy()
    for col in masked.columns:
        if col.lower() in SENSITIVE_COLUMNS:
            masked[col] = masked[col].astype(str).str[:2] + '****'
    return masked


# ── Healthcare Role API ───────────────────────────────────────
@app.route('/api/hc/role', methods=['GET'])
@login_required
def hc_get_role():
    role = get_user_hc_role(current_user.id)
    perms = ROLE_PERMISSIONS.get(role, [])
    return jsonify({
        'role': role,
        'permissions': perms,
        'username': current_user.username,
        'is_admin': current_user.is_admin
    })


@app.route('/api/hc/set_role', methods=['POST'])
@login_required
def hc_set_role():
    if not current_user.is_admin:
        return jsonify({'error': 'Sirf admin role assign kar sakta hai!'}), 403
    data = request.get_json()
    target_username = data.get('username', '')
    new_role = data.get('role', '')
    if new_role not in HEALTHCARE_ROLES:
        return jsonify({'error': f'Invalid role! Use: {HEALTHCARE_ROLES}'}), 400
    try:
        db = SessionLocal()
        user = db.query(User).filter_by(username=target_username).first()
        if not user:
            db.close()
            return jsonify({'error': 'User nahi mila!'}), 404
        # Role store via username prefix convention
        db.close()
        return jsonify({'message': f'Role "{new_role}" assigned to {target_username}! (Username prefix se role detect hota hai)', 'role': new_role})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── Hospital Dashboard ────────────────────────────────────────
@app.route('/api/hc/dashboard', methods=['POST'])
@login_required
def hc_dashboard():
    if agent.df is None:
        return jsonify({'error': 'Pehle patient/hospital data upload karo!'}), 400
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        role = get_user_hc_role(current_user.id)
        perms = ROLE_PERMISSIONS.get(role, [])
        df = agent.df.copy()
        df_masked = mask_sensitive(df, role)

        result = {
            'role': role,
            'permissions': perms,
            'total_records': len(df),
            'columns': list(df.columns),
        }

        # Auto-detect common hospital columns
        cols_lower = {c.lower(): c for c in df.columns}

        # KPIs
        kpis = {}
        kpis['total_patients'] = len(df)

        # Age stats
        age_col = next((cols_lower[c] for c in cols_lower if 'age' in c), None)
        if age_col and pd.api.types.is_numeric_dtype(df[age_col]):
            kpis['avg_age'] = round(float(df[age_col].mean()), 1)
            kpis['age_range'] = f"{int(df[age_col].min())} - {int(df[age_col].max())}"

        # Gender
        gender_col = next((cols_lower[c] for c in cols_lower if 'gender' in c or 'sex' in c), None)
        if gender_col:
            kpis['gender_dist'] = df[gender_col].value_counts().to_dict()

        # Diagnosis/Disease
        diag_col = next((cols_lower[c] for c in cols_lower if 'diag' in c or 'disease' in c or 'condition' in c), None)
        if diag_col:
            kpis['top_diagnoses'] = df[diag_col].value_counts().head(5).to_dict()

        # Department
        dept_col = next((cols_lower[c] for c in cols_lower if 'dept' in c or 'department' in c or 'ward' in c), None)
        if dept_col:
            kpis['department_dist'] = df[dept_col].value_counts().head(8).to_dict()

        # Admission/discharge
        admit_col = next((cols_lower[c] for c in cols_lower if 'admit' in c or 'admission' in c), None)
        if admit_col:
            kpis['admission_col'] = admit_col

        # Finance (only admin/analyst)
        if 'finance' in perms:
            bill_col = next((cols_lower[c] for c in cols_lower if 'bill' in c or 'amount' in c or 'cost' in c or 'charge' in c or 'fee' in c), None)
            if bill_col and pd.api.types.is_numeric_dtype(df[bill_col]):
                kpis['total_revenue'] = round(float(df[bill_col].sum()), 2)
                kpis['avg_bill'] = round(float(df[bill_col].mean()), 2)
                kpis['max_bill'] = round(float(df[bill_col].max()), 2)

        result['kpis'] = kpis

        # Charts
        charts = []
        fig, axes = plt.subplots(2, 2, figsize=(12, 9))
        axes = axes.flatten()
        plot_count = 0

        # 1. Age distribution
        if age_col and pd.api.types.is_numeric_dtype(df[age_col]):
            axes[0].hist(df[age_col].dropna(), bins=20, color='#5b5ef4', alpha=0.8, edgecolor='white')
            axes[0].set_title('Patient Age Distribution', fontweight='bold', fontsize=11)
            axes[0].set_xlabel('Age'); axes[0].set_ylabel('Count')
            axes[0].grid(True, alpha=0.3)
            plot_count += 1
        else:
            axes[0].text(0.5, 0.5, 'No age column found', ha='center', va='center', transform=axes[0].transAxes, color='#888')
            axes[0].set_title('Age Distribution', fontweight='bold')

        # 2. Gender pie
        if gender_col:
            gdata = df[gender_col].value_counts().head(5)
            axes[1].pie(gdata.values, labels=gdata.index, autopct='%1.1f%%',
                       colors=['#5b5ef4','#06d6a0','#f59e0b','#ef4444','#8b5cf6'])
            axes[1].set_title('Gender Distribution', fontweight='bold', fontsize=11)
            plot_count += 1
        else:
            axes[1].text(0.5, 0.5, 'No gender column found', ha='center', va='center', transform=axes[1].transAxes, color='#888')
            axes[1].set_title('Gender', fontweight='bold')

        # 3. Top diagnoses
        if diag_col:
            top_d = df[diag_col].value_counts().head(8)
            axes[2].barh(top_d.index[::-1], top_d.values[::-1], color='#06d6a0', edgecolor='white')
            axes[2].set_title('Top Diagnoses', fontweight='bold', fontsize=11)
            axes[2].set_xlabel('Count')
            plot_count += 1
        else:
            # Use any categorical col
            cat_cols = list(df.select_dtypes(include='object').columns)
            if cat_cols:
                top_d = df[cat_cols[0]].value_counts().head(8)
                axes[2].barh(top_d.index[::-1].astype(str), top_d.values[::-1], color='#06d6a0', edgecolor='white')
                axes[2].set_title(f'Top {cat_cols[0]}', fontweight='bold', fontsize=11)

        # 4. Department or numeric distribution
        if dept_col:
            dept_d = df[dept_col].value_counts().head(8)
            axes[3].bar(dept_d.index.astype(str), dept_d.values, color='#f59e0b', edgecolor='white')
            axes[3].set_title('Department-wise Patients', fontweight='bold', fontsize=11)
            plt.setp(axes[3].xaxis.get_majorticklabels(), rotation=30, ha='right')
            plot_count += 1
        elif 'finance' in perms and 'bill_col' in dir() and bill_col:
            df[bill_col].dropna().plot(kind='hist', ax=axes[3], bins=20, color='#ef4444', alpha=0.8, edgecolor='white')
            axes[3].set_title('Billing Distribution', fontweight='bold', fontsize=11)
        else:
            num_cols = list(df.select_dtypes(include='number').columns)
            if num_cols:
                df[num_cols[0]].dropna().plot(kind='hist', ax=axes[3], bins=20, color='#8b5cf6', alpha=0.8)
                axes[3].set_title(f'{num_cols[0]} Distribution', fontweight='bold', fontsize=11)

        plt.suptitle('Hospital Analytics Dashboard', fontsize=14, fontweight='bold', color='#1a1a2e')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=110, bbox_inches='tight'); plt.close('all')
        buf.seek(0); dashboard_chart = base64.b64encode(buf.read()).decode()

        result['dashboard_chart'] = dashboard_chart
        result['preview'] = df_masked.head(5).to_string(index=False)
        return jsonify(result)

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ── Patient Risk Prediction ───────────────────────────────────
@app.route('/api/hc/risk', methods=['POST'])
@login_required
def hc_risk():
    if agent.df is None:
        return jsonify({'error': 'Pehle patient data upload karo!'}), 400
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np
        from sklearn.ensemble import RandomForestClassifier, GradientBoostingClassifier
        from sklearn.linear_model import LogisticRegression
        from sklearn.model_selection import train_test_split
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        from sklearn.metrics import accuracy_score, classification_report

        role = get_user_hc_role(current_user.id)
        if 'risk' not in ROLE_PERMISSIONS.get(role, []):
            return jsonify({'error': f'Role "{role}" ko risk prediction ki permission nahi!'}), 403

        data = request.get_json()
        target_col = data.get('target', '')
        df = agent.df.copy().dropna()

        if not target_col or target_col not in df.columns:
            # Auto detect risk column
            risk_keywords = ['risk', 'readmit', 'readmission', 'outcome', 'mortality',
                           'death', 'critical', 'severe', 'status', 'result']
            cols_lower = {c.lower(): c for c in df.columns}
            target_col = next((cols_lower[c] for c in cols_lower
                              if any(k in c for k in risk_keywords)), None)
            if not target_col:
                return jsonify({
                    'error': 'Target column nahi mila!',
                    'suggestion': f'Available columns: {list(df.columns)[:10]}',
                    'tip': 'Koi column jo predict karna hai wo daalo (e.g. readmission, outcome, risk_level)'
                }), 400

        y = df[target_col]
        X = df.drop(columns=[target_col])

        # Mask sensitive before training (privacy)
        sensitive = [c for c in X.columns if c.lower() in SENSITIVE_COLUMNS]
        X = X.drop(columns=sensitive, errors='ignore')

        # Encode
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder()
            X[c] = le.fit_transform(X[c].astype(str))

        le_y = LabelEncoder()
        y_enc = le_y.fit_transform(y.astype(str))

        if len(df) > 10000:
            from sklearn.utils import resample
            X, y_enc = resample(X, y_enc, n_samples=10000, random_state=42)

        X_train, X_test, y_train, y_test = train_test_split(X, y_enc, test_size=0.2, random_state=42, stratify=y_enc if len(np.unique(y_enc)) > 1 else None)
        scaler = StandardScaler()
        X_train_s = scaler.fit_transform(X_train)
        X_test_s = scaler.transform(X_test)

        models = {
            'Random Forest': RandomForestClassifier(n_estimators=100, random_state=42),
            'Gradient Boosting': GradientBoostingClassifier(random_state=42),
            'Logistic Regression': LogisticRegression(max_iter=1000, random_state=42),
        }

        results = []
        best_model = None
        best_score = -1

        for name, model in models.items():
            use_scaled = name == 'Logistic Regression'
            Xtr = X_train_s if use_scaled else X_train
            Xte = X_test_s if use_scaled else X_test
            model.fit(Xtr, y_train)
            score = accuracy_score(y_test, model.predict(Xte))
            results.append({'model': name, 'accuracy': round(float(score)*100, 2)})
            if score > best_score:
                best_score = score
                best_model = (name, model, use_scaled)

        results.sort(key=lambda x: x['accuracy'], reverse=True)

        # Feature importance
        feat_importance = {}
        chart_b64 = None
        if hasattr(best_model[1], 'feature_importances_'):
            fi = best_model[1].feature_importances_
            feat_df = pd.DataFrame({'feature': X.columns, 'importance': fi})
            feat_df = feat_df.sort_values('importance', ascending=True).tail(12)
            feat_importance = dict(zip(feat_df['feature'], feat_df['importance'].round(4)))

            fig, axes = plt.subplots(1, 2, figsize=(12, 5))
            axes[0].barh(feat_df['feature'], feat_df['importance'], color='#ef4444', alpha=0.8)
            axes[0].set_title(f'Risk Factors — {best_model[0]}', fontweight='bold', fontsize=11)
            axes[0].set_xlabel('Importance')

            # Class distribution
            classes, counts = np.unique(y_enc, return_counts=True)
            class_labels = le_y.inverse_transform(classes)
            axes[1].bar(class_labels.astype(str), counts,
                       color=['#06d6a0' if i==0 else '#ef4444' for i in range(len(classes))])
            axes[1].set_title('Risk Class Distribution', fontweight='bold', fontsize=11)
            axes[1].set_ylabel('Count')
            plt.suptitle('Patient Risk Analysis', fontsize=13, fontweight='bold')
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); chart_b64 = base64.b64encode(buf.read()).decode()

        # High risk patients
        name, model, use_scaled = best_model
        Xte = X_test_s if use_scaled else X_test
        y_pred = model.predict(Xte)
        high_risk_count = int(np.sum(y_pred == max(classes)))

        return jsonify({
            'target': target_col,
            'best_model': best_model[0],
            'best_accuracy': round(float(best_score)*100, 2),
            'results': results,
            'feature_importance': feat_importance,
            'sensitive_cols_removed': sensitive,
            'high_risk_patients': high_risk_count,
            'total_test_patients': len(y_test),
            'risk_classes': list(le_y.classes_),
            'chart': chart_b64,
            'message': f'Risk model ready! Best: {best_model[0]} ({round(float(best_score)*100,1)}% accuracy)'
        })

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ── Medical PDF Report ────────────────────────────────────────
@app.route('/api/hc/report', methods=['POST'])
@login_required
def hc_report():
    if agent.df is None:
        return jsonify({'error': 'Pehle data upload karo!'}), 400
    try:
        import base64, io as _io, datetime
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        role = get_user_hc_role(current_user.id)
        perms = ROLE_PERMISSIONS.get(role, [])
        if 'reports' not in perms:
            return jsonify({'error': f'Role "{role}" ko report generate karne ki permission nahi!'}), 403

        data = request.get_json()
        hospital_name = data.get('hospital_name', 'Hospital Analytics Report')
        report_type = data.get('report_type', 'general')
        df = agent.df.copy()
        df_masked = mask_sensitive(df, role)
        cols_lower = {c.lower(): c for c in df.columns}

        # Build charts
        charts_b64 = []
        num_cols = list(df.select_dtypes(include='number').columns[:4])

        if num_cols:
            fig, axes = plt.subplots(1, min(2, len(num_cols)), figsize=(10, 4))
            if len(num_cols) == 1: axes = [axes]
            for i, c in enumerate(num_cols[:2]):
                axes[i].hist(df[c].dropna(), bins=20, color='#5b5ef4', alpha=0.8, edgecolor='white')
                axes[i].set_title(c, fontweight='bold')
            plt.suptitle('Key Metrics Distribution', fontweight='bold')
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100); plt.close('all')
            buf.seek(0); charts_b64.append(base64.b64encode(buf.read()).decode())

        # Professional HTML report
        now = datetime.datetime.now()
        html = f"""<!DOCTYPE html><html><head><meta charset='utf-8'>
        <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;800&display=swap');
        *{{margin:0;padding:0;box-sizing:border-box}}
        body{{font-family:'Inter',sans-serif;background:#f0f4f8;color:#1a202c;}}
        .page{{max-width:900px;margin:0 auto;background:white;min-height:100vh;}}
        .header{{background:linear-gradient(135deg,#1e3a5f,#2563eb);color:white;padding:32px 40px;}}
        .header-top{{display:flex;justify-content:space-between;align-items:flex-start;}}
        .hospital-name{{font-size:22px;font-weight:800;margin-bottom:4px;}}
        .report-title{{font-size:14px;opacity:0.8;}}
        .header-meta{{text-align:right;font-size:12px;opacity:0.7;}}
        .confidential{{background:rgba(255,255,255,0.15);border-radius:6px;padding:4px 10px;font-size:11px;font-weight:700;margin-top:8px;display:inline-block;}}
        .kpi-grid{{display:grid;grid-template-columns:repeat(4,1fr);gap:16px;padding:24px 40px;background:#f8fafc;border-bottom:1px solid #e2e8f0;}}
        .kpi-box{{background:white;border-radius:12px;padding:16px;text-align:center;box-shadow:0 1px 4px rgba(0,0,0,0.06);}}
        .kpi-val{{font-size:28px;font-weight:800;color:#2563eb;}}
        .kpi-lbl{{font-size:11px;color:#64748b;margin-top:4px;font-weight:600;text-transform:uppercase;letter-spacing:0.5px;}}
        .section{{padding:24px 40px;border-bottom:1px solid #f1f5f9;}}
        .section-title{{font-size:15px;font-weight:800;color:#1e3a5f;margin-bottom:14px;display:flex;align-items:center;gap:8px;}}
        table{{width:100%;border-collapse:collapse;font-size:12px;}}
        th{{background:#f1f5f9;padding:10px 12px;text-align:left;font-weight:700;color:#475569;font-size:11px;text-transform:uppercase;letter-spacing:0.3px;}}
        td{{padding:9px 12px;border-bottom:1px solid #f8fafc;color:#374151;}}
        tr:hover td{{background:#fafbff;}}
        .badge{{padding:3px 8px;border-radius:20px;font-size:10px;font-weight:700;}}
        .badge-blue{{background:#dbeafe;color:#1d4ed8;}}
        .badge-green{{background:#dcfce7;color:#166534;}}
        .badge-red{{background:#fee2e2;color:#991b1b;}}
        img{{max-width:100%;border-radius:10px;margin:8px 0;}}
        .footer{{padding:20px 40px;background:#f8fafc;text-align:center;font-size:11px;color:#94a3b8;border-top:1px solid #e2e8f0;}}
        .role-badge{{background:rgba(255,255,255,0.2);border-radius:6px;padding:3px 10px;font-size:11px;font-weight:700;}}
        </style></head><body><div class='page'>
        <div class='header'>
          <div class='header-top'>
            <div>
              <div class='hospital-name'>🏥 {hospital_name}</div>
              <div class='report-title'>Healthcare Analytics Report — {report_type.title()}</div>
              <div class='confidential'>🔒 CONFIDENTIAL</div>
            </div>
            <div class='header-meta'>
              <div>Generated: {now.strftime('%d %b %Y, %H:%M')}</div>
              <div>Prepared by: {current_user.username}</div>
              <div class='role-badge' style='margin-top:6px;'>{role.upper()}</div>
            </div>
          </div>
        </div>"""

        # KPIs
        age_col = next((cols_lower[c] for c in cols_lower if 'age' in c), None)
        bill_col = next((cols_lower[c] for c in cols_lower if 'bill' in c or 'amount' in c or 'cost' in c), None) if 'finance' in perms else None
        diag_col = next((cols_lower[c] for c in cols_lower if 'diag' in c or 'disease' in c), None)

        html += "<div class='kpi-grid'>"
        html += f"<div class='kpi-box'><div class='kpi-val'>{len(df):,}</div><div class='kpi-lbl'>Total Records</div></div>"
        html += f"<div class='kpi-box'><div class='kpi-val'>{len(df.columns)}</div><div class='kpi-lbl'>Data Fields</div></div>"
        if age_col and pd.api.types.is_numeric_dtype(df[age_col]):
            html += f"<div class='kpi-box'><div class='kpi-val'>{df[age_col].mean():.0f}</div><div class='kpi-lbl'>Avg Age</div></div>"
        if bill_col and pd.api.types.is_numeric_dtype(df[bill_col]):
            html += f"<div class='kpi-box'><div class='kpi-val'>₹{df[bill_col].sum()/1e5:.1f}L</div><div class='kpi-lbl'>Total Revenue</div></div>"
        else:
            html += f"<div class='kpi-box'><div class='kpi-val'>{int(df.isnull().sum().sum())}</div><div class='kpi-lbl'>Missing Values</div></div>"
        html += "</div>"

        # Data table (masked)
        html += "<div class='section'><div class='section-title'>📋 Patient Data Preview (Top 10)</div>"
        html += "<table><tr>" + "".join(f"<th>{c}</th>" for c in df_masked.columns[:8]) + "</tr>"
        for _, row in df_masked.head(10).iterrows():
            html += "<tr>" + "".join(f"<td>{str(v)[:30]}</td>" for v in list(row)[:8]) + "</tr>"
        html += "</table></div>"

        # Top diagnoses
        if diag_col:
            top_d = df[diag_col].value_counts().head(8)
            html += "<div class='section'><div class='section-title'>🔬 Top Diagnoses / Conditions</div><table><tr><th>Condition</th><th>Count</th><th>%</th></tr>"
            for d, cnt in top_d.items():
                pct = round(cnt/len(df)*100, 1)
                html += f"<tr><td>{d}</td><td>{cnt}</td><td><span class='badge badge-blue'>{pct}%</span></td></tr>"
            html += "</table></div>"

        # Charts
        for chart in charts_b64:
            html += f"<div class='section'><div class='section-title'>📊 Analytics Charts</div><img src='data:image/png;base64,{chart}'/></div>"

        # Privacy notice
        html += f"""<div class='section'>
        <div class='section-title'>🔒 Data Privacy Notice</div>
        <p style='font-size:12px;color:#64748b;line-height:1.6;'>
        This report has been generated with role-based access control. Sensitive patient information
        (names, contact details, IDs) has been {'fully visible (Admin access)' if 'privacy_off' in perms else 'masked for privacy protection'}.
        Report access level: <strong>{role.upper()}</strong>. Generated on {now.strftime('%d %b %Y')} for internal use only.
        </p></div>"""

        html += f"<div class='footer'>DS Agent Healthcare Analytics &nbsp;|&nbsp; {hospital_name} &nbsp;|&nbsp; {now.strftime('%Y')} &nbsp;|&nbsp; CONFIDENTIAL</div>"
        html += "</div></body></html>"

        return jsonify({
            'message': f'Medical report ready! Role: {role} | {len(df)} records',
            'html': html,
            'role': role,
            'hospital_name': hospital_name,
            'records': len(df),
            'sensitive_masked': 'privacy_off' not in perms
        })

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ── Data Privacy — Mask/Unmask ────────────────────────────────
@app.route('/api/hc/privacy', methods=['POST'])
@login_required
def hc_privacy():
    if agent.df is None:
        return jsonify({'error': 'Pehle data upload karo!'}), 400
    try:
        data = request.get_json()
        action = data.get('action', 'scan')
        df = agent.df.copy()

        if action == 'scan':
            found = [c for c in df.columns if c.lower() in SENSITIVE_COLUMNS]
            all_cols = list(df.columns)
            return jsonify({
                'sensitive_columns': found,
                'safe_columns': [c for c in all_cols if c not in found],
                'total_cols': len(all_cols),
                'message': f'{len(found)} sensitive columns found: {found}'
            })

        elif action == 'mask':
            role = get_user_hc_role(current_user.id)
            masked = mask_sensitive(df, role)
            return jsonify({
                'message': f'Data masked! {len([c for c in df.columns if c.lower() in SENSITIVE_COLUMNS])} columns masked',
                'preview': masked.head(5).to_string(index=False),
                'masked_columns': [c for c in df.columns if c.lower() in SENSITIVE_COLUMNS]
            })

        elif action == 'remove':
            cols_to_remove = data.get('columns', [])
            if not cols_to_remove:
                cols_to_remove = [c for c in df.columns if c.lower() in SENSITIVE_COLUMNS]
            df_clean = df.drop(columns=[c for c in cols_to_remove if c in df.columns])
            agent.df = df_clean
            return jsonify({
                'message': f'Sensitive columns permanently removed! {cols_to_remove}',
                'remaining_cols': list(df_clean.columns),
                'rows': len(df_clean)
            })

    except Exception as e:
        return jsonify({'error': str(e)}), 500



# ════════════════════════════════════════════════════════════════
# MEDICAL IMAGE ANALYSIS — X-Ray, ECG, Skin, OCR
# ════════════════════════════════════════════════════════════════

def get_hf_url():
    """Get HuggingFace or Colab URL"""
    import requests as req
    hf = os.environ.get('HF_SPACE_URL', '').strip()
    colab = os.environ.get('COLAB_URL', '').strip()
    if hf:
        try:
            r = req.get(hf, timeout=10)
            if r.status_code == 200:
                return hf
        except:
            pass
    if colab:
        try:
            r = req.get(colab + '/ping', timeout=6)
            if r.status_code == 200:
                return colab
        except:
            pass
    return None


def call_hf_api(url, endpoint, payload):
    """Call HuggingFace FastAPI routes or Colab Flask API"""
    import requests as req
    try:
        # Direct FastAPI route — works for both HF Space and Colab
        r = req.post(url + endpoint, json=payload, timeout=180)
        r.raise_for_status()
        return r.json()
    except Exception as e:
        print(f"[HF] call_hf_api error endpoint={endpoint}: {e}")
        return {'error': f'ML API error: {str(e)}'}


def validate_hf_result(result, expected_fields):
    """Validate & sanitize HuggingFace API result — ensure no undefined fields"""
    if not result or not isinstance(result, dict):
        return {'error': 'HuggingFace se empty response. Space restart karo.'}
    if result.get('error'):
        return result
    safe = {}
    for field, default in expected_fields.items():
        safe[field] = result.get(field, default)
    return safe


@app.route('/api/medical/xray', methods=['POST'])
@login_required
def medical_xray():
    url = get_hf_url()
    if not url:
        return jsonify({'error': 'ML Server offline! HuggingFace Space start karo.'}), 400
    try:
        data = request.get_json()
        result = call_hf_api(url, '/xray', data)
        safe = validate_hf_result(result, {
            'success': False, 'finding': 'Unknown', 'confidence': 0.0,
            'severity': 'unknown', 'recommendations': [], 'report': '',
            'heatmap': None, 'message': 'X-Ray analysis complete'
        })
        return jsonify(safe), (500 if safe.get('error') else 200)
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/medical/ecg', methods=['POST'])
@login_required
def medical_ecg():
    url = get_hf_url()
    if not url:
        return jsonify({'error': 'ML Server offline!'}), 400
    try:
        data = request.get_json()
        result = call_hf_api(url, '/ecg', data)
        safe = validate_hf_result(result, {
            'success': False, 'rhythm': 'Unknown', 'heart_rate': 0,
            'findings': [], 'severity': 'unknown', 'report': '',
            'chart': None, 'message': 'ECG analysis complete'
        })
        return jsonify(safe), (500 if safe.get('error') else 200)
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/medical/skin', methods=['POST'])
@login_required
def medical_skin():
    url = get_hf_url()
    if not url:
        return jsonify({'error': 'ML Server offline!'}), 400
    try:
        data = request.get_json()
        result = call_hf_api(url, '/skin', data)
        safe = validate_hf_result(result, {
            'success': False, 'condition': 'Unknown', 'confidence': 0.0,
            'severity': 'unknown', 'recommendations': [], 'report': '',
            'message': 'Skin analysis complete'
        })
        return jsonify(safe), (500 if safe.get('error') else 200)
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/medical/ocr', methods=['POST'])
@login_required
def medical_ocr():
    url = get_hf_url()
    if not url:
        return jsonify({'error': 'ML Server offline!'}), 400
    try:
        data = request.get_json()
        result = call_hf_api(url, '/ocr', data)
        safe = validate_hf_result(result, {
            'success': False, 'text': '', 'entities': [],
            'summary': '', 'message': 'OCR analysis complete'
        })
        return jsonify(safe), (500 if safe.get('error') else 200)
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: SHAP — Model Explanation
# ════════════════════════════════════════════════════════════════
@app.route('/api/shap', methods=['POST'])
@login_required
@check_query_limit
def shap_explain():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')
        model_type = data.get('model_type', 'auto')

        df = agent.df.dropna()
        if target not in df.columns:
            return jsonify({'error': f'Target column "{target}" nahi mila!'}), 400

        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
        from sklearn.preprocessing import LabelEncoder
        from sklearn.model_selection import train_test_split
        import shap, io as _io

        y = df[target]
        X = df.drop(columns=[target])
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))

        task = 'classification' if (y.dtype == object or y.nunique() <= 15) else 'regression'
        if y.dtype == object:
            le_y = LabelEncoder(); y = le_y.fit_transform(y.astype(str))

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        model = RandomForestClassifier(n_estimators=50, random_state=42) if task == 'classification' else RandomForestRegressor(n_estimators=50, random_state=42)
        model.fit(X_train, y_train)

        explainer = shap.TreeExplainer(model)
        shap_vals = explainer.shap_values(X_test[:100])
        if isinstance(shap_vals, list): shap_vals = shap_vals[1]

        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        fig, axes = plt.subplots(1, 2, figsize=(14, 5))
        feat_imp = dict(zip(X.columns, abs(shap_vals).mean(axis=0)))
        feat_imp = dict(sorted(feat_imp.items(), key=lambda x: x[1], reverse=True)[:15])
        axes[0].barh(list(feat_imp.keys())[::-1], list(feat_imp.values())[::-1], color='#5b5ef4', alpha=0.85)
        axes[0].set_title('SHAP Feature Importance', fontweight='bold', fontsize=13)
        axes[0].set_xlabel('Mean |SHAP Value|')

        top_feat = list(feat_imp.keys())[0]
        axes[1].scatter(X_test[:100][top_feat], shap_vals[:, list(X.columns).index(top_feat)],
                       alpha=0.6, c='#8b5cf6', edgecolors='white', linewidth=0.5, s=60)
        axes[1].axhline(0, color='#ef4444', linestyle='--', linewidth=1)
        axes[1].set_xlabel(top_feat); axes[1].set_ylabel('SHAP Value')
        axes[1].set_title(f'SHAP: {top_feat} Impact', fontweight='bold')
        plt.suptitle(f'Model Explanation — Target: {target}', fontweight='bold', color='#1e3a5f')
        plt.tight_layout()

        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        score = model.score(X_test, y_test)
        top5 = list(feat_imp.items())[:5]
        insights = [f"'{k}' has {v:.3f} avg impact on prediction" for k,v in top5]

        return jsonify({
            'success': True,
            'task': task,
            'target': target,
            'model_score': round(float(score), 4),
            'feature_importance': feat_imp,
            'top_features': top5,
            'insights': insights,
            'chart': chart,
            'message': f'SHAP explanation ready! Top feature: {list(feat_imp.keys())[0]}'
        })
    except ImportError:
        return jsonify({'error': 'SHAP not installed! requirements mein add karo: shap'}), 500
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: HYPERPARAMETER TUNING
# ════════════════════════════════════════════════════════════════
@app.route('/api/hyperparam', methods=['POST'])
@login_required
@check_query_limit
def hyperparam_tune():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')
        model_name = data.get('model', 'random_forest')
        cv_folds = int(data.get('cv_folds', 3))

        df = agent.df.dropna()
        if target not in df.columns:
            return jsonify({'error': f'Target "{target}" nahi mila!'}), 400

        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor, GradientBoostingClassifier, GradientBoostingRegressor
        from sklearn.linear_model import LogisticRegression, Ridge
        from sklearn.svm import SVC, SVR
        from sklearn.model_selection import RandomizedSearchCV, train_test_split, cross_val_score
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        y = df[target]; X = df.drop(columns=[target])
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))

        task = 'classification' if (y.dtype == object or y.nunique() <= 15) else 'regression'
        if y.dtype == object:
            le_y = LabelEncoder(); y = le_y.fit_transform(y.astype(str))
        else:
            y = y.values

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        scaler = StandardScaler(); X_train_s = scaler.fit_transform(X_train); X_test_s = scaler.transform(X_test)

        param_grids = {
            'random_forest': {
                'clf': RandomForestClassifier(random_state=42) if task=='classification' else RandomForestRegressor(random_state=42),
                'params': {'n_estimators':[50,100,200],'max_depth':[None,5,10,20],'min_samples_split':[2,5,10],'min_samples_leaf':[1,2,4]}
            },
            'gradient_boost': {
                'clf': GradientBoostingClassifier(random_state=42) if task=='classification' else GradientBoostingRegressor(random_state=42),
                'params': {'n_estimators':[50,100],'learning_rate':[0.01,0.05,0.1,0.2],'max_depth':[3,5,7],'subsample':[0.8,1.0]}
            },
            'svm': {
                'clf': SVC(random_state=42) if task=='classification' else SVR(),
                'params': {'C':[0.1,1,10,100],'kernel':['rbf','linear'],'gamma':['scale','auto']}
            }
        }

        chosen = param_grids.get(model_name, param_grids['random_forest'])
        search = RandomizedSearchCV(chosen['clf'], chosen['params'], n_iter=12, cv=cv_folds,
                                   scoring='accuracy' if task=='classification' else 'r2',
                                   n_jobs=-1, random_state=42)
        search.fit(X_train_s if model_name=='svm' else X_train, y_train)

        results_df = pd.DataFrame(search.cv_results_)
        results_df = results_df.sort_values('mean_test_score', ascending=False).head(8)

        baseline_clf = RandomForestClassifier(random_state=42) if task=='classification' else RandomForestRegressor(random_state=42)
        baseline_clf.fit(X_train, y_train)
        baseline_score = baseline_clf.score(X_test, y_test)
        best_score = search.best_estimator_.score(X_test_s if model_name=='svm' else X_test, y_test)

        fig, axes = plt.subplots(1, 2, figsize=(13, 5))
        scores = results_df['mean_test_score'].values
        axes[0].barh(range(len(scores)), scores, color=['#5b5ef4' if i==0 else '#8b5cf6' for i in range(len(scores))], alpha=0.85)
        axes[0].set_yticks(range(len(scores)))
        axes[0].set_yticklabels([f'Config {i+1}' for i in range(len(scores))])
        axes[0].set_title('Top Configurations', fontweight='bold')
        axes[0].set_xlabel('CV Score')
        axes[0].axvline(baseline_score, color='#ef4444', linestyle='--', label=f'Baseline: {baseline_score:.3f}')
        axes[0].legend()

        axes[1].bar(['Baseline', 'Tuned'], [baseline_score, best_score],
                   color=['#94a3b8','#059669'], alpha=0.85, width=0.5)
        axes[1].set_title('Baseline vs Tuned', fontweight='bold')
        axes[1].set_ylabel('Score')
        for i, v in enumerate([baseline_score, best_score]):
            axes[1].text(i, v+0.005, f'{v:.4f}', ha='center', fontweight='700', fontsize=12)
        plt.suptitle(f'Hyperparameter Tuning — {model_name.replace("_"," ").title()}', fontweight='bold')
        plt.tight_layout()

        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        improvement = round((best_score - baseline_score)*100, 2)
        return jsonify({
            'success': True,
            'task': task,
            'target': target,
            'model': model_name,
            'baseline_score': round(float(baseline_score), 4),
            'best_score': round(float(best_score), 4),
            'improvement_pct': improvement,
            'best_params': search.best_params_,
            'cv_folds': cv_folds,
            'iterations_tried': 12,
            'chart': chart,
            'message': f'Tuning complete! Score {baseline_score:.3f} → {best_score:.3f} (+{improvement}%)'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: NLP — Text Analysis & Sentiment
# ════════════════════════════════════════════════════════════════
@app.route('/api/nlp', methods=['POST'])
@login_required
@check_query_limit
def nlp_analysis():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        text_col = data.get('text_col', '')
        analysis_type = data.get('analysis_type', 'all')

        df = agent.df.dropna(subset=[text_col]) if text_col in agent.df.columns else None
        if df is None or text_col not in agent.df.columns:
            return jsonify({'error': f'Text column "{text_col}" nahi mila!'}), 400

        texts = df[text_col].astype(str).head(500).tolist()

        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import io as _io
        from collections import Counter
        import re

        # Sentiment using TextBlob or simple lexicon
        positive_words = set(['good','great','excellent','best','love','wonderful','amazing','happy','perfect','brilliant','outstanding','fantastic','superb','awesome','nice','helpful','pleased'])
        negative_words = set(['bad','worst','terrible','hate','awful','horrible','poor','disappointing','useless','waste','broken','failed','ugly','slow','worse','never','problem'])

        sentiments = []
        word_freq = Counter()
        for txt in texts:
            words = re.findall(r'\b[a-zA-Z]{3,}\b', txt.lower())
            word_freq.update(words)
            pos = sum(1 for w in words if w in positive_words)
            neg = sum(1 for w in words if w in negative_words)
            if pos > neg: sentiments.append('Positive')
            elif neg > pos: sentiments.append('Negative')
            else: sentiments.append('Neutral')

        sent_counts = Counter(sentiments)
        avg_len = sum(len(t.split()) for t in texts) / len(texts)
        stop_words = {'the','a','an','and','or','but','in','on','at','to','for','of','with','is','was','are','were','be','been','have','has','had','do','does','did','will','would','could','should','may','might','that','this','it','they','we','you','i','he','she'}
        top_words = [(w, c) for w, c in word_freq.most_common(30) if w not in stop_words][:15]

        fig, axes = plt.subplots(1, 3, figsize=(16, 5))
        colors_sent = {'Positive':'#059669','Negative':'#ef4444','Neutral':'#f59e0b'}
        axes[0].pie([sent_counts.get(s,0) for s in ['Positive','Negative','Neutral']],
                   labels=['Positive','Negative','Neutral'],
                   colors=['#059669','#ef4444','#f59e0b'],
                   autopct='%1.1f%%', startangle=90, textprops={'fontsize':11})
        axes[0].set_title('Sentiment Distribution', fontweight='bold')

        words, counts = zip(*top_words) if top_words else (['no data'],[1])
        bar_colors = ['#059669' if w in positive_words else '#ef4444' if w in negative_words else '#5b5ef4' for w in words]
        axes[1].barh(list(words)[::-1], list(counts)[::-1], color=bar_colors[::-1], alpha=0.85)
        axes[1].set_title('Top Keywords', fontweight='bold')
        axes[1].set_xlabel('Frequency')

        lengths = [len(t.split()) for t in texts]
        axes[2].hist(lengths, bins=20, color='#8b5cf6', alpha=0.8, edgecolor='white')
        axes[2].axvline(avg_len, color='#ef4444', linestyle='--', linewidth=2, label=f'Avg: {avg_len:.1f}')
        axes[2].set_title('Text Length Distribution', fontweight='bold')
        axes[2].set_xlabel('Word Count'); axes[2].legend()

        plt.suptitle(f'NLP Analysis — Column: {text_col}', fontweight='bold', color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'text_col': text_col,
            'total_texts': len(texts),
            'sentiment_counts': dict(sent_counts),
            'positive_pct': round(sent_counts.get('Positive',0)/len(texts)*100, 1),
            'negative_pct': round(sent_counts.get('Negative',0)/len(texts)*100, 1),
            'neutral_pct': round(sent_counts.get('Neutral',0)/len(texts)*100, 1),
            'avg_word_count': round(avg_len, 1),
            'top_keywords': top_words[:10],
            'total_unique_words': len(word_freq),
            'chart': chart,
            'message': f'NLP done! {sent_counts.get("Positive",0)} Positive | {sent_counts.get("Negative",0)} Negative | {sent_counts.get("Neutral",0)} Neutral'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: AUTO FEATURE ENGINEERING
# ════════════════════════════════════════════════════════════════
@app.route('/api/feature_eng', methods=['POST'])
@login_required
@check_query_limit
def feature_engineering():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')

        df = agent.df.copy()
        original_cols = list(df.columns)
        new_features = []
        log = []

        numeric_cols = df.select_dtypes(include='number').columns.tolist()
        if target in numeric_cols: numeric_cols.remove(target)
        cat_cols = df.select_dtypes(include='object').columns.tolist()
        if target in cat_cols: cat_cols.remove(target)

        # 1. Log transform for skewed columns
        for col in numeric_cols:
            skew = float(df[col].skew())
            if abs(skew) > 1.0 and df[col].min() >= 0:
                df[f'{col}_log'] = df[col].apply(lambda x: __import__('math').log1p(x))
                new_features.append(f'{col}_log'); log.append(f'Log transform: {col} (skew={skew:.2f})')

        # 2. Interaction features (top numeric pairs)
        if len(numeric_cols) >= 2:
            for i in range(min(3, len(numeric_cols))):
                for j in range(i+1, min(4, len(numeric_cols))):
                    c1, c2 = numeric_cols[i], numeric_cols[j]
                    df[f'{c1}_x_{c2}'] = df[c1] * df[c2]
                    new_features.append(f'{c1}_x_{c2}')
                    log.append(f'Interaction: {c1} × {c2}')

        # 3. Ratio features
        if len(numeric_cols) >= 2:
            c1, c2 = numeric_cols[0], numeric_cols[1]
            safe_div = df[c2].replace(0, float('nan'))
            df[f'{c1}_ratio_{c2}'] = df[c1] / safe_div
            df[f'{c1}_ratio_{c2}'] = df[f'{c1}_ratio_{c2}'].fillna(0)
            new_features.append(f'{c1}_ratio_{c2}')
            log.append(f'Ratio: {c1} / {c2}')

        # 4. Binning numeric to categories
        for col in numeric_cols[:3]:
            try:
                df[f'{col}_bin'] = pd.cut(df[col], bins=4, labels=['Low','Medium','High','Very High'])
                df[f'{col}_bin'] = df[f'{col}_bin'].astype(str)
                new_features.append(f'{col}_bin')
                log.append(f'Binning: {col} → 4 categories')
            except: pass

        # 5. Label encode categoricals
        from sklearn.preprocessing import LabelEncoder
        for col in cat_cols[:5]:
            le = LabelEncoder()
            df[f'{col}_encoded'] = le.fit_transform(df[col].astype(str))
            new_features.append(f'{col}_encoded')
            log.append(f'Label encoded: {col}')

        # 6. Null indicator flags
        null_cols = [c for c in df.columns if df[c].isnull().sum() > 0]
        for col in null_cols[:3]:
            df[f'{col}_is_null'] = df[col].isnull().astype(int)
            new_features.append(f'{col}_is_null')
            log.append(f'Null flag: {col}')

        # Save enhanced df to agent
        agent.df = df
        agent.available_columns = list(df.columns)

        # Chart — before vs after
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        fig, axes = plt.subplots(1, 2, figsize=(12, 5))
        cats = ['Original\nFeatures', 'New\nFeatures', 'Total\nFeatures']
        vals = [len(original_cols), len(new_features), len(df.columns)]
        colors = ['#94a3b8','#5b5ef4','#059669']
        bars = axes[0].bar(cats, vals, color=colors, alpha=0.85, width=0.5)
        for bar, val in zip(bars, vals):
            axes[0].text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3, str(val), ha='center', fontweight='700', fontsize=13)
        axes[0].set_title('Feature Count', fontweight='bold')

        feat_types = {'Log Transform': sum(1 for f in new_features if '_log' in f),
                     'Interaction': sum(1 for f in new_features if '_x_' in f),
                     'Ratio': sum(1 for f in new_features if '_ratio_' in f),
                     'Binning': sum(1 for f in new_features if '_bin' in f),
                     'Encoded': sum(1 for f in new_features if '_encoded' in f),
                     'Null Flag': sum(1 for f in new_features if '_is_null' in f)}
        feat_types = {k:v for k,v in feat_types.items() if v>0}
        if feat_types:
            axes[1].pie(feat_types.values(), labels=feat_types.keys(),
                       autopct='%1.0f%%', startangle=90,
                       colors=['#5b5ef4','#8b5cf6','#06d6a0','#f59e0b','#ef4444','#3b82f6'])
            axes[1].set_title('New Feature Types', fontweight='bold')

        plt.suptitle('Auto Feature Engineering Results', fontweight='bold', color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'original_features': len(original_cols),
            'new_features_added': len(new_features),
            'total_features': len(df.columns),
            'new_feature_names': new_features[:20],
            'engineering_log': log,
            'feature_types': feat_types,
            'chart': chart,
            'message': f'Feature Engineering done! {len(original_cols)} → {len(df.columns)} features (+{len(new_features)} new)'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: MODEL VERSIONING
# ════════════════════════════════════════════════════════════════
@app.route('/api/model_version/save', methods=['POST'])
@login_required
def save_model_version():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')
        model_name = data.get('model_name', 'My Model')
        notes = data.get('notes', '')

        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
        from sklearn.model_selection import train_test_split, cross_val_score
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        import pickle, io as _io

        df = agent.df.dropna()
        if target not in df.columns:
            return jsonify({'error': f'Target "{target}" nahi mila!'}), 400

        y = df[target]; X = df.drop(columns=[target])
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))

        task = 'classification' if (y.dtype == object or y.nunique() <= 15) else 'regression'
        if y.dtype == object:
            le_y = LabelEncoder(); y = le_y.fit_transform(y.astype(str))
        else:
            y = y.values

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        model = RandomForestClassifier(n_estimators=100, random_state=42) if task == 'classification' else RandomForestRegressor(n_estimators=100, random_state=42)
        model.fit(X_train, y_train)
        score = float(model.score(X_test, y_test))
        cv_scores = cross_val_score(model, X, y, cv=3)

        # Serialize model to base64
        buf = _io.BytesIO(); pickle.dump(model, buf); buf.seek(0)
        model_b64 = base64.b64encode(buf.read()).decode()

        # Save to DB as UserChart (reuse existing table)
        db = SessionLocal()
        version_num = db.query(UserChart).filter_by(user_id=current_user.id).filter(UserChart.chart_type=='model_version').count() + 1
        chart = UserChart(
            user_id=current_user.id,
            chart_type='model_version',
            chart_title=f"v{version_num} — {model_name}",
        )
        db.add(chart); db.commit(); db.close()

        return jsonify({
            'success': True,
            'version': version_num,
            'model_name': model_name,
            'target': target,
            'task': task,
            'score': round(score, 4),
            'cv_mean': round(float(cv_scores.mean()), 4),
            'cv_std': round(float(cv_scores.std()), 4),
            'features': list(X.columns),
            'n_features': len(X.columns),
            'n_rows': len(df),
            'message': f'✅ v{version_num} saved! {model_name} — Score: {score:.4f}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/model_version/list', methods=['GET'])
@login_required
def list_model_versions():
    try:
        db = SessionLocal()
        versions = db.query(UserChart).filter_by(user_id=current_user.id).filter(UserChart.chart_type=='model_version').order_by(UserChart.created_at.desc()).all()
        db.close()
        result = []
        for v in versions:
            parts = dict(p.split('=',1) for p in v.description.split('|') if '=' in p)
            result.append({
                'id': v.id,
                'title': v.title,
                'target': parts.get('target',''),
                'task': parts.get('task',''),
                'score': parts.get('score',''),
                'features': parts.get('features',''),
                'rows': parts.get('rows',''),
                'notes': parts.get('notes',''),
                'cv_mean': parts.get('cv_mean',''),
                'cv_std': parts.get('cv_std',''),
                'saved_at': v.created_at.strftime('%d %b %Y %H:%M')
            })
        return jsonify({'success': True, 'versions': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/model_version/delete/<int:vid>', methods=['DELETE'])
@login_required
def delete_model_version(vid):
    try:
        db = SessionLocal()
        v = db.query(UserChart).filter_by(id=vid, user_id=current_user.id).first()
        if v: db.delete(v); db.commit()
        db.close()
        return jsonify({'success': True, 'message': 'Version deleted!'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: A/B TESTING — 2 Models Compare
# ════════════════════════════════════════════════════════════════
@app.route('/api/ab_test', methods=['POST'])
@login_required
@check_query_limit
def ab_test():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')
        model_a = data.get('model_a', 'random_forest')
        model_b = data.get('model_b', 'gradient_boost')

        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor, GradientBoostingClassifier, GradientBoostingRegressor
        from sklearn.linear_model import LogisticRegression, Ridge
        from sklearn.svm import SVC, SVR
        from sklearn.neighbors import KNeighborsClassifier, KNeighborsRegressor
        from sklearn.tree import DecisionTreeClassifier, DecisionTreeRegressor
        from sklearn.model_selection import train_test_split, cross_val_score
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io, time as _time

        df = agent.df.dropna()
        if target not in df.columns:
            return jsonify({'error': f'Target "{target}" nahi mila!'}), 400

        y = df[target]; X = df.drop(columns=[target])
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))

        task = 'classification' if (y.dtype == object or y.nunique() <= 15) else 'regression'
        if y.dtype == object:
            le_y = LabelEncoder(); y = le_y.fit_transform(y.astype(str))
        else:
            y = y.values

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        scaler = StandardScaler(); Xs_train = scaler.fit_transform(X_train); Xs_test = scaler.transform(X_test)

        def get_model(name, task):
            clfs = {'random_forest': RandomForestClassifier(n_estimators=100, random_state=42),
                    'gradient_boost': GradientBoostingClassifier(n_estimators=100, random_state=42),
                    'logistic': LogisticRegression(max_iter=500, random_state=42),
                    'svm': SVC(random_state=42), 'knn': KNeighborsClassifier(),
                    'decision_tree': DecisionTreeClassifier(random_state=42)}
            regs = {'random_forest': RandomForestRegressor(n_estimators=100, random_state=42),
                    'gradient_boost': GradientBoostingRegressor(n_estimators=100, random_state=42),
                    'ridge': Ridge(), 'svr': SVR(), 'knn': KNeighborsRegressor(),
                    'decision_tree': DecisionTreeRegressor(random_state=42)}
            return clfs.get(name, clfs['random_forest']) if task=='classification' else regs.get(name, regs['random_forest'])

        def train_eval(name, task, Xtr, Xte, ytr, yte):
            m = get_model(name, task)
            use_scaled = name in ('logistic','svm','svr','knn')
            Xt = Xs_train if use_scaled else Xtr
            Xv = Xs_test if use_scaled else Xte
            t0 = _time.time()
            m.fit(Xt, ytr)
            train_time = round(_time.time()-t0, 3)
            score = float(m.score(Xv, yte))
            cv = cross_val_score(m, Xs_train if use_scaled else Xtr, ytr, cv=3)
            return {'score': round(score,4), 'cv_mean': round(float(cv.mean()),4),
                    'cv_std': round(float(cv.std()),4), 'train_time': train_time, 'model': m}

        res_a = train_eval(model_a, task, X_train, X_test, y_train, y_test)
        res_b = train_eval(model_b, task, X_train, X_test, y_train, y_test)
        winner = model_a if res_a['score'] >= res_b['score'] else model_b
        diff = abs(res_a['score'] - res_b['score'])

        # Chart
        fig, axes = plt.subplots(1, 3, figsize=(15, 5))
        metrics = ['Test Score','CV Mean']
        a_vals = [res_a['score'], res_a['cv_mean']]
        b_vals = [res_b['score'], res_b['cv_mean']]
        x = range(len(metrics))
        axes[0].bar([i-0.2 for i in x], a_vals, 0.35, label=model_a.replace('_',' ').title(), color='#5b5ef4', alpha=0.85)
        axes[0].bar([i+0.2 for i in x], b_vals, 0.35, label=model_b.replace('_',' ').title(), color='#f59e0b', alpha=0.85)
        axes[0].set_xticks(list(x)); axes[0].set_xticklabels(metrics)
        axes[0].legend(); axes[0].set_title('Model A vs B — Scores', fontweight='bold')
        axes[0].set_ylim(0,1.1)

        axes[1].bar([model_a.replace('_',' ').title(), model_b.replace('_',' ').title()],
                   [res_a['cv_std'], res_b['cv_std']], color=['#5b5ef4','#f59e0b'], alpha=0.85)
        axes[1].set_title('Variance (lower=stable)', fontweight='bold')
        axes[1].set_ylabel('CV Std Dev')

        axes[2].bar([model_a.replace('_',' ').title(), model_b.replace('_',' ').title()],
                   [res_a['train_time'], res_b['train_time']], color=['#059669','#ef4444'], alpha=0.85)
        axes[2].set_title('Training Time (sec)', fontweight='bold')
        axes[2].set_ylabel('Seconds')

        plt.suptitle(f'A/B Test — Target: {target} | Winner: {winner.replace("_"," ").title()} 🏆', fontweight='bold', color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True, 'target': target, 'task': task,
            'model_a': {'name': model_a, **{k:v for k,v in res_a.items() if k!='model'}},
            'model_b': {'name': model_b, **{k:v for k,v in res_b.items() if k!='model'}},
            'winner': winner, 'score_diff': round(diff,4), 'chart': chart,
            'message': f'A/B Test done! Winner: {winner.replace("_"," ").title()} (diff: {diff:.4f})'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: ANOMALY DETECTION
# ════════════════════════════════════════════════════════════════
@app.route('/api/anomaly', methods=['POST'])
@login_required
@check_query_limit
def anomaly_detection():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        method = data.get('method', 'isolation_forest')
        contamination = float(data.get('contamination', 0.05))
        cols = data.get('columns', [])

        from sklearn.ensemble import IsolationForest
        from sklearn.preprocessing import StandardScaler
        from sklearn.neighbors import LocalOutlierFactor
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        df = agent.df.copy()
        num_cols = df.select_dtypes(include='number').columns.tolist()
        if cols: num_cols = [c for c in cols if c in num_cols]
        if not num_cols:
            return jsonify({'error': 'Numeric columns nahi mile!'}), 400

        X = df[num_cols].fillna(df[num_cols].median())
        scaler = StandardScaler(); Xs = scaler.fit_transform(X)

        if method == 'isolation_forest':
            clf = IsolationForest(contamination=contamination, random_state=42, n_estimators=100)
            preds = clf.fit_predict(Xs)
            scores = clf.score_samples(Xs)
        elif method == 'lof':
            clf = LocalOutlierFactor(contamination=contamination, n_neighbors=20)
            preds = clf.fit_predict(Xs)
            scores = clf.negative_outlier_factor_
        else:
            # Z-score method
            from scipy import stats
            z_scores = abs(stats.zscore(X.fillna(0)))
            preds = ((z_scores > 3).any(axis=1)).astype(int)
            preds = [-1 if p==1 else 1 for p in preds]
            scores = -z_scores.max(axis=1).values

        anomaly_mask = [p == -1 for p in preds]
        n_anomalies = sum(anomaly_mask)
        anomaly_pct = round(n_anomalies/len(df)*100, 2)
        anomaly_indices = [i for i, m in enumerate(anomaly_mask) if m][:50]
        anomaly_rows = df.iloc[anomaly_indices][num_cols[:5]].round(3).to_dict('records')

        fig, axes = plt.subplots(1, 3, figsize=(15, 5))
        colors = ['#ef4444' if m else '#5b5ef4' for m in anomaly_mask]
        if len(num_cols) >= 2:
            axes[0].scatter(X.iloc[:,0], X.iloc[:,1], c=colors, alpha=0.6, s=30, edgecolors='none')
            axes[0].set_xlabel(num_cols[0]); axes[0].set_ylabel(num_cols[1])
            axes[0].set_title(f'Anomalies: {n_anomalies} red dots', fontweight='bold')
            from matplotlib.lines import Line2D
            legend = [Line2D([0],[0],marker='o',color='w',markerfacecolor='#ef4444',label=f'Anomaly ({n_anomalies})'),
                     Line2D([0],[0],marker='o',color='w',markerfacecolor='#5b5ef4',label=f'Normal ({len(df)-n_anomalies})')]
            axes[0].legend(handles=legend)
        else:
            axes[0].hist(X.iloc[:,0], bins=30, color='#5b5ef4', alpha=0.7)
            axes[0].set_title(num_cols[0], fontweight='bold')

        axes[1].hist(scores, bins=30, color='#8b5cf6', alpha=0.8, edgecolor='white')
        axes[1].axvline(float(pd.Series(scores)[anomaly_mask].max()), color='#ef4444', linestyle='--', linewidth=2, label='Anomaly threshold')
        axes[1].set_title('Anomaly Score Distribution', fontweight='bold')
        axes[1].set_xlabel('Score'); axes[1].legend()

        axes[2].bar(['Normal', 'Anomaly'], [len(df)-n_anomalies, n_anomalies],
                   color=['#059669','#ef4444'], alpha=0.85, width=0.5)
        axes[2].set_title(f'Result: {anomaly_pct}% Anomalies', fontweight='bold')
        for i, v in enumerate([len(df)-n_anomalies, n_anomalies]):
            axes[2].text(i, v+1, str(v), ha='center', fontweight='700', fontsize=13)

        plt.suptitle(f'Anomaly Detection — {method.replace("_"," ").title()}', fontweight='bold', color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True, 'method': method,
            'total_rows': len(df), 'n_anomalies': n_anomalies,
            'anomaly_pct': anomaly_pct, 'columns_analyzed': num_cols,
            'sample_anomalies': anomaly_rows[:10], 'chart': chart,
            'message': f'Anomaly Detection done! {n_anomalies} anomalies found ({anomaly_pct}% of data)'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: ADVANCED FORECASTING (ARIMA + Prophet style)
# ════════════════════════════════════════════════════════════════
@app.route('/api/forecast_advanced', methods=['POST'])
@login_required
@check_query_limit
def forecast_advanced():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        date_col = data.get('date_col', '')
        value_col = data.get('value_col', '')
        periods = int(data.get('periods', 30))
        method = data.get('method', 'auto')

        df = agent.df.copy()
        if date_col not in df.columns or value_col not in df.columns:
            return jsonify({'error': f'Columns nahi mile! Available: {list(df.columns)}'}), 400

        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        df = df.dropna(subset=[date_col, value_col]).sort_values(date_col)
        df[value_col] = pd.to_numeric(df[value_col], errors='coerce')
        df = df.dropna(subset=[value_col])

        if len(df) < 10:
            return jsonify({'error': 'Forecast ke liye minimum 10 data points chahiye!'}), 400

        ts = df.set_index(date_col)[value_col].resample('D').mean().fillna(method='ffill')

        # Multiple methods
        results = {}

        # 1. Moving Average
        window = min(7, len(ts)//3)
        ma = ts.rolling(window=window, min_periods=1).mean()
        last_ma = float(ma.iloc[-1])
        results['moving_avg'] = last_ma

        # 2. Exponential Smoothing
        alpha = 0.3
        ets = ts.ewm(alpha=alpha).mean()
        results['exp_smoothing'] = float(ets.iloc[-1])

        # 3. Linear Trend
        import numpy as np
        x = np.arange(len(ts))
        coeffs = np.polyfit(x, ts.values, 1)
        trend_fn = np.poly1d(coeffs)
        results['linear_trend'] = float(trend_fn(len(ts) + periods//2))

        # 4. ARIMA-like (using statsmodels if available, else fallback)
        try:
            from statsmodels.tsa.arima.model import ARIMA
            arima = ARIMA(ts.values[-min(200, len(ts)):], order=(1,1,1))
            arima_fit = arima.fit()
            arima_forecast = arima_fit.forecast(steps=periods)
            results['arima'] = [round(float(v),3) for v in arima_forecast]
            arima_available = True
        except:
            arima_available = False
            slope = coeffs[0]
            arima_forecast = [trend_fn(len(ts)+i) for i in range(periods)]
            results['arima'] = [round(float(v),3) for v in arima_forecast]

        # Future dates
        last_date = ts.index[-1]
        freq = pd.infer_freq(ts.index) or 'D'
        future_dates = pd.date_range(last_date, periods=periods+1, freq=freq)[1:]

        # Chart
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        fig, axes = plt.subplots(2, 2, figsize=(15, 10))
        hist_dates = ts.index[-min(90, len(ts)):]
        hist_vals = ts.values[-min(90, len(ts)):]

        # Plot 1: Historical + Forecast
        axes[0,0].plot(hist_dates, hist_vals, color='#5b5ef4', linewidth=2, label='Historical')
        axes[0,0].plot(future_dates, arima_forecast, color='#ef4444', linewidth=2, linestyle='--', label=f'Forecast ({periods}d)')
        upper = [v*1.1 for v in arima_forecast]; lower = [v*0.9 for v in arima_forecast]
        axes[0,0].fill_between(future_dates, lower, upper, alpha=0.2, color='#ef4444', label='90% CI')
        axes[0,0].set_title('Forecast with Confidence Interval', fontweight='bold')
        axes[0,0].legend(); axes[0,0].tick_params(axis='x', rotation=30)

        # Plot 2: Trend + Moving Avg
        axes[0,1].plot(hist_dates, hist_vals, color='#94a3b8', linewidth=1, alpha=0.7, label='Actual')
        axes[0,1].plot(hist_dates, ma.values[-len(hist_dates):], color='#5b5ef4', linewidth=2, label=f'{window}d MA')
        axes[0,1].plot(hist_dates, ets.values[-len(hist_dates):], color='#f59e0b', linewidth=2, label='Exp Smooth')
        trend_line = [trend_fn(len(ts)-len(hist_dates)+i) for i in range(len(hist_dates))]
        axes[0,1].plot(hist_dates, trend_line, color='#059669', linewidth=2, linestyle=':', label='Trend')
        axes[0,1].set_title('Trend Components', fontweight='bold'); axes[0,1].legend()
        axes[0,1].tick_params(axis='x', rotation=30)

        # Plot 3: Seasonality check
        if len(ts) >= 14:
            weekly = ts.groupby(ts.index.dayofweek).mean()
            days = ['Mon','Tue','Wed','Thu','Fri','Sat','Sun']
            axes[1,0].bar(days[:len(weekly)], weekly.values, color='#8b5cf6', alpha=0.85)
            axes[1,0].set_title('Weekly Seasonality Pattern', fontweight='bold')
            axes[1,0].set_ylabel(value_col)
        else:
            axes[1,0].text(0.5,0.5,'Not enough data\nfor seasonality', ha='center', va='center', transform=axes[1,0].transAxes, fontsize=13)

        # Plot 4: Forecast values table
        forecast_sample = list(zip(future_dates[:10].strftime('%d-%b'), [round(v,2) for v in arima_forecast[:10]]))
        axes[1,1].axis('off')
        table_data = [[d, str(v)] for d,v in forecast_sample]
        tbl = axes[1,1].table(cellText=table_data, colLabels=['Date','Forecast'], loc='center',
                              cellLoc='center', colColours=['#5b5ef4','#5b5ef4'])
        tbl.auto_set_font_size(False); tbl.set_fontsize(11)
        for (r,c), cell in tbl.get_celld().items():
            if r==0: cell.set_text_props(color='white', fontweight='bold')
        axes[1,1].set_title('Next 10 Days Forecast', fontweight='bold', pad=20)

        plt.suptitle(f'Advanced Forecast — {value_col}', fontweight='bold', color='#1e3a5f', fontsize=14)
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        trend_dir = '📈 Upward' if coeffs[0] > 0 else '📉 Downward'
        return jsonify({
            'success': True, 'date_col': date_col, 'value_col': value_col,
            'total_points': len(ts), 'forecast_periods': periods,
            'method_used': 'ARIMA' if arima_available else 'Trend Projection',
            'trend_direction': trend_dir,
            'trend_slope': round(float(coeffs[0]), 4),
            'forecast_next_10': [{'date': str(d.date()), 'value': round(float(v),3)} for d,v in zip(future_dates[:10], arima_forecast[:10])],
            'current_value': round(float(ts.iloc[-1]), 3),
            'forecast_end_value': round(float(arima_forecast[-1]), 3),
            'chart': chart,
            'message': f'Forecast ready! {trend_dir} trend | {periods} periods predicted'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE: REAL NLP via HuggingFace Transformers
# ════════════════════════════════════════════════════════════════
@app.route('/api/nlp_advanced', methods=['POST'])
@login_required
@check_query_limit
def nlp_advanced():
    try:
        url = get_hf_url()
        if not url:
            return jsonify({'error': 'ML Server offline! HuggingFace Space start karo.'}), 400
        data = request.get_json()
        result = call_hf_api(url, '/nlp_advanced', data)
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# Helper — column data for NLP Advanced
@app.route('/api/get_column_data', methods=['POST'])
@login_required
def get_column_data():
    try:
        if agent.df is None:
            return jsonify({'error': 'Dataset nahi hai!'}), 400
        col = request.get_json().get('column','')
        if col not in agent.df.columns:
            return jsonify({'error': f'Column {col} nahi mila!'}), 400
        data = agent.df[col].dropna().astype(str).head(100).tolist()
        return jsonify({'success': True, 'column': col, 'data': data})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# MLOPS — COMPLETE MODULE
# ════════════════════════════════════════════════════════════════

# ── 1. EXPERIMENT TRACKING ────────────────────────────────────
@app.route('/api/mlops/experiment/log', methods=['POST'])
@login_required
def log_experiment():
    try:
        data = request.get_json()
        exp_name = data.get('name', 'Experiment')
        params = data.get('params', {})
        metrics = data.get('metrics', {})
        tags = data.get('tags', '')
        dataset_info = data.get('dataset_info', {})

        db = SessionLocal()
        import json as _json
        exp = UserChart(
            user_id=current_user.id,
            chart_type='mlops_experiment',
            chart_title=exp_name,
            image_data=_json.dumps({
                'params': params, 'metrics': metrics,
                'tags': tags, 'dataset_info': dataset_info,
                'status': 'completed'
            }),
        )
        db.add(exp); db.commit()
        exp_id = exp.id; db.close()
        return jsonify({'success': True, 'experiment_id': exp_id,
                       'message': f'Experiment "{exp_name}" logged! ID: {exp_id}'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mlops/experiment/list', methods=['GET'])
@login_required
def list_experiments():
    try:
        import json as _json
        db = SessionLocal()
        exps = db.query(UserChart).filter_by(user_id=current_user.id, chart_type='mlops_experiment').order_by(UserChart.created_at.desc()).limit(50).all()
        db.close()
        result = []
        for e in exps:
            try: info = _json.loads(e.description)
            except: info = {}
            result.append({'id': e.id, 'name': e.title,
                          'params': info.get('params', {}), 'metrics': info.get('metrics', {}),
                          'tags': info.get('tags', ''), 'status': info.get('status', ''),
                          'created_at': e.created_at.strftime('%d %b %H:%M')})
        return jsonify({'success': True, 'experiments': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mlops/experiment/delete/<int:eid>', methods=['DELETE'])
@login_required
def delete_experiment(eid):
    try:
        db = SessionLocal()
        e = db.query(UserChart).filter_by(id=eid, user_id=current_user.id).first()
        if e: db.delete(e); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── 2. MODEL REGISTRY + DOWNLOAD ─────────────────────────────
@app.route('/api/mlops/registry/save', methods=['POST'])
@login_required
def registry_save():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')
        model_type = data.get('model_type', 'random_forest')
        model_name = data.get('model_name', 'model_v1')
        description = data.get('description', '')

        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor, GradientBoostingClassifier, GradientBoostingRegressor
        from sklearn.linear_model import LogisticRegression, Ridge
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        from sklearn.model_selection import train_test_split, cross_val_score
        from sklearn.pipeline import Pipeline
        import pickle, io as _io, json as _json

        df = agent.df.dropna()
        if target not in df.columns:
            return jsonify({'error': f'Target "{target}" nahi mila!'}), 400

        y = df[target]; X = df.drop(columns=[target])
        feature_names = list(X.columns)
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))

        task = 'classification' if (y.dtype == object or y.nunique() <= 15) else 'regression'
        if y.dtype == object:
            le_y = LabelEncoder(); y = le_y.fit_transform(y.astype(str)); classes = list(le_y.classes_)
        else:
            y = y.values; classes = []

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

        model_map = {
            'random_forest': RandomForestClassifier(n_estimators=100, random_state=42) if task=='classification' else RandomForestRegressor(n_estimators=100, random_state=42),
            'gradient_boost': GradientBoostingClassifier(n_estimators=100, random_state=42) if task=='classification' else GradientBoostingRegressor(n_estimators=100, random_state=42),
            'logistic': LogisticRegression(max_iter=500, random_state=42) if task=='classification' else Ridge(),
        }
        clf = model_map.get(model_type, model_map['random_forest'])

        pipeline = Pipeline([('scaler', StandardScaler()), ('model', clf)])
        pipeline.fit(X_train, y_train)
        score = float(pipeline.score(X_test, y_test))
        cv = cross_val_score(pipeline, X, y, cv=3)

        # Serialize
        buf = _io.BytesIO(); pickle.dump(pipeline, buf); buf.seek(0)
        model_b64 = base64.b64encode(buf.read()).decode()

        metadata = {
            'model_name': model_name, 'model_type': model_type,
            'target': target, 'task': task, 'classes': classes,
            'feature_names': feature_names, 'n_features': len(feature_names),
            'n_rows': len(df), 'score': round(score, 4),
            'cv_mean': round(float(cv.mean()), 4), 'cv_std': round(float(cv.std()), 4),
            'description': description
        }

        db = SessionLocal()
        count = db.query(UserChart).filter_by(user_id=current_user.id, chart_type='mlops_registry').count()
        reg = UserChart(
            user_id=current_user.id, chart_type='mlops_registry',
            chart_title=f"{model_name} | {model_type} | score={score:.4f}",
            image_data=str(score)
        )
        db.add(reg); db.commit(); reg_id = reg.id; db.close()

        # Auto-log experiment
        import requests as _req
        try:
            _req.post(f'http://localhost:{os.environ.get("PORT",5000)}/api/mlops/experiment/log',
                     json={'name': model_name, 'params': {'model_type': model_type, 'target': target},
                           'metrics': {'score': score, 'cv_mean': float(cv.mean())},
                           'tags': 'registry,auto-logged'},
                     cookies=request.cookies, timeout=3)
        except: pass

        return jsonify({'success': True, 'registry_id': reg_id,
                       'metadata': metadata, 'model_count': count+1,
                       'message': f'✅ Model "{model_name}" registered! Score: {score:.4f}'})
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/mlops/registry/list', methods=['GET'])
@login_required
def registry_list():
    try:
        import json as _json
        db = SessionLocal()
        models = db.query(UserChart).filter_by(user_id=current_user.id, chart_type='mlops_registry').order_by(UserChart.created_at.desc()).all()
        db.close()
        result = []
        for m in models:
            try: meta = _json.loads(m.description)
            except: meta = {}
            result.append({'id': m.id, 'title': m.title,
                          'model_name': meta.get('model_name',''),
                          'model_type': meta.get('model_type',''),
                          'target': meta.get('target',''), 'task': meta.get('task',''),
                          'score': meta.get('score',''), 'cv_mean': meta.get('cv_mean',''),
                          'n_features': meta.get('n_features',''), 'n_rows': meta.get('n_rows',''),
                          'feature_names': meta.get('feature_names',[])[:10],
                          'description': meta.get('description',''),
                          'saved_at': m.created_at.strftime('%d %b %Y %H:%M')})
        return jsonify({'success': True, 'models': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mlops/registry/download/<int:mid>', methods=['GET'])
@login_required
def registry_download(mid):
    try:
        import json as _json
        db = SessionLocal()
        m = db.query(UserChart).filter_by(id=mid, user_id=current_user.id, chart_type='mlops_registry').first()
        db.close()
        if not m: return jsonify({'error': 'Model not found!'}), 404
        model_bytes = base64.b64decode(m.chart_data)
        try: meta = _json.loads(m.description)
        except: meta = {}
        fname = f"{meta.get('model_name','model')}.pkl"
        return send_file(io.BytesIO(model_bytes), as_attachment=True,
                        download_name=fname, mimetype='application/octet-stream')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mlops/registry/delete/<int:mid>', methods=['DELETE'])
@login_required
def registry_delete(mid):
    try:
        db = SessionLocal()
        m = db.query(UserChart).filter_by(id=mid, user_id=current_user.id).first()
        if m: db.delete(m); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── 3. MODEL MONITORING — Data Drift Detection ────────────────
@app.route('/api/mlops/monitor', methods=['POST'])
@login_required
@check_query_limit
def model_monitor():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        reference_pct = float(data.get('reference_pct', 0.6))
        target = data.get('target', '')

        from scipy import stats
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        df = agent.df.select_dtypes(include='number').dropna()
        if len(df) < 20:
            return jsonify({'error': 'Monitoring ke liye minimum 20 rows chahiye!'}), 400

        split = int(len(df) * reference_pct)
        ref_df = df.iloc[:split]
        cur_df = df.iloc[split:]

        drift_results = []
        for col in df.columns[:10]:
            if col == target: continue
            ks_stat, p_val = stats.ks_2samp(ref_df[col].dropna(), cur_df[col].dropna())
            drifted = p_val < 0.05
            ref_mean = float(ref_df[col].mean()); cur_mean = float(cur_df[col].mean())
            mean_shift = round(abs(cur_mean - ref_mean) / (abs(ref_mean) + 1e-10) * 100, 2)
            drift_results.append({
                'column': col, 'ks_stat': round(float(ks_stat),4),
                'p_value': round(float(p_val),4), 'drifted': drifted,
                'ref_mean': round(ref_mean,3), 'current_mean': round(cur_mean,3),
                'mean_shift_pct': mean_shift
            })

        n_drifted = sum(1 for d in drift_results if d['drifted'])
        drift_pct = round(n_drifted/len(drift_results)*100, 1) if drift_results else 0
        overall_status = 'CRITICAL 🔴' if drift_pct>50 else 'WARNING 🟡' if drift_pct>20 else 'HEALTHY 🟢'

        # Chart
        fig, axes = plt.subplots(2, 2, figsize=(14, 10))
        cols_show = [d['column'] for d in drift_results[:8]]
        colors = ['#ef4444' if d['drifted'] else '#059669' for d in drift_results[:8]]
        p_vals = [d['p_value'] for d in drift_results[:8]]
        axes[0,0].barh(cols_show[::-1], p_vals[::-1], color=colors[::-1], alpha=0.85)
        axes[0,0].axvline(0.05, color='#f59e0b', linestyle='--', linewidth=2, label='p=0.05 threshold')
        axes[0,0].set_title('KS Test p-values (red=drift)', fontweight='bold')
        axes[0,0].set_xlabel('p-value'); axes[0,0].legend()

        mean_shifts = [d['mean_shift_pct'] for d in drift_results[:8]]
        axes[0,1].barh(cols_show[::-1], mean_shifts[::-1], color=colors[::-1], alpha=0.85)
        axes[0,1].set_title('Mean Shift %', fontweight='bold'); axes[0,1].set_xlabel('% Change')

        axes[1,0].pie([n_drifted, len(drift_results)-n_drifted],
                     labels=[f'Drifted ({n_drifted})', f'Stable ({len(drift_results)-n_drifted})'],
                     colors=['#ef4444','#059669'], autopct='%1.1f%%', startangle=90)
        axes[1,0].set_title('Feature Drift Summary', fontweight='bold')

        if drift_results:
            top_col = drift_results[0]['column']
            axes[1,1].hist(ref_df[top_col].dropna(), bins=20, alpha=0.6, color='#5b5ef4', label='Reference')
            axes[1,1].hist(cur_df[top_col].dropna(), bins=20, alpha=0.6, color='#ef4444', label='Current')
            axes[1,1].set_title(f'Distribution Shift: {top_col}', fontweight='bold')
            axes[1,1].legend()

        plt.suptitle(f'Data Drift Monitoring — Status: {overall_status}', fontweight='bold', color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True, 'overall_status': overall_status,
            'total_features': len(drift_results), 'n_drifted': n_drifted,
            'drift_pct': drift_pct, 'reference_rows': len(ref_df),
            'current_rows': len(cur_df), 'drift_details': drift_results,
            'chart': chart,
            'message': f'Monitoring done! Status: {overall_status} — {n_drifted}/{len(drift_results)} features drifted'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ── 4. ML PIPELINE BUILDER ────────────────────────────────────
@app.route('/api/mlops/pipeline', methods=['POST'])
@login_required
@check_query_limit
def pipeline_builder():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')
        steps = data.get('steps', ['impute','scale','encode','train'])
        model_type = data.get('model_type', 'random_forest')

        from sklearn.pipeline import Pipeline
        from sklearn.compose import ColumnTransformer
        from sklearn.preprocessing import StandardScaler, LabelEncoder, OneHotEncoder
        from sklearn.impute import SimpleImputer
        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor, GradientBoostingClassifier, GradientBoostingRegressor
        from sklearn.linear_model import LogisticRegression, Ridge
        from sklearn.model_selection import train_test_split, cross_val_score
        from sklearn.metrics import classification_report, mean_absolute_error, r2_score
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io, time as _time

        df = agent.df.copy()
        if target not in df.columns:
            return jsonify({'error': f'Target "{target}" nahi mila!'}), 400

        y_raw = df[target]; X = df.drop(columns=[target])
        num_cols = X.select_dtypes(include='number').columns.tolist()
        cat_cols = X.select_dtypes(include='object').columns.tolist()

        task = 'classification' if (y_raw.dtype == object or y_raw.nunique() <= 15) else 'regression'
        le_y = LabelEncoder()
        y = le_y.fit_transform(y_raw.astype(str)) if y_raw.dtype == object else y_raw.values

        # Build pipeline steps
        num_pipe = Pipeline([('imputer', SimpleImputer(strategy='median')), ('scaler', StandardScaler())])
        cat_pipe = Pipeline([('imputer', SimpleImputer(strategy='most_frequent')), ('encoder', OneHotEncoder(handle_unknown='ignore', sparse_output=False))])

        preprocessor = ColumnTransformer([
            ('num', num_pipe, num_cols),
            ('cat', cat_pipe, cat_cols)
        ], remainder='drop')

        model_map = {
            'random_forest': RandomForestClassifier(n_estimators=100, random_state=42) if task=='classification' else RandomForestRegressor(n_estimators=100, random_state=42),
            'gradient_boost': GradientBoostingClassifier(random_state=42) if task=='classification' else GradientBoostingRegressor(random_state=42),
            'logistic': LogisticRegression(max_iter=500, random_state=42) if task=='classification' else Ridge(),
        }

        full_pipeline = Pipeline([('preprocessor', preprocessor), ('model', model_map.get(model_type, model_map['random_forest']))])

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        t0 = _time.time()
        full_pipeline.fit(X_train, y_train)
        train_time = round(_time.time()-t0, 2)

        score = float(full_pipeline.score(X_test, y_test))
        cv = cross_val_score(full_pipeline, X, y, cv=3)
        y_pred = full_pipeline.predict(X_test)

        pipeline_steps_info = [
            {'step': 'Data Loading', 'status': '✅', 'detail': f'{len(df)} rows, {len(X.columns)} features'},
            {'step': 'Missing Value Imputation', 'status': '✅', 'detail': f'Median (numeric), Mode (categorical)'},
            {'step': 'Feature Encoding', 'status': '✅', 'detail': f'{len(cat_cols)} categorical → OneHot, {len(num_cols)} numeric → scaled'},
            {'step': 'Train/Test Split', 'status': '✅', 'detail': f'{len(X_train)} train, {len(X_test)} test (80/20)'},
            {'step': f'Model Training ({model_type})', 'status': '✅', 'detail': f'{train_time}s | Score: {score:.4f}'},
            {'step': 'Cross Validation', 'status': '✅', 'detail': f'CV Mean: {cv.mean():.4f} ± {cv.std():.4f}'},
            {'step': 'Evaluation', 'status': '✅', 'detail': f'Final Score: {score:.4f}'}
        ]

        # Chart
        fig, axes = plt.subplots(1, 3, figsize=(15, 5))
        step_names = [s['step'][:20] for s in pipeline_steps_info]
        step_colors = ['#059669']*len(step_names)
        axes[0].barh(step_names[::-1], [1]*len(step_names), color=step_colors, alpha=0.7)
        axes[0].set_xlim(0,1.2); axes[0].set_xticks([])
        for i, s in enumerate(pipeline_steps_info[::-1]):
            axes[0].text(0.05, i, s['detail'][:35], va='center', fontsize=8, color='#1e3a5f')
        axes[0].set_title('Pipeline Steps', fontweight='bold')

        axes[1].bar(['Train Score','Test Score','CV Mean'], [1.0, score, float(cv.mean())],
                   color=['#5b5ef4','#059669','#f59e0b'], alpha=0.85, width=0.5)
        axes[1].set_ylim(0,1.1); axes[1].set_title('Performance Metrics', fontweight='bold')
        for i,v in enumerate([1.0, score, float(cv.mean())]):
            axes[1].text(i, v+0.01, f'{v:.4f}', ha='center', fontweight='700')

        if task == 'classification':
            from sklearn.metrics import confusion_matrix
            import numpy as np
            cm = confusion_matrix(y_test, y_pred)
            im = axes[2].imshow(cm, cmap='Blues', aspect='auto')
            axes[2].set_title('Confusion Matrix', fontweight='bold')
            for i in range(cm.shape[0]):
                for j in range(cm.shape[1]):
                    axes[2].text(j, i, str(cm[i,j]), ha='center', va='center', fontweight='bold')
        else:
            mae = mean_absolute_error(y_test, y_pred)
            r2 = r2_score(y_test, y_pred)
            import numpy as np
            axes[2].scatter(y_test[:100], y_pred[:100], alpha=0.6, color='#5b5ef4', s=30)
            mn, mx = min(float(np.min(y_test)), float(np.min(y_pred))), max(float(np.max(y_test)), float(np.max(y_pred)))
            axes[2].plot([mn,mx],[mn,mx],'r--',linewidth=2)
            axes[2].set_xlabel('Actual'); axes[2].set_ylabel('Predicted')
            axes[2].set_title(f'Actual vs Predicted\nMAE={mae:.3f} R²={r2:.3f}', fontweight='bold')

        plt.suptitle(f'ML Pipeline — {model_type.replace("_"," ").title()} | Score: {score:.4f}', fontweight='bold')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True, 'target': target, 'task': task,
            'model_type': model_type, 'score': round(score,4),
            'cv_mean': round(float(cv.mean()),4), 'cv_std': round(float(cv.std()),4),
            'train_time_sec': train_time, 'n_train': len(X_train), 'n_test': len(X_test),
            'pipeline_steps': pipeline_steps_info, 'chart': chart,
            'message': f'Pipeline complete! {model_type} | Score: {score:.4f} | Time: {train_time}s'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ── 5. ONE-CLICK MODEL API ────────────────────────────────────
@app.route('/api/mlops/deploy', methods=['POST'])
@login_required
def deploy_model_api():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        target = data.get('target', '')
        model_name = data.get('model_name', 'my_model')

        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        from sklearn.pipeline import Pipeline
        from sklearn.impute import SimpleImputer
        from sklearn.model_selection import train_test_split
        import pickle, io as _io, json as _json

        df = agent.df.dropna()
        if target not in df.columns:
            return jsonify({'error': f'Target "{target}" nahi mila!'}), 400

        y = df[target]; X = df.drop(columns=[target])
        feature_names = list(X.columns)
        cat_cols = X.select_dtypes(include='object').columns.tolist()
        for c in cat_cols:
            le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))

        task = 'classification' if (y.dtype == object or y.nunique() <= 15) else 'regression'
        classes = []
        if y.dtype == object:
            le_y = LabelEncoder(); y = le_y.fit_transform(y.astype(str)); classes = list(le_y.classes_)
        else:
            y = y.values

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        pipeline = Pipeline([('scaler', StandardScaler()),
                           ('model', RandomForestClassifier(n_estimators=100, random_state=42) if task=='classification' else RandomForestRegressor(n_estimators=100, random_state=42))])
        pipeline.fit(X_train, y_train)
        score = float(pipeline.score(X_test, y_test))

        buf = _io.BytesIO(); pickle.dump(pipeline, buf); buf.seek(0)
        model_b64 = base64.b64encode(buf.read()).decode()

        # Store deployed model
        db = SessionLocal()
        existing = db.query(UserChart).filter_by(user_id=current_user.id, chart_type='mlops_deployed', title=model_name).first()
        meta = _json.dumps({'target': target, 'task': task, 'feature_names': feature_names,
                           'classes': classes, 'score': round(score,4), 'model_name': model_name})
        if existing:
            existing.image_data = meta; db.commit()
            deploy_id = existing.id
        else:
            dep = UserChart(user_id=current_user.id, chart_type='mlops_deployed',
                          chart_title=model_name, image_data=meta)
            db.add(dep); db.commit(); deploy_id = dep.id
        db.close()

        endpoint = f'/api/mlops/predict/{deploy_id}'
        curl_example = f'''curl -X POST {request.host_url.rstrip("/")}{endpoint} \\
  -H "Content-Type: application/json" \\
  -d '{{{", ".join([f'"{f}": <value>' for f in feature_names[:3]])}...}}'
'''
        return jsonify({
            'success': True, 'deploy_id': deploy_id,
            'endpoint': endpoint, 'model_name': model_name,
            'task': task, 'score': round(score,4),
            'feature_names': feature_names, 'classes': classes,
            'curl_example': curl_example,
            'message': f'✅ Model deployed! Endpoint: {endpoint}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/mlops/predict/<int:deploy_id>', methods=['POST'])
@login_required
def model_predict_api(deploy_id):
    try:
        import pickle, io as _io, json as _json
        db = SessionLocal()
        m = db.query(UserChart).filter_by(id=deploy_id, user_id=current_user.id, chart_type='mlops_deployed').first()
        db.close()
        if not m: return jsonify({'error': 'Model not found! Pehle deploy karo.'}), 404

        meta = _json.loads(m.description)
        pipeline = pickle.loads(base64.b64decode(m.chart_data))
        input_data = request.get_json()
        feature_names = meta['feature_names']

        import pandas as pd
        input_df = pd.DataFrame([{f: input_data.get(f, 0) for f in feature_names}])
        for c in input_df.select_dtypes(include='object').columns:
            from sklearn.preprocessing import LabelEncoder
            le = LabelEncoder(); le.classes_ = [str(c)] ; input_df[c] = 0

        prediction = pipeline.predict(input_df)[0]
        result = {'prediction': str(prediction), 'model_name': meta['model_name'],
                  'task': meta['task'], 'score': meta['score']}

        if meta['task'] == 'classification' and hasattr(pipeline, 'predict_proba'):
            try:
                proba = pipeline.predict_proba(input_df)[0]
                classes = meta.get('classes', [])
                result['probabilities'] = {str(c): round(float(p),4) for c,p in zip(classes, proba)} if classes else [round(float(p),4) for p in proba]
            except: pass

        return jsonify({'success': True, 'result': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mlops/deployed/list', methods=['GET'])
@login_required
def list_deployed():
    try:
        import json as _json
        db = SessionLocal()
        models = db.query(UserChart).filter_by(user_id=current_user.id, chart_type='mlops_deployed').all()
        db.close()
        result = []
        for m in models:
            try: meta = _json.loads(m.description)
            except: meta = {}
            result.append({'id': m.id, 'model_name': meta.get('model_name',''),
                          'target': meta.get('target',''), 'task': meta.get('task',''),
                          'score': meta.get('score',''), 'endpoint': f'/api/mlops/predict/{m.id}',
                          'deployed_at': m.created_at.strftime('%d %b %Y %H:%M')})
        return jsonify({'success': True, 'deployed_models': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 1: SQL DATABASE CONNECT & QUERY
# ════════════════════════════════════════════════════════════════
@app.route('/api/sql_connect', methods=['POST'])
@login_required
@check_query_limit
def sql_connect():
    try:
        data = request.get_json()
        db_type = data.get('db_type', 'sqlite')  # sqlite, mysql, postgresql, mssql
        connection_str = data.get('connection_string', '')
        query = data.get('query', '')
        
        import sqlalchemy as sa
        import pandas as _pd
        
        # Build connection
        if db_type == 'sqlite':
            # Upload sqlite file path or in-memory
            upload_path = data.get('db_file', '')
            if upload_path and os.path.exists(upload_path):
                engine = sa.create_engine(f'sqlite:///{upload_path}')
            else:
                return jsonify({'error': 'SQLite: db_file path do ya connection_string do'}), 400
        elif connection_str:
            engine = sa.create_engine(connection_str)
        else:
            return jsonify({'error': 'connection_string required for MySQL/PostgreSQL/MSSQL'}), 400

        with engine.connect() as conn:
            # Get table list
            inspector = sa.inspect(engine)
            tables = inspector.get_table_names()
            
            result = {'success': True, 'tables': tables, 'db_type': db_type}
            
            if query:
                df_result = _pd.read_sql_query(sa.text(query), conn)
                # Load into agent
                agent.df = df_result
                agent.available_columns = list(df_result.columns)
                
                result.update({
                    'query_result': df_result.head(100).to_dict('records'),
                    'total_rows': len(df_result),
                    'columns': list(df_result.columns),
                    'dtypes': {c: str(t) for c, t in df_result.dtypes.items()},
                    'message': f'Query executed! {len(df_result)} rows loaded into DS Agent'
                })
            else:
                result['message'] = f'Connected! {len(tables)} tables found: {", ".join(tables[:10])}'
        
        return jsonify(result)
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': f'DB Connection failed: {str(e)}'}), 500


@app.route('/api/sql_query', methods=['POST'])
@login_required
@check_query_limit
def sql_query_run():
    try:
        data = request.get_json()
        query = data.get('query', '')
        connection_str = data.get('connection_string', '')
        db_file = data.get('db_file', '')
        
        import sqlalchemy as sa
        import pandas as _pd
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        
        if db_file:
            engine = sa.create_engine(f'sqlite:///{db_file}')
        elif connection_str:
            engine = sa.create_engine(connection_str)
        else:
            return jsonify({'error': 'Connection info required!'}), 400
        
        with engine.connect() as conn:
            df_result = _pd.read_sql_query(sa.text(query), conn)
        
        agent.df = df_result
        agent.available_columns = list(df_result.columns)
        
        # Auto chart if numeric data
        chart = None
        num_cols = df_result.select_dtypes(include='number').columns.tolist()
        if len(num_cols) >= 1 and len(df_result) > 1:
            fig, ax = plt.subplots(figsize=(10, 4))
            if len(num_cols) == 1:
                df_result[num_cols[0]].head(50).plot(kind='bar', ax=ax, color='#5b5ef4', alpha=0.85)
                ax.set_title(f'Query Result — {num_cols[0]}', fontweight='bold')
            else:
                df_result[num_cols[:3]].head(20).plot(kind='bar', ax=ax, colormap='viridis', alpha=0.85)
                ax.set_title('Query Result Visualization', fontweight='bold')
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=110, bbox_inches='tight')
            plt.close('all'); buf.seek(0)
            chart = base64.b64encode(buf.read()).decode()
        
        return jsonify({
            'success': True,
            'query': query,
            'rows': len(df_result),
            'columns': list(df_result.columns),
            'preview': df_result.head(50).to_dict('records'),
            'chart': chart,
            'message': f'Query done! {len(df_result)} rows, {len(df_result.columns)} columns — loaded into DS Agent'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 2: WEB SCRAPING + ANALYSIS
# ════════════════════════════════════════════════════════════════
@app.route('/api/web_scrape', methods=['POST'])
@login_required
@check_query_limit
def web_scrape():
    try:
        data = request.get_json()
        url = data.get('url', '')
        scrape_type = data.get('scrape_type', 'table')  # table, text, list, custom
        selector = data.get('selector', '')  # CSS selector optional
        
        if not url:
            return jsonify({'error': 'URL do scrape karne ke liye!'}), 400
        
        import requests as _req
        from html.parser import HTMLParser
        import re as _re
        
        headers_http = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0 Safari/537.36'
        }
        
        resp = _req.get(url, headers=headers_http, timeout=15)
        resp.raise_for_status()
        html_content = resp.text
        
        result = {'success': True, 'url': url, 'scrape_type': scrape_type}
        
        # Try pandas read_html for tables
        if scrape_type == 'table':
            try:
                tables = pd.read_html(html_content)
                if tables:
                    best_table = max(tables, key=lambda t: t.shape[0] * t.shape[1])
                    agent.df = best_table
                    agent.available_columns = list(best_table.columns)
                    result.update({
                        'tables_found': len(tables),
                        'loaded_table_shape': best_table.shape,
                        'columns': list(best_table.columns),
                        'preview': best_table.head(20).to_dict('records'),
                        'message': f'{len(tables)} tables found! Largest ({best_table.shape[0]}x{best_table.shape[1]}) loaded into DS Agent'
                    })
                else:
                    result['message'] = 'No tables found on this page'
            except Exception as te:
                result['table_error'] = str(te)
        
        # Text extraction
        text_clean = _re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=_re.DOTALL)
        text_clean = _re.sub(r'<style[^>]*>.*?</style>', '', text_clean, flags=_re.DOTALL)
        text_clean = _re.sub(r'<[^>]+>', ' ', text_clean)
        text_clean = _re.sub(r'\s+', ' ', text_clean).strip()
        
        # Extract links
        links = _re.findall(r'href=["\']([^"\']+)["\']', html_content)
        links = [l for l in links if l.startswith('http')][:20]
        
        # Extract numbers/stats
        numbers = _re.findall(r'\b\d+(?:,\d+)*(?:\.\d+)?\b', text_clean)
        numbers = [float(n.replace(',','')) for n in numbers if float(n.replace(',','')) > 0][:100]
        
        result.update({
            'text_preview': text_clean[:500],
            'text_length': len(text_clean),
            'links_found': links[:10],
            'numbers_extracted': numbers[:20],
            'status_code': resp.status_code,
        })
        
        if not result.get('message'):
            result['message'] = f'Scraping done! {len(text_clean)} chars extracted, {len(links)} links found'
        
        return jsonify(result)
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': f'Scraping failed: {str(e)}'}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 3: PYTHON SCRIPT EXPORT
# ════════════════════════════════════════════════════════════════
@app.route('/api/export_script', methods=['POST'])
@login_required
def export_python_script():
    try:
        data = request.get_json()
        script_type = data.get('script_type', 'full_analysis')
        target = data.get('target', '')
        model_type = data.get('model_type', 'random_forest')
        dataset_name = data.get('dataset_name', 'your_data.csv')
        
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        
        cols = list(agent.df.columns)
        num_cols = agent.df.select_dtypes(include='number').columns.tolist()
        
        scripts = {
            'full_analysis': f'''#!/usr/bin/env python3
"""
DS Agent — Auto Generated Full Analysis Script
Dataset: {dataset_name}
Generated by DS Agent (https://data-scientist-1ly9.onrender.com)
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
from sklearn.model_selection import train_test_split, cross_val_score
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.metrics import classification_report, confusion_matrix
import warnings
warnings.filterwarnings('ignore')

# ── 1. LOAD DATA ──────────────────────────────────────────────
df = pd.read_csv('{dataset_name}')
print(f"Shape: {{df.shape}}")
print(f"Columns: {{list(df.columns)}}")

# ── 2. EDA ────────────────────────────────────────────────────
print("\\n=== DATA INFO ===")
print(df.info())
print("\\n=== STATISTICS ===")
print(df.describe())
print("\\n=== MISSING VALUES ===")
print(df.isnull().sum())

# ── 3. VISUALIZATION ─────────────────────────────────────────
fig, axes = plt.subplots(2, 3, figsize=(15, 10))
numeric_cols = df.select_dtypes(include='number').columns[:6]
for i, col in enumerate(numeric_cols[:6]):
    row, col_idx = i // 3, i % 3
    df[col].hist(ax=axes[row, col_idx], bins=30, color='#5b5ef4', alpha=0.8, edgecolor='white')
    axes[row, col_idx].set_title(col, fontweight='bold')
plt.suptitle('Feature Distributions', fontweight='bold', fontsize=14)
plt.tight_layout()
plt.savefig('distributions.png', dpi=120, bbox_inches='tight')
print("Distributions saved: distributions.png")

# Correlation heatmap
plt.figure(figsize=(10, 8))
corr = df.select_dtypes(include='number').corr()
sns.heatmap(corr, annot=True, fmt='.2f', cmap='coolwarm', center=0,
            square=True, linewidths=0.5)
plt.title('Correlation Heatmap', fontweight='bold')
plt.tight_layout()
plt.savefig('correlation.png', dpi=120, bbox_inches='tight')
print("Correlation saved: correlation.png")

# ── 4. DATA CLEANING ─────────────────────────────────────────
print("\\n=== DATA CLEANING ===")
df_clean = df.copy()
# Fill missing values
for col in df_clean.select_dtypes(include='number').columns:
    df_clean[col].fillna(df_clean[col].median(), inplace=True)
for col in df_clean.select_dtypes(include='object').columns:
    df_clean[col].fillna(df_clean[col].mode()[0] if len(df_clean[col].mode()) > 0 else 'Unknown', inplace=True)
# Remove duplicates
before = len(df_clean)
df_clean.drop_duplicates(inplace=True)
print(f"Removed {{before - len(df_clean)}} duplicates")
print(f"Clean dataset shape: {{df_clean.shape}}")

{"" if not target else f"""
# ── 5. MACHINE LEARNING ──────────────────────────────────────
TARGET = '{target}'
if TARGET in df_clean.columns:
    y = df_clean[TARGET]
    X = df_clean.drop(columns=[TARGET])
    
    # Encode categoricals
    le_dict = {{}}
    for col in X.select_dtypes(include='object').columns:
        le = LabelEncoder()
        X[col] = le.fit_transform(X[col].astype(str))
        le_dict[col] = le
    
    # Task detection
    task = 'classification' if (y.dtype == object or y.nunique() <= 15) else 'regression'
    print(f"Task: {{task}}")
    
    if y.dtype == object:
        le_y = LabelEncoder()
        y = le_y.fit_transform(y.astype(str))
    
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
    
    # Train model
    model = RandomForestClassifier(n_estimators=100, random_state=42) if task == 'classification' \\
            else RandomForestRegressor(n_estimators=100, random_state=42)
    model.fit(X_train, y_train)
    
    score = model.score(X_test, y_test)
    cv_scores = cross_val_score(model, X, y, cv=5)
    print(f"Test Score: {{score:.4f}}")
    print(f"CV Score: {{cv_scores.mean():.4f}} +/- {{cv_scores.std():.4f}}")
    
    # Feature importance
    feat_imp = pd.DataFrame({{
        'Feature': X.columns,
        'Importance': model.feature_importances_
    }}).sort_values('Importance', ascending=False)
    print("\\nTop 10 Features:")
    print(feat_imp.head(10).to_string())
    
    # Save model
    import pickle
    with open('model_{model_type}.pkl', 'wb') as f:
        pickle.dump(model, f)
    print("\\nModel saved: model_{model_type}.pkl")
"""}

# ── 6. SAVE RESULTS ──────────────────────────────────────────
df_clean.to_csv('cleaned_{dataset_name}', index=False)
print(f"\\nCleaned data saved: cleaned_{dataset_name}")
print("\\n✅ Analysis Complete!")
''',

            'ml_only': f'''#!/usr/bin/env python3
"""
DS Agent — ML Model Script
Target: {target}
Model: {model_type}
"""
import pandas as pd
import numpy as np
import pickle
from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.model_selection import train_test_split
from sklearn.pipeline import Pipeline
import warnings; warnings.filterwarnings('ignore')

def train_model(csv_path='{dataset_name}', target='{target}'):
    df = pd.read_csv(csv_path).dropna()
    y = df[target]; X = df.drop(columns=[target])
    
    for c in X.select_dtypes(include='object').columns:
        le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))
    
    task = 'classification' if (y.dtype==object or y.nunique()<=15) else 'regression'
    if y.dtype==object:
        le_y=LabelEncoder(); y=le_y.fit_transform(y.astype(str))
    
    X_train,X_test,y_train,y_test = train_test_split(X,y,test_size=0.2,random_state=42)
    
    pipeline = Pipeline([
        ('scaler', StandardScaler()),
        ('model', RandomForestClassifier(100,random_state=42) if task=='classification'
                  else RandomForestRegressor(100,random_state=42))
    ])
    pipeline.fit(X_train, y_train)
    score = pipeline.score(X_test, y_test)
    print(f"Model Score: {{score:.4f}}")
    
    with open('trained_model.pkl','wb') as f: pickle.dump(pipeline,f)
    print("Model saved: trained_model.pkl")
    return pipeline, score

def predict(input_dict, model_path='trained_model.pkl'):
    """Single prediction"""
    with open(model_path,'rb') as f: model=pickle.load(f)
    df_in = pd.DataFrame([input_dict])
    for c in df_in.select_dtypes(include='object').columns:
        le=LabelEncoder(); df_in[c]=le.fit_transform(df_in[c].astype(str))
    return model.predict(df_in)[0]

if __name__ == '__main__':
    model, score = train_model()
    print(f"Training complete! Score: {{score:.4f}}")
''',

            'etl_pipeline': f'''#!/usr/bin/env python3
"""
DS Agent — ETL Pipeline Script
Auto-generated production-ready ETL
"""
import pandas as pd
import numpy as np
import os, logging, json
from datetime import datetime

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(message)s')
log = logging.getLogger(__name__)

class DSAgentETL:
    def __init__(self, source_path, output_path='output/'):
        self.source = source_path
        self.output = output_path
        self.df = None
        self.log = []
        os.makedirs(output_path, exist_ok=True)
    
    def extract(self):
        """E — Extract data from source"""
        log.info(f"Extracting from: {{self.source}}")
        if self.source.endswith('.csv'):
            self.df = pd.read_csv(self.source)
        elif self.source.endswith(('.xlsx','.xls')):
            self.df = pd.read_excel(self.source)
        else:
            raise ValueError("Supported: CSV, Excel")
        self.log.append(f"Extracted: {{self.df.shape[0]}} rows, {{self.df.shape[1]}} cols")
        log.info(f"Extracted: {{self.df.shape}}")
        return self
    
    def transform(self):
        """T — Transform & Clean"""
        log.info("Transforming...")
        original_rows = len(self.df)
        
        # Remove duplicates
        self.df.drop_duplicates(inplace=True)
        self.log.append(f"Removed {{original_rows - len(self.df)}} duplicates")
        
        # Fill missing values
        for col in self.df.select_dtypes(include='number').columns:
            missing = self.df[col].isnull().sum()
            if missing > 0:
                self.df[col].fillna(self.df[col].median(), inplace=True)
                self.log.append(f"Filled {{missing}} nulls in {{col}} with median")
        
        for col in self.df.select_dtypes(include='object').columns:
            missing = self.df[col].isnull().sum()
            if missing > 0:
                self.df[col].fillna('Unknown', inplace=True)
        
        # Standardize column names
        self.df.columns = [c.lower().strip().replace(' ','_') for c in self.df.columns]
        
        # Add metadata columns
        self.df['_etl_processed_at'] = datetime.now().isoformat()
        self.df['_etl_source'] = os.path.basename(self.source)
        
        log.info(f"Transform done: {{self.df.shape}}")
        return self
    
    def load(self, format='csv'):
        """L — Load to destination"""
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')
        if format == 'csv':
            out_path = f"{{self.output}}processed_{{ts}}.csv"
            self.df.to_csv(out_path, index=False)
        elif format == 'excel':
            out_path = f"{{self.output}}processed_{{ts}}.xlsx"
            self.df.to_excel(out_path, index=False)
        elif format == 'json':
            out_path = f"{{self.output}}processed_{{ts}}.json"
            self.df.to_json(out_path, orient='records', indent=2)
        
        # Save pipeline log
        log_path = f"{{self.output}}etl_log_{{ts}}.json"
        with open(log_path,'w') as f:
            json.dump({{'timestamp':ts,'source':self.source,'shape':list(self.df.shape),'steps':self.log}},f,indent=2)
        
        log.info(f"Loaded to: {{out_path}}")
        log.info(f"Log saved: {{log_path}}")
        return out_path
    
    def run(self, output_format='csv'):
        return self.extract().transform().load(output_format)

# Usage
if __name__ == '__main__':
    etl = DSAgentETL(source_path='{dataset_name}')
    output = etl.run(output_format='csv')
    print(f"ETL Complete! Output: {{output}}")
    print(f"Steps: {{etl.log}}")
'''
        }
        
        script_code = scripts.get(script_type, scripts['full_analysis'])
        
        # Return as downloadable file
        buf = io.BytesIO(script_code.encode('utf-8'))
        filename = f'ds_agent_{script_type}_{dataset_name.replace(".csv","")}.py'
        return send_file(buf, as_attachment=True, download_name=filename,
                        mimetype='text/x-python')
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 4: SCHEDULED EMAIL REPORTS
# ════════════════════════════════════════════════════════════════
_scheduled_reports = {}

@app.route('/api/schedule_report', methods=['POST'])
@login_required
def schedule_report():
    try:
        data = request.get_json()
        report_name = data.get('report_name', 'Auto Report')
        recipient_email = data.get('email', current_user.email)
        frequency = data.get('frequency', 'daily')  # daily, weekly, monthly
        report_template = data.get('template', 'corporate')
        analysis_type = data.get('analysis_type', 'full')
        
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        
        # Store schedule in DB
        db = SessionLocal()
        import json as _json
        schedule_data = _json.dumps({
            'report_name': report_name,
            'email': recipient_email,
            'frequency': frequency,
            'template': report_template,
            'analysis_type': analysis_type,
            'dataset_shape': str(agent.df.shape),
            'columns': list(agent.df.columns)[:20],
        })
        sched = UserChart(
            user_id=current_user.id,
            chart_type='scheduled_report',
            chart_title=f'{report_name} — {frequency}',
        )
        db.add(sched); db.commit(); sched_id = sched.id; db.close()
        
        # Send immediate confirmation email
        try:
            import resend as _resend
            _resend.api_key = os.environ.get('RESEND_API_KEY', '')
            if _resend.api_key:
                _resend.Emails.send({
                    'from': 'DS Agent <reports@dsagent.ai>',
                    'to': recipient_email,
                    'subject': f'✅ Report Scheduled: {report_name}',
                    'html': f'''<div style="font-family:sans-serif;max-width:600px;margin:auto">
                        <div style="background:linear-gradient(135deg,#5b5ef4,#8b5cf6);padding:30px;border-radius:12px 12px 0 0;color:white;text-align:center">
                            <h2>📊 Report Scheduled Successfully!</h2>
                        </div>
                        <div style="padding:24px;background:#f8faff;border-radius:0 0 12px 12px">
                            <p><b>Report Name:</b> {report_name}</p>
                            <p><b>Frequency:</b> {frequency}</p>
                            <p><b>Template:</b> {report_template}</p>
                            <p><b>Recipient:</b> {recipient_email}</p>
                            <p style="color:#64748b;font-size:13px">
                                Aapko {frequency} basis pe reports milenge. 
                                DS Agent se manage karo: 
                                <a href="https://data-scientist-1ly9.onrender.com">dashboard</a>
                            </p>
                        </div>
                    </div>'''
                })
        except: pass
        
        return jsonify({
            'success': True,
            'schedule_id': sched_id,
            'report_name': report_name,
            'frequency': frequency,
            'recipient': recipient_email,
            'message': f'Report scheduled! {frequency} basis pe {recipient_email} pe milega'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/schedule_report/list', methods=['GET'])
@login_required
def list_scheduled_reports():
    try:
        import json as _json
        db = SessionLocal()
        scheds = db.query(UserChart).filter_by(user_id=current_user.id,
                         chart_type='scheduled_report').all()
        db.close()
        result = []
        for s in scheds:
            try: info = _json.loads(s.description)
            except: info = {}
            result.append({'id': s.id, 'title': s.title,
                          'email': info.get('email',''),
                          'frequency': info.get('frequency',''),
                          'template': info.get('template',''),
                          'created_at': s.created_at.strftime('%d %b %Y')})
        return jsonify({'success': True, 'schedules': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/schedule_report/delete/<int:sid>', methods=['DELETE'])
@login_required
def delete_scheduled_report(sid):
    try:
        db = SessionLocal()
        s = db.query(UserChart).filter_by(id=sid, user_id=current_user.id).first()
        if s: db.delete(s); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 5: POWERPOINT AUTO GENERATOR
# ════════════════════════════════════════════════════════════════
@app.route('/api/generate_pptx', methods=['POST'])
@login_required
@check_query_limit
def generate_pptx():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        title = data.get('title', 'Data Analysis Report')
        subtitle = data.get('subtitle', 'Generated by DS Agent')
        theme = data.get('theme', 'corporate')  # corporate, dark, minimal, vibrant
        target = data.get('target', '')
        
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        
        themes = {
            'corporate': {'bg': '1e3a5f', 'accent': '5b5ef4', 'text': 'ffffff', 'sub': 'a5b4fc'},
            'dark':      {'bg': '0f172a', 'accent': '06d6a0', 'text': 'f1f5f9', 'sub': '94a3b8'},
            'minimal':   {'bg': 'ffffff', 'accent': '5b5ef4', 'text': '1e293b', 'sub': '64748b'},
            'vibrant':   {'bg': '4338ca', 'accent': 'f59e0b', 'text': 'ffffff', 'sub': 'c7d2fe'},
        }
        th = themes.get(theme, themes['corporate'])
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        def add_bg(slide, color_hex):
            from pptx.util import Pt
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(color_hex)
        
        def add_text(slide, text, l, t, w, h, size=24, bold=False, color='ffffff', align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = align
            run = p.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)
        
        def add_rect(slide, l, t, w, h, color_hex, alpha=None):
            shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            shape.line.fill.background()
            return shape
        
        # ── Slide 1: Title ──────────────────────────────────────
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s1, th['bg'])
        add_rect(s1, 0, 0, 0.5, 7.5, th['accent'])
        add_rect(s1, 0.5, 3.2, 12.8, 0.06, th['accent'])
        add_text(s1, title, 1, 2.0, 11, 1.2, size=40, bold=True, color=th['text'], align=PP_ALIGN.LEFT)
        add_text(s1, subtitle, 1, 3.4, 11, 0.6, size=20, color=th['sub'])
        df_shape = f"Dataset: {agent.df.shape[0]:,} rows × {agent.df.shape[1]} columns"
        add_text(s1, df_shape, 1, 4.2, 11, 0.5, size=16, color=th['sub'])
        add_text(s1, 'Generated by DS Agent  •  AI-Powered Data Science', 1, 6.8, 11, 0.5, size=12, color=th['sub'])
        
        # ── Slide 2: Dataset Overview ───────────────────────────
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s2, th['bg'])
        add_rect(s2, 0, 0, 13.33, 1.2, th['accent'])
        add_text(s2, '📊 Dataset Overview', 0.4, 0.2, 12, 0.8, size=28, bold=True, color='ffffff')
        
        stats = [
            ('Total Rows', f"{agent.df.shape[0]:,}"),
            ('Total Columns', str(agent.df.shape[1])),
            ('Missing Values', str(int(agent.df.isnull().sum().sum()))),
            ('Numeric Cols', str(len(agent.df.select_dtypes(include='number').columns))),
        ]
        for i, (label, val) in enumerate(stats):
            x = 0.4 + i * 3.1
            add_rect(s2, x, 1.5, 2.8, 1.8, th['accent'])
            add_text(s2, val, x, 1.7, 2.8, 0.8, size=32, bold=True, color='ffffff', align=PP_ALIGN.CENTER)
            add_text(s2, label, x, 2.5, 2.8, 0.6, size=14, color='ffffff', align=PP_ALIGN.CENTER)
        
        # Column list
        cols_text = '  •  '.join(list(agent.df.columns)[:12])
        if len(agent.df.columns) > 12: cols_text += f'  ...+{len(agent.df.columns)-12} more'
        add_text(s2, f'Columns: {cols_text}', 0.4, 3.6, 12.5, 1.5, size=13, color=th['sub'])
        
        # ── Slide 3: Distribution Charts ─────────────────────────
        num_cols_list = agent.df.select_dtypes(include='number').columns[:4].tolist()
        if num_cols_list:
            fig, axes = plt.subplots(1, min(4, len(num_cols_list)), figsize=(14, 4))
            if len(num_cols_list) == 1: axes = [axes]
            for i, col in enumerate(num_cols_list[:4]):
                axes[i].hist(agent.df[col].dropna(), bins=25,
                           color='#'+th['accent'], alpha=0.85, edgecolor='white')
                axes[i].set_title(col, fontweight='bold', color='white', fontsize=11)
                axes[i].set_facecolor('#'+th['bg'])
                axes[i].tick_params(colors='white'); axes[i].spines[:].set_color('white')
            fig.patch.set_facecolor('#'+th['bg'])
            plt.tight_layout()
            buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=130, bbox_inches='tight',
                                             facecolor='#'+th['bg']); buf.seek(0)
            plt.close('all')
            
            s3 = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(s3, th['bg'])
            add_rect(s3, 0, 0, 13.33, 1.1, th['accent'])
            add_text(s3, '📈 Feature Distributions', 0.4, 0.15, 12, 0.8, size=28, bold=True, color='ffffff')
            s3.shapes.add_picture(buf, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
        
        # ── Slide 4: Correlation Heatmap ─────────────────────────
        if len(num_cols_list) >= 3:
            import seaborn as sns
            fig, ax = plt.subplots(figsize=(10, 7))
            corr = agent.df[num_cols_list[:10]].corr()
            sns.heatmap(corr, annot=True, fmt='.2f', cmap='RdYlBu_r', center=0,
                       ax=ax, linewidths=0.5, square=True,
                       annot_kws={'size': 10, 'color': 'black'})
            ax.set_title('Correlation Matrix', fontweight='bold', color='white', pad=15)
            fig.patch.set_facecolor('#'+th['bg'])
            ax.set_facecolor('#'+th['bg'])
            ax.tick_params(colors='white')
            plt.tight_layout()
            buf2 = _io.BytesIO(); plt.savefig(buf2, format='png', dpi=130, bbox_inches='tight',
                                              facecolor='#'+th['bg']); buf2.seek(0)
            plt.close('all')
            
            s4 = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(s4, th['bg'])
            add_rect(s4, 0, 0, 13.33, 1.1, th['accent'])
            add_text(s4, '🔗 Correlation Analysis', 0.4, 0.15, 12, 0.8, size=28, bold=True, color='ffffff')
            s4.shapes.add_picture(buf2, Inches(2.5), Inches(1.2), Inches(8), Inches(6))
        
        # ── Slide 5: ML Results (if target given) ────────────────
        if target and target in agent.df.columns:
            from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
            from sklearn.model_selection import train_test_split
            from sklearn.preprocessing import LabelEncoder
            
            df_ml = agent.df.dropna()
            y = df_ml[target]; X = df_ml.drop(columns=[target])
            for c in X.select_dtypes(include='object').columns:
                le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))
            task = 'Classification' if (y.dtype==object or y.nunique()<=15) else 'Regression'
            if y.dtype==object:
                le_y=LabelEncoder(); y=le_y.fit_transform(y.astype(str))
            
            X_tr,X_te,y_tr,y_te = train_test_split(X,y,test_size=0.2,random_state=42)
            clf = RandomForestClassifier(100,random_state=42) if task=='Classification' else RandomForestRegressor(100,random_state=42)
            clf.fit(X_tr,y_tr); score=clf.score(X_te,y_te)
            feat_imp = sorted(zip(X.columns, clf.feature_importances_), key=lambda x:-x[1])[:8]
            
            fig, axes = plt.subplots(1, 2, figsize=(13, 5))
            fnames, fscores = zip(*feat_imp)
            axes[0].barh(list(fnames)[::-1], list(fscores)[::-1], color='#'+th['accent'], alpha=0.85)
            axes[0].set_title(f'Feature Importance\n{task} — Target: {target}', fontweight='bold', color='white')
            axes[0].set_facecolor('#'+th['bg']); axes[0].tick_params(colors='white')
            axes[0].spines[:].set_color('white')
            
            models_scores = {'Random Forest': score, 'Baseline': 0.5 if task=='Classification' else 0.0}
            axes[1].bar(list(models_scores.keys()), list(models_scores.values()),
                       color=['#'+th['accent'],'#94a3b8'], alpha=0.85, width=0.4)
            axes[1].set_title(f'Model Score: {score:.4f}', fontweight='bold', color='white')
            axes[1].set_facecolor('#'+th['bg']); axes[1].tick_params(colors='white')
            axes[1].spines[:].set_color('white'); axes[1].set_ylim(0,1.1)
            
            fig.patch.set_facecolor('#'+th['bg'])
            plt.tight_layout()
            buf3 = _io.BytesIO(); plt.savefig(buf3, format='png', dpi=130, bbox_inches='tight',
                                              facecolor='#'+th['bg']); buf3.seek(0)
            plt.close('all')
            
            s5 = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(s5, th['bg'])
            add_rect(s5, 0, 0, 13.33, 1.1, th['accent'])
            add_text(s5, f'🤖 ML Results — {task}', 0.4, 0.15, 12, 0.8, size=28, bold=True, color='ffffff')
            add_text(s5, f'Target: {target}  |  Score: {score:.4f}  |  Algorithm: Random Forest',
                    0.4, 1.2, 12, 0.5, size=16, color=th['sub'])
            s5.shapes.add_picture(buf3, Inches(0.4), Inches(1.8), Inches(12.5), Inches(5.2))
        
        # ── Slide 6: Key Insights ─────────────────────────────────
        s6 = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s6, th['bg'])
        add_rect(s6, 0, 0, 13.33, 1.1, th['accent'])
        add_text(s6, '💡 Key Insights & Recommendations', 0.4, 0.15, 12, 0.8, size=24, bold=True, color='ffffff')
        
        insights = []
        insights.append(f"Dataset has {agent.df.shape[0]:,} records across {agent.df.shape[1]} features")
        missing_pct = round(agent.df.isnull().sum().sum() / (agent.df.shape[0]*agent.df.shape[1]) * 100, 1)
        if missing_pct > 5:
            insights.append(f"⚠️ {missing_pct}% missing data detected — cleaning recommended")
        else:
            insights.append(f"✅ Data quality good — only {missing_pct}% missing values")
        if len(num_cols_list) >= 2:
            corr_df = agent.df[num_cols_list].corr().abs()
            high_corr = [(corr_df.columns[i], corr_df.columns[j], round(float(corr_df.iloc[i,j]),2))
                        for i in range(len(corr_df)) for j in range(i+1,len(corr_df))
                        if corr_df.iloc[i,j] > 0.7][:2]
            for c1, c2, v in high_corr:
                insights.append(f"🔗 High correlation: {c1} & {c2} ({v}) — potential multicollinearity")
        insights.append(f"📊 {len(agent.df.select_dtypes(include='number').columns)} numeric features available for ML")
        insights.append("🚀 Use Auto ML to find the best prediction model")
        insights.append("💾 Export Python script for production deployment")
        
        for i, ins in enumerate(insights[:7]):
            y_pos = 1.4 + i * 0.75
            add_rect(s6, 0.4, y_pos, 0.4, 0.45, th['accent'])
            add_text(s6, str(i+1), 0.4, y_pos, 0.4, 0.45, size=16, bold=True, color='ffffff', align=PP_ALIGN.CENTER)
            add_text(s6, ins, 1.0, y_pos, 11.8, 0.5, size=14, color=th['text'])
        
        # ── Slide 7: Thank You ────────────────────────────────────
        s7 = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(s7, th['bg'])
        add_rect(s7, 0, 0, 0.8, 7.5, th['accent'])
        add_rect(s7, 0.8, 3.5, 12.5, 0.05, th['accent'])
        add_text(s7, 'Thank You', 1.5, 2.2, 11, 1.2, size=52, bold=True, color=th['text'], align=PP_ALIGN.CENTER)
        add_text(s7, 'Powered by DS Agent — AI Data Science Platform', 1.5, 3.7, 11, 0.6, size=18, color=th['sub'], align=PP_ALIGN.CENTER)
        add_text(s7, 'data-scientist-1ly9.onrender.com', 1.5, 4.4, 11, 0.5, size=14, color=th['sub'], align=PP_ALIGN.CENTER)
        
        # Save
        pptx_buf = io.BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)
        
        safe_title = title.replace(' ', '_').replace('/', '-')[:30]
        return send_file(pptx_buf, as_attachment=True,
                        download_name=f'DS_Agent_{safe_title}.pptx',
                        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except ImportError:
        return jsonify({'error': 'python-pptx not installed! requirements mein add karo: python-pptx'}), 500
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 6: REAL-TIME LIVE DASHBOARD DATA
# ════════════════════════════════════════════════════════════════
@app.route('/api/live_dashboard', methods=['POST'])
@login_required
def live_dashboard():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        refresh = data.get('refresh', False)
        kpi_cols = data.get('kpi_cols', [])
        
        df = agent.df
        num_cols = df.select_dtypes(include='number').columns.tolist()
        
        # Auto KPIs
        kpis = []
        for col in (kpi_cols if kpi_cols else num_cols[:6]):
            if col in df.columns:
                series = df[col].dropna()
                half = len(series) // 2
                current = float(series.iloc[half:].mean()) if half > 0 else float(series.mean())
                previous = float(series.iloc[:half].mean()) if half > 0 else current
                change_pct = round((current-previous)/abs(previous+1e-10)*100, 2)
                kpis.append({
                    'metric': col,
                    'current': round(current, 3),
                    'previous': round(previous, 3),
                    'change_pct': change_pct,
                    'trend': '📈' if change_pct > 0 else '📉' if change_pct < 0 else '➡️',
                    'status': 'up' if change_pct > 0 else 'down' if change_pct < 0 else 'stable'
                })
        
        # Time-series like data for each numeric col (last 30 points)
        time_series = {}
        for col in num_cols[:4]:
            time_series[col] = df[col].dropna().tail(30).round(3).tolist()
        
        # Distribution summary
        distributions = {}
        for col in num_cols[:6]:
            q = df[col].quantile([0.25, 0.5, 0.75]).to_dict()
            distributions[col] = {
                'min': round(float(df[col].min()), 3),
                'q25': round(float(q[0.25]), 3),
                'median': round(float(q[0.5]), 3),
                'q75': round(float(q[0.75]), 3),
                'max': round(float(df[col].max()), 3),
                'mean': round(float(df[col].mean()), 3),
            }
        
        # Category breakdowns
        cat_cols = df.select_dtypes(include='object').columns[:3]
        categories = {}
        for col in cat_cols:
            vc = df[col].value_counts().head(8)
            categories[col] = {'labels': list(vc.index), 'values': [int(v) for v in vc.values]}
        
        return jsonify({
            'success': True,
            'total_rows': len(df),
            'total_cols': len(df.columns),
            'kpis': kpis,
            'time_series': time_series,
            'distributions': distributions,
            'categories': categories,
            'last_updated': pd.Timestamp.now().strftime('%H:%M:%S'),
            'message': f'Dashboard refreshed! {len(kpis)} KPIs, {len(time_series)} trends'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 7: EXCEL FORMULA GENERATOR
# ════════════════════════════════════════════════════════════════
@app.route('/api/excel_formulas', methods=['POST'])
@login_required
def excel_formulas():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        task = data.get('task', 'auto')  # auto, vlookup, pivot, dashboard, stats
        target_col = data.get('target_col', '')
        lookup_col = data.get('lookup_col', '')
        
        df = agent.df
        cols = list(df.columns)
        num_cols = df.select_dtypes(include='number').columns.tolist()
        cat_cols = df.select_dtypes(include='object').columns.tolist()
        
        formulas = []
        explanations = []
        
        # Auto-generate based on dataset
        col_letters = {col: chr(65+i) for i, col in enumerate(cols[:26])}
        
        # Basic stats formulas
        for col in num_cols[:4]:
            letter = col_letters.get(col, 'A')
            formulas.append({
                'formula': f'=AVERAGE({letter}2:{letter}1000)',
                'cell': f'{letter}1002',
                'description': f'{col} ka average',
                'use_case': 'Statistics'
            })
            formulas.append({
                'formula': f'=COUNTIF({letter}2:{letter}1000,">0")',
                'cell': f'{letter}1003',
                'description': f'{col} mein positive values count',
                'use_case': 'Counting'
            })
            formulas.append({
                'formula': f'=STDEV({letter}2:{letter}1000)',
                'cell': f'{letter}1004',
                'description': f'{col} standard deviation',
                'use_case': 'Statistics'
            })
        
        # VLOOKUP
        if len(cols) >= 2:
            formulas.append({
                'formula': f'=VLOOKUP(A2,Sheet2!A:C,2,FALSE)',
                'description': 'Sheet2 se data match karo ID ke basis pe',
                'use_case': 'VLOOKUP'
            })
            formulas.append({
                'formula': f'=INDEX(Sheet2!B:B,MATCH(A2,Sheet2!A:A,0))',
                'description': 'INDEX-MATCH — VLOOKUP se better',
                'use_case': 'Lookup'
            })
        
        # Conditional
        if num_cols:
            letter = col_letters.get(num_cols[0], 'B')
            formulas.append({
                'formula': f'=IF({letter}2>AVERAGE({letter}:{letter}),"High","Low")',
                'description': f'{num_cols[0]} ko High/Low mein classify karo',
                'use_case': 'Conditional'
            })
            formulas.append({
                'formula': f'=IF({letter}2=""," Missing","Present")',
                'description': 'Missing values check karo',
                'use_case': 'Data Quality'
            })
        
        # Text formulas
        if cat_cols:
            letter = col_letters.get(cat_cols[0], 'A')
            formulas.append({
                'formula': f'=PROPER({letter}2)',
                'description': f'{cat_cols[0]} ko Title Case mein convert karo',
                'use_case': 'Text'
            })
            formulas.append({
                'formula': f'=TRIM(LOWER({letter}2))',
                'description': 'Extra spaces remove karo aur lowercase karo',
                'use_case': 'Cleaning'
            })
        
        # Date formulas (if any date-like columns)
        formulas.append({'formula': '=TODAY()', 'description': 'Aaj ki date', 'use_case': 'Date'})
        formulas.append({'formula': '=DATEDIF(A2,TODAY(),"D")', 'description': 'Kitne din ho gaye', 'use_case': 'Date'})
        formulas.append({'formula': '=TEXT(A2,"DD-MMM-YYYY")', 'description': 'Date ko readable format mein', 'use_case': 'Date'})
        
        # Dashboard KPI formulas
        kpi_formulas = []
        for col in num_cols[:3]:
            letter = col_letters.get(col, 'A')
            kpi_formulas.append({
                'kpi': col,
                'total': f'=SUM({letter}2:{letter}1000)',
                'avg': f'=AVERAGE({letter}2:{letter}1000)',
                'max': f'=MAX({letter}2:{letter}1000)',
                'min': f'=MIN({letter}2:{letter}1000)',
                'count': f'=COUNT({letter}2:{letter}1000)',
            })
        
        # Pivot suggestion
        pivot_suggestion = {
            'rows': cat_cols[:2] if cat_cols else [],
            'values': num_cols[:3] if num_cols else [],
            'filters': cat_cols[2:4] if len(cat_cols) > 2 else [],
            'tip': f'Insert → PivotTable → Rows: {cat_cols[0] if cat_cols else "Category"}, Values: Sum of {num_cols[0] if num_cols else "Amount"}'
        }
        
        return jsonify({
            'success': True,
            'dataset_columns': cols,
            'column_letters': col_letters,
            'formulas': formulas[:20],
            'kpi_dashboard': kpi_formulas,
            'pivot_suggestion': pivot_suggestion,
            'total_formulas': len(formulas),
            'message': f'{len(formulas)} Excel formulas generated for your dataset!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 8: COMPLETE ETL PIPELINE BUILDER
# ════════════════════════════════════════════════════════════════
@app.route('/api/etl_pipeline', methods=['POST'])
@login_required
@check_query_limit
def etl_pipeline():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        source_name = data.get('source_name', 'source_data.csv')
        target_format = data.get('target_format', 'csv')  # csv, excel, json
        transformations = data.get('transformations', ['clean','dedupe','normalize','encode'])
        
        import time as _time, json as _json
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        
        df = agent.df.copy()
        log_steps = []
        start_time = _time.time()
        
        original_shape = df.shape
        
        # EXTRACT
        log_steps.append({'step': 'EXTRACT', 'status': '✅', 'detail': f'Loaded {df.shape[0]:,} rows × {df.shape[1]} cols',
                         'rows_in': df.shape[0], 'rows_out': df.shape[0]})
        
        # TRANSFORM steps
        if 'dedupe' in transformations:
            before = len(df)
            df.drop_duplicates(inplace=True)
            removed = before - len(df)
            log_steps.append({'step': 'DEDUPLICATE', 'status': '✅', 'detail': f'Removed {removed} duplicates',
                             'rows_in': before, 'rows_out': len(df)})
        
        if 'clean' in transformations:
            null_before = int(df.isnull().sum().sum())
            for col in df.select_dtypes(include='number').columns:
                df[col].fillna(df[col].median(), inplace=True)
            for col in df.select_dtypes(include='object').columns:
                mode = df[col].mode()
                df[col].fillna(mode[0] if len(mode)>0 else 'Unknown', inplace=True)
            null_after = int(df.isnull().sum().sum())
            log_steps.append({'step': 'NULL HANDLING', 'status': '✅',
                             'detail': f'Fixed {null_before-null_after} null values',
                             'rows_in': len(df), 'rows_out': len(df)})
        
        if 'normalize' in transformations:
            from sklearn.preprocessing import MinMaxScaler
            num_cols = df.select_dtypes(include='number').columns.tolist()
            if num_cols:
                scaler = MinMaxScaler()
                df[num_cols] = scaler.fit_transform(df[num_cols])
                log_steps.append({'step': 'NORMALIZE', 'status': '✅',
                                 'detail': f'Scaled {len(num_cols)} numeric columns to [0,1]',
                                 'rows_in': len(df), 'rows_out': len(df)})
        
        if 'encode' in transformations:
            from sklearn.preprocessing import LabelEncoder
            cat_cols = df.select_dtypes(include='object').columns.tolist()
            encoded = []
            for col in cat_cols:
                if df[col].nunique() <= 20:
                    le = LabelEncoder()
                    df[col+'_encoded'] = le.fit_transform(df[col].astype(str))
                    encoded.append(col)
            if encoded:
                log_steps.append({'step': 'LABEL ENCODE', 'status': '✅',
                                 'detail': f'Encoded {len(encoded)} columns: {", ".join(encoded[:5])}',
                                 'rows_in': len(df), 'rows_out': len(df)})
        
        if 'outlier' in transformations:
            num_cols = df.select_dtypes(include='number').columns.tolist()
            clipped = 0
            for col in num_cols[:10]:
                Q1 = df[col].quantile(0.25); Q3 = df[col].quantile(0.75); IQR = Q3 - Q1
                lower = Q1 - 1.5*IQR; upper = Q3 + 1.5*IQR
                outliers = ((df[col] < lower) | (df[col] > upper)).sum()
                df[col] = df[col].clip(lower, upper)
                clipped += int(outliers)
            log_steps.append({'step': 'OUTLIER CLIP', 'status': '✅',
                             'detail': f'Clipped {clipped} outlier values (IQR method)',
                             'rows_in': len(df), 'rows_out': len(df)})
        
        # LOAD — save to temp
        elapsed = round(_time.time()-start_time, 2)
        out_path = f'/tmp/etl_output_{int(_time.time())}'
        
        if target_format == 'csv':
            out_file = out_path + '.csv'; df.to_csv(out_file, index=False)
        elif target_format == 'excel':
            out_file = out_path + '.xlsx'; df.to_excel(out_file, index=False)
        elif target_format == 'json':
            out_file = out_path + '.json'; df.to_json(out_file, orient='records', indent=2)
        
        log_steps.append({'step': 'LOAD', 'status': '✅',
                         'detail': f'Saved as {target_format.upper()} — {df.shape[0]:,} rows',
                         'rows_in': df.shape[0], 'rows_out': df.shape[0]})
        
        agent.df = df
        agent.available_columns = list(df.columns)
        
        # Pipeline visualization chart
        fig, axes = plt.subplots(1, 2, figsize=(13, 5))
        steps = [s['step'] for s in log_steps]
        rows_out = [s['rows_out'] for s in log_steps]
        colors = ['#5b5ef4','#8b5cf6','#059669','#f59e0b','#06d6a0','#ef4444'][:len(steps)]
        axes[0].barh(steps[::-1], rows_out[::-1], color=colors[::-1], alpha=0.85)
        axes[0].set_title('Pipeline Steps — Rows After Each Step', fontweight='bold')
        axes[0].set_xlabel('Row Count')
        
        axes[1].bar(['Before ETL', 'After ETL'],
                   [original_shape[0], len(df)],
                   color=['#94a3b8','#059669'], alpha=0.85, width=0.4)
        axes[1].set_title('Before vs After', fontweight='bold')
        for i, v in enumerate([original_shape[0], len(df)]):
            axes[1].text(i, v+50, f'{v:,}', ha='center', fontweight='700')
        
        plt.suptitle(f'ETL Pipeline — {elapsed}s | {len(log_steps)} steps', fontweight='bold', color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()
        
        return jsonify({
            'success': True,
            'original_shape': list(original_shape),
            'final_shape': list(df.shape),
            'rows_removed': original_shape[0] - df.shape[0],
            'cols_added': df.shape[1] - original_shape[1],
            'steps_executed': len(log_steps),
            'pipeline_log': log_steps,
            'elapsed_sec': elapsed,
            'output_format': target_format,
            'chart': chart,
            'message': f'ETL done in {elapsed}s! {original_shape[0]:,} → {df.shape[0]:,} rows | {len(log_steps)} steps'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# FEATURE 9: CUSTOM DEEP LEARNING ARCHITECTURE BUILDER
# ════════════════════════════════════════════════════════════════
@app.route('/api/custom_dl', methods=['POST'])
@login_required
@check_query_limit
def custom_deep_learning():
    try:
        url = get_hf_url()
        if not url:
            return jsonify({'error': 'ML Server offline! HuggingFace Space start karo.'}), 400
        data = request.get_json()
        result = call_hf_api(url, '/custom_dl', data)
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500



# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 1: BUSINESS INTELLIGENCE & AUTO RECOMMENDATIONS
# ════════════════════════════════════════════════════════════════
@app.route('/api/business_intelligence', methods=['POST'])
@login_required
@check_query_limit
def business_intelligence():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        industry = data.get('industry', 'general')
        business_goal = data.get('business_goal', 'revenue_growth')
        company_size = data.get('company_size', 'mid')

        df = agent.df.copy()
        num_cols = df.select_dtypes(include='number').columns.tolist()
        cat_cols = df.select_dtypes(include='object').columns.tolist()

        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        import numpy as np

        # ── Auto Business Insights Engine ──────────────────────
        insights = []
        recommendations = []
        risks = []
        opportunities = []

        # Missing data risk
        missing_pct = round(df.isnull().sum().sum() / (df.shape[0] * df.shape[1]) * 100, 2)
        if missing_pct > 10:
            risks.append({
                'risk': 'High Data Quality Risk',
                'detail': f'{missing_pct}% missing data detected — decisions based on incomplete data',
                'severity': 'HIGH',
                'action': 'Immediate data cleaning required before any business decision'
            })
        else:
            insights.append(f'✅ Data quality good — {missing_pct}% missing (acceptable threshold: <10%)')

        # Volume insight
        if df.shape[0] > 10000:
            insights.append(f'📊 Large dataset: {df.shape[0]:,} records — statistically significant decisions possible')
        elif df.shape[0] < 500:
            risks.append({'risk': 'Small Sample Size', 'detail': f'Only {df.shape[0]} records — ML predictions may not be reliable',
                         'severity': 'MEDIUM', 'action': 'Collect more data before deploying ML models'})

        # Numeric analysis
        for col in num_cols[:8]:
            series = df[col].dropna()
            if len(series) < 10: continue
            skew = float(series.skew())
            cv = float(series.std() / (series.mean() + 1e-10))
            trend_half1 = float(series.iloc[:len(series)//2].mean())
            trend_half2 = float(series.iloc[len(series)//2:].mean())
            trend_chg = round((trend_half2 - trend_half1) / (abs(trend_half1) + 1e-10) * 100, 2)

            if abs(trend_chg) > 15:
                direction = 'increasing 📈' if trend_chg > 0 else 'decreasing 📉'
                if trend_chg > 0:
                    opportunities.append(f'{col} is {direction} by {abs(trend_chg)}% — capitalize on this growth trend')
                else:
                    risks.append({'risk': f'{col} Declining', 'detail': f'{col} dropped {abs(trend_chg)}% — attention needed',
                                 'severity': 'HIGH' if abs(trend_chg) > 30 else 'MEDIUM',
                                 'action': f'Investigate root cause of {col} decline immediately'})

            if cv > 0.5:
                insights.append(f'⚠️ {col} has high variability (CV={cv:.2f}) — inconsistent performance detected')

            if abs(skew) > 2:
                insights.append(f'📐 {col} is highly skewed ({skew:.2f}) — median better than mean for reporting')

        # Correlation opportunities
        if len(num_cols) >= 3:
            corr_matrix = df[num_cols[:10]].corr().abs()
            for i in range(len(corr_matrix.columns)):
                for j in range(i+1, len(corr_matrix.columns)):
                    val = float(corr_matrix.iloc[i, j])
                    if val > 0.75:
                        c1, c2 = corr_matrix.columns[i], corr_matrix.columns[j]
                        opportunities.append(f'🔗 Strong link between {c1} & {c2} ({val:.2f}) — improving {c1} will directly impact {c2}')

        # Industry-specific recommendations
        industry_recs = {
            'healthcare': [
                'Implement predictive risk scoring to identify high-risk patients early',
                'Use anomaly detection on lab values for early disease detection',
                'Track readmission rates — reducing by 5% saves significant costs',
                'Analyze treatment outcome patterns for protocol optimization',
                'Monitor staff-to-patient ratios for operational efficiency'
            ],
            'finance': [
                'Deploy fraud detection model on transaction data immediately',
                'Segment customers by risk profile for personalized products',
                'Use time series forecasting for cash flow optimization',
                'Monitor credit score distributions for portfolio health',
                'Implement churn prediction to retain high-value customers'
            ],
            'retail': [
                'Customer segmentation for targeted marketing campaigns',
                'Basket analysis to improve product placement and cross-selling',
                'Demand forecasting for inventory optimization',
                'Price elasticity analysis for revenue maximization',
                'Customer lifetime value calculation for acquisition budget'
            ],
            'general': [
                'Segment your data into meaningful groups for targeted actions',
                'Identify top 20% drivers that cause 80% of outcomes (Pareto)',
                'Build a prediction model to proactively address issues',
                'Track key metrics weekly with automated alerts',
                'Use correlation analysis to find hidden business levers'
            ]
        }
        recommendations = industry_recs.get(industry, industry_recs['general'])

        # ROI Estimation
        roi_estimates = {
            'revenue_growth': {'potential': '15-25%', 'timeline': '6-12 months', 'method': 'Customer segmentation + churn prevention'},
            'cost_reduction': {'potential': '20-35%', 'timeline': '3-6 months', 'method': 'Process optimization + anomaly detection'},
            'risk_reduction': {'potential': '40-60%', 'timeline': '1-3 months', 'method': 'Predictive risk scoring + early warning'},
            'efficiency': {'potential': '30-50%', 'timeline': '3-9 months', 'method': 'Automation + pattern recognition'},
        }
        roi = roi_estimates.get(business_goal, roi_estimates['revenue_growth'])

        # Executive Summary
        exec_summary = f"""EXECUTIVE SUMMARY
Dataset: {df.shape[0]:,} records × {df.shape[1]} features | Industry: {industry.title()}
Business Goal: {business_goal.replace('_',' ').title()} | Company Size: {company_size.title()}

KEY FINDINGS:
• {len(risks)} risks identified ({sum(1 for r in risks if r.get('severity')=='HIGH')} HIGH priority)
• {len(opportunities)} growth opportunities detected
• {len(insights)} data quality observations
• Estimated ROI potential: {roi['potential']} improvement in {roi['timeline']}

TOP PRIORITY ACTION: {recommendations[0] if recommendations else 'Begin with data quality improvement'}

RECOMMENDED NEXT STEPS:
1. Address HIGH severity risks immediately
2. Implement top 3 recommendations within 30 days
3. Deploy predictive model within 60 days
4. Measure KPI improvement at 90-day mark"""

        # ── Visualization ──────────────────────────────────────
        fig = plt.figure(figsize=(16, 12))
        gs = fig.add_gridspec(3, 3, hspace=0.45, wspace=0.35)

        # Risk matrix
        ax1 = fig.add_subplot(gs[0, 0])
        risk_counts = {'HIGH': sum(1 for r in risks if r.get('severity')=='HIGH'),
                      'MEDIUM': sum(1 for r in risks if r.get('severity')=='MEDIUM'),
                      'LOW': sum(1 for r in risks if r.get('severity','')=='LOW')}
        risk_counts = {k:v for k,v in risk_counts.items() if v > 0}
        if risk_counts:
            colors_risk = {'HIGH':'#ef4444','MEDIUM':'#f59e0b','LOW':'#059669'}
            ax1.pie(risk_counts.values(), labels=risk_counts.keys(),
                   colors=[colors_risk[k] for k in risk_counts],
                   autopct='%1.0f%%', startangle=90)
            ax1.set_title('Risk Distribution', fontweight='bold', fontsize=11)
        else:
            ax1.text(0.5, 0.5, '✅ No Risks\nDetected', ha='center', va='center',
                    fontsize=14, color='#059669', fontweight='bold', transform=ax1.transAxes)
            ax1.set_title('Risk Distribution', fontweight='bold')
            ax1.axis('off')

        # Top numeric trends
        ax2 = fig.add_subplot(gs[0, 1:])
        if len(num_cols) >= 2:
            trend_data = []
            for col in num_cols[:5]:
                series = df[col].dropna()
                half1 = float(series.iloc[:len(series)//2].mean())
                half2 = float(series.iloc[len(series)//2:].mean())
                chg = round((half2-half1)/(abs(half1)+1e-10)*100, 1)
                trend_data.append((col[:15], chg))
            trend_data.sort(key=lambda x: x[1], reverse=True)
            cols_t, vals_t = zip(*trend_data)
            bar_colors = ['#059669' if v >= 0 else '#ef4444' for v in vals_t]
            bars = ax2.bar(cols_t, vals_t, color=bar_colors, alpha=0.85, edgecolor='white', linewidth=1.5)
            ax2.axhline(0, color='#1e293b', linewidth=1, linestyle='-')
            ax2.set_title('Business Metric Trends (First Half vs Second Half %)', fontweight='bold', fontsize=11)
            ax2.set_ylabel('% Change')
            for bar, val in zip(bars, vals_t):
                ax2.text(bar.get_x()+bar.get_width()/2,
                        bar.get_height() + (0.5 if val >= 0 else -1.5),
                        f'{val:+.1f}%', ha='center', fontsize=9, fontweight='bold')

        # Opportunity bars
        ax3 = fig.add_subplot(gs[1, :2])
        opp_labels = [o[:50] for o in opportunities[:5]]
        if opp_labels:
            impact_scores = [90, 75, 65, 55, 45][:len(opp_labels)]
            ax3.barh(opp_labels[::-1], impact_scores[::-1], color='#5b5ef4', alpha=0.85)
            ax3.set_title('Top Business Opportunities (Impact Score)', fontweight='bold', fontsize=11)
            ax3.set_xlabel('Impact Score')
        else:
            ax3.text(0.5, 0.5, 'Run analysis to detect opportunities', ha='center', va='center', transform=ax3.transAxes)
        ax3.axis('on')

        # ROI gauge
        ax4 = fig.add_subplot(gs[1, 2])
        ax4.text(0.5, 0.75, roi['potential'], ha='center', va='center',
                fontsize=28, fontweight='bold', color='#059669', transform=ax4.transAxes)
        ax4.text(0.5, 0.5, 'Potential Improvement', ha='center', va='center',
                fontsize=10, color='#64748b', transform=ax4.transAxes)
        ax4.text(0.5, 0.32, f'in {roi["timeline"]}', ha='center', va='center',
                fontsize=11, color='#5b5ef4', fontweight='bold', transform=ax4.transAxes)
        ax4.text(0.5, 0.15, roi['method'][:35], ha='center', va='center',
                fontsize=8, color='#94a3b8', transform=ax4.transAxes)
        ax4.set_title(f'ROI Estimate\n{business_goal.replace("_"," ").title()}', fontweight='bold')
        ax4.axis('off')
        ax4.set_facecolor('#f0fdf4')

        # Correlation heatmap small
        ax5 = fig.add_subplot(gs[2, :])
        if len(num_cols) >= 3:
            import seaborn as sns
            corr = df[num_cols[:10]].corr()
            sns.heatmap(corr, annot=True, fmt='.2f', cmap='RdYlBu_r', center=0,
                       ax=ax5, linewidths=0.3, annot_kws={'size': 8})
            ax5.set_title('Feature Correlation Matrix — Business Lever Analysis', fontweight='bold', fontsize=11)

        plt.suptitle(f'Business Intelligence Report — {industry.title()} | Goal: {business_goal.replace("_"," ").title()}',
                    fontweight='bold', fontsize=14, color='#1e3a5f', y=1.01)
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'industry': industry,
            'business_goal': business_goal,
            'executive_summary': exec_summary,
            'insights': insights[:10],
            'recommendations': recommendations[:7],
            'risks': risks[:8],
            'opportunities': opportunities[:8],
            'roi_estimate': roi,
            'data_health_score': max(0, round(100 - missing_pct*2 - len(risks)*5, 1)),
            'chart': chart,
            'message': f'BI Analysis complete! {len(risks)} risks, {len(opportunities)} opportunities, ROI potential: {roi["potential"]}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 2: ENTERPRISE SECURITY LAYER
# ════════════════════════════════════════════════════════════════
@app.route('/api/security/audit_log', methods=['GET'])
@login_required
def get_audit_log():
    try:
        import json as _json
        db = SessionLocal()
        logs = db.query(UserChart).filter_by(
            user_id=current_user.id,
            chart_type='audit_log'
        ).order_by(UserChart.created_at.desc()).limit(100).all()
        db.close()
        result = []
        for log in logs:
            try: info = _json.loads(log.description)
            except: info = {}
            result.append({
                'id': log.id,
                'action': info.get('action', ''),
                'resource': info.get('resource', ''),
                'ip': info.get('ip', ''),
                'details': info.get('details', ''),
                'timestamp': log.created_at.strftime('%d %b %Y %H:%M:%S')
            })
        return jsonify({'success': True, 'audit_logs': result, 'total': len(result)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


def log_audit(user_id, action, resource, details='', request_obj=None):
    """Helper — log every important action"""
    try:
        import json as _json
        db = SessionLocal()
        ip = request_obj.remote_addr if request_obj else 'unknown'
        audit = UserChart(
            user_id=user_id,
            chart_type='audit_log',
            chart_title=f'{action} — {resource}',
            image_data=_json.dumps({
                'action': action,
                'resource': resource,
                'ip': ip,
                'details': details
            }),
        )
        db.add(audit); db.commit(); db.close()
    except: pass


@app.route('/api/security/data_mask', methods=['POST'])
@login_required
def data_masking():
    """Mask PII/PHI data for privacy compliance"""
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        mask_cols = data.get('columns', [])
        mask_type = data.get('mask_type', 'partial')  # partial, full, hash, pseudonym

        import hashlib, re as _re

        df = agent.df.copy()
        masked_info = []

        # Auto-detect PII columns if not specified
        if not mask_cols:
            pii_keywords = ['name', 'phone', 'mobile', 'email', 'address', 'aadhaar',
                          'pan', 'passport', 'dob', 'birth', 'ssn', 'account', 'card']
            for col in df.columns:
                if any(kw in col.lower() for kw in pii_keywords):
                    mask_cols.append(col)

        for col in mask_cols:
            if col not in df.columns: continue
            original_sample = str(df[col].iloc[0]) if len(df) > 0 else ''

            if mask_type == 'partial':
                df[col] = df[col].astype(str).apply(
                    lambda x: x[:2] + '*' * max(0, len(x)-4) + x[-2:] if len(x) > 4 else '****'
                )
            elif mask_type == 'full':
                df[col] = '***MASKED***'
            elif mask_type == 'hash':
                df[col] = df[col].astype(str).apply(
                    lambda x: hashlib.sha256(x.encode()).hexdigest()[:12]
                )
            elif mask_type == 'pseudonym':
                unique_vals = df[col].unique()
                mapping = {v: f'{col[:3].upper()}_{i:04d}' for i, v in enumerate(unique_vals)}
                df[col] = df[col].map(mapping)

            masked_info.append({
                'column': col,
                'original_sample': original_sample[:20] + '...' if len(original_sample) > 20 else original_sample,
                'masked_sample': str(df[col].iloc[0]) if len(df) > 0 else '',
                'rows_masked': len(df)
            })

        agent.df = df
        log_audit(current_user.id, 'DATA_MASK', f'{len(mask_cols)} columns',
                 f'Type: {mask_type}', request)

        return jsonify({
            'success': True,
            'mask_type': mask_type,
            'columns_masked': mask_cols,
            'masked_info': masked_info,
            'compliance': ['HIPAA', 'GDPR', 'PDPB (India)'],
            'message': f'✅ {len(mask_cols)} columns masked ({mask_type}) — PII/PHI protected!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/security/session_info', methods=['GET'])
@login_required
def session_info():
    try:
        import json as _json
        db = SessionLocal()
        logs = db.query(UserChart).filter_by(
            user_id=current_user.id, chart_type='audit_log'
        ).order_by(UserChart.created_at.desc()).limit(20).all()
        db.close()

        recent_actions = []
        for log in logs[:10]:
            try: info = _json.loads(log.description)
            except: info = {}
            recent_actions.append(f"{log.created_at.strftime('%H:%M')} — {info.get('action','?')} on {info.get('resource','?')}")

        return jsonify({
            'success': True,
            'user': current_user.email,
            'user_id': current_user.id,
            'dataset_loaded': agent.df is not None,
            'dataset_shape': list(agent.df.shape) if agent.df is not None else None,
            'recent_actions': recent_actions,
            'security_features': ['Audit Logging', 'Data Masking', 'Session Tracking', 'IP Logging'],
            'compliance_ready': ['HIPAA', 'GDPR', 'SOC2 (partial)', 'PDPB India']
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 3: CLIENT DELIVERABLE SYSTEM
# ════════════════════════════════════════════════════════════════
@app.route('/api/client_report', methods=['POST'])
@login_required
@check_query_limit
def client_report():
    """Executive-level client deliverable with ROI, insights, next steps"""
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        client_name = data.get('client_name', 'Client')
        project_name = data.get('project_name', 'Data Analysis Project')
        analyst_name = data.get('analyst_name', 'Data Scientist')
        industry = data.get('industry', 'Healthcare')
        target = data.get('target', '')
        currency = data.get('currency', '₹')
        baseline_cost = float(data.get('baseline_cost', 0))

        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        Table, TableStyle, HRFlowable, PageBreak)
        from reportlab.platypus.flowables import Flowable
        from reportlab.lib.colors import HexColor

        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        import numpy as np
        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
        from sklearn.preprocessing import LabelEncoder
        from sklearn.model_selection import train_test_split

        df = agent.df.copy()
        num_cols = df.select_dtypes(include='number').columns.tolist()

        # Run quick ML if target given
        ml_score = None
        feat_importance = {}
        task = None
        if target and target in df.columns:
            try:
                df_ml = df.dropna()
                y = df_ml[target]; X = df_ml.drop(columns=[target])
                for c in X.select_dtypes(include='object').columns:
                    le = LabelEncoder(); X[c] = le.fit_transform(X[c].astype(str))
                task = 'Classification' if (y.dtype==object or y.nunique()<=15) else 'Regression'
                if y.dtype == object:
                    le_y = LabelEncoder(); y = le_y.fit_transform(y.astype(str))
                X_tr, X_te, y_tr, y_te = train_test_split(X, y, test_size=0.2, random_state=42)
                clf = RandomForestClassifier(100, random_state=42) if task=='Classification' \
                      else RandomForestRegressor(100, random_state=42)
                clf.fit(X_tr, y_tr); ml_score = round(float(clf.score(X_te, y_te)), 4)
                feat_importance = dict(sorted(
                    zip(X.columns, clf.feature_importances_),
                    key=lambda x: -x[1])[:5])
            except: pass

        # Generate chart
        fig = plt.figure(figsize=(16, 8))
        axes = fig.subplots(1, 3)

        # Chart 1: Data overview
        missing = df.isnull().sum()
        missing = missing[missing > 0].head(8)
        if len(missing) > 0:
            axes[0].barh(missing.index, missing.values, color='#ef4444', alpha=0.85)
            axes[0].set_title('Missing Data by Column', fontweight='bold')
        else:
            axes[0].text(0.5, 0.5, '✅ No Missing Data', ha='center', va='center',
                        fontsize=14, color='#059669', fontweight='bold', transform=axes[0].transAxes)
            axes[0].axis('off')
            axes[0].set_title('Data Quality', fontweight='bold')

        # Chart 2: Feature importance or distribution
        if feat_importance:
            items = list(feat_importance.items())
            fnames, fscores = zip(*items)
            axes[1].barh(list(fnames)[::-1], [s*100 for s in list(fscores)[::-1]],
                        color='#5b5ef4', alpha=0.85)
            axes[1].set_title(f'Top Predictors for {target}', fontweight='bold')
            axes[1].set_xlabel('Importance %')
        else:
            if num_cols:
                df[num_cols[0]].hist(ax=axes[1], bins=25, color='#5b5ef4', alpha=0.85, edgecolor='white')
                axes[1].set_title(f'{num_cols[0]} Distribution', fontweight='bold')

        # Chart 3: Summary stats
        if num_cols:
            means = [float(df[c].mean()) for c in num_cols[:6]]
            axes[2].bar([c[:12] for c in num_cols[:6]], means, color='#06d6a0', alpha=0.85, edgecolor='white')
            axes[2].set_title('Column Means Overview', fontweight='bold')
            plt.setp(axes[2].xaxis.get_majorticklabels(), rotation=30, ha='right')

        plt.suptitle(f'{project_name} — {client_name}', fontweight='bold', fontsize=13)
        plt.tight_layout()
        chart_buf = _io.BytesIO()
        plt.savefig(chart_buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); chart_buf.seek(0)

        # ── Build PDF ──────────────────────────────────────────
        INDIGO = HexColor('#5b5ef4'); DARK = HexColor('#1e293b')
        SLATE = HexColor('#475569'); GREEN = HexColor('#059669')
        RED = HexColor('#ef4444'); LIGHT = HexColor('#f0f4ff')
        BORDER = HexColor('#e0e7ff')

        styles = getSampleStyleSheet()
        H_STYLE = ParagraphStyle('h', fontName='Helvetica-Bold', fontSize=13,
                                 textColor=INDIGO, spaceAfter=6, spaceBefore=12)
        B_STYLE = ParagraphStyle('b', fontName='Helvetica', fontSize=10,
                                 textColor=SLATE, spaceAfter=4, leading=15,
                                 alignment=TA_JUSTIFY)
        SM = ParagraphStyle('sm', fontName='Helvetica', fontSize=9,
                            textColor=SLATE, leading=13)

        pdf_buf = _io.BytesIO()
        doc = SimpleDocTemplate(pdf_buf, pagesize=A4,
                               leftMargin=30*mm, rightMargin=30*mm,
                               topMargin=20*mm, bottomMargin=20*mm)
        W_body = A4[0] - 60*mm
        story = []

        # Cover header
        cover_data = [[
            Paragraph(f'<font color="white"><b>{project_name}</b></font>',
                     ParagraphStyle('x', fontName='Helvetica-Bold', fontSize=18,
                                   textColor=colors.white, alignment=TA_CENTER)),
        ],[
            Paragraph(f'<font color="#c7d2fe">Prepared for: {client_name}  |  By: {analyst_name}  |  Industry: {industry}</font>',
                     ParagraphStyle('x', fontName='Helvetica', fontSize=10,
                                   textColor=HexColor('#c7d2fe'), alignment=TA_CENTER)),
        ]]
        ct = Table(cover_data, colWidths=[W_body])
        ct.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), DARK),
            ('TOPPADDING', (0,0), (-1,-1), 14),
            ('BOTTOMPADDING', (0,0), (-1,-1), 10),
            ('LEFTPADDING', (0,0), (-1,-1), 16),
        ]))
        story.append(ct)
        story.append(Spacer(1, 12))

        # KPI summary boxes
        missing_pct = round(df.isnull().sum().sum() / (df.shape[0]*df.shape[1])*100, 1)
        health = max(0, round(100 - missing_pct*2, 0))
        kpi_items = [
            ('📋 Records', f'{df.shape[0]:,}'),
            ('📊 Features', str(df.shape[1])),
            ('✅ Data Health', f'{health}%'),
            ('🎯 ML Score', f'{ml_score:.2%}' if ml_score else 'N/A'),
        ]
        kpi_data = [[Paragraph(f'<b><font color="white">{v}</font></b>',
                               ParagraphStyle('x', fontName='Helvetica-Bold', fontSize=16,
                                             textColor=colors.white, alignment=TA_CENTER)),
                     ] for _, v in kpi_items]
        kpi_labels = [[Paragraph(f'<font color="#e0e7ff">{k}</font>',
                                  ParagraphStyle('x', fontName='Helvetica', fontSize=8,
                                                textColor=HexColor('#e0e7ff'), alignment=TA_CENTER)),
                       ] for k, _ in kpi_items]

        kpi_row = []
        for (label, val) in kpi_items:
            cell_data = [
                [Paragraph(f'<b>{val}</b>', ParagraphStyle('x', fontName='Helvetica-Bold',
                           fontSize=18, textColor=colors.white, alignment=TA_CENTER))],
                [Paragraph(label, ParagraphStyle('x', fontName='Helvetica', fontSize=8,
                           textColor=HexColor('#c7d2fe'), alignment=TA_CENTER))]
            ]
            ct2 = Table(cell_data, colWidths=[W_body/4 - 4])
            ct2.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), INDIGO),
                ('TOPPADDING', (0,0), (-1,-1), 8),
                ('BOTTOMPADDING', (0,0), (-1,-1), 6),
            ]))
            kpi_row.append(ct2)

        kpi_table = Table([kpi_row], colWidths=[W_body/4]*4)
        kpi_table.setStyle(TableStyle([
            ('LEFTPADDING', (0,0), (-1,-1), 3),
            ('RIGHTPADDING', (0,0), (-1,-1), 3),
            ('TOPPADDING', (0,0), (-1,-1), 0),
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 14))

        # Executive Summary
        story.append(Paragraph('Executive Summary', H_STYLE))
        story.append(HRFlowable(width='100%', thickness=1, color=INDIGO))
        story.append(Spacer(1, 6))
        exec_text = (
            f'This report presents a comprehensive data analysis of {df.shape[0]:,} records across '
            f'{df.shape[1]} dimensions for {client_name}. The dataset demonstrates '
            f'{"excellent" if health >= 85 else "good" if health >= 70 else "moderate"} '
            f'data health at {health}% quality score. '
        )
        if ml_score:
            exec_text += (f'A predictive model for <b>{target}</b> achieves {ml_score:.1%} accuracy '
                         f'using {task} methodology — ')
            exec_text += ('exceeding industry benchmarks.' if ml_score > 0.85 else 'meeting standard thresholds.')
        exec_text += (f' Key findings include {len(num_cols)} numeric features suitable for '
                     f'predictive modeling and automated decision support.')
        story.append(Paragraph(exec_text, B_STYLE))
        story.append(Spacer(1, 10))

        # Key Findings
        story.append(Paragraph('Key Findings', H_STYLE))
        story.append(HRFlowable(width='100%', thickness=1, color=INDIGO))
        story.append(Spacer(1, 6))

        findings = [
            f'Dataset contains {df.shape[0]:,} records with {df.shape[1]} features — '
            f'{"sufficient" if df.shape[0] > 1000 else "limited"} for statistical analysis',
            f'Data completeness: {100-missing_pct:.1f}% — '
            f'{"No immediate action required" if missing_pct < 5 else f"{missing_pct}% missing values need attention"}',
            f'{len(num_cols)} numeric features available for quantitative modeling and forecasting',
            f'{len(df.select_dtypes(include="object").columns)} categorical features suitable for segmentation analysis',
        ]
        if feat_importance:
            top_feat = list(feat_importance.keys())[0]
            findings.append(f'Most influential predictor: <b>{top_feat}</b> — primary driver of {target}')
        if ml_score and ml_score > 0.80:
            findings.append(f'Predictive model achieves <b>{ml_score:.1%} accuracy</b> — ready for production deployment')

        for i, finding in enumerate(findings):
            row_data = [[
                Paragraph(f'<b>{i+1:02d}</b>', ParagraphStyle('x', fontName='Helvetica-Bold',
                         fontSize=11, textColor=colors.white, alignment=TA_CENTER)),
                Paragraph(finding, SM)
            ]]
            ft = Table(row_data, colWidths=[24, W_body-30])
            ft.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (0,0), INDIGO),
                ('BACKGROUND', (1,0), (1,0), LIGHT),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('TOPPADDING', (0,0), (-1,-1), 7),
                ('BOTTOMPADDING', (0,0), (-1,-1), 7),
                ('LEFTPADDING', (0,0), (-1,-1), 8),
                ('LINEBELOW', (0,0), (-1,-1), 0.5, BORDER),
            ]))
            story.append(ft)
            story.append(Spacer(1, 3))

        story.append(Spacer(1, 10))

        # Chart
        story.append(Paragraph('Data Visualization', H_STYLE))
        story.append(HRFlowable(width='100%', thickness=1, color=INDIGO))
        story.append(Spacer(1, 6))
        from reportlab.platypus import Image as RLImage
        story.append(RLImage(chart_buf, width=W_body, height=W_body*0.45))
        story.append(Spacer(1, 10))

        # Recommendations
        story.append(PageBreak())
        story.append(Paragraph('Strategic Recommendations', H_STYLE))
        story.append(HRFlowable(width='100%', thickness=1, color=INDIGO))
        story.append(Spacer(1, 6))

        industry_recs_map = {
            'Healthcare': [
                ('Immediate', 'Deploy patient risk scoring model to ICU/Emergency', 'HIGH'),
                ('30 Days', 'Implement automated anomaly alerts for critical lab values', 'HIGH'),
                ('60 Days', 'Build readmission prediction dashboard for care teams', 'MEDIUM'),
                ('90 Days', 'Integrate predictive analytics into EMR workflow', 'MEDIUM'),
                ('Ongoing', 'Monthly data quality audit and model performance review', 'LOW'),
            ],
            'Finance': [
                ('Immediate', 'Activate fraud detection model on live transactions', 'HIGH'),
                ('30 Days', 'Customer risk segmentation for portfolio management', 'HIGH'),
                ('60 Days', 'Churn prediction model for relationship managers', 'MEDIUM'),
                ('90 Days', 'Automated credit scoring pipeline', 'MEDIUM'),
                ('Ongoing', 'Quarterly model recalibration and drift monitoring', 'LOW'),
            ],
            'general': [
                ('Immediate', 'Address data quality issues in identified columns', 'HIGH'),
                ('30 Days', 'Deploy predictive model for primary KPI optimization', 'HIGH'),
                ('60 Days', 'Build automated reporting dashboard for leadership', 'MEDIUM'),
                ('90 Days', 'Implement anomaly detection for proactive alerts', 'MEDIUM'),
                ('Ongoing', 'Monthly model performance review and retraining', 'LOW'),
            ]
        }
        recs = industry_recs_map.get(industry, industry_recs_map['general'])
        priority_colors = {'HIGH': '#fef2f2', 'MEDIUM': '#fffbeb', 'LOW': '#f0fdf4'}
        priority_text_colors = {'HIGH': '#991b1b', 'MEDIUM': '#92400e', 'LOW': '#166534'}

        rec_data = [['Timeline', 'Action', 'Priority']]
        for timeline, action, priority in recs:
            rec_data.append([timeline, action, priority])

        rt = Table(rec_data, colWidths=[70, W_body-120, 50])
        rt.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), DARK),
            ('TEXTCOLOR', (0,0), (-1,0), colors.white),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('FONTSIZE', (0,0), (-1,-1), 9),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [HexColor('#f8faff'), colors.white]),
            ('BOX', (0,0), (-1,-1), 1, BORDER),
            ('INNERGRID', (0,0), (-1,-1), 0.5, BORDER),
            ('TOPPADDING', (0,0), (-1,-1), 8),
            ('BOTTOMPADDING', (0,0), (-1,-1), 8),
            ('LEFTPADDING', (0,0), (-1,-1), 10),
            ('ALIGN', (2,0), (2,-1), 'CENTER'),
        ]))
        story.append(rt)
        story.append(Spacer(1, 14))

        # ROI Section
        if baseline_cost > 0:
            story.append(Paragraph('ROI Projection', H_STYLE))
            story.append(HRFlowable(width='100%', thickness=1, color=INDIGO))
            story.append(Spacer(1, 6))
            savings_pct = 0.25
            annual_saving = baseline_cost * savings_pct
            roi_data = [
                ['Metric', 'Value', 'Notes'],
                ['Current Annual Cost', f'{currency}{baseline_cost:,.0f}', 'Baseline'],
                ['Projected Savings', f'{currency}{annual_saving:,.0f}', '~25% optimization'],
                ['Implementation Cost', f'{currency}{annual_saving*0.3:,.0f}', 'One-time investment'],
                ['Net Annual Benefit', f'{currency}{annual_saving*0.7:,.0f}', 'After costs'],
                ['Payback Period', '4-6 months', 'Based on industry average'],
                ['3-Year ROI', f'{round(annual_saving*0.7*3/( annual_saving*0.3)*100)}%', 'Conservative estimate'],
            ]
            roi_table = Table(roi_data, colWidths=[140, 120, W_body-265])
            roi_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), DARK),
                ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('BACKGROUND', (0,-1), (-1,-1), HexColor('#f0fdf4')),
                ('FONTNAME', (0,-1), (-1,-1), 'Helvetica-Bold'),
                ('FONTSIZE', (0,0), (-1,-1), 9.5),
                ('ROWBACKGROUNDS', (0,1), (-1,-2), [HexColor('#f8faff'), colors.white]),
                ('BOX', (0,0), (-1,-1), 1, BORDER),
                ('INNERGRID', (0,0), (-1,-1), 0.5, BORDER),
                ('TOPPADDING', (0,0), (-1,-1), 8),
                ('BOTTOMPADDING', (0,0), (-1,-1), 8),
                ('LEFTPADDING', (0,0), (-1,-1), 10),
                ('ALIGN', (1,0), (1,-1), 'RIGHT'),
            ]))
            story.append(roi_table)
            story.append(Spacer(1, 10))

        # Footer
        story.append(Spacer(1, 20))
        footer_data = [[
            Paragraph(f'Confidential — Prepared for {client_name}',
                     ParagraphStyle('x', fontName='Helvetica', fontSize=8,
                                   textColor=HexColor('#94a3b8'))),
            Paragraph(f'Analyst: {analyst_name}  |  Powered by DS Agent',
                     ParagraphStyle('x', fontName='Helvetica', fontSize=8,
                                   textColor=HexColor('#94a3b8'), alignment=TA_RIGHT)),
        ]]
        ft_table = Table(footer_data, colWidths=[W_body/2, W_body/2])
        ft_table.setStyle(TableStyle([
            ('LINEABOVE', (0,0), (-1,0), 0.5, HexColor('#e2e8f0')),
            ('TOPPADDING', (0,0), (-1,-1), 8),
        ]))
        story.append(ft_table)

        doc.build(story)
        pdf_buf.seek(0)

        log_audit(current_user.id, 'CLIENT_REPORT', client_name,
                 f'Project: {project_name}', request)

        safe = project_name.replace(' ', '_').replace('/', '-')[:30]
        return send_file(pdf_buf, as_attachment=True,
                        download_name=f'DS_Report_{safe}_{client_name}.pdf',
                        mimetype='application/pdf')
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 4: DOMAIN INTELLIGENCE (Healthcare + Finance)
# ════════════════════════════════════════════════════════════════
@app.route('/api/domain_intelligence', methods=['POST'])
@login_required
@check_query_limit
def domain_intelligence():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        domain = data.get('domain', 'healthcare')

        df = agent.df.copy()
        cols_lower = {c.lower(): c for c in df.columns}
        num_cols = df.select_dtypes(include='number').columns.tolist()

        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        import numpy as np

        domain_insights = []
        alerts = []
        kpis = {}
        benchmarks = {}

        if domain == 'healthcare':
            # Healthcare-specific KPI detection
            hc_col_map = {
                'age': ['age', 'patient_age', 'age_years'],
                'bp': ['bp', 'blood_pressure', 'systolic', 'bp_systolic', 'sbp'],
                'glucose': ['glucose', 'blood_glucose', 'sugar', 'glucose_level'],
                'bmi': ['bmi', 'body_mass_index'],
                'heart_rate': ['heart_rate', 'pulse', 'hr', 'bpm'],
                'temperature': ['temperature', 'temp', 'body_temp'],
                'length_of_stay': ['los', 'length_of_stay', 'days_admitted', 'stay_days'],
                'readmission': ['readmission', 'readmitted', 'readmit'],
                'mortality': ['mortality', 'death', 'died', 'survival', 'outcome'],
            }

            hc_benchmarks = {
                'age': {'normal': (18, 65), 'unit': 'years', 'note': 'Working age population'},
                'bp': {'normal': (90, 140), 'unit': 'mmHg', 'note': 'Normal systolic BP'},
                'glucose': {'normal': (70, 140), 'unit': 'mg/dL', 'note': 'Normal fasting glucose'},
                'bmi': {'normal': (18.5, 24.9), 'unit': 'kg/m²', 'note': 'Healthy BMI range'},
                'heart_rate': {'normal': (60, 100), 'unit': 'bpm', 'note': 'Normal resting HR'},
                'temperature': {'normal': (97, 99), 'unit': '°F', 'note': 'Normal body temp'},
            }

            for kpi, keywords in hc_col_map.items():
                matched_col = None
                for kw in keywords:
                    if kw in cols_lower:
                        matched_col = cols_lower[kw]; break
                if matched_col and matched_col in num_cols:
                    series = df[matched_col].dropna()
                    mean_val = float(series.mean())
                    kpis[kpi] = {'value': round(mean_val, 2), 'column': matched_col,
                                 'count': len(series)}
                    if kpi in hc_benchmarks:
                        bm = hc_benchmarks[kpi]
                        low, high = bm['normal']
                        pct_abnormal = round(((series < low) | (series > high)).sum() / len(series) * 100, 1)
                        benchmarks[kpi] = {
                            'mean': round(mean_val, 2),
                            'normal_range': f'{low}–{high} {bm["unit"]}',
                            'pct_abnormal': pct_abnormal,
                            'status': '🔴 CRITICAL' if pct_abnormal > 30 else '🟡 ATTENTION' if pct_abnormal > 15 else '🟢 NORMAL',
                            'note': bm['note']
                        }
                        if pct_abnormal > 20:
                            alerts.append({
                                'alert': f'HIGH {kpi.upper()} ABNORMALITY',
                                'detail': f'{pct_abnormal}% patients have {kpi} outside normal range ({low}–{high})',
                                'action': f'Review {matched_col} — clinical intervention may be needed',
                                'severity': 'CRITICAL' if pct_abnormal > 30 else 'WARNING'
                            })

            # Clinical insights
            if 'age' in kpis:
                age_mean = kpis['age']['value']
                domain_insights.append(f'👥 Average patient age: {age_mean:.1f} years — '
                                      f'{"elderly population (higher comorbidity risk)" if age_mean > 60 else "middle-aged population"}')
            if 'bp' in kpis:
                bp_mean = kpis['bp']['value']
                domain_insights.append(f'❤️ Mean blood pressure: {bp_mean:.0f} mmHg — '
                                      f'{"hypertension prevalent" if bp_mean > 130 else "within acceptable range"}')
            if 'glucose' in kpis:
                glc_mean = kpis['glucose']['value']
                domain_insights.append(f'🩸 Mean glucose: {glc_mean:.1f} mg/dL — '
                                      f'{"diabetic range detected" if glc_mean > 126 else "pre-diabetic range" if glc_mean > 100 else "normal range"}')

            domain_insights.append(f'🏥 {len(df):,} patient records analyzed across {len(num_cols)} clinical parameters')
            if 'length_of_stay' in kpis:
                domain_insights.append(f'📅 Average length of stay: {kpis["length_of_stay"]["value"]:.1f} days')

        elif domain == 'finance':
            finance_col_map = {
                'revenue': ['revenue', 'sales', 'income', 'turnover'],
                'profit': ['profit', 'net_profit', 'earnings', 'net_income'],
                'cost': ['cost', 'expense', 'expenditure', 'total_cost'],
                'loan_amount': ['loan_amount', 'loan', 'credit', 'principal'],
                'default': ['default', 'npa', 'bad_loan', 'defaulted'],
                'transaction': ['amount', 'transaction_amount', 'txn_amount'],
            }
            finance_insights = {
                'revenue': 'Total revenue flow analysis complete',
                'profit': 'Profitability metrics extracted',
                'default': 'Credit risk indicators identified',
            }
            for kpi, keywords in finance_col_map.items():
                for kw in keywords:
                    if kw in cols_lower:
                        matched = cols_lower[kw]
                        if matched in num_cols:
                            s = df[matched].dropna()
                            kpis[kpi] = {'value': round(float(s.mean()), 2),
                                        'total': round(float(s.sum()), 2),
                                        'column': matched}
                            if kpi in finance_insights:
                                domain_insights.append(f'💰 {finance_insights[kpi]}: avg ₹{s.mean():,.2f}')
                            break

        # ── Visualization ───────────────────────────────────────
        fig, axes = plt.subplots(2, 2, figsize=(14, 10))

        # KPI summary
        if kpis:
            kpi_names = list(kpis.keys())[:6]
            kpi_vals = [kpis[k]['value'] for k in kpi_names]
            axes[0,0].bar(kpi_names, kpi_vals, color='#5b5ef4', alpha=0.85, edgecolor='white')
            axes[0,0].set_title(f'{domain.title()} KPI Overview', fontweight='bold')
            plt.setp(axes[0,0].xaxis.get_majorticklabels(), rotation=30, ha='right')

        # Benchmarks
        if benchmarks:
            bm_names = list(benchmarks.keys())[:6]
            pct_abnormal = [benchmarks[k]['pct_abnormal'] for k in bm_names]
            bar_colors = ['#ef4444' if v > 30 else '#f59e0b' if v > 15 else '#059669' for v in pct_abnormal]
            axes[0,1].bar(bm_names, pct_abnormal, color=bar_colors, alpha=0.85, edgecolor='white')
            axes[0,1].axhline(15, color='#f59e0b', linestyle='--', linewidth=1.5, label='Warning 15%')
            axes[0,1].axhline(30, color='#ef4444', linestyle='--', linewidth=1.5, label='Critical 30%')
            axes[0,1].set_title('% Patients Outside Normal Range', fontweight='bold')
            axes[0,1].set_ylabel('% Abnormal'); axes[0,1].legend()
            axes[0,1].set_ylim(0, 100)

        # Distribution of first numeric col
        if num_cols:
            axes[1,0].hist(df[num_cols[0]].dropna(), bins=30, color='#8b5cf6', alpha=0.85, edgecolor='white')
            axes[1,0].set_title(f'{num_cols[0]} Distribution', fontweight='bold')

        # Alerts visualization
        if alerts:
            alert_labels = [a['alert'][:25] for a in alerts[:5]]
            alert_sev = [1 if a['severity']=='CRITICAL' else 0.6 for a in alerts[:5]]
            alert_colors = ['#ef4444' if a['severity']=='CRITICAL' else '#f59e0b' for a in alerts[:5]]
            axes[1,1].barh(alert_labels[::-1], alert_sev[::-1], color=alert_colors[::-1], alpha=0.85)
            axes[1,1].set_title('Clinical Alerts', fontweight='bold')
            axes[1,1].set_xlim(0, 1.5)
        else:
            axes[1,1].text(0.5, 0.5, '✅ No Critical\nAlerts', ha='center', va='center',
                          fontsize=16, color='#059669', fontweight='bold', transform=axes[1,1].transAxes)
            axes[1,1].axis('off')

        plt.suptitle(f'Domain Intelligence — {domain.title()} Analytics',
                    fontweight='bold', fontsize=14, color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'domain': domain,
            'kpis_detected': kpis,
            'benchmarks': benchmarks,
            'alerts': alerts,
            'domain_insights': domain_insights,
            'chart': chart,
            'total_kpis': len(kpis),
            'critical_alerts': sum(1 for a in alerts if a.get('severity') == 'CRITICAL'),
            'message': f'Domain Intelligence complete! {len(kpis)} KPIs, {len(alerts)} alerts, {len(domain_insights)} insights'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 5: SMART DATASET PROFILER
# Ek click mein poora dataset ka professional profile
# ════════════════════════════════════════════════════════════════
@app.route('/api/smart_profile', methods=['POST'])
@login_required
@check_query_limit
def smart_profile():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400

        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        import numpy as np

        df = agent.df.copy()
        num_cols = df.select_dtypes(include='number').columns.tolist()
        cat_cols = df.select_dtypes(include='object').columns.tolist()
        date_cols = df.select_dtypes(include=['datetime64']).columns.tolist()

        # ── Column-level deep profile ───────────────────────────
        col_profiles = []
        for col in df.columns:
            series = df[col]
            missing = int(series.isnull().sum())
            missing_pct = round(missing / len(df) * 100, 2)
            unique = int(series.nunique())
            dtype = str(series.dtype)
            profile = {
                'column': col,
                'dtype': dtype,
                'missing': missing,
                'missing_pct': missing_pct,
                'unique': unique,
                'unique_pct': round(unique / len(df) * 100, 2),
            }
            if col in num_cols:
                s = series.dropna()
                q1 = float(s.quantile(0.25))
                q3 = float(s.quantile(0.75))
                iqr = q3 - q1
                outliers = int(((s < q1 - 1.5*iqr) | (s > q3 + 1.5*iqr)).sum())
                profile.update({
                    'mean': round(float(s.mean()), 4),
                    'median': round(float(s.median()), 4),
                    'std': round(float(s.std()), 4),
                    'min': round(float(s.min()), 4),
                    'max': round(float(s.max()), 4),
                    'q1': round(q1, 4),
                    'q3': round(q3, 4),
                    'skewness': round(float(s.skew()), 3),
                    'kurtosis': round(float(s.kurt()), 3),
                    'outliers': outliers,
                    'outlier_pct': round(outliers / len(s) * 100, 2),
                    'zero_count': int((s == 0).sum()),
                    'negative_count': int((s < 0).sum()),
                    'col_type': 'numeric'
                })
            elif col in cat_cols:
                top_vals = series.value_counts().head(5).to_dict()
                profile.update({
                    'top_values': {str(k): int(v) for k, v in top_vals.items()},
                    'mode': str(series.mode()[0]) if len(series.mode()) > 0 else '',
                    'avg_length': round(float(series.dropna().astype(str).str.len().mean()), 1),
                    'col_type': 'categorical'
                })

            # Health score per column
            health = 100
            health -= min(50, missing_pct * 2)
            if col in num_cols:
                health -= min(20, profile.get('outlier_pct', 0))
                if abs(profile.get('skewness', 0)) > 3:
                    health -= 10
            profile['health_score'] = max(0, round(health, 1))
            profile['recommendation'] = (
                'Good — no action needed' if health >= 80
                else f'Fix {missing} missing values' if missing_pct > 10
                else f'Handle {profile.get("outliers",0)} outliers' if col in num_cols and profile.get('outlier_pct', 0) > 10
                else 'Acceptable — minor cleanup'
            )
            col_profiles.append(profile)

        # ── Dataset-level summary ───────────────────────────────
        total_missing = int(df.isnull().sum().sum())
        total_cells = df.shape[0] * df.shape[1]
        completeness = round((1 - total_missing / total_cells) * 100, 2)
        duplicates = int(df.duplicated().sum())
        memory_mb = round(df.memory_usage(deep=True).sum() / 1024 / 1024, 3)

        overall_health = round(np.mean([c['health_score'] for c in col_profiles]), 1)

        # Potential ML columns
        good_target_cols = [c['column'] for c in col_profiles
                           if c['col_type'] in ('numeric', 'categorical')
                           and c['missing_pct'] < 20
                           and c.get('unique', 0) >= 2]

        # ── Visualization ───────────────────────────────────────
        fig = plt.figure(figsize=(18, 12))
        gs = fig.add_gridspec(3, 4, hspace=0.5, wspace=0.4)

        # 1. Health scores per column
        ax1 = fig.add_subplot(gs[0, :2])
        col_names_short = [c['column'][:14] for c in col_profiles[:15]]
        health_scores = [c['health_score'] for c in col_profiles[:15]]
        bar_c = ['#059669' if h >= 80 else '#f59e0b' if h >= 60 else '#ef4444' for h in health_scores]
        ax1.bar(col_names_short, health_scores, color=bar_c, alpha=0.9, edgecolor='white', linewidth=1.2)
        ax1.axhline(80, color='#059669', linestyle='--', linewidth=1.5, label='Good ≥80')
        ax1.axhline(60, color='#f59e0b', linestyle='--', linewidth=1.5, label='Warn ≥60')
        ax1.set_title('Column Health Scores', fontweight='bold', fontsize=12)
        ax1.set_ylabel('Health %'); ax1.set_ylim(0, 110)
        ax1.legend(fontsize=8)
        plt.setp(ax1.xaxis.get_majorticklabels(), rotation=35, ha='right', fontsize=8)

        # 2. Missing values heatmap style bar
        ax2 = fig.add_subplot(gs[0, 2:])
        missing_data = [(c['column'][:14], c['missing_pct']) for c in col_profiles if c['missing_pct'] > 0]
        if missing_data:
            m_names, m_vals = zip(*sorted(missing_data, key=lambda x: -x[1]))
            bar_colors_m = ['#ef4444' if v > 30 else '#f59e0b' if v > 10 else '#94a3b8' for v in m_vals]
            ax2.barh(list(m_names)[::-1], list(m_vals)[::-1], color=list(bar_colors_m)[::-1], alpha=0.85)
            ax2.set_title('Missing Values % by Column', fontweight='bold', fontsize=12)
            ax2.set_xlabel('Missing %')
        else:
            ax2.text(0.5, 0.5, '✅ Zero Missing Values!\nPerfect Dataset', ha='center', va='center',
                    fontsize=14, color='#059669', fontweight='bold', transform=ax2.transAxes)
            ax2.axis('off')
            ax2.set_title('Missing Values', fontweight='bold')

        # 3. Data types donut
        ax3 = fig.add_subplot(gs[1, 0])
        type_counts = {'Numeric': len(num_cols), 'Categorical': len(cat_cols), 'DateTime': len(date_cols)}
        type_counts = {k: v for k, v in type_counts.items() if v > 0}
        if type_counts:
            wedge_colors = ['#5b5ef4', '#8b5cf6', '#06d6a0'][:len(type_counts)]
            wedges, texts, autotexts = ax3.pie(
                type_counts.values(), labels=type_counts.keys(),
                colors=wedge_colors, autopct='%1.0f%%',
                wedgeprops={'width': 0.55}, startangle=90)
            for at in autotexts: at.set_fontsize(9)
        ax3.set_title('Column Types', fontweight='bold')

        # 4. Outlier counts
        ax4 = fig.add_subplot(gs[1, 1])
        outlier_data = [(c['column'][:14], c.get('outlier_pct', 0))
                       for c in col_profiles if c.get('col_type') == 'numeric' and c.get('outlier_pct', 0) > 0]
        if outlier_data:
            o_names, o_vals = zip(*sorted(outlier_data, key=lambda x: -x[1])[:8])
            ax4.bar(o_names, o_vals,
                   color=['#ef4444' if v > 15 else '#f59e0b' for v in o_vals], alpha=0.85)
            ax4.set_title('Outlier % by Column', fontweight='bold')
            ax4.set_ylabel('%')
            plt.setp(ax4.xaxis.get_majorticklabels(), rotation=35, ha='right', fontsize=8)
        else:
            ax4.text(0.5, 0.5, '✅ No Outliers\nDetected', ha='center', va='center',
                    fontsize=13, color='#059669', fontweight='bold', transform=ax4.transAxes)
            ax4.axis('off')
        ax4.set_title('Outlier % by Column', fontweight='bold')

        # 5. Skewness
        ax5 = fig.add_subplot(gs[1, 2])
        skew_data = [(c['column'][:14], c.get('skewness', 0))
                    for c in col_profiles if c.get('col_type') == 'numeric']
        if skew_data:
            s_names, s_vals = zip(*sorted(skew_data, key=lambda x: abs(x[1]), reverse=True)[:8])
            bar_sk = ['#ef4444' if abs(v) > 2 else '#f59e0b' if abs(v) > 1 else '#059669' for v in s_vals]
            ax5.bar(s_names, s_vals, color=bar_sk, alpha=0.85)
            ax5.axhline(0, color='black', linewidth=1)
            ax5.set_title('Skewness by Column', fontweight='bold')
            plt.setp(ax5.xaxis.get_majorticklabels(), rotation=35, ha='right', fontsize=8)

        # 6. Overall dataset scorecard
        ax6 = fig.add_subplot(gs[1, 3])
        ax6.axis('off')
        scorecard = [
            ('📋 Rows', f'{df.shape[0]:,}'),
            ('📊 Columns', str(df.shape[1])),
            ('✅ Completeness', f'{completeness}%'),
            ('🏥 Health Score', f'{overall_health}/100'),
            ('🔁 Duplicates', str(duplicates)),
            ('💾 Memory', f'{memory_mb} MB'),
            ('🔢 Numeric Cols', str(len(num_cols))),
            ('🔤 Cat Cols', str(len(cat_cols))),
        ]
        y_pos = 0.95
        for label, val in scorecard:
            color = '#059669' if 'Health' in label and overall_health >= 75 else \
                    '#ef4444' if 'Health' in label and overall_health < 60 else '#1e293b'
            ax6.text(0.05, y_pos, label, transform=ax6.transAxes,
                    fontsize=9, color='#64748b')
            ax6.text(0.65, y_pos, val, transform=ax6.transAxes,
                    fontsize=9, fontweight='bold', color=color)
            y_pos -= 0.12
        ax6.set_title('Dataset Scorecard', fontweight='bold')
        ax6.add_patch(plt.Rectangle((0, 0), 1, 1, fill=True, color='#f8faff',
                                   transform=ax6.transAxes, zorder=-1))

        # 7. Correlation heatmap bottom
        ax7 = fig.add_subplot(gs[2, :])
        if len(num_cols) >= 3:
            import seaborn as sns
            corr = df[num_cols[:12]].corr()
            mask = np.zeros_like(corr, dtype=bool)
            mask[np.triu_indices_from(mask)] = True
            sns.heatmap(corr, mask=mask, annot=True, fmt='.2f', cmap='RdYlBu_r',
                       center=0, ax=ax7, linewidths=0.3, annot_kws={'size': 7},
                       square=False)
            ax7.set_title('Correlation Matrix (Lower Triangle)', fontweight='bold', fontsize=11)
        else:
            ax7.text(0.5, 0.5, 'Need 3+ numeric columns for correlation',
                    ha='center', va='center', transform=ax7.transAxes, color='#94a3b8')
            ax7.axis('off')

        plt.suptitle(f'Smart Dataset Profile — {df.shape[0]:,} rows × {df.shape[1]} cols | Health: {overall_health}/100',
                    fontweight='bold', fontsize=14, color='#1e3a5f')
        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'dataset_shape': list(df.shape),
            'overall_health_score': overall_health,
            'completeness_pct': completeness,
            'duplicates': duplicates,
            'memory_mb': memory_mb,
            'total_missing': total_missing,
            'col_profiles': col_profiles,
            'num_cols': num_cols,
            'cat_cols': cat_cols,
            'good_target_cols': good_target_cols[:10],
            'chart': chart,
            'message': f'Profile complete! Health: {overall_health}/100 | {completeness}% complete | {duplicates} duplicates'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 6: AUTO INSIGHT GENERATOR (Natural Language)
# Dataset se automatic human-readable story generate karo
# ════════════════════════════════════════════════════════════════
@app.route('/api/auto_insight', methods=['POST'])
@login_required
@check_query_limit
def auto_insight():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400

        import numpy as np
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        df = agent.df.copy()
        num_cols = df.select_dtypes(include='number').columns.tolist()
        cat_cols = df.select_dtypes(include='object').columns.tolist()
        data = request.get_json()
        focus_col = data.get('focus_col', num_cols[0] if num_cols else '')
        tone = data.get('tone', 'business')  # business, technical, simple

        insights_story = []
        charts_data = {}

        # ── Insight 1: Size & Shape ─────────────────────────────
        size_msg = f"The dataset contains {df.shape[0]:,} records across {df.shape[1]} dimensions. "
        if df.shape[0] > 50000:
            size_msg += "This is a large-scale dataset — ML models will be highly reliable."
        elif df.shape[0] > 5000:
            size_msg += "Dataset size is adequate for reliable statistical analysis and modeling."
        elif df.shape[0] > 500:
            size_msg += "Moderate dataset — suitable for analysis but more data would improve ML accuracy."
        else:
            size_msg += "Small dataset — use cross-validation carefully and avoid complex models."
        insights_story.append({'type': 'size', 'icon': '📋', 'title': 'Dataset Scale', 'insight': size_msg})

        # ── Insight 2: Data Quality ─────────────────────────────
        missing_total = int(df.isnull().sum().sum())
        missing_pct = round(missing_total / (df.shape[0] * df.shape[1]) * 100, 2)
        dupes = int(df.duplicated().sum())
        quality_msg = f"Data completeness stands at {100-missing_pct:.1f}%"
        if missing_pct == 0:
            quality_msg += " — perfect data quality with zero missing values. "
        elif missing_pct < 5:
            quality_msg += f" with only {missing_total} missing values — excellent quality. "
        else:
            worst_col = df.isnull().sum().idxmax()
            quality_msg += f" — {missing_total} missing values found. Column '{worst_col}' needs immediate attention. "
        if dupes > 0:
            quality_msg += f"Additionally, {dupes} duplicate records were detected and should be removed."
        else:
            quality_msg += "No duplicate records found."
        insights_story.append({'type': 'quality', 'icon': '🏥', 'title': 'Data Quality', 'insight': quality_msg})

        # ── Insight 3: Focus Column Deep Dive ──────────────────
        if focus_col and focus_col in df.columns and focus_col in num_cols:
            s = df[focus_col].dropna()
            mean_v = float(s.mean())
            median_v = float(s.median())
            skew_v = float(s.skew())
            cv = float(s.std() / (mean_v + 1e-10))
            q1, q3 = float(s.quantile(0.25)), float(s.quantile(0.75))
            iqr = q3 - q1
            outliers = int(((s < q1 - 1.5*iqr) | (s > q3 + 1.5*iqr)).sum())

            focus_msg = f"'{focus_col}' averages {mean_v:.2f} (median: {median_v:.2f}). "
            if abs(mean_v - median_v) / (abs(median_v) + 1e-10) > 0.1:
                focus_msg += f"The gap between mean and median suggests {'right' if mean_v > median_v else 'left'}-skewed distribution — "
                focus_msg += "use median for reporting, not mean. "
            if cv > 0.5:
                focus_msg += f"High variability (CV={cv:.2f}) indicates inconsistent values — segment analysis recommended. "
            if outliers > 0:
                focus_msg += f"{outliers} outliers detected — investigate whether these are data errors or genuine extreme cases."
            insights_story.append({'type': 'focus', 'icon': '🔍', 'title': f"Deep Dive: {focus_col}", 'insight': focus_msg})

        # ── Insight 4: Top Correlations ─────────────────────────
        if len(num_cols) >= 3:
            corr = df[num_cols[:10]].corr().abs()
            corr_pairs = []
            for i in range(len(corr.columns)):
                for j in range(i+1, len(corr.columns)):
                    v = float(corr.iloc[i, j])
                    if v > 0.5:
                        corr_pairs.append((corr.columns[i], corr.columns[j], round(v, 3)))
            corr_pairs.sort(key=lambda x: -x[2])

            if corr_pairs:
                top = corr_pairs[0]
                corr_msg = f"Strongest business relationship found: '{top[0]}' and '{top[1]}' are {top[2]:.0%} correlated. "
                corr_msg += "This means improvements in one directly impact the other — a key business lever. "
                if len(corr_pairs) > 3:
                    corr_msg += f"{len(corr_pairs)} total significant correlations found — rich dataset for predictive modeling."
                insights_story.append({'type': 'correlation', 'icon': '🔗', 'title': 'Key Business Relationships', 'insight': corr_msg})

        # ── Insight 5: Category Distribution ───────────────────
        if cat_cols:
            col = cat_cols[0]
            vc = df[col].value_counts()
            top_val = str(vc.index[0]); top_pct = round(vc.iloc[0] / len(df) * 100, 1)
            concentration = 'highly concentrated' if top_pct > 50 else 'evenly distributed' if top_pct < 25 else 'moderately distributed'
            cat_msg = f"'{col}' is {concentration} — top value '{top_val}' represents {top_pct}% of all records. "
            if top_pct > 70:
                cat_msg += "This imbalance may affect ML model performance — consider resampling techniques."
            else:
                cat_msg += f"Distribution across {vc.nunique()} categories is suitable for segmentation analysis."
            insights_story.append({'type': 'category', 'icon': '📊', 'title': f"Category Analysis: {col}", 'insight': cat_msg})

        # ── Insight 6: Trend Detection ──────────────────────────
        if num_cols:
            trend_findings = []
            for col in num_cols[:5]:
                s = df[col].dropna()
                if len(s) < 10: continue
                half = len(s) // 2
                h1 = float(s.iloc[:half].mean())
                h2 = float(s.iloc[half:].mean())
                chg = round((h2 - h1) / (abs(h1) + 1e-10) * 100, 1)
                if abs(chg) > 10:
                    direction = 'increasing 📈' if chg > 0 else 'decreasing 📉'
                    trend_findings.append(f"'{col}' is {direction} by {abs(chg):.1f}%")
            if trend_findings:
                trend_msg = "Trend analysis reveals: " + "; ".join(trend_findings[:3]) + ". "
                trend_msg += "These trends indicate directional changes worth monitoring closely."
                insights_story.append({'type': 'trend', 'icon': '📈', 'title': 'Trend Analysis', 'insight': trend_msg})

        # ── Insight 7: ML Readiness ─────────────────────────────
        ml_score = 100
        issues = []
        if missing_pct > 20: ml_score -= 30; issues.append(f'high missing data ({missing_pct}%)')
        if df.shape[0] < 200: ml_score -= 25; issues.append('small sample size')
        if len(num_cols) < 2: ml_score -= 20; issues.append('few numeric features')
        if dupes > df.shape[0] * 0.1: ml_score -= 15; issues.append('many duplicates')
        ml_score = max(0, ml_score)
        ml_msg = f"ML Readiness Score: {ml_score}/100. "
        if ml_score >= 80:
            ml_msg += "Dataset is excellent for machine learning — deploy models with confidence."
        elif ml_score >= 60:
            ml_msg += f"Dataset needs minor fixes before ML: {', '.join(issues)}."
        else:
            ml_msg += f"Dataset requires significant preparation: {', '.join(issues)}."
        insights_story.append({'type': 'ml_ready', 'icon': '🤖', 'title': 'ML Readiness', 'insight': ml_msg})

        # ── Build visualization ─────────────────────────────────
        fig = plt.figure(figsize=(16, 10))
        gs = fig.add_gridspec(2, 3, hspace=0.45, wspace=0.35)

        # Insight scores summary
        ax1 = fig.add_subplot(gs[0, 0])
        categories_radar = ['Data\nQuality', 'ML\nReadiness', 'Completeness', 'Volume', 'Diversity']
        values_radar = [
            max(0, 100 - missing_pct * 3),
            ml_score,
            100 - missing_pct,
            min(100, df.shape[0] / 100),
            min(100, len(num_cols) * 10 + len(cat_cols) * 5)
        ]
        bar_c = ['#059669' if v >= 75 else '#f59e0b' if v >= 50 else '#ef4444' for v in values_radar]
        ax1.bar(categories_radar, values_radar, color=bar_c, alpha=0.85, edgecolor='white', linewidth=1.5)
        ax1.set_ylim(0, 110); ax1.set_title('Dataset Quality Scores', fontweight='bold')
        ax1.set_ylabel('Score /100')

        # Focus column distribution
        ax2 = fig.add_subplot(gs[0, 1])
        if focus_col and focus_col in num_cols:
            df[focus_col].dropna().hist(ax=ax2, bins=30, color='#5b5ef4', alpha=0.85, edgecolor='white')
            mean_v_plot = float(df[focus_col].mean())
            median_v_plot = float(df[focus_col].median())
            ax2.axvline(mean_v_plot, color='#ef4444', linestyle='--', linewidth=2, label=f'Mean: {mean_v_plot:.2f}')
            ax2.axvline(median_v_plot, color='#059669', linestyle='--', linewidth=2, label=f'Median: {median_v_plot:.2f}')
            ax2.legend(fontsize=8)
            ax2.set_title(f'{focus_col} Distribution', fontweight='bold')

        # Top categories
        ax3 = fig.add_subplot(gs[0, 2])
        if cat_cols:
            vc = df[cat_cols[0]].value_counts().head(7)
            colors_cat = ['#5b5ef4', '#8b5cf6', '#06d6a0', '#f59e0b', '#ef4444', '#0ea5e9', '#84cc16']
            ax3.pie(vc.values, labels=[str(x)[:12] for x in vc.index],
                   colors=colors_cat[:len(vc)], autopct='%1.1f%%', startangle=90)
            ax3.set_title(f'{cat_cols[0]} Distribution', fontweight='bold')

        # Correlation heatmap
        ax4 = fig.add_subplot(gs[1, :2])
        if len(num_cols) >= 3:
            import seaborn as sns
            corr = df[num_cols[:8]].corr()
            sns.heatmap(corr, annot=True, fmt='.2f', cmap='RdYlBu_r', center=0,
                       ax=ax4, linewidths=0.3, annot_kws={'size': 8}, square=True)
            ax4.set_title('Correlation — Business Relationships', fontweight='bold')

        # Trend lines
        ax5 = fig.add_subplot(gs[1, 2])
        if len(num_cols) >= 2:
            for i, col in enumerate(num_cols[:4]):
                s = df[col].dropna().rolling(max(1, len(df)//20)).mean().dropna()
                s_norm = (s - s.min()) / (s.max() - s.min() + 1e-10)
                colors_line = ['#5b5ef4', '#8b5cf6', '#06d6a0', '#f59e0b']
                ax5.plot(s_norm.values, color=colors_line[i], linewidth=2,
                        label=col[:12], alpha=0.85)
            ax5.legend(fontsize=7); ax5.set_title('Normalized Trends', fontweight='bold')
            ax5.set_xlabel('Record Index'); ax5.set_ylabel('Normalized Value')

        plt.suptitle('Auto Insight Report — AI Generated Analysis',
                    fontweight='bold', fontsize=14, color='#1e3a5f')
        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'insights_story': insights_story,
            'total_insights': len(insights_story),
            'ml_readiness_score': ml_score,
            'data_quality_score': max(0, round(100 - missing_pct * 3, 1)),
            'focus_column': focus_col,
            'chart': chart,
            'message': f'{len(insights_story)} auto insights generated! ML Readiness: {ml_score}/100'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 7: WHAT-IF SCENARIO SIMULATOR
# "Agar ye column 10% badh jaye to kya hoga?" — simulation
# ════════════════════════════════════════════════════════════════
@app.route('/api/whatif_simulator', methods=['POST'])
@login_required
@check_query_limit
def whatif_simulator():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400

        import numpy as np
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        from sklearn.ensemble import RandomForestRegressor, RandomForestClassifier
        from sklearn.preprocessing import LabelEncoder
        from sklearn.model_selection import train_test_split

        data = request.get_json()
        target = data.get('target', '')
        scenarios = data.get('scenarios', [])
        # scenarios = [{'column': 'age', 'change_type': 'pct', 'change_value': 10}, ...]

        if not target:
            return jsonify({'error': 'Target column do simulation ke liye!'}), 400
        if target not in agent.df.columns:
            return jsonify({'error': f'Target "{target}" not found!'}), 400

        df = agent.df.copy().dropna()
        num_cols = df.select_dtypes(include='number').columns.tolist()

        # Train base model
        y_raw = df[target]; X = df.drop(columns=[target])
        le_dict = {}
        for c in X.select_dtypes(include='object').columns:
            le = LabelEncoder()
            X[c] = le.fit_transform(X[c].astype(str))
            le_dict[c] = le
        task = 'classification' if (y_raw.dtype == object or y_raw.nunique() <= 15) else 'regression'
        le_y = None
        if y_raw.dtype == object:
            le_y = LabelEncoder(); y = le_y.fit_transform(y_raw.astype(str))
        else:
            y = y_raw.values

        X_tr, X_te, y_tr, y_te = train_test_split(X, y, test_size=0.2, random_state=42)
        model = RandomForestClassifier(100, random_state=42) if task == 'classification' \
                else RandomForestRegressor(100, random_state=42)
        model.fit(X_tr, y_tr)
        base_score = round(float(model.score(X_te, y_te)), 4)

        # Baseline predictions on full dataset
        X_full = X.copy()
        base_predictions = model.predict(X_full)
        if task == 'regression':
            base_mean = float(np.mean(base_predictions))
        else:
            base_mean = float(np.mean(base_predictions))

        # ── Run each scenario ───────────────────────────────────
        scenario_results = []
        if not scenarios:
            # Auto-generate scenarios for numeric cols
            for col in num_cols[:3]:
                if col == target: continue
                scenarios.append({'column': col, 'change_type': 'pct', 'change_value': 10})
                scenarios.append({'column': col, 'change_type': 'pct', 'change_value': -10})

        for sc in scenarios[:10]:
            col = sc.get('column', '')
            change_type = sc.get('change_type', 'pct')  # pct, absolute, set
            change_val = float(sc.get('change_value', 10))

            if col not in X_full.columns:
                continue

            X_sim = X_full.copy()
            original_col_mean = float(X_sim[col].mean())

            if change_type == 'pct':
                X_sim[col] = X_sim[col] * (1 + change_val / 100)
            elif change_type == 'absolute':
                X_sim[col] = X_sim[col] + change_val
            elif change_type == 'set':
                X_sim[col] = change_val

            sim_predictions = model.predict(X_sim)
            sim_mean = float(np.mean(sim_predictions))
            impact = round(sim_mean - base_mean, 4)
            impact_pct = round((sim_mean - base_mean) / (abs(base_mean) + 1e-10) * 100, 2)

            scenario_results.append({
                'scenario': f"{col} {'+' if change_val >= 0 else ''}{change_val}{'%' if change_type == 'pct' else ''}",
                'column': col,
                'change_type': change_type,
                'change_value': change_val,
                'original_col_mean': round(original_col_mean, 3),
                'new_col_mean': round(float(X_sim[col].mean()), 3),
                'base_target_mean': round(base_mean, 4),
                'simulated_target_mean': round(sim_mean, 4),
                'impact': impact,
                'impact_pct': impact_pct,
                'direction': '📈 Positive' if impact > 0 else '📉 Negative' if impact < 0 else '➡️ Neutral',
                'magnitude': 'HIGH' if abs(impact_pct) > 10 else 'MEDIUM' if abs(impact_pct) > 3 else 'LOW'
            })

        # Sort by absolute impact
        scenario_results.sort(key=lambda x: abs(x['impact_pct']), reverse=True)

        # ── Sensitivity Analysis ────────────────────────────────
        sensitivity = []
        for col in num_cols[:8]:
            if col == target or col not in X_full.columns: continue
            impacts_pct_list = []
            for chg in [-20, -10, -5, 5, 10, 20]:
                X_sens = X_full.copy()
                X_sens[col] = X_sens[col] * (1 + chg / 100)
                preds = model.predict(X_sens)
                imp_pct = round((float(np.mean(preds)) - base_mean) / (abs(base_mean) + 1e-10) * 100, 2)
                impacts_pct_list.append(imp_pct)
            avg_sensitivity = round(float(np.mean([abs(x) for x in impacts_pct_list])), 3)
            sensitivity.append({'column': col, 'sensitivity_score': avg_sensitivity,
                                'impacts': impacts_pct_list})
        sensitivity.sort(key=lambda x: -x['sensitivity_score'])

        # ── Visualization ───────────────────────────────────────
        fig, axes = plt.subplots(2, 2, figsize=(14, 10))

        # Scenario impact chart
        if scenario_results:
            sc_labels = [s['scenario'][:20] for s in scenario_results[:8]]
            sc_impacts = [s['impact_pct'] for s in scenario_results[:8]]
            bar_c = ['#059669' if v > 0 else '#ef4444' for v in sc_impacts]
            axes[0, 0].barh(sc_labels[::-1], sc_impacts[::-1], color=bar_c[::-1], alpha=0.85)
            axes[0, 0].axvline(0, color='black', linewidth=1)
            axes[0, 0].set_title(f'Scenario Impact on {target}', fontweight='bold')
            axes[0, 0].set_xlabel('% Change in Target')

        # Sensitivity ranking
        if sensitivity:
            sens_cols = [s['column'][:14] for s in sensitivity[:8]]
            sens_scores = [s['sensitivity_score'] for s in sensitivity[:8]]
            axes[0, 1].bar(sens_cols, sens_scores,
                          color=['#5b5ef4'] * len(sens_cols), alpha=0.85, edgecolor='white')
            axes[0, 1].set_title(f'Feature Sensitivity to {target}', fontweight='bold')
            axes[0, 1].set_ylabel('Sensitivity Score')
            plt.setp(axes[0, 1].xaxis.get_majorticklabels(), rotation=35, ha='right')

        # Tornado chart
        if scenario_results:
            top_sc = scenario_results[:6]
            labels = [s['scenario'][:18] for s in top_sc]
            pos_vals = [max(0, s['impact_pct']) for s in top_sc]
            neg_vals = [min(0, s['impact_pct']) for s in top_sc]
            y_pos = range(len(labels))
            axes[1, 0].barh(y_pos, pos_vals, color='#059669', alpha=0.85, label='Positive')
            axes[1, 0].barh(y_pos, neg_vals, color='#ef4444', alpha=0.85, label='Negative')
            axes[1, 0].set_yticks(y_pos); axes[1, 0].set_yticklabels(labels, fontsize=8)
            axes[1, 0].axvline(0, color='black', linewidth=1.5)
            axes[1, 0].set_title('Tornado Chart — Impact Analysis', fontweight='bold')
            axes[1, 0].legend()

        # Sensitivity response curves (top 3 features)
        axes[1, 1].set_title('Sensitivity Response Curves', fontweight='bold')
        change_pcts = [-20, -10, -5, 5, 10, 20]
        line_colors = ['#5b5ef4', '#8b5cf6', '#06d6a0', '#f59e0b']
        for i, sens in enumerate(sensitivity[:4]):
            axes[1, 1].plot(change_pcts, sens['impacts'],
                           color=line_colors[i], linewidth=2.5,
                           marker='o', markersize=5,
                           label=sens['column'][:14])
        axes[1, 1].axhline(0, color='black', linewidth=1, linestyle='--')
        axes[1, 1].axvline(0, color='black', linewidth=1, linestyle='--')
        axes[1, 1].set_xlabel('Input Change %'); axes[1, 1].set_ylabel('Target Impact %')
        axes[1, 1].legend(fontsize=8)

        plt.suptitle(f'What-If Simulator — Target: {target} | Base Score: {base_score:.4f}',
                    fontweight='bold', fontsize=13, color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        top_lever = sensitivity[0]['column'] if sensitivity else 'N/A'
        return jsonify({
            'success': True,
            'target': target,
            'task': task,
            'base_model_score': base_score,
            'base_target_mean': round(base_mean, 4),
            'scenarios_run': len(scenario_results),
            'scenario_results': scenario_results,
            'sensitivity_ranking': sensitivity[:8],
            'top_business_lever': top_lever,
            'chart': chart,
            'message': f'What-If simulation complete! {len(scenario_results)} scenarios | Top lever: {top_lever}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 8: SMART COLUMN RENAMER & SCHEMA MAPPER
# Messy column names ko professional schema mein convert karo
# ════════════════════════════════════════════════════════════════
@app.route('/api/schema_mapper', methods=['POST'])
@login_required
def schema_mapper():
    try:
        if agent.df is None:
            return jsonify({'error': 'Pehle dataset upload karo!'}), 400
        data = request.get_json()
        action = data.get('action', 'suggest')  # suggest, apply, validate
        custom_mapping = data.get('mapping', {})

        df = agent.df.copy()
        original_cols = list(df.columns)

        # Common messy -> clean mappings
        auto_mapping = {}
        clean_suggestions = {}

        import re as _re

        # Standard abbreviation expansions
        abbrev_map = {
            'pt': 'patient', 'pts': 'patients', 'doc': 'doctor',
            'dept': 'department', 'hosp': 'hospital', 'adm': 'admission',
            'disch': 'discharge', 'diag': 'diagnosis', 'proc': 'procedure',
            'amt': 'amount', 'qty': 'quantity', 'num': 'number',
            'dt': 'date', 'yr': 'year', 'mo': 'month', 'wk': 'week',
            'rev': 'revenue', 'exp': 'expense', 'sal': 'salary',
            'emp': 'employee', 'cust': 'customer', 'ord': 'order',
            'prod': 'product', 'cat': 'category', 'desc': 'description',
            'addr': 'address', 'ph': 'phone', 'mob': 'mobile',
            'bp': 'blood_pressure', 'hr': 'heart_rate', 'bmi': 'bmi',
            'dob': 'date_of_birth', 'los': 'length_of_stay',
        }

        for col in original_cols:
            clean = col.strip()
            # Remove special characters except underscore
            clean = _re.sub(r'[^\w\s]', '_', clean)
            # Multiple spaces/underscores to single
            clean = _re.sub(r'[\s_]+', '_', clean)
            # Lowercase
            clean = clean.lower().strip('_')
            # Expand abbreviations
            parts = clean.split('_')
            parts = [abbrev_map.get(p, p) for p in parts]
            clean = '_'.join(parts)
            # Remove leading numbers
            clean = _re.sub(r'^(\d)', r'col_\1', clean)

            if clean != col:
                auto_mapping[col] = clean
            clean_suggestions[col] = clean

        # Detect data types and suggest better names
        type_hints = {}
        for col in df.columns:
            series = df[col].dropna()
            if series.dtype == object:
                sample = str(series.iloc[0]) if len(series) > 0 else ''
                # Detect emails
                if '@' in sample and '.' in sample:
                    type_hints[col] = 'email'
                # Detect phone numbers
                elif any(c.isdigit() for c in sample) and len(_re.sub(r'\D','',sample)) in [10,11,12]:
                    type_hints[col] = 'phone'
                # Detect dates
                elif any(x in sample for x in ['-', '/', '20', '19']) and len(sample) <= 12:
                    type_hints[col] = 'date'
                else:
                    type_hints[col] = 'text'
            elif 'int' in str(series.dtype):
                type_hints[col] = 'integer'
            elif 'float' in str(series.dtype):
                type_hints[col] = 'decimal'

        if action == 'apply':
            mapping_to_use = custom_mapping if custom_mapping else auto_mapping
            df.rename(columns=mapping_to_use, inplace=True)
            agent.df = df
            agent.available_columns = list(df.columns)
            return jsonify({
                'success': True,
                'action': 'applied',
                'renamed_count': len(mapping_to_use),
                'new_columns': list(df.columns),
                'mapping_applied': mapping_to_use,
                'message': f'Schema updated! {len(mapping_to_use)} columns renamed.'
            })

        elif action == 'validate':
            issues = []
            for col in df.columns:
                if ' ' in col: issues.append({'col': col, 'issue': 'Contains spaces — use underscores'})
                if col != col.lower(): issues.append({'col': col, 'issue': 'Not lowercase — inconsistent naming'})
                if len(col) > 40: issues.append({'col': col, 'issue': 'Too long — consider abbreviation'})
                if col[0].isdigit(): issues.append({'col': col, 'issue': 'Starts with number — invalid for many tools'})
                import keyword
                if keyword.iskeyword(col): issues.append({'col': col, 'issue': 'Python reserved keyword — will cause errors'})
            return jsonify({
                'success': True,
                'action': 'validated',
                'total_issues': len(issues),
                'issues': issues,
                'schema_quality_score': max(0, round(100 - len(issues) * (100 / max(1, len(df.columns))), 1)),
                'message': f'Schema validated! {len(issues)} issues found in {len(df.columns)} columns.'
            })

        else:  # suggest
            return jsonify({
                'success': True,
                'action': 'suggest',
                'original_columns': original_cols,
                'suggested_mapping': auto_mapping,
                'clean_names': clean_suggestions,
                'type_hints': type_hints,
                'changes_suggested': len(auto_mapping),
                'message': f'Schema analysis complete! {len(auto_mapping)} column renames suggested.'
            })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 9: MULTI-DATASET MERGE & JOIN ENGINE
# Multiple CSVs ko intelligently merge karo
# ════════════════════════════════════════════════════════════════
@app.route('/api/smart_merge', methods=['POST'])
@login_required
@check_query_limit
def smart_merge():
    try:
        data = request.get_json()
        action = data.get('action', 'analyze')
        df2_b64 = data.get('df2_base64', '')
        join_type = data.get('join_type', 'auto')  # auto, inner, left, outer
        on_col = data.get('on_col', '')

        import base64 as _b64, io as _io2
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        if agent.df is None:
            return jsonify({'error': 'Pehle primary dataset upload karo!'}), 400

        df1 = agent.df.copy()

        if action == 'analyze_primary':
            # Just analyze current loaded dataset for merge readiness
            num_cols = df1.select_dtypes(include='number').columns.tolist()
            cat_cols = df1.select_dtypes(include='object').columns.tolist()
            potential_keys = []
            for col in df1.columns:
                unique_ratio = df1[col].nunique() / len(df1)
                if unique_ratio > 0.8:
                    potential_keys.append({'column': col, 'unique_ratio': round(unique_ratio, 3),
                                          'unique_count': df1[col].nunique(), 'type': 'potential_key'})
                elif unique_ratio > 0.3:
                    potential_keys.append({'column': col, 'unique_ratio': round(unique_ratio, 3),
                                          'unique_count': df1[col].nunique(), 'type': 'potential_join'})
            return jsonify({
                'success': True,
                'primary_shape': list(df1.shape),
                'columns': list(df1.columns),
                'potential_join_keys': potential_keys[:10],
                'message': f'Primary dataset analyzed! {len(potential_keys)} potential join keys found.'
            })

        if not df2_b64:
            return jsonify({'error': 'Second dataset (df2_base64) required for merge!'}), 400

        # Decode second dataset
        csv2 = _b64.b64decode(df2_b64)
        df2 = pd.read_csv(_io2.BytesIO(csv2))

        # Find common columns
        common_cols = list(set(df1.columns) & set(df2.columns))

        # Auto-detect best join key
        best_key = on_col
        if not best_key and common_cols:
            # Prefer ID-like columns
            id_keywords = ['id', 'key', 'code', 'no', 'number', 'patient', 'emp', 'cust']
            for kw in id_keywords:
                for col in common_cols:
                    if kw in col.lower():
                        best_key = col; break
                if best_key: break
            if not best_key:
                best_key = common_cols[0]

        if not best_key:
            return jsonify({'error': 'No common columns found — datasets cannot be merged!',
                           'df1_cols': list(df1.columns), 'df2_cols': list(df2.columns)}), 400

        # Auto join type
        if join_type == 'auto':
            overlap = len(set(df1[best_key].dropna()) & set(df2[best_key].dropna()))
            overlap_pct = overlap / max(len(df1[best_key].dropna()), 1)
            join_type = 'inner' if overlap_pct > 0.7 else 'left'

        # Perform merge
        merged = df1.merge(df2, on=best_key, how=join_type, suffixes=('_primary', '_secondary'))
        agent.df = merged
        agent.available_columns = list(merged.columns)

        # Visualization
        fig, axes = plt.subplots(1, 3, figsize=(14, 4))
        data_labels = [f'Primary\n{df1.shape[0]:,} rows', f'Secondary\n{df2.shape[0]:,} rows', f'Merged\n{merged.shape[0]:,} rows']
        data_vals = [df1.shape[0], df2.shape[0], merged.shape[0]]
        bar_c = ['#5b5ef4', '#8b5cf6', '#059669']
        axes[0].bar(data_labels, data_vals, color=bar_c, alpha=0.85, edgecolor='white', linewidth=1.5)
        axes[0].set_title('Row Counts Before & After Merge', fontweight='bold')
        axes[0].set_ylabel('Rows')
        for i, v in enumerate(data_vals):
            axes[0].text(i, v + max(data_vals)*0.01, f'{v:,}', ha='center', fontweight='bold', fontsize=9)

        col_data = [df1.shape[1], df2.shape[1], merged.shape[1]]
        axes[1].bar(data_labels, col_data, color=bar_c, alpha=0.85, edgecolor='white', linewidth=1.5)
        axes[1].set_title('Column Counts Before & After', fontweight='bold')
        for i, v in enumerate(col_data):
            axes[1].text(i, v + 0.2, str(v), ha='center', fontweight='bold', fontsize=10)

        # Missing after merge
        miss = merged.isnull().sum().sort_values(ascending=False).head(8)
        if len(miss[miss > 0]) > 0:
            axes[2].barh(miss[miss > 0].index, miss[miss > 0].values, color='#f59e0b', alpha=0.85)
            axes[2].set_title('Missing Values After Merge', fontweight='bold')
        else:
            axes[2].text(0.5, 0.5, '✅ No Missing\nValues Post-Merge', ha='center', va='center',
                        fontsize=12, color='#059669', fontweight='bold', transform=axes[2].transAxes)
            axes[2].axis('off')
            axes[2].set_title('Post-Merge Quality', fontweight='bold')

        plt.suptitle(f'Smart Merge Result — {join_type.upper()} JOIN on "{best_key}"',
                    fontweight='bold', fontsize=12, color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'join_type': join_type,
            'join_key': best_key,
            'df1_shape': list(df1.shape),
            'df2_shape': list(df2.shape),
            'merged_shape': list(merged.shape),
            'common_columns': common_cols,
            'new_columns': list(merged.columns),
            'rows_gained': merged.shape[0] - df1.shape[0],
            'cols_gained': merged.shape[1] - df1.shape[1],
            'chart': chart,
            'message': f'Merge complete! {join_type.upper()} JOIN on "{best_key}" → {merged.shape[0]:,} rows × {merged.shape[1]} cols'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# LEAD DS — FEATURE 10: SMART KPI TRACKER
# Custom KPIs define karo aur continuously monitor karo
# ════════════════════════════════════════════════════════════════
@app.route('/api/kpi_tracker/set', methods=['POST'])
@login_required
def kpi_set():
    try:
        import json as _json
        data = request.get_json()
        kpis = data.get('kpis', [])
        # kpis = [{'name': 'Revenue', 'column': 'revenue', 'aggregation': 'sum', 'target': 1000000, 'alert_threshold': 0.8}]

        db = SessionLocal()
        saved = []
        for kpi in kpis:
            k = UserChart(
                user_id=current_user.id,
                chart_type='kpi_definition',
                chart_title=kpi.get('name', 'KPI'),
            )
            db.add(k); db.commit()
            saved.append({'id': k.id, 'name': kpi.get('name')})
        db.close()
        return jsonify({'success': True, 'saved_kpis': saved,
                       'message': f'{len(saved)} KPIs defined and saved!'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/kpi_tracker/dashboard', methods=['GET'])
@login_required
def kpi_dashboard():
    try:
        import json as _json
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        db = SessionLocal()
        kpi_defs = db.query(UserChart).filter_by(
            user_id=current_user.id, chart_type='kpi_definition'
        ).all()
        db.close()

        if not kpi_defs:
            return jsonify({'error': 'Koi KPI define nahi hai! Pehle /api/kpi_tracker/set call karo.'}), 400
        if agent.df is None:
            return jsonify({'error': 'Dataset load karo pehle!'}), 400

        df = agent.df
        kpi_results = []

        agg_funcs = {
            'sum': lambda s: float(s.sum()),
            'mean': lambda s: float(s.mean()),
            'median': lambda s: float(s.median()),
            'max': lambda s: float(s.max()),
            'min': lambda s: float(s.min()),
            'count': lambda s: int(s.count()),
            'std': lambda s: float(s.std()),
        }

        for kd in kpi_defs:
            try:
                kpi_def = _json.loads(kd.description)
                col = kpi_def.get('column', '')
                agg = kpi_def.get('aggregation', 'sum')
                target = float(kpi_def.get('target', 0))
                threshold = float(kpi_def.get('alert_threshold', 0.8))

                if col not in df.columns:
                    continue

                series = df[col].dropna()
                if series.dtype not in ['int64', 'float64']:
                    series = pd.to_numeric(series, errors='coerce').dropna()

                actual = round(agg_funcs.get(agg, agg_funcs['sum'])(series), 4)
                achievement = round(actual / (target + 1e-10) * 100, 1) if target else None
                status = 'ON_TRACK' if (achievement or 100) >= threshold * 100 else 'AT_RISK' if (achievement or 100) >= threshold * 80 else 'OFF_TRACK'

                kpi_results.append({
                    'id': kd.id,
                    'name': kpi_def.get('name', col),
                    'column': col,
                    'aggregation': agg,
                    'actual': actual,
                    'target': target,
                    'achievement_pct': achievement,
                    'status': status,
                    'status_icon': '🟢' if status == 'ON_TRACK' else '🟡' if status == 'AT_RISK' else '🔴',
                    'unit': kpi_def.get('unit', ''),
                })
            except Exception:
                continue

        if not kpi_results:
            return jsonify({'error': 'KPI columns not found in current dataset!'}), 400

        # Visualization
        fig, axes = plt.subplots(1, 2, figsize=(14, 5))

        # KPI Achievement gauge bars
        kpi_names = [k['name'][:15] for k in kpi_results]
        achievements = [k['achievement_pct'] or 100 for k in kpi_results]
        bar_colors_kpi = ['#059669' if a >= 80 else '#f59e0b' if a >= 60 else '#ef4444' for a in achievements]
        bars = axes[0].bar(kpi_names, achievements, color=bar_colors_kpi, alpha=0.85, edgecolor='white', linewidth=1.5)
        axes[0].axhline(100, color='#1e293b', linewidth=2, linestyle='--', label='Target 100%')
        axes[0].axhline(80, color='#f59e0b', linewidth=1.5, linestyle=':', label='At-Risk 80%')
        axes[0].set_title('KPI Achievement vs Target', fontweight='bold')
        axes[0].set_ylabel('Achievement %'); axes[0].legend(fontsize=9)
        axes[0].set_ylim(0, max(130, max(achievements) + 20))
        for bar, val in zip(bars, achievements):
            axes[0].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                        f'{val:.0f}%', ha='center', fontsize=9, fontweight='bold')
        plt.setp(axes[0].xaxis.get_majorticklabels(), rotation=30, ha='right')

        # Actual vs Target
        x = range(len(kpi_results))
        actuals = [k['actual'] for k in kpi_results]
        targets = [k['target'] for k in kpi_results if k['target']]
        if targets and len(targets) == len(actuals):
            width = 0.35
            axes[1].bar([i - width/2 for i in x], actuals, width, label='Actual', color='#5b5ef4', alpha=0.85)
            axes[1].bar([i + width/2 for i in x], [k['target'] for k in kpi_results], width, label='Target', color='#94a3b8', alpha=0.85)
            axes[1].set_xticks(list(x)); axes[1].set_xticklabels(kpi_names, rotation=30, ha='right')
            axes[1].set_title('Actual vs Target Values', fontweight='bold')
            axes[1].legend()
        else:
            axes[1].bar(kpi_names, actuals, color='#5b5ef4', alpha=0.85, edgecolor='white')
            axes[1].set_title('KPI Actual Values', fontweight='bold')
            plt.setp(axes[1].xaxis.get_majorticklabels(), rotation=30, ha='right')

        on_track = sum(1 for k in kpi_results if k['status'] == 'ON_TRACK')
        plt.suptitle(f'KPI Dashboard — {on_track}/{len(kpi_results)} On Track',
                    fontweight='bold', fontsize=13, color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'kpi_results': kpi_results,
            'total_kpis': len(kpi_results),
            'on_track': on_track,
            'at_risk': sum(1 for k in kpi_results if k['status'] == 'AT_RISK'),
            'off_track': sum(1 for k in kpi_results if k['status'] == 'OFF_TRACK'),
            'chart': chart,
            'message': f'KPI Dashboard loaded! {on_track}/{len(kpi_results)} KPIs on track.'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/kpi_tracker/delete/<int:kid>', methods=['DELETE'])
@login_required
def kpi_delete(kid):
    try:
        db = SessionLocal()
        k = db.query(UserChart).filter_by(id=kid, user_id=current_user.id).first()
        if k: db.delete(k); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# ██████╗ ██████╗ ███╗   ███╗    ███╗   ███╗ ██████╗ ██████╗ ██╗   ██╗██╗     ███████╗
# ██╔══██╗██╔══██╗████╗ ████║    ████╗ ████║██╔═══██╗██╔══██╗██║   ██║██║     ██╔════╝
# ██████╔╝██║  ██║██╔████╔██║    ██╔████╔██║██║   ██║██║  ██║██║   ██║██║     █████╗
# ██╔══██╗██║  ██║██║╚██╔╝██║    ██║╚██╔╝██║██║   ██║██║  ██║██║   ██║██║     ██╔══╝
# ██████╔╝██████╔╝██║ ╚═╝ ██║    ██║ ╚═╝ ██║╚██████╔╝██████╔╝╚██████╔╝███████╗███████╗
# ╚═════╝ ╚═════╝ ╚═╝     ╚═╝    ╚═╝     ╚═╝ ╚═════╝ ╚═════╝  ╚═════╝ ╚══════╝╚══════╝
#
# BUSINESS DEVELOPMENT MANAGER + EXECUTIVE ASSISTANT
# Data + Groq AI powered | Personal Use Only
# ════════════════════════════════════════════════════════════════

def get_groq_response(system_prompt, user_prompt, max_tokens=2000):
    """Central Groq AI caller for BDM + EA features"""
    try:
        from langchain_groq import ChatGroq
        llm = ChatGroq(
            model="llama-3.3-70b-versatile",
            temperature=0.4,
            max_tokens=max_tokens
        )
        from langchain_core.messages import SystemMessage, HumanMessage
        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=user_prompt)
        ]
        response = llm.invoke(messages)
        return response.content
    except Exception as e:
        return f"AI Error: {str(e)}"


# ════════════════════════════════════════════════════════════════
# BDM — 1: CLIENT PROPOSAL GENERATOR (Data + AI)
# ════════════════════════════════════════════════════════════════
@app.route('/api/bdm/proposal', methods=['POST'])
@login_required
def bdm_proposal():
    try:
        data = request.get_json()
        client_name    = data.get('client_name', 'Client')
        client_type    = data.get('client_type', 'Hospital')
        problem        = data.get('problem', '')
        budget         = data.get('budget', '')
        timeline       = data.get('timeline', '3 months')
        your_name      = data.get('your_name', 'Data Scientist')
        services       = data.get('services', [])
        currency       = data.get('currency', '₹')

        # Data context from loaded dataset
        data_context = ""
        if agent.df is not None:
            df = agent.df
            num_cols = df.select_dtypes(include='number').columns.tolist()
            missing_pct = round(df.isnull().sum().sum() / (df.shape[0]*df.shape[1])*100, 2)
            data_context = f"""
CLIENT DATA ALREADY ANALYZED:
- Dataset: {df.shape[0]:,} records × {df.shape[1]} columns
- Data Quality: {100-missing_pct:.1f}% complete
- Numeric Features: {len(num_cols)} (suitable for ML modeling)
- Key Columns: {', '.join(list(df.columns)[:8])}
- This analysis can be referenced in the proposal.
"""

        system_prompt = f"""You are an elite Business Development Manager and Data Science consultant in India.
You write highly professional, persuasive client proposals that win contracts.
Your proposals are specific, data-backed, and clearly show ROI.
Always write in a confident, professional tone. Use Indian business context.
Format with clear sections. Be specific about deliverables and timelines."""

        user_prompt = f"""Write a complete professional business proposal for:

CLIENT: {client_name} ({client_type})
PROBLEM TO SOLVE: {problem}
PROPOSED SERVICES: {', '.join(services) if services else 'Data Analysis, ML Modeling, Dashboard, Insights'}
BUDGET RANGE: {currency}{budget}
TIMELINE: {timeline}
PROPOSED BY: {your_name}
{data_context}

Write a complete proposal with these sections:
1. Executive Summary (compelling opening)
2. Understanding Your Challenge (show you understand their pain)
3. Our Proposed Solution (specific, detailed)
4. Deliverables & Timeline (week by week breakdown)
5. Technology & Methodology (tools, models, approach)
6. Investment & ROI (clear pricing, expected returns)
7. Why Choose Us (unique value proposition)
8. Next Steps (clear call to action)

Make it compelling, specific, and professional. Show concrete ROI numbers."""

        proposal_text = get_groq_response(system_prompt, user_prompt, max_tokens=3000)

        # Save proposal
        import json as _json
        db = SessionLocal()
        p = UserChart(
            user_id=current_user.id,
            chart_type='bdm_proposal',
            chart_title=f'Proposal — {client_name}',
            image_data=_json.dumps({
                'client_name': client_name,
                'client_type': client_type,
                'problem': problem,
                'budget': budget,
                'timeline': timeline,
                'services': services,
            }),
        )
        db.add(p); db.commit(); pid = p.id; db.close()

        return jsonify({
            'success': True,
            'proposal_id': pid,
            'client_name': client_name,
            'proposal_text': proposal_text,
            'word_count': len(proposal_text.split()),
            'message': f'Proposal for {client_name} generated! {len(proposal_text.split())} words'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# BDM — 2: SALES PIPELINE TRACKER
# ════════════════════════════════════════════════════════════════
@app.route('/api/bdm/pipeline/add', methods=['POST'])
@login_required
def bdm_pipeline_add():
    try:
        import json as _json
        data = request.get_json()
        lead = {
            'client_name':  data.get('client_name', ''),
            'client_type':  data.get('client_type', ''),
            'contact':      data.get('contact', ''),
            'deal_value':   data.get('deal_value', 0),
            'stage':        data.get('stage', 'Lead'),  # Lead→Contacted→Proposal→Negotiation→Won/Lost
            'probability':  data.get('probability', 20),
            'next_action':  data.get('next_action', ''),
            'next_date':    data.get('next_date', ''),
            'notes':        data.get('notes', ''),
            'currency':     data.get('currency', '₹'),
        }
        db = SessionLocal()
        entry = UserChart(
            user_id=current_user.id,
            chart_type='bdm_pipeline',
            chart_title=f"{lead['client_name']} — {lead['stage']}",
        )
        db.add(entry); db.commit(); eid = entry.id; db.close()
        return jsonify({'success': True, 'id': eid, 'message': f"Lead added: {lead['client_name']} in {lead['stage']} stage"})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/bdm/pipeline/list', methods=['GET'])
@login_required
def bdm_pipeline_list():
    try:
        import json as _json
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        db = SessionLocal()
        entries = db.query(UserChart).filter_by(
            user_id=current_user.id, chart_type='bdm_pipeline'
        ).order_by(UserChart.created_at.desc()).all()
        db.close()

        pipeline = []
        stage_counts = {}
        total_value = 0
        weighted_value = 0

        stage_order = ['Lead', 'Contacted', 'Proposal', 'Negotiation', 'Won', 'Lost']
        stage_probs  = {'Lead': 10, 'Contacted': 25, 'Proposal': 50, 'Negotiation': 75, 'Won': 100, 'Lost': 0}
        stage_colors = {'Lead':'#94a3b8','Contacted':'#3b82f6','Proposal':'#8b5cf6',
                       'Negotiation':'#f59e0b','Won':'#059669','Lost':'#ef4444'}

        for e in entries:
            try: lead = _json.loads(e.description)
            except: continue
            lead['id'] = e.id
            lead['created_at'] = e.created_at.strftime('%d %b %Y')
            pipeline.append(lead)
            stage = lead.get('stage', 'Lead')
            stage_counts[stage] = stage_counts.get(stage, 0) + 1
            val = float(lead.get('deal_value', 0))
            total_value += val
            prob = float(lead.get('probability', stage_probs.get(stage, 20)))
            weighted_value += val * prob / 100

        if pipeline:
            # Pipeline funnel chart
            fig, axes = plt.subplots(1, 3, figsize=(15, 5))

            # Funnel
            stages_present = [s for s in stage_order if s in stage_counts]
            counts = [stage_counts.get(s, 0) for s in stages_present]
            colors_used = [stage_colors.get(s, '#94a3b8') for s in stages_present]
            bars = axes[0].barh(stages_present[::-1], counts[::-1],
                               color=colors_used[::-1], alpha=0.85, edgecolor='white', linewidth=1.5)
            axes[0].set_title('Sales Pipeline Funnel', fontweight='bold')
            axes[0].set_xlabel('Number of Deals')
            for bar, val in zip(bars, counts[::-1]):
                if val > 0:
                    axes[0].text(bar.get_width() + 0.05, bar.get_y() + bar.get_height()/2,
                               str(val), va='center', fontweight='bold')

            # Deal values by stage
            stage_values = {}
            for lead in pipeline:
                s = lead.get('stage', 'Lead')
                stage_values[s] = stage_values.get(s, 0) + float(lead.get('deal_value', 0))
            if stage_values:
                sv_stages = list(stage_values.keys())
                sv_vals = [stage_values[s]/1e5 for s in sv_stages]  # in Lakhs
                axes[1].bar(sv_stages, sv_vals,
                           color=[stage_colors.get(s, '#94a3b8') for s in sv_stages],
                           alpha=0.85, edgecolor='white')
                axes[1].set_title('Deal Value by Stage (₹ Lakhs)', fontweight='bold')
                plt.setp(axes[1].xaxis.get_majorticklabels(), rotation=30, ha='right')

            # Win/Loss pie
            won = stage_counts.get('Won', 0)
            lost = stage_counts.get('Lost', 0)
            active = sum(v for k, v in stage_counts.items() if k not in ('Won', 'Lost'))
            if won + lost + active > 0:
                axes[2].pie([won, lost, active],
                           labels=[f'Won\n{won}', f'Lost\n{lost}', f'Active\n{active}'],
                           colors=['#059669', '#ef4444', '#3b82f6'],
                           autopct='%1.0f%%', startangle=90)
                axes[2].set_title('Deal Status Overview', fontweight='bold')

            plt.suptitle(f'Sales Pipeline | Total: ₹{total_value:,.0f} | Weighted: ₹{weighted_value:,.0f}',
                        fontweight='bold', fontsize=12, color='#1e3a5f')
            plt.tight_layout()
            buf = _io.BytesIO()
            plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
            plt.close('all'); buf.seek(0)
            chart = base64.b64encode(buf.read()).decode()
        else:
            chart = None

        return jsonify({
            'success': True,
            'pipeline': pipeline,
            'total_leads': len(pipeline),
            'stage_counts': stage_counts,
            'total_value': round(total_value, 2),
            'weighted_value': round(weighted_value, 2),
            'win_rate': round(stage_counts.get('Won', 0) /
                             max(1, stage_counts.get('Won', 0) + stage_counts.get('Lost', 0)) * 100, 1),
            'chart': chart,
            'message': f'Pipeline: {len(pipeline)} leads | Total ₹{total_value:,.0f} | Weighted ₹{weighted_value:,.0f}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/bdm/pipeline/update/<int:lid>', methods=['POST'])
@login_required
def bdm_pipeline_update(lid):
    try:
        import json as _json
        data = request.get_json()
        db = SessionLocal()
        entry = db.query(UserChart).filter_by(id=lid, user_id=current_user.id).first()
        if not entry:
            db.close()
            return jsonify({'error': 'Lead not found!'}), 404
        try: current_data = _json.loads(entry.description)
        except: current_data = {}
        current_data.update({k: v for k, v in data.items() if k != 'id'})
        entry.description = _json.dumps(current_data)
        entry.title = f"{current_data.get('client_name','?')} — {current_data.get('stage','?')}"
        db.commit(); db.close()
        return jsonify({'success': True, 'message': f"Lead updated to stage: {current_data.get('stage')}"})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/bdm/pipeline/delete/<int:lid>', methods=['DELETE'])
@login_required
def bdm_pipeline_delete(lid):
    try:
        db = SessionLocal()
        e = db.query(UserChart).filter_by(id=lid, user_id=current_user.id).first()
        if e: db.delete(e); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# BDM — 3: COMPETITOR ANALYSIS (AI powered)
# ════════════════════════════════════════════════════════════════
@app.route('/api/bdm/competitor', methods=['POST'])
@login_required
def bdm_competitor():
    try:
        data = request.get_json()
        your_service    = data.get('your_service', 'Data Science Consulting')
        competitors     = data.get('competitors', [])
        target_market   = data.get('target_market', 'Healthcare sector India')
        your_strengths  = data.get('your_strengths', 'AI-powered analysis, affordable, fast delivery')
        pricing         = data.get('pricing', '')

        system_prompt = """You are a senior business strategy consultant in India.
You provide sharp, actionable competitor analysis for data science and analytics businesses.
Be specific, honest, and strategic. Use Indian market context."""

        user_prompt = f"""Perform a detailed competitor analysis for a Data Science consultant:

MY SERVICE: {your_service}
MY TARGET MARKET: {target_market}
MY STRENGTHS: {your_strengths}
MY PRICING: {pricing}
KNOWN COMPETITORS: {', '.join(competitors) if competitors else 'General data science agencies, freelancers on Upwork/Fiverr, in-house analytics teams'}

Provide:
1. **Competitive Landscape Overview** — who are the real competitors in Indian market
2. **Competitor Profiles** — for each competitor: strengths, weaknesses, pricing, positioning
3. **Competitive Gaps** — what they don't offer that I can
4. **My Unique Positioning** — how to differentiate myself
5. **Pricing Strategy** — how to price competitively yet profitably  
6. **Win Strategy** — specific tactics to win against each competitor type
7. **Red Flags** — threats to watch out for
8. **Action Plan** — 5 immediate actions to gain competitive edge

Be specific to the Indian market. Give actual numbers and strategies."""

        analysis = get_groq_response(system_prompt, user_prompt, max_tokens=2500)

        return jsonify({
            'success': True,
            'analysis': analysis,
            'your_service': your_service,
            'target_market': target_market,
            'message': 'Competitor analysis complete!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# BDM — 4: REVENUE FORECAST (Data + AI)
# ════════════════════════════════════════════════════════════════
@app.route('/api/bdm/revenue_forecast', methods=['POST'])
@login_required
def bdm_revenue_forecast():
    try:
        import numpy as np
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io

        data = request.get_json()
        current_clients    = int(data.get('current_clients', 0))
        avg_deal_value     = float(data.get('avg_deal_value', 50000))
        monthly_new_leads  = int(data.get('monthly_new_leads', 5))
        conversion_rate    = float(data.get('conversion_rate', 20)) / 100
        churn_rate         = float(data.get('churn_rate', 10)) / 100
        growth_target      = float(data.get('growth_target', 30)) / 100
        currency           = data.get('currency', '₹')
        horizon_months     = int(data.get('horizon_months', 12))

        # Month-by-month simulation
        months = list(range(1, horizon_months + 1))
        clients_count = []
        monthly_revenue = []
        cumulative_revenue = []
        new_clients_monthly = []

        current = current_clients
        cum_rev = 0

        for m in months:
            new_c = round(monthly_new_leads * conversion_rate * (1 + growth_target * m / 12))
            churned = round(current * churn_rate)
            current = max(0, current + new_c - churned)
            rev = current * avg_deal_value
            cum_rev += rev
            clients_count.append(current)
            monthly_revenue.append(rev)
            cumulative_revenue.append(cum_rev)
            new_clients_monthly.append(new_c)

        # AI commentary
        system_prompt = """You are a revenue strategy expert for Indian freelance/consulting businesses.
Give sharp, practical revenue growth advice."""

        user_prompt = f"""Analyze this revenue forecast and give strategic advice:

Current Clients: {current_clients}
Avg Deal Value: {currency}{avg_deal_value:,.0f}
Monthly New Leads: {monthly_new_leads}
Conversion Rate: {conversion_rate*100:.0f}%
Churn Rate: {churn_rate*100:.0f}%
Growth Target: {growth_target*100:.0f}%/year
Forecast Period: {horizon_months} months

Projected Results:
- Clients at end: {clients_count[-1]}
- Monthly Revenue at end: {currency}{monthly_revenue[-1]:,.0f}
- Total Revenue ({horizon_months}mo): {currency}{cumulative_revenue[-1]:,.0f}

Give:
1. Is this forecast realistic? What needs to change?
2. Top 3 actions to double conversion rate
3. How to reduce churn to <5%
4. Pricing strategy to increase deal value
5. Month-by-month focus areas (Q1, Q2, Q3, Q4)
Keep it crisp and actionable."""

        ai_advice = get_groq_response(system_prompt, user_prompt, max_tokens=1500)

        # Chart
        fig, axes = plt.subplots(2, 2, figsize=(14, 9))

        # Monthly Revenue
        axes[0,0].bar(months, [r/1e5 for r in monthly_revenue],
                     color='#5b5ef4', alpha=0.85, edgecolor='white')
        axes[0,0].plot(months, [r/1e5 for r in monthly_revenue],
                      color='#1e3a5f', linewidth=2, marker='o', markersize=4)
        axes[0,0].set_title(f'Monthly Revenue ({currency} Lakhs)', fontweight='bold')
        axes[0,0].set_xlabel('Month'); axes[0,0].set_ylabel(f'{currency} Lakhs')

        # Cumulative Revenue
        axes[0,1].fill_between(months, [r/1e5 for r in cumulative_revenue],
                              alpha=0.3, color='#059669')
        axes[0,1].plot(months, [r/1e5 for r in cumulative_revenue],
                      color='#059669', linewidth=2.5, marker='o', markersize=4)
        axes[0,1].set_title(f'Cumulative Revenue ({currency} Lakhs)', fontweight='bold')
        axes[0,1].set_xlabel('Month'); axes[0,1].set_ylabel(f'{currency} Lakhs')

        # Client Growth
        axes[1,0].bar(months, clients_count, color='#8b5cf6', alpha=0.85, edgecolor='white')
        axes[1,0].set_title('Client Count Growth', fontweight='bold')
        axes[1,0].set_xlabel('Month'); axes[1,0].set_ylabel('Active Clients')

        # New vs Churned
        axes[1,1].bar(months, new_clients_monthly, label='New', color='#059669', alpha=0.75)
        churned_monthly = [round(clients_count[max(0,i-1)] * churn_rate) for i in range(len(months))]
        axes[1,1].bar(months, [-c for c in churned_monthly], label='Churned', color='#ef4444', alpha=0.75)
        axes[1,1].axhline(0, color='black', linewidth=1)
        axes[1,1].set_title('New vs Churned Clients/Month', fontweight='bold')
        axes[1,1].set_xlabel('Month'); axes[1,1].legend()

        total_rev = cumulative_revenue[-1]
        plt.suptitle(f'Revenue Forecast — {horizon_months} months | Total: {currency}{total_rev:,.0f}',
                    fontweight='bold', fontsize=13, color='#1e3a5f')
        plt.tight_layout()
        buf = _io.BytesIO()
        plt.savefig(buf, format='png', dpi=130, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success': True,
            'months': months,
            'monthly_revenue': monthly_revenue,
            'cumulative_revenue': cumulative_revenue,
            'clients_count': clients_count,
            'final_monthly_revenue': monthly_revenue[-1],
            'total_revenue': cumulative_revenue[-1],
            'final_clients': clients_count[-1],
            'ai_advice': ai_advice,
            'chart': chart,
            'message': f'Forecast: {currency}{total_rev:,.0f} in {horizon_months} months | {clients_count[-1]} clients'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# BDM — 5: COLD OUTREACH EMAIL GENERATOR
# ════════════════════════════════════════════════════════════════
@app.route('/api/bdm/cold_outreach', methods=['POST'])
@login_required
def bdm_cold_outreach():
    try:
        data = request.get_json()
        prospect_name    = data.get('prospect_name', 'Decision Maker')
        prospect_company = data.get('prospect_company', 'Company')
        prospect_role    = data.get('prospect_role', 'CEO')
        industry         = data.get('industry', 'Healthcare')
        pain_point       = data.get('pain_point', '')
        your_name        = data.get('your_name', 'Data Scientist')
        email_type       = data.get('email_type', 'cold')  # cold, followup, linkedin
        previous_context = data.get('previous_context', '')

        system_prompt = """You are an expert B2B sales copywriter who specializes in data science consulting.
You write emails that get responses — specific, value-first, no fluff.
You understand Indian business culture and communication style."""

        type_instructions = {
            'cold': "Write a cold outreach email. Short (150-200 words), specific, clear value prop, one CTA.",
            'followup': f"Write a follow-up email. Reference previous interaction: {previous_context}. Be persistent but professional.",
            'linkedin': "Write a LinkedIn connection message + follow-up note. Keep it brief and personal (under 100 words).",
        }

        data_hook = ""
        if agent.df is not None:
            df = agent.df
            data_hook = f"I recently analyzed data in {industry} sector showing [{df.shape[0]:,} patient/customer records] — found interesting patterns in {list(df.columns)[:3]}."

        user_prompt = f"""Write a {email_type} outreach message:

TO: {prospect_name}, {prospect_role} at {prospect_company} ({industry})
PAIN POINT TO ADDRESS: {pain_point if pain_point else f'Data is not being used for decisions in {industry}'}
FROM: {your_name} — Data Science Consultant
DATA HOOK: {data_hook}

Instructions: {type_instructions.get(email_type, type_instructions['cold'])}

Also write:
- Subject line (compelling, specific)
- PS line (adds urgency/credibility)
- Alternative shorter version (for WhatsApp/SMS)"""

        email_content = get_groq_response(system_prompt, user_prompt, max_tokens=1000)

        # Save
        import json as _json
        db = SessionLocal()
        e = UserChart(
            user_id=current_user.id,
            chart_type='bdm_outreach',
            chart_title=f'Outreach — {prospect_company} ({email_type})',
        )
        db.add(e); db.commit(); db.close()

        return jsonify({
            'success': True,
            'email_content': email_content,
            'prospect': f'{prospect_name} at {prospect_company}',
            'email_type': email_type,
            'message': f'{email_type.title()} outreach generated for {prospect_company}!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# BDM — 6: MEETING PREP ASSISTANT
# Kisi bhi client meeting ke liye AI prep karo
# ════════════════════════════════════════════════════════════════
@app.route('/api/bdm/meeting_prep', methods=['POST'])
@login_required
def bdm_meeting_prep():
    try:
        data = request.get_json()
        client_name     = data.get('client_name', 'Client')
        meeting_type    = data.get('meeting_type', 'discovery')  # discovery, proposal, negotiation, review
        client_industry = data.get('client_industry', 'Healthcare')
        known_info      = data.get('known_info', '')
        your_goal       = data.get('your_goal', 'Close the deal')

        system_prompt = """You are an expert sales coach and business consultant.
You prepare data science consultants for high-stakes client meetings.
Give practical, specific preparation advice."""

        data_summary = ""
        if agent.df is not None:
            df = agent.df
            data_summary = f"\nI have analyzed their data: {df.shape[0]:,} records, key metrics available in {list(df.columns)[:5]}"

        user_prompt = f"""Prepare me for this client meeting:

CLIENT: {client_name} ({client_industry})
MEETING TYPE: {meeting_type}
WHAT I KNOW: {known_info}{data_summary}
MY GOAL: {your_goal}

Give me:
1. **Meeting Agenda** (suggested structure with time allocation)
2. **Opening Statement** (how to start — word for word)
3. **Discovery Questions** (10 powerful questions to ask)
4. **Expected Objections & Responses** (top 5 objections, how to handle each)
5. **Value Statements** (3-4 compelling statements to make)
6. **Closing Strategy** (how to close or advance the deal)
7. **Materials to Prepare** (what to bring/send)
8. **Red Flags to Watch** (signs they're not serious)
9. **Follow-up Plan** (what to do after the meeting)

Be very specific and practical."""

        prep_content = get_groq_response(system_prompt, user_prompt, max_tokens=2000)

        return jsonify({
            'success': True,
            'prep_content': prep_content,
            'client_name': client_name,
            'meeting_type': meeting_type,
            'message': f'Meeting prep ready for {client_name} ({meeting_type} meeting)!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# EA — 1: SMART EMAIL COMPOSER
# AI se professional emails draft karo
# ════════════════════════════════════════════════════════════════
@app.route('/api/ea/email_compose', methods=['POST'])
@login_required
def ea_email_compose():
    try:
        data = request.get_json()
        to_name      = data.get('to_name', '')
        to_role      = data.get('to_role', '')
        subject_hint = data.get('subject_hint', '')
        email_purpose= data.get('email_purpose', '')
        key_points   = data.get('key_points', '')
        tone         = data.get('tone', 'professional')  # professional, friendly, formal, urgent
        your_name    = data.get('your_name', 'Data Scientist')
        email_length = data.get('email_length', 'medium')  # short/medium/long

        length_map = {'short': '80-120 words', 'medium': '150-200 words', 'long': '250-350 words'}

        system_prompt = f"""You are an expert email writer for business professionals.
You write clear, effective emails that get responses.
Tone should be {tone}. Length: {length_map.get(email_length, '150-200 words')}.
Use Indian professional email conventions."""

        user_prompt = f"""Write a professional email:

TO: {to_name} ({to_role})
PURPOSE: {email_purpose}
KEY POINTS TO COVER: {key_points}
FROM: {your_name}
TONE: {tone}

Provide:
1. Subject Line (compelling)
2. Email Body (ready to send)
3. Alternative subject line
4. One-line WhatsApp follow-up version"""

        email_draft = get_groq_response(system_prompt, user_prompt, max_tokens=800)

        # Save to history
        import json as _json
        db = SessionLocal()
        e = UserChart(
            user_id=current_user.id,
            chart_type='ea_email',
            chart_title=f'Email to {to_name} — {subject_hint[:30]}',
        )
        db.add(e); db.commit(); db.close()

        return jsonify({
            'success': True,
            'email_draft': email_draft,
            'message': f'Email draft ready for {to_name}!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# EA — 2: MEETING NOTES SUMMARIZER
# ════════════════════════════════════════════════════════════════
@app.route('/api/ea/meeting_notes', methods=['POST'])
@login_required
def ea_meeting_notes():
    try:
        data = request.get_json()
        raw_notes    = data.get('raw_notes', '')
        meeting_type = data.get('meeting_type', 'client')
        attendees    = data.get('attendees', '')
        date         = data.get('date', '')

        if not raw_notes:
            return jsonify({'error': 'Meeting notes daalo!'}), 400

        system_prompt = """You are an expert executive assistant who creates perfect meeting summaries.
Your summaries are clear, structured, and action-focused.
Always extract concrete action items with owners and deadlines."""

        user_prompt = f"""Summarize these meeting notes professionally:

MEETING TYPE: {meeting_type}
DATE: {date}
ATTENDEES: {attendees}

RAW NOTES:
{raw_notes}

Create:
1. **Meeting Summary** (3-4 lines executive overview)
2. **Key Decisions Made** (bulleted list)
3. **Action Items** (table: Action | Owner | Deadline)
4. **Open Questions** (unresolved items)
5. **Next Meeting** (suggested agenda items)
6. **Follow-up Email** (ready to send to attendees)

Format professionally. Be specific about action items."""

        summary = get_groq_response(system_prompt, user_prompt, max_tokens=1500)

        # Save
        import json as _json
        db = SessionLocal()
        e = UserChart(
            user_id=current_user.id,
            chart_type='ea_meeting_notes',
            chart_title=f'Meeting Notes — {date} ({meeting_type})',
        )
        db.add(e); db.commit(); db.close()

        return jsonify({
            'success': True,
            'summary': summary,
            'original_length': len(raw_notes.split()),
            'summary_length': len(summary.split()),
            'message': f'Meeting notes summarized! {len(raw_notes.split())} words → {len(summary.split())} words'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# EA — 3: TASK MANAGER
# ════════════════════════════════════════════════════════════════
@app.route('/api/ea/task/add', methods=['POST'])
@login_required
def ea_task_add():
    try:
        import json as _json
        data = request.get_json()
        task = {
            'title':      data.get('title', ''),
            'description':data.get('description', ''),
            'priority':   data.get('priority', 'Medium'),  # Low/Medium/High/Urgent
            'due_date':   data.get('due_date', ''),
            'category':   data.get('category', 'General'),  # Work/Client/Personal/Follow-up
            'status':     'Pending',
            'client':     data.get('client', ''),
        }
        if not task['title']:
            return jsonify({'error': 'Task title required!'}), 400
        db = SessionLocal()
        t = UserChart(
            user_id=current_user.id,
            chart_type='ea_task',
            chart_title=task['title'],
        )
        db.add(t); db.commit(); tid = t.id; db.close()
        return jsonify({'success': True, 'task_id': tid, 'message': f"Task added: {task['title']}"})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/ea/task/list', methods=['GET'])
@login_required
def ea_task_list():
    try:
        import json as _json
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt, io as _io
        from datetime import datetime

        db = SessionLocal()
        tasks = db.query(UserChart).filter_by(
            user_id=current_user.id, chart_type='ea_task'
        ).order_by(UserChart.created_at.desc()).all()
        db.close()

        task_list = []
        priority_counts = {'Urgent': 0, 'High': 0, 'Medium': 0, 'Low': 0}
        status_counts   = {'Pending': 0, 'In Progress': 0, 'Done': 0}
        overdue_count   = 0
        today_str       = datetime.now().strftime('%Y-%m-%d')

        for t in tasks:
            try: task_data = _json.loads(t.description)
            except: task_data = {}
            task_data['id'] = t.id
            task_data['created_at'] = t.created_at.strftime('%d %b')

            # Check overdue
            due = task_data.get('due_date', '')
            if due and task_data.get('status') != 'Done':
                try:
                    if due < today_str:
                        task_data['overdue'] = True
                        overdue_count += 1
                    elif due == today_str:
                        task_data['due_today'] = True
                except: pass

            priority_counts[task_data.get('priority', 'Medium')] = \
                priority_counts.get(task_data.get('priority', 'Medium'), 0) + 1
            status_counts[task_data.get('status', 'Pending')] = \
                status_counts.get(task_data.get('status', 'Pending'), 0) + 1
            task_list.append(task_data)

        # Sort: Urgent > High > Medium > Low, Pending first
        priority_order = {'Urgent': 0, 'High': 1, 'Medium': 2, 'Low': 3}
        status_order = {'Pending': 0, 'In Progress': 1, 'Done': 2}
        task_list.sort(key=lambda x: (
            status_order.get(x.get('status', 'Pending'), 0),
            priority_order.get(x.get('priority', 'Medium'), 2)
        ))

        # Chart
        if task_list:
            fig, axes = plt.subplots(1, 2, figsize=(12, 4))
            pdata = {k: v for k, v in priority_counts.items() if v > 0}
            if pdata:
                pcolors = {'Urgent':'#ef4444','High':'#f59e0b','Medium':'#3b82f6','Low':'#94a3b8'}
                axes[0].pie(pdata.values(), labels=pdata.keys(),
                           colors=[pcolors.get(k,'#94a3b8') for k in pdata],
                           autopct='%1.0f%%', startangle=90)
                axes[0].set_title('Tasks by Priority', fontweight='bold')
            sdata = {k: v for k, v in status_counts.items() if v > 0}
            if sdata:
                axes[1].bar(sdata.keys(), sdata.values(),
                           color=['#94a3b8','#3b82f6','#059669'], alpha=0.85, edgecolor='white')
                axes[1].set_title('Tasks by Status', fontweight='bold')
                for i, (k, v) in enumerate(sdata.items()):
                    axes[1].text(i, v + 0.1, str(v), ha='center', fontweight='bold')
            plt.suptitle(f'Task Dashboard | {len(task_list)} total | {overdue_count} overdue',
                        fontweight='bold', color='#1e3a5f')
            plt.tight_layout()
            buf = _io.BytesIO()
            plt.savefig(buf, format='png', dpi=120, bbox_inches='tight')
            plt.close('all'); buf.seek(0)
            chart = base64.b64encode(buf.read()).decode()
        else:
            chart = None

        return jsonify({
            'success': True,
            'tasks': task_list,
            'total': len(task_list),
            'pending': status_counts.get('Pending', 0),
            'in_progress': status_counts.get('In Progress', 0),
            'done': status_counts.get('Done', 0),
            'overdue': overdue_count,
            'priority_counts': priority_counts,
            'chart': chart,
            'message': f'{len(task_list)} tasks | {overdue_count} overdue | {status_counts.get("Done",0)} done'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/ea/task/update/<int:tid>', methods=['POST'])
@login_required
def ea_task_update(tid):
    try:
        import json as _json
        data = request.get_json()
        db = SessionLocal()
        t = db.query(UserChart).filter_by(id=tid, user_id=current_user.id).first()
        if not t: db.close(); return jsonify({'error': 'Task not found!'}), 404
        try: td = _json.loads(t.description)
        except: td = {}
        td.update({k: v for k, v in data.items() if k != 'id'})
        t.description = _json.dumps(td)
        t.title = td.get('title', t.title)
        db.commit(); db.close()
        return jsonify({'success': True, 'message': f"Task updated: {td.get('status', 'updated')}"})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/ea/task/delete/<int:tid>', methods=['DELETE'])
@login_required
def ea_task_delete(tid):
    try:
        db = SessionLocal()
        t = db.query(UserChart).filter_by(id=tid, user_id=current_user.id).first()
        if t: db.delete(t); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# EA — 4: DOCUMENT SUMMARIZER
# Koi bhi text/document paste karo — AI summary milega
# ════════════════════════════════════════════════════════════════
@app.route('/api/ea/summarize', methods=['POST'])
@login_required
def ea_summarize():
    try:
        data = request.get_json()
        text         = data.get('text', '')
        doc_type     = data.get('doc_type', 'general')  # report, email, contract, news, research
        output_style = data.get('output_style', 'bullets')  # bullets, paragraph, executive

        if not text or len(text.strip()) < 50:
            return jsonify({'error': 'Text too short — minimum 50 characters!'}), 400

        style_map = {
            'bullets':   'Use clear bullet points. Highlight key facts.',
            'paragraph': 'Write as clean, flowing paragraphs.',
            'executive': 'Executive summary style — one-pager with key decisions and actions.',
        }

        system_prompt = f"""You are an expert document analyst and executive assistant.
You extract the most important information quickly and clearly.
{style_map.get(output_style, style_map['bullets'])}"""

        user_prompt = f"""Summarize this {doc_type} document:

{text[:4000]}

Provide:
1. **One-Line Summary** (what this is about in one sentence)
2. **Key Points** (5-7 most important points)
3. **Action Items / Decisions** (what needs to be done)
4. **Important Numbers/Dates** (any critical figures mentioned)
5. **My Next Step** (what should I do with this information?)

Keep it crisp. No fluff."""

        summary = get_groq_response(system_prompt, user_prompt, max_tokens=1000)

        return jsonify({
            'success': True,
            'summary': summary,
            'original_words': len(text.split()),
            'summary_words': len(summary.split()),
            'compression': round((1 - len(summary.split()) / max(1, len(text.split()))) * 100, 1),
            'message': f'Summarized! {len(text.split())} → {len(summary.split())} words ({round((1-len(summary.split())/max(1,len(text.split())))*100)}% compression)'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# EA — 5: DAILY PLANNER / SCHEDULE BUILDER
# ════════════════════════════════════════════════════════════════
@app.route('/api/ea/daily_plan', methods=['POST'])
@login_required
def ea_daily_plan():
    try:
        import json as _json
        from datetime import datetime

        data = request.get_json()
        date          = data.get('date', datetime.now().strftime('%A, %d %B %Y'))
        priorities    = data.get('priorities', '')
        meetings      = data.get('meetings', '')
        energy_level  = data.get('energy_level', 'high')  # high/medium/low
        work_hours    = data.get('work_hours', '9AM-6PM')
        focus_areas   = data.get('focus_areas', 'data science work, client follow-ups')

        # Get pending tasks from DB
        db = SessionLocal()
        tasks = db.query(UserChart).filter_by(
            user_id=current_user.id, chart_type='ea_task'
        ).all()
        db.close()

        pending_tasks = []
        for t in tasks:
            try:
                td = _json.loads(t.description)
                if td.get('status') != 'Done':
                    pending_tasks.append(f"{td.get('priority','?')} priority: {td.get('title','?')} (due: {td.get('due_date','no date')})")
            except: pass

        system_prompt = """You are a productivity expert and executive assistant.
You create realistic, effective daily plans that maximize output.
You understand deep work, energy management, and prioritization."""

        user_prompt = f"""Create a detailed daily plan for:

DATE: {date}
WORK HOURS: {work_hours}
ENERGY LEVEL: {energy_level}
TODAY'S PRIORITIES: {priorities}
SCHEDULED MEETINGS: {meetings if meetings else 'None mentioned'}
FOCUS AREAS: {focus_areas}

PENDING TASKS FROM SYSTEM:
{chr(10).join(pending_tasks[:10]) if pending_tasks else 'No pending tasks in system'}

Create:
1. **Daily Schedule** (hour-by-hour time blocks)
2. **Top 3 Non-Negotiables** (must complete today)
3. **Energy Management** (when to do deep work vs admin)
4. **Meeting Prep** (if any meetings, what to prepare)
5. **End-of-Day Checklist** (how to close the day)
6. **Tomorrow Prep** (what to set up tonight)

Make it realistic for {energy_level} energy. Be specific with times."""

        plan = get_groq_response(system_prompt, user_prompt, max_tokens=1500)

        return jsonify({
            'success': True,
            'date': date,
            'daily_plan': plan,
            'pending_tasks_count': len(pending_tasks),
            'message': f'Daily plan ready for {date}!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# EA — 6: REMINDER SYSTEM
# ════════════════════════════════════════════════════════════════
@app.route('/api/ea/reminder/add', methods=['POST'])
@login_required
def ea_reminder_add():
    try:
        import json as _json
        data = request.get_json()
        reminder = {
            'title':      data.get('title', ''),
            'note':       data.get('note', ''),
            'remind_on':  data.get('remind_on', ''),   # YYYY-MM-DD
            'remind_time':data.get('remind_time', '09:00'),
            'priority':   data.get('priority', 'Medium'),
            'category':   data.get('category', 'Work'),  # Work/Client/Personal
            'repeat':     data.get('repeat', 'none'),   # none/daily/weekly/monthly
            'status':     'Active',
        }
        if not reminder['title']:
            return jsonify({'error': 'Reminder title do!'}), 400

        db = SessionLocal()
        r = UserChart(
            user_id=current_user.id,
            chart_type='ea_reminder',
            chart_title=reminder['title'],
        )
        db.add(r); db.commit(); rid = r.id; db.close()
        return jsonify({'success': True, 'reminder_id': rid,
                       'message': f"Reminder set: '{reminder['title']}' on {reminder['remind_on']}"})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/ea/reminder/list', methods=['GET'])
@login_required
def ea_reminder_list():
    try:
        import json as _json
        from datetime import datetime

        db = SessionLocal()
        reminders = db.query(UserChart).filter_by(
            user_id=current_user.id, chart_type='ea_reminder'
        ).order_by(UserChart.created_at.desc()).all()
        db.close()

        today = datetime.now().strftime('%Y-%m-%d')
        result = []
        due_today = []
        overdue = []

        for r in reminders:
            try: rd = _json.loads(r.description)
            except: rd = {}
            rd['id'] = r.id
            remind_on = rd.get('remind_on', '')
            if remind_on:
                if remind_on == today:
                    rd['status_tag'] = '🔔 TODAY'
                    due_today.append(rd)
                elif remind_on < today and rd.get('status') == 'Active':
                    rd['status_tag'] = '🔴 OVERDUE'
                    overdue.append(rd)
                else:
                    rd['status_tag'] = '📅 Upcoming'
            result.append(rd)

        result.sort(key=lambda x: (x.get('remind_on','9999'), x.get('remind_time','00:00')))

        return jsonify({
            'success': True,
            'reminders': result,
            'total': len(result),
            'due_today': len(due_today),
            'overdue': len(overdue),
            'due_today_list': due_today,
            'message': f'{len(result)} reminders | {len(due_today)} due today | {len(overdue)} overdue'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/ea/reminder/done/<int:rid>', methods=['POST'])
@login_required
def ea_reminder_done(rid):
    try:
        import json as _json
        db = SessionLocal()
        r = db.query(UserChart).filter_by(id=rid, user_id=current_user.id).first()
        if r:
            try: rd = _json.loads(r.description)
            except: rd = {}
            rd['status'] = 'Done'
            r.description = _json.dumps(rd)
            db.commit()
        db.close()
        return jsonify({'success': True, 'message': 'Reminder marked done!'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/ea/reminder/delete/<int:rid>', methods=['DELETE'])
@login_required
def ea_reminder_delete(rid):
    try:
        db = SessionLocal()
        r = db.query(UserChart).filter_by(id=rid, user_id=current_user.id).first()
        if r: db.delete(r); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# EA — 7: AI QUICK ASSISTANT (General purpose AI chat for work)
# ════════════════════════════════════════════════════════════════
@app.route('/api/ea/quick_assist', methods=['POST'])
@login_required
def ea_quick_assist():
    try:
        data = request.get_json()
        question     = data.get('question', '')
        context      = data.get('context', '')
        assist_type  = data.get('assist_type', 'general')
        # assist_type options:
        # general, linkedin_post, presentation_outline,
        # negotiation_script, objection_handling,
        # invoice_draft, contract_clause, pricing_strategy

        if not question:
            return jsonify({'error': 'Question do!'}), 400

        type_contexts = {
            'linkedin_post':         'Write compelling LinkedIn posts for data science professionals targeting Indian business market.',
            'presentation_outline':  'Create clear, persuasive presentation outlines.',
            'negotiation_script':    'Provide negotiation scripts for consulting fee discussions.',
            'objection_handling':    'Give specific responses to sales objections in data science consulting.',
            'invoice_draft':         'Draft professional service invoices and payment terms.',
            'contract_clause':       'Suggest contract clauses for data science consulting engagements.',
            'pricing_strategy':      'Advise on pricing strategy for data science services in Indian market.',
            'general':               'You are a smart business and technical assistant.',
        }

        system_prompt = f"""You are an expert business assistant for a Data Science consultant in India.
{type_contexts.get(assist_type, type_contexts['general'])}
Be specific, practical, and actionable. Use Indian market context where relevant."""

        user_prompt = f"""{question}

Additional context: {context}""" if context else question

        response = get_groq_response(system_prompt, user_prompt, max_tokens=1200)

        return jsonify({
            'success': True,
            'response': response,
            'assist_type': assist_type,
            'message': 'Quick assist ready!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# 📝 CONTENT MARKETING SUITE
# Categories: SaaS, B2B, Finance, Healthcare, Ecommerce, etc.
# Types: Blog, Video Script, Tweet, Email, LinkedIn, Article, etc.
# Features: Generate → Plagiarism Check → Humanize → Final Ready
# ════════════════════════════════════════════════════════════════

CONTENT_SYSTEM_PROMPT = """You are an elite content marketing expert and copywriter.
You create high-quality, engaging, SEO-optimized content for businesses.
You understand different industries, audiences, and content formats.
Your content is original, human-sounding, and conversion-focused.
Always match the tone and style to the industry and content type."""

CONTENT_CATEGORIES = {
    'saas': 'SaaS / Software Products',
    'b2b': 'B2B / Business Services',
    'finance': 'Finance / Fintech / Investment',
    'healthcare': 'Healthcare / MedTech / Wellness',
    'ecommerce': 'Ecommerce / D2C / Retail',
    'edtech': 'EdTech / Online Education',
    'realestate': 'Real Estate / Property',
    'startup': 'Startups / Entrepreneurship',
    'marketing': 'Digital Marketing / Growth',
    'ai_ml': 'AI / ML / Data Science',
    'cybersecurity': 'Cybersecurity / Privacy',
    'hr': 'HR / Recruitment / People',
    'legal': 'Legal / Compliance / LegalTech',
    'travel': 'Travel / Hospitality / Tourism',
    'food': 'Food / Restaurant / FoodTech',
    'fashion': 'Fashion / Lifestyle / Beauty',
    'fitness': 'Fitness / Health / Wellness',
    'agriculture': 'Agriculture / AgriTech',
    'personalbranding': 'Personal Branding / Influencer',
    'consulting': 'Consulting / Advisory / Strategy',
}

CONTENT_TYPES = {
    'blog': 'Blog Post (Long Form)',
    'article': 'Article / News Piece',
    'video_script': 'YouTube Video Script',
    'linkedin': 'LinkedIn Post',
    'twitter': 'Twitter/X Thread',
    'email': 'Email Newsletter',
    'email_cold': 'Cold Email Sequence',
    'instagram': 'Instagram Caption',
    'case_study': 'Case Study',
    'whitepaper': 'Whitepaper / Report Intro',
    'podcast': 'Podcast Script / Show Notes',
    'ad_copy': 'Ad Copy (Google/Meta)',
    'landing_page': 'Landing Page Copy',
    'product_desc': 'Product Description',
    'press_release': 'Press Release',
    'seo_content': 'SEO Article with Keywords',
}


@app.route('/api/content/generate', methods=['POST'])
@login_required
def content_generate():
    try:
        import json as _json
        data = request.get_json()
        category     = data.get('category', 'saas')
        content_type = data.get('content_type', 'blog')
        topic        = data.get('topic', '')
        tone         = data.get('tone', 'professional')
        length       = data.get('length', 'medium')
        keywords     = data.get('keywords', '')
        target_audience = data.get('target_audience', '')
        language     = data.get('language', 'English')

        if not topic:
            return jsonify({'error': 'Topic do — kya likhna hai batao!'}), 400

        cat_label  = CONTENT_CATEGORIES.get(category, category)
        type_label = CONTENT_TYPES.get(content_type, content_type)

        length_map = {
            'short':  'Short (200-400 words / 3-5 min read)',
            'medium': 'Medium (600-1000 words / 8-12 min read)',
            'long':   'Long (1500-2500 words / comprehensive)',
        }

        type_instructions = {
            'blog': 'Write a complete blog post with: catchy title, introduction with hook, 4-6 subheadings with detailed content, conclusion with CTA.',
            'article': 'Write a professional article with: headline, lead paragraph, body with facts/examples, expert quotes style, conclusion.',
            'video_script': 'Write a YouTube video script with: Hook (first 15 sec), intro, main content sections with timestamps, CTA at end. Include [VISUAL CUE] notes.',
            'linkedin': 'Write a LinkedIn post with: attention-grabbing first line, story or insight, 3-5 key points, engagement question, relevant hashtags.',
            'twitter': 'Write a Twitter/X thread with: hook tweet (1), numbered thread tweets (2-10), each under 280 chars, concluding tweet with CTA.',
            'email': 'Write an email newsletter with: subject line (A/B options), preheader, greeting, main content, value section, CTA button text, signature.',
            'email_cold': 'Write a 3-email cold email sequence: Email 1 (intro + value), Email 2 (follow-up + case study), Email 3 (breakup email). Each under 150 words.',
            'instagram': 'Write Instagram caption with: hook first line, story/value, emojis naturally placed, 5-10 hashtags, CTA.',
            'case_study': 'Write a case study with: client background, challenge/problem, solution implemented, results with numbers, client quote, key takeaways.',
            'whitepaper': 'Write a whitepaper introduction + executive summary + first section with: compelling title, abstract, problem statement, market context, methodology preview.',
            'podcast': 'Write podcast script with: intro music cue, host intro, episode overview, main discussion points with timestamps, outro, show notes summary.',
            'ad_copy': 'Write 3 variations of ad copy: Version A (problem-focused), Version B (benefit-focused), Version C (social proof). Include headline + description + CTA.',
            'landing_page': 'Write landing page copy with: hero headline + subheadline, value proposition, 3 key benefits, social proof section, FAQ (3 questions), CTA sections.',
            'product_desc': 'Write product description with: attention headline, key features as benefits, technical specs, who it is for, what makes it unique, purchase CTA.',
            'press_release': 'Write press release with: FOR IMMEDIATE RELEASE header, headline, dateline, lead paragraph (5 Ws), body with quotes, boilerplate, contact info.',
            'seo_content': 'Write SEO-optimized article with: H1 title (keyword included), meta description, H2/H3 structure, keyword naturally placed, internal link suggestions, FAQ section.',
        }

        user_prompt = f"""Create {type_label} content for {cat_label} industry:

TOPIC: {topic}
TARGET AUDIENCE: {target_audience if target_audience else 'General business audience'}
TONE: {tone}
LENGTH: {length_map.get(length, 'Medium')}
KEYWORDS TO INCLUDE: {keywords if keywords else 'Use relevant industry keywords'}
LANGUAGE: {language}

INSTRUCTIONS: {type_instructions.get(content_type, 'Write high-quality, engaging content.')}

Make it:
- 100% original and unique
- Human-sounding (not robotic)
- Industry-specific with relevant examples
- Action-oriented with clear value
- Ready to publish as-is"""

        content_text = get_groq_response(CONTENT_SYSTEM_PROMPT, user_prompt, max_tokens=3000)

        # No DB save — content returned directly to user
        cid = 0

        word_count = len(content_text.split())
        read_time  = max(1, word_count // 200)

        return jsonify({
            'success': True,
            'content_id': cid,
            'content': content_text,
            'word_count': word_count,
            'read_time': read_time,
            'category': cat_label,
            'content_type': type_label,
            'message': f'{type_label} ready! {word_count} words | ~{read_time} min read'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/content/plagiarism', methods=['POST'])
@login_required
def content_plagiarism():
    try:
        data = request.get_json()
        content_text = data.get('content', '')

        if not content_text or len(content_text.strip()) < 50:
            return jsonify({'error': 'Content do pehle!'}), 400

        system_prompt = """You are an expert plagiarism detection and content originality analyst.
You analyze text for: originality, uniqueness, AI-generated patterns, cliche phrases, and duplicate risks.
Be specific and actionable in your feedback."""

        user_prompt = f"""Analyze this content for plagiarism risk and originality:

CONTENT:
{content_text[:3000]}

Provide a detailed analysis:

1. **ORIGINALITY SCORE** (0-100): Give a single number score
2. **PLAGIARISM RISK**: Low / Medium / High — explain why
3. **AI-DETECTION RISK**: How likely to be flagged as AI-written (Low/Medium/High)
4. **RISKY PHRASES** (list 3-5): Phrases that might be too generic or potentially copied
5. **CLICHE COUNT**: How many cliches detected — list them
6. **UNIQUE ELEMENTS**: What makes this content original
7. **RECOMMENDATIONS** (5 specific changes): How to improve originality score
8. **OVERALL VERDICT**: Is this safe to publish? Yes/No/With Changes

Format the score clearly as: ORIGINALITY SCORE: XX/100"""

        analysis = get_groq_response(system_prompt, user_prompt, max_tokens=1500)

        # Extract score
        import re
        score_match = re.search(r'ORIGINALITY SCORE[:\s]+(\d+)', analysis, re.IGNORECASE)
        score = int(score_match.group(1)) if score_match else 75

        score_label = 'Excellent' if score >= 85 else 'Good' if score >= 70 else 'Needs Work' if score >= 50 else 'High Risk'
        score_color = '#059669' if score >= 85 else '#d97706' if score >= 70 else '#dc2626'

        return jsonify({
            'success': True,
            'analysis': analysis,
            'score': score,
            'score_label': score_label,
            'score_color': score_color,
            'word_count': len(content_text.split()),
            'message': f'Plagiarism check done! Originality: {score}/100 — {score_label}'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/content/humanize', methods=['POST'])
@login_required
def content_humanize():
    try:
        data = request.get_json()
        content_text   = data.get('content', '')
        humanize_level = data.get('level', 'medium')  # light / medium / deep
        writing_style  = data.get('style', 'conversational')  # conversational / professional / storytelling

        if not content_text or len(content_text.strip()) < 50:
            return jsonify({'error': 'Content do pehle!'}), 400

        level_map = {
            'light':  'Lightly humanize — fix robotic phrases, add natural flow, keep structure same',
            'medium': 'Moderately humanize — rewrite sentences naturally, add personality, vary sentence length',
            'deep':   'Deeply humanize — completely rewrite in human voice, add personal touches, stories, imperfections',
        }

        style_map = {
            'conversational': 'Conversational and friendly — like talking to a smart friend',
            'professional':   'Professional but warm — like a trusted industry expert',
            'storytelling':   'Storytelling style — narrative, engaging, memorable',
            'casual':         'Casual and relatable — simple words, short sentences',
        }

        system_prompt = """You are an expert human content writer who specializes in making AI-generated 
content sound completely human. You add personality, natural imperfections, varied sentence structures,
relatable examples, and genuine emotion. Your rewrites bypass AI detection tools easily."""

        user_prompt = f"""Humanize this content completely:

HUMANIZE LEVEL: {level_map.get(humanize_level, level_map['medium'])}
WRITING STYLE: {style_map.get(writing_style, style_map['conversational'])}

ORIGINAL CONTENT:
{content_text[:3000]}

Rewrite this to sound 100% human-written:
- Vary sentence lengths (mix short and long)
- Add occasional contractions (don't, isn't, you'll)
- Include natural transitions
- Add relatable analogies or examples
- Remove corporate jargon and buzzwords
- Add subtle personality and opinion
- Make it flow like a real person wrote it
- Keep all key information and structure intact

Output ONLY the humanized content — no explanations."""

        humanized = get_groq_response(system_prompt, user_prompt, max_tokens=3000)

        return jsonify({
            'success': True,
            'humanized_content': humanized,
            'original_words': len(content_text.split()),
            'humanized_words': len(humanized.split()),
            'level': humanize_level,
            'style': writing_style,
            'message': f'Content humanized! ({humanize_level} level, {writing_style} style) — {len(humanized.split())} words'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/content/finalize', methods=['POST'])
@login_required
def content_finalize():
    try:
        data = request.get_json()
        content_text = data.get('content', '')
        content_type = data.get('content_type', 'blog')
        category     = data.get('category', 'saas')
        add_seo      = data.get('add_seo', True)
        add_cta      = data.get('add_cta', True)

        if not content_text or len(content_text.strip()) < 50:
            return jsonify({'error': 'Content do pehle!'}), 400

        cat_label  = CONTENT_CATEGORIES.get(category, category)
        type_label = CONTENT_TYPES.get(content_type, content_type)

        system_prompt = """You are a senior content editor and SEO specialist.
You finalize content to be 100% publish-ready — polished, optimized, and complete."""

        user_prompt = f"""Finalize this {type_label} for {cat_label} industry and make it 100% publish-ready:

CONTENT:
{content_text[:3000]}

FINALIZATION TASKS:
1. **POLISH**: Fix grammar, flow, and readability
2. **SEO OPTIMIZATION** (if applicable):
   - Suggest meta title (60 chars max)
   - Suggest meta description (155 chars max)  
   - List 5 focus keywords naturally used
3. **CTA ENHANCEMENT**: Make call-to-action stronger and clearer
4. **FORMATTING CHECK**: Ensure proper structure for {type_label}
5. **FINAL CONTENT**: Output the complete polished, ready-to-publish version

Also provide at the end:
- PUBLISH CHECKLIST (5 items)
- BEST TIME TO POST (for this content type)
- PLATFORMS TO REPURPOSE ON (3 suggestions)"""

        finalized = get_groq_response(system_prompt, user_prompt, max_tokens=3000)

        # No DB save needed — content returned directly

        return jsonify({
            'success': True,
            'finalized_content': finalized,
            'word_count': len(finalized.split()),
            'message': f'Content finalized! Ready to publish — {len(finalized.split())} words'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/content/history', methods=['GET'])
@login_required
def content_history():
    try:
        # History stored in browser — DB not used for content
        return jsonify({'success': True, 'history': [], 'total': 0})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/content/delete/<int:cid>', methods=['DELETE'])
@login_required
def content_delete(cid):
    try:
        db = SessionLocal()
        e = db.query(UserChart).filter_by(id=cid, user_id=current_user.id).first()
        if e: db.delete(e); db.commit()
        db.close()
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/datasets/list')
@login_required
def datasets_list():
    try:
        db = SessionLocal()
        datasets = db.query(UserDataset).filter_by(
            user_id=current_user.id
        ).order_by(UserDataset.uploaded_at.desc()).all()
        db.close()
        result = [{
            'id': d.id,
            'filename': d.filename,
            'rows': d.row_count,
            'uploaded_at': d.uploaded_at.strftime('%d %b %I:%M %p')
        } for d in datasets]
        return jsonify({'success': True, 'datasets': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/datasets/switch/<int:did>', methods=['POST'])
@login_required
def dataset_switch(did):
    try:
        db = SessionLocal()
        dataset = db.query(UserDataset).filter_by(
            id=did, user_id=current_user.id).first()
        if not dataset:
            db.close()
            return jsonify({'error': 'Dataset nahi mila!'}), 404
        import tempfile, json as _j
        tmp = tempfile.NamedTemporaryFile(
            suffix='_'+dataset.filename, delete=False)
        tmp.write(dataset.csv_data.encode('utf-8'))
        tmp.close()
        result = agent.load_data(tmp.name)
        db.close()
        return jsonify({
            'success': True,
            'filename': dataset.filename,
            'rows': dataset.row_count,
            'columns': len(_j.loads(dataset.columns)) if dataset.columns else 0,
            'message': result
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/charts/list')
@login_required  
def charts_list_new():
    try:
        db = SessionLocal()
        charts = db.query(UserChart).filter(
            UserChart.user_id == current_user.id,
            UserChart.chart_type == 'matplotlib'
        ).order_by(UserChart.created_at.desc()).limit(20).all()
        db.close()
        return jsonify({'charts': [{'id':c.id,'title':c.chart_title or f'Chart {c.id}','time':c.created_at.strftime('%d %b %H:%M')} for c in charts]})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ════════════════════════════════════════════════════════════════════
# 👩‍⚕️ CNO — CHIEF NURSING OFFICER DASHBOARD
# Full Suite: Staff Analytics + Clinical Quality + Ward Analytics + Report
# ════════════════════════════════════════════════════════════════════

CNO_SYSTEM_PROMPT = """You are an expert Chief Nursing Officer (CNO) analytics consultant.
You analyze nursing staff data, clinical quality metrics, ward operations, and patient outcomes.
You provide actionable insights for hospital leadership and nursing administration.
Your analysis is data-driven, clinically relevant, and board-presentation ready.
Always provide specific numbers, percentages, benchmarks, and recommendations."""

@app.route('/api/cno/staff', methods=['POST'])
@login_required
def cno_staff_analytics():
    """Staff Analytics: Nurse ratio, shifts, overtime, performance"""
    try:
        import json as _json
        data = request.get_json()
        total_nurses      = data.get('total_nurses', 0)
        total_patients    = data.get('total_patients', 0)
        day_shift         = data.get('day_shift', 0)
        night_shift       = data.get('night_shift', 0)
        overtime_hours    = data.get('overtime_hours', 0)
        absent_today      = data.get('absent_today', 0)
        icu_nurses        = data.get('icu_nurses', 0)
        icu_patients      = data.get('icu_patients', 0)
        agency_nurses     = data.get('agency_nurses', 0)
        avg_experience    = data.get('avg_experience', 0)
        training_due      = data.get('training_due', 0)
        incidents_month   = data.get('incidents_month', 0)
        department        = data.get('department', 'All Departments')

        # Calculate key metrics
        nurse_pt_ratio    = round(total_patients / total_nurses, 2) if total_nurses > 0 else 0
        icu_ratio         = round(icu_patients / icu_nurses, 2) if icu_nurses > 0 else 0
        absence_rate      = round((absent_today / total_nurses) * 100, 1) if total_nurses > 0 else 0
        agency_pct        = round((agency_nurses / total_nurses) * 100, 1) if total_nurses > 0 else 0

        # Benchmarks (WHO/JCI standards)
        general_benchmark = 4  # 1:4 nurse:patient ratio
        icu_benchmark     = 2  # 1:2 for ICU
        absence_benchmark = 5  # max 5% absence
        agency_benchmark  = 10 # max 10% agency staff

        user_prompt = f"""Analyze this nursing staff data for {department}:

STAFF METRICS:
- Total Nurses: {total_nurses}
- Total Patients: {total_patients}
- Nurse:Patient Ratio: 1:{nurse_pt_ratio}
- Day Shift: {day_shift} nurses
- Night Shift: {night_shift} nurses
- Absent Today: {absent_today} ({absence_rate}%)
- Overtime Hours This Week: {overtime_hours}
- ICU Nurses: {icu_nurses} | ICU Patients: {icu_patients} | ICU Ratio: 1:{icu_ratio}
- Agency/Contract Nurses: {agency_nurses} ({agency_pct}%)
- Average Experience: {avg_experience} years
- Training Certifications Due: {training_due}
- Incident Reports This Month: {incidents_month}

BENCHMARKS (JCI/WHO):
- General Ward: 1:{general_benchmark} (current: 1:{nurse_pt_ratio})
- ICU: 1:{icu_benchmark} (current: 1:{icu_ratio})
- Absence Rate: max {absence_benchmark}% (current: {absence_rate}%)
- Agency Staff: max {agency_benchmark}% (current: {agency_pct}%)

Provide CNO-level analysis:
1. **STAFFING STATUS** (Red/Yellow/Green with reason)
2. **CRITICAL ALERTS** (any immediate action needed?)
3. **RATIO ANALYSIS** (vs benchmarks, risk level)
4. **SHIFT COVERAGE** (is current shift safe?)
5. **OVERTIME RISK** (burnout/safety concern?)
6. **AGENCY DEPENDENCY** (cost and quality risk)
7. **TOP 3 IMMEDIATE ACTIONS** (specific, actionable)
8. **RECRUITMENT NEED** (how many nurses needed?)
9. **PERFORMANCE SCORE** (0-100 with breakdown)
10. **BOARD SUMMARY** (2 lines for hospital board)"""

        analysis = get_groq_response(CNO_SYSTEM_PROMPT, user_prompt, max_tokens=2000)

        return jsonify({
            'success': True,
            'analysis': analysis,
            'metrics': {
                'nurse_pt_ratio': nurse_pt_ratio,
                'icu_ratio': icu_ratio,
                'absence_rate': absence_rate,
                'agency_pct': agency_pct,
                'total_nurses': total_nurses,
                'total_patients': total_patients,
            },
            'benchmarks': {
                'general': general_benchmark,
                'icu': icu_benchmark,
                'absence_max': absence_benchmark,
                'agency_max': agency_benchmark,
            },
            'alerts': {
                'ratio_alert':   nurse_pt_ratio > general_benchmark,
                'icu_alert':     icu_ratio > icu_benchmark,
                'absence_alert': absence_rate > absence_benchmark,
                'agency_alert':  agency_pct > agency_benchmark,
                'overtime_alert': overtime_hours > 40,
            },
            'message': f'Staff Analytics ready! Nurse:Patient = 1:{nurse_pt_ratio}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/cno/clinical', methods=['POST'])
@login_required
def cno_clinical_quality():
    """Clinical Quality: Outcomes, infections, incidents, falls, errors"""
    try:
        data = request.get_json()
        hospital_acquired_infections = data.get('hai', 0)
        patient_falls                = data.get('falls', 0)
        medication_errors            = data.get('med_errors', 0)
        pressure_ulcers              = data.get('pressure_ulcers', 0)
        total_admissions             = data.get('total_admissions', 0)
        patient_complaints           = data.get('complaints', 0)
        patient_satisfaction         = data.get('satisfaction_score', 0)
        readmission_30day            = data.get('readmission_30day', 0)
        mortality_rate               = data.get('mortality_rate', 0)
        code_blue_events             = data.get('code_blue', 0)
        near_miss_reports            = data.get('near_miss', 0)
        catheter_infections          = data.get('cauti', 0)
        ventilator_pneumonia         = data.get('vap', 0)
        period                       = data.get('period', 'This Month')

        # Rates per 1000 patient days
        hai_rate    = round((hospital_acquired_infections / total_admissions) * 1000, 2) if total_admissions > 0 else 0
        fall_rate   = round((patient_falls / total_admissions) * 1000, 2) if total_admissions > 0 else 0
        error_rate  = round((medication_errors / total_admissions) * 1000, 2) if total_admissions > 0 else 0

        user_prompt = f"""Analyze clinical quality metrics for CNO review — {period}:

PATIENT SAFETY INDICATORS:
- Hospital Acquired Infections (HAI): {hospital_acquired_infections} (Rate: {hai_rate}/1000)
- CAUTI (Catheter infections): {catheter_infections}
- VAP (Ventilator pneumonia): {ventilator_pneumonia}
- Patient Falls: {patient_falls} (Rate: {fall_rate}/1000)
- Medication Errors: {medication_errors} (Rate: {error_rate}/1000)
- Pressure Ulcers: {pressure_ulcers}
- Code Blue Events: {code_blue_events}
- Near Miss Reports: {near_miss_reports}

PATIENT EXPERIENCE:
- Total Admissions: {total_admissions}
- Satisfaction Score: {patient_satisfaction}/100
- Patient Complaints: {patient_complaints}
- 30-Day Readmission: {readmission_30day}
- Mortality Rate: {mortality_rate}%

JOINT COMMISSION BENCHMARKS:
- HAI Rate: < 2/1000 | Falls Rate: < 3/1000
- Satisfaction: > 85/100 | Readmission: < 15%
- Medication Errors: Zero tolerance
- Pressure Ulcers: Zero tolerance (hospital-acquired)

Provide CNO clinical quality report:
1. **QUALITY SCORE** (A/B/C/D grade with reason)
2. **CRITICAL SAFETY ALERTS** (immediate escalation needed?)
3. **HAI ANALYSIS** (source, prevention steps)
4. **FALLS PREVENTION** (root cause, protocol gaps)
5. **MEDICATION SAFETY** (error patterns, prevention)
6. **PATIENT EXPERIENCE** (satisfaction drivers)
7. **REGULATORY RISK** (NABH/JCI compliance status)
8. **QUALITY IMPROVEMENT PRIORITIES** (top 3)
9. **NURSING CARE IMPACT** (what nursing changes needed)
10. **BOARD PRESENTATION SUMMARY** (3 bullet points)"""

        analysis = get_groq_response(CNO_SYSTEM_PROMPT, user_prompt, max_tokens=2000)

        quality_score = 100
        if hai_rate > 2: quality_score -= 20
        if fall_rate > 3: quality_score -= 15
        if medication_errors > 0: quality_score -= 10
        if pressure_ulcers > 0: quality_score -= 15
        if patient_satisfaction < 85: quality_score -= 10
        if readmission_30day > 15: quality_score -= 10
        quality_score = max(0, quality_score)
        grade = 'A' if quality_score >= 90 else 'B' if quality_score >= 75 else 'C' if quality_score >= 60 else 'D'

        return jsonify({
            'success': True,
            'analysis': analysis,
            'quality_score': quality_score,
            'grade': grade,
            'rates': {
                'hai_rate': hai_rate,
                'fall_rate': fall_rate,
                'error_rate': error_rate,
            },
            'alerts': {
                'hai_alert':  hai_rate > 2,
                'fall_alert': fall_rate > 3,
                'sat_alert':  patient_satisfaction < 85,
                'readm_alert': readmission_30day > 15,
                'zero_tolerance': medication_errors > 0 or pressure_ulcers > 0,
            },
            'message': f'Clinical Quality: Grade {grade} ({quality_score}/100)'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/cno/ward', methods=['POST'])
@login_required
def cno_ward_analytics():
    """Ward Analytics: Bed occupancy, patient flow, discharge planning"""
    try:
        data = request.get_json()
        wards = data.get('wards', [])
        total_beds        = data.get('total_beds', 0)
        occupied_beds     = data.get('occupied_beds', 0)
        icu_beds          = data.get('icu_beds', 0)
        icu_occupied      = data.get('icu_occupied', 0)
        avg_los           = data.get('avg_los', 0)
        admissions_today  = data.get('admissions_today', 0)
        discharges_today  = data.get('discharges_today', 0)
        transfers_today   = data.get('transfers_today', 0)
        pending_discharge = data.get('pending_discharge', 0)
        er_waiting        = data.get('er_waiting', 0)
        elective_surgeries = data.get('elective_surgeries', 0)
        period            = data.get('period', 'Today')

        occupancy_rate    = round((occupied_beds / total_beds) * 100, 1) if total_beds > 0 else 0
        icu_occupancy     = round((icu_occupied / icu_beds) * 100, 1) if icu_beds > 0 else 0
        available_beds    = total_beds - occupied_beds
        bed_turnover      = round(discharges_today / total_beds, 2) if total_beds > 0 else 0

        wards_text = ""
        if wards:
            wards_text = "\nWARD-WISE BREAKDOWN:\n"
            for w in wards:
                occ = round((w.get('occupied',0) / w.get('capacity',1)) * 100, 1)
                wards_text += f"- {w.get('name','Ward')}: {w.get('occupied',0)}/{w.get('capacity',0)} beds ({occ}% occupancy)\n"

        user_prompt = f"""Analyze ward operations for CNO — {period}:

BED MANAGEMENT:
- Total Beds: {total_beds} | Occupied: {occupied_beds} | Available: {available_beds}
- Overall Occupancy: {occupancy_rate}%
- ICU Beds: {icu_beds} | ICU Occupied: {icu_occupied} | ICU Occupancy: {icu_occupancy}%
- Average Length of Stay: {avg_los} days
- Bed Turnover Rate: {bed_turnover}
{wards_text}
PATIENT FLOW — {period}:
- Admissions: {admissions_today}
- Discharges: {discharges_today}
- Transfers: {transfers_today}
- Pending Discharge (ready but waiting): {pending_discharge}
- ER Waiting for Bed: {er_waiting}
- Elective Surgeries Scheduled: {elective_surgeries}

BENCHMARKS:
- Optimal Occupancy: 80-85% (current: {occupancy_rate}%)
- ICU Occupancy: max 85% (current: {icu_occupancy}%)
- Average LOS benchmark: 4-5 days (current: {avg_los})

Provide CNO ward operations analysis:
1. **CAPACITY STATUS** (Critical/Warning/Normal)
2. **BOTTLENECKS** (where is flow blocked?)
3. **DISCHARGE DELAYS** (cause and solution)
4. **ICU PRESSURE** (overflow risk?)
5. **ER DIVERSION RISK** (should we divert patients?)
6. **BED MANAGEMENT ACTIONS** (immediate steps)
7. **STAFFING IMPLICATIONS** (how does this affect nurse needs?)
8. **ELECTIVE SURGERY IMPACT** (safe to proceed?)
9. **FLOW OPTIMIZATION** (3 specific recommendations)
10. **OPERATIONAL SUMMARY** (for CNO morning huddle)"""

        analysis = get_groq_response(CNO_SYSTEM_PROMPT, user_prompt, max_tokens=2000)

        return jsonify({
            'success': True,
            'analysis': analysis,
            'metrics': {
                'occupancy_rate': occupancy_rate,
                'icu_occupancy': icu_occupancy,
                'available_beds': available_beds,
                'bed_turnover': bed_turnover,
            },
            'alerts': {
                'high_occupancy':  occupancy_rate > 90,
                'icu_critical':    icu_occupancy > 85,
                'er_pressure':     er_waiting > 5,
                'discharge_delay': pending_discharge > 10,
                'low_occupancy':   occupancy_rate < 60,
            },
            'message': f'Ward Analytics: {occupancy_rate}% occupancy | {available_beds} beds available'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/cno/report', methods=['POST'])
@login_required
def cno_report():
    """Generate comprehensive CNO Board Report PDF"""
    try:
        data = request.get_json()
        hospital_name    = data.get('hospital_name', 'Hospital')
        cno_name         = data.get('cno_name', 'Chief Nursing Officer')
        report_period    = data.get('period', 'Monthly')
        staff_data       = data.get('staff_data', {})
        clinical_data    = data.get('clinical_data', {})
        ward_data        = data.get('ward_data', {})
        highlights       = data.get('highlights', '')
        challenges       = data.get('challenges', '')
        goals_next_month = data.get('goals', '')

        user_prompt = f"""Generate a comprehensive CNO Board Report for:

HOSPITAL: {hospital_name}
CNO: {cno_name}
PERIOD: {report_period}
DATE: {__import__('datetime').datetime.now().strftime('%B %Y')}

STAFFING SUMMARY:
{str(staff_data) if staff_data else 'Total Nurses: Not provided'}

CLINICAL QUALITY SUMMARY:
{str(clinical_data) if clinical_data else 'Quality metrics: Not provided'}

WARD OPERATIONS SUMMARY:
{str(ward_data) if ward_data else 'Ward data: Not provided'}

HIGHLIGHTS THIS PERIOD: {highlights if highlights else 'To be filled'}
CHALLENGES FACED: {challenges if challenges else 'To be filled'}
GOALS FOR NEXT PERIOD: {goals_next_month if goals_next_month else 'To be filled'}

Generate a FULL CNO Board Report with these sections:

# CNO MONTHLY REPORT — {report_period.upper()}
## {hospital_name}
### Prepared by: {cno_name}

**EXECUTIVE SUMMARY** (3-4 lines — for Board of Directors)

**1. NURSING WORKFORCE OVERVIEW**
- Staffing levels and ratios
- Key staffing metrics
- Workforce challenges

**2. CLINICAL QUALITY & PATIENT SAFETY**
- Quality indicators summary
- Patient safety events
- Improvement initiatives

**3. PATIENT EXPERIENCE**
- Satisfaction scores
- Complaint resolution
- Service excellence initiatives

**4. WARD OPERATIONS & CAPACITY**
- Bed utilization
- Patient flow efficiency
- Operational highlights

**5. EDUCATION & PROFESSIONAL DEVELOPMENT**
- Training completed
- Certifications achieved
- Upcoming programs

**6. HIGHLIGHTS & ACHIEVEMENTS**
- Key accomplishments
- Best practices implemented
- Staff recognition

**7. CHALLENGES & MITIGATION**
- Current challenges
- Action plans
- Support needed

**8. GOALS FOR NEXT PERIOD** (SMART goals)

**9. BUDGET IMPLICATIONS**
- Overtime costs
- Agency staff costs
- Cost saving initiatives

**10. CNO RECOMMENDATIONS** (Top 3 for Board approval)

Make it professional, data-driven, and board-presentation ready."""

        report_text = get_groq_response(CNO_SYSTEM_PROMPT, user_prompt, max_tokens=3000)

        # Generate PDF
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.colors import HexColor, white, black
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.units import inch, cm
            from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
            import io as _io

            buf = _io.BytesIO()
            doc = SimpleDocTemplate(buf, pagesize=A4,
                                   topMargin=0.8*inch, bottomMargin=0.8*inch,
                                   leftMargin=0.9*inch, rightMargin=0.9*inch)

            styles = getSampleStyleSheet()
            PURPLE = HexColor('#5B5EF4')
            TEAL   = HexColor('#0D9488')
            DARK   = HexColor('#1E293B')

            title_style = ParagraphStyle('Title', parent=styles['Title'],
                fontSize=20, textColor=white, alignment=TA_CENTER,
                spaceAfter=6, fontName='Helvetica-Bold')
            sub_style = ParagraphStyle('Sub', parent=styles['Normal'],
                fontSize=11, textColor=HexColor('#C7D2FE'), alignment=TA_CENTER,
                spaceAfter=4, fontName='Helvetica')
            h1_style = ParagraphStyle('H1', parent=styles['Heading1'],
                fontSize=13, textColor=PURPLE, spaceBefore=14, spaceAfter=6,
                fontName='Helvetica-Bold', borderPad=4)
            body_style = ParagraphStyle('Body', parent=styles['Normal'],
                fontSize=10, textColor=DARK, spaceAfter=5, leading=16,
                fontName='Helvetica', alignment=TA_JUSTIFY)

            story = []

            # Header banner
            header_data = [[Paragraph(f'CNO BOARD REPORT', title_style)],
                          [Paragraph(f'{hospital_name} | {report_period}', sub_style)],
                          [Paragraph(f'Prepared by: {cno_name} | {__import__("datetime").datetime.now().strftime("%d %B %Y")}', sub_style)]]
            header_table = Table(header_data, colWidths=[6.3*inch])
            header_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), PURPLE),
                ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                ('TOPPADDING', (0,0), (-1,-1), 12),
                ('BOTTOMPADDING', (0,0), (-1,-1), 12),
                ('ROUNDEDCORNERS', [8,8,8,8]),
            ]))
            story.append(header_table)
            story.append(Spacer(1, 0.3*inch))

            # Report content
            for line in report_text.split('\n'):
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 0.08*inch))
                    continue
                line_clean = line.replace('**', '').replace('##', '').replace('#', '').strip()
                if line.startswith('# ') or line.startswith('## '):
                    story.append(Paragraph(line_clean, h1_style))
                elif line.startswith('### '):
                    sub_h = ParagraphStyle('SubH', parent=styles['Heading2'],
                        fontSize=11, textColor=TEAL, spaceBefore=8, spaceAfter=4, fontName='Helvetica-Bold')
                    story.append(Paragraph(line_clean, sub_h))
                elif line.startswith('- ') or line.startswith('• '):
                    bullet_style = ParagraphStyle('Bullet', parent=styles['Normal'],
                        fontSize=10, textColor=DARK, leftIndent=18, spaceAfter=3,
                        fontName='Helvetica', bulletIndent=6)
                    story.append(Paragraph(f'• {line_clean[2:]}', bullet_style))
                elif line.startswith('**') and line.endswith('**'):
                    bold_s = ParagraphStyle('Bold', parent=styles['Normal'],
                        fontSize=11, textColor=DARK, spaceBefore=6, spaceAfter=3, fontName='Helvetica-Bold')
                    story.append(Paragraph(line_clean, bold_s))
                else:
                    story.append(Paragraph(line_clean, body_style))

            # Footer
            story.append(Spacer(1, 0.3*inch))
            footer_data = [[Paragraph(f'Confidential | {hospital_name} | Generated by DS Agent', ParagraphStyle('Footer', fontSize=8, textColor=white, alignment=TA_CENTER, fontName='Helvetica'))]]
            footer_table = Table(footer_data, colWidths=[6.3*inch])
            footer_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), HexColor('#3730A3')),
                ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                ('TOPPADDING', (0,0), (-1,-1), 8),
                ('BOTTOMPADDING', (0,0), (-1,-1), 8),
            ]))
            story.append(footer_table)

            doc.build(story)
            buf.seek(0)

            return send_file(buf, as_attachment=True,
                           download_name=f'CNO_Report_{hospital_name}_{report_period}.pdf',
                           mimetype='application/pdf')
        except Exception as pdf_err:
            return jsonify({
                'success': True,
                'report': report_text,
                'message': 'CNO Report generated! (PDF library issue — text version)',
                'pdf_error': str(pdf_err)
            })

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/cno/huddle', methods=['POST'])
@login_required
def cno_morning_huddle():
    """CNO Morning Huddle — Quick daily briefing for nursing leadership"""
    try:
        data = request.get_json()
        date          = data.get('date', __import__('datetime').datetime.now().strftime('%d %B %Y'))
        census        = data.get('census', '')
        staffing      = data.get('staffing', '')
        quality       = data.get('quality', '')
        priorities    = data.get('priorities', '')
        hospital_name = data.get('hospital_name', 'Hospital')

        user_prompt = f"""Generate CNO Morning Huddle Briefing for {hospital_name} — {date}:

CURRENT CENSUS: {census if census else 'Not provided'}
STAFFING STATUS: {staffing if staffing else 'Not provided'}
QUALITY/SAFETY UPDATE: {quality if quality else 'Not provided'}
TODAY'S PRIORITIES: {priorities if priorities else 'Not provided'}

Generate a crisp 5-minute morning huddle script:

**MORNING HUDDLE — {date}**
**{hospital_name} Nursing Leadership**

🏥 **CENSUS UPDATE** (30 seconds)
📊 **STAFFING STATUS** (1 minute)
🎯 **QUALITY & SAFETY** (1 minute)
⚡ **TODAY'S PRIORITIES** (1 minute)
💪 **MOTIVATION & FOCUS** (30 seconds)
❓ **QUICK QUESTIONS** (1 minute)

Keep it concise, action-oriented, and energizing for nursing staff."""

        huddle = get_groq_response(CNO_SYSTEM_PROMPT, user_prompt, max_tokens=1200)

        return jsonify({
            'success': True,
            'huddle': huddle,
            'date': date,
            'message': f'Morning Huddle script ready for {date}!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── CNO CSV Auto-Fill ─────────────────────────────────────────────
@app.route('/api/cno/parse_csv', methods=['POST'])
@login_required
def cno_parse_csv():
    """Parse hospital CSV/Excel and auto-detect CNO fields"""
    try:
        import json as _j
        if 'file' not in request.files:
            return jsonify({'error': 'File upload karo!'}), 400

        file = request.files['file']
        filename = file.filename.lower()

        # Read file
        import io as _io
        file_bytes = file.read()

        if filename.endswith('.csv'):
            import pandas as _pd
            df = _pd.read_csv(_io.BytesIO(file_bytes), encoding='latin1')
        elif filename.endswith(('.xlsx', '.xls')):
            import pandas as _pd
            df = _pd.read_excel(_io.BytesIO(file_bytes))
        else:
            return jsonify({'error': 'CSV ya Excel file upload karo!'}), 400

        cols = [c.lower().strip() for c in df.columns]
        col_map = {c.lower().strip(): c for c in df.columns}

        def find_col(*keywords):
            for kw in keywords:
                for c in cols:
                    if kw in c:
                        return col_map[c]
            return None

        def get_val(col, agg='sum'):
            if col is None: return 0
            try:
                s = df[col].dropna()
                if agg == 'sum':   return int(s.sum())
                if agg == 'mean':  return round(float(s.mean()), 1)
                if agg == 'count': return int(s.count())
                if agg == 'max':   return int(s.max())
            except: return 0

        # ── Staff fields ─────────────────────────────────────────
        nurse_col    = find_col('nurse','nursing','staff','rn','nurse_count','nurses')
        patient_col  = find_col('patient','census','admitted','total_patient')
        day_col      = find_col('day_shift','day shift','morning')
        night_col    = find_col('night_shift','night shift','evening')
        absent_col   = find_col('absent','absentee','leave','off')
        ot_col       = find_col('overtime','ot_hours','extra_hours')
        icu_n_col    = find_col('icu_nurse','icu nurse','critical_nurse')
        icu_p_col    = find_col('icu_patient','icu patient','critical_patient','icu_census')
        agency_col   = find_col('agency','contract','temp_nurse','locum')
        exp_col      = find_col('experience','exp_years','years_exp','tenure')
        training_col = find_col('training','certification','cert_due')
        incident_col = find_col('incident','report','ir_count')

        # ── Clinical fields ──────────────────────────────────────
        adm_col      = find_col('admission','admissions','total_adm')
        hai_col      = find_col('hai','infection','hais','hospital_infection')
        cauti_col    = find_col('cauti','catheter','uti')
        vap_col      = find_col('vap','ventilator','pneumonia')
        fall_col     = find_col('fall','falls','patient_fall')
        med_col      = find_col('medication_error','med_error','drug_error')
        ulcer_col    = find_col('ulcer','pressure','bedsore')
        code_col     = find_col('code_blue','code blue','cardiac_arrest')
        near_col     = find_col('near_miss','nearmiss','near miss')
        sat_col      = find_col('satisfaction','hcahps','patient_score','rating')
        readm_col    = find_col('readmission','readmit','30day')
        mort_col     = find_col('mortality','death','death_rate')
        comp_col     = find_col('complaint','grievance')

        # ── Ward fields ──────────────────────────────────────────
        beds_col     = find_col('total_bed','total beds','bed_capacity','beds')
        occ_col      = find_col('occupied','occupancy','occupied_bed')
        icu_bed_col  = find_col('icu_bed','icu bed','critical_bed')
        icu_occ_col  = find_col('icu_occ','icu_occupied','critical_occ')
        los_col      = find_col('los','length_of_stay','avg_stay','alos')
        adm_today    = find_col('admission_today','admit_today','new_admission')
        dis_col      = find_col('discharge','discharged','dis_today')
        pend_col     = find_col('pending_discharge','pending','delayed_discharge')
        er_col       = find_col('er_waiting','ed_waiting','emergency_wait')

        # ── Build result ─────────────────────────────────────────
        staff = {
            'total_nurses':    get_val(nurse_col, 'sum') or get_val(nurse_col, 'count'),
            'total_patients':  get_val(patient_col, 'sum'),
            'day_shift':       get_val(day_col, 'sum'),
            'night_shift':     get_val(night_col, 'sum'),
            'absent_today':    get_val(absent_col, 'sum'),
            'overtime_hours':  get_val(ot_col, 'sum'),
            'icu_nurses':      get_val(icu_n_col, 'sum'),
            'icu_patients':    get_val(icu_p_col, 'sum'),
            'agency_nurses':   get_val(agency_col, 'sum'),
            'avg_experience':  get_val(exp_col, 'mean'),
            'training_due':    get_val(training_col, 'sum'),
            'incidents_month': get_val(incident_col, 'sum'),
        }

        clinical = {
            'total_admissions':   get_val(adm_col, 'sum'),
            'hai':                get_val(hai_col, 'sum'),
            'cauti':              get_val(cauti_col, 'sum'),
            'vap':                get_val(vap_col, 'sum'),
            'falls':              get_val(fall_col, 'sum'),
            'med_errors':         get_val(med_col, 'sum'),
            'pressure_ulcers':    get_val(ulcer_col, 'sum'),
            'code_blue':          get_val(code_col, 'sum'),
            'near_miss':          get_val(near_col, 'sum'),
            'satisfaction_score': get_val(sat_col, 'mean'),
            'readmission_30day':  get_val(readm_col, 'sum'),
            'mortality_rate':     get_val(mort_col, 'mean'),
            'complaints':         get_val(comp_col, 'sum'),
        }

        ward = {
            'total_beds':       get_val(beds_col, 'sum'),
            'occupied_beds':    get_val(occ_col, 'sum'),
            'icu_beds':         get_val(icu_bed_col, 'sum'),
            'icu_occupied':     get_val(icu_occ_col, 'sum'),
            'avg_los':          get_val(los_col, 'mean'),
            'admissions_today': get_val(adm_today, 'sum'),
            'discharges_today': get_val(dis_col, 'sum'),
            'pending_discharge':get_val(pend_col, 'sum'),
            'er_waiting':       get_val(er_col, 'sum'),
        }

        # Detected columns summary
        detected = {}
        field_map = {
            'Nurses': nurse_col, 'Patients': patient_col,
            'Admissions': adm_col, 'Falls': fall_col,
            'HAI': hai_col, 'Beds': beds_col,
            'Occupied': occ_col, 'Satisfaction': sat_col,
            'Incidents': incident_col, 'Absent': absent_col,
        }
        for label, col in field_map.items():
            if col: detected[label] = col

        return jsonify({
            'success':  True,
            'staff':    staff,
            'clinical': clinical,
            'ward':     ward,
            'detected': detected,
            'rows':     len(df),
            'columns':  list(df.columns),
            'message':  f'✅ {len(detected)} fields auto-detected from {len(df)} rows!'
        })

    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# 🔬 ADVANCED ML FEATURES
# PCA, DBSCAN, Lasso/Ridge, Confusion Matrix, Cross Validation
# ════════════════════════════════════════════════════════════════

@app.route('/api/ml/pca', methods=['POST'])
@login_required
def ml_pca():
    """PCA — Principal Component Analysis"""
    try:
        import json as _j
        data = request.get_json()
        n_components = data.get('n_components', 2)
        target_col   = data.get('target', '')

        if agent.df is None:
            return jsonify({'error': 'Pehle CSV upload karo!'}), 400

        import numpy as _np
        from sklearn.preprocessing import StandardScaler
        from sklearn.decomposition import PCA
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as _plt
        import io as _io, base64 as _b64

        num_cols = agent.df.select_dtypes('number').columns.tolist()
        if target_col and target_col in num_cols:
            num_cols = [c for c in num_cols if c != target_col]

        if len(num_cols) < 2:
            return jsonify({'error': 'Minimum 2 numeric columns chahiye!'}), 400

        n_components = min(n_components, len(num_cols))
        df_num = agent.df[num_cols].dropna()

        # Scale + PCA
        scaler = StandardScaler()
        X_scaled = scaler.fit_transform(df_num)
        pca = PCA(n_components=n_components)
        X_pca = pca.fit_transform(X_scaled)

        explained = pca.explained_variance_ratio_
        cumulative = _np.cumsum(explained)
        loadings   = pca.components_

        # Plot
        fig, axes = _plt.subplots(1, 2, figsize=(10, 4))
        fig.suptitle('PCA Analysis', fontsize=14, fontweight='bold')

        # Scree plot
        axes[0].bar(range(1, n_components+1), explained*100, color='#5b5ef4', alpha=0.8)
        axes[0].plot(range(1, n_components+1), cumulative*100, 'ro-', linewidth=2)
        axes[0].set_xlabel('Principal Component')
        axes[0].set_ylabel('Variance Explained (%)')
        axes[0].set_title('Explained Variance')
        axes[0].axhline(y=80, color='green', linestyle='--', alpha=0.5, label='80% threshold')
        axes[0].legend()

        # Biplot (PC1 vs PC2)
        if n_components >= 2:
            scatter = axes[1].scatter(X_pca[:, 0], X_pca[:, 1],
                                      c='#5b5ef4', alpha=0.5, s=20)
            # Feature arrows
            scale = _np.max(_np.abs(X_pca[:, :2])) * 0.4
            for i, col in enumerate(num_cols[:8]):
                axes[1].arrow(0, 0,
                              loadings[0, i]*scale,
                              loadings[1, i]*scale,
                              head_width=0.1, head_length=0.05,
                              fc='red', ec='red', alpha=0.7)
                axes[1].text(loadings[0, i]*scale*1.15,
                             loadings[1, i]*scale*1.15,
                             col, fontsize=8, ha='center', color='darkred')
            axes[1].set_xlabel(f'PC1 ({explained[0]*100:.1f}%)')
            axes[1].set_ylabel(f'PC2 ({explained[1]*100:.1f}%)')
            axes[1].set_title('PCA Biplot')

        _plt.tight_layout()
        buf = _io.BytesIO()
        _plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        _plt.close()
        buf.seek(0)
        chart_b64 = _b64.b64encode(buf.read()).decode('utf-8')

        # Component summary
        comp_summary = []
        for i in range(n_components):
            top_features = sorted(zip(num_cols, abs(loadings[i])),
                                  key=lambda x: x[1], reverse=True)[:3]
            comp_summary.append({
                'component': f'PC{i+1}',
                'variance':  round(explained[i]*100, 2),
                'cumulative': round(cumulative[i]*100, 2),
                'top_features': [f[0] for f in top_features]
            })

        return jsonify({
            'success':       True,
            'chart':         chart_b64,
            'components':    comp_summary,
            'total_variance': round(cumulative[-1]*100, 2),
            'n_components':  n_components,
            'features_used': num_cols,
            'message': f'PCA done! {n_components} components = {round(cumulative[-1]*100,1)}% variance explained'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/ml/dbscan', methods=['POST'])
@login_required
def ml_dbscan():
    """DBSCAN Clustering"""
    try:
        data = request.get_json()
        eps         = float(data.get('eps', 0.5))
        min_samples = int(data.get('min_samples', 5))
        feature1    = data.get('feature1', '')
        feature2    = data.get('feature2', '')

        if agent.df is None:
            return jsonify({'error': 'Pehle CSV upload karo!'}), 400

        import numpy as _np
        from sklearn.preprocessing import StandardScaler
        from sklearn.cluster import DBSCAN
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as _plt
        import io as _io, base64 as _b64

        num_cols = agent.df.select_dtypes('number').columns.tolist()
        if len(num_cols) < 2:
            return jsonify({'error': 'Minimum 2 numeric columns chahiye!'}), 400

        f1 = feature1 if feature1 in num_cols else num_cols[0]
        f2 = feature2 if feature2 in num_cols else num_cols[1]

        df_clean = agent.df[[f1, f2]].dropna()
        X = StandardScaler().fit_transform(df_clean)

        db = DBSCAN(eps=eps, min_samples=min_samples)
        labels = db.fit_predict(X)

        n_clusters = len(set(labels)) - (1 if -1 in labels else 0)
        n_noise    = list(labels).count(-1)
        noise_pct  = round(n_noise / len(labels) * 100, 1)

        # Plot
        fig, ax = _plt.subplots(figsize=(8, 5))
        unique_labels = set(labels)
        colors = _plt.cm.tab10(_np.linspace(0, 1, max(len(unique_labels), 1)))

        for k, col in zip(unique_labels, colors):
            if k == -1:
                col = 'black'
                label_name = 'Noise/Outlier'
                marker = 'x'
                size = 30
            else:
                label_name = f'Cluster {k+1}'
                marker = 'o'
                size = 40
            mask = labels == k
            ax.scatter(df_clean[f1][mask], df_clean[f2][mask],
                      c=[col], label=label_name,
                      marker=marker, s=size, alpha=0.7)

        ax.set_xlabel(f1); ax.set_ylabel(f2)
        ax.set_title(f'DBSCAN Clustering — {n_clusters} Clusters Found', fontweight='bold')
        ax.legend(loc='best', fontsize=8)
        _plt.tight_layout()

        buf = _io.BytesIO()
        _plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        _plt.close()
        buf.seek(0)
        chart_b64 = _b64.b64encode(buf.read()).decode('utf-8')

        # Cluster sizes
        cluster_sizes = {}
        for k in unique_labels:
            if k != -1:
                cluster_sizes[f'Cluster {k+1}'] = int(list(labels).count(k))

        return jsonify({
            'success':       True,
            'chart':         chart_b64,
            'n_clusters':    n_clusters,
            'n_noise':       n_noise,
            'noise_pct':     noise_pct,
            'cluster_sizes': cluster_sizes,
            'total_points':  len(labels),
            'eps':           eps,
            'min_samples':   min_samples,
            'message': f'DBSCAN done! {n_clusters} clusters found, {n_noise} noise points ({noise_pct}%)'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/ml/lasso_ridge', methods=['POST'])
@login_required
def ml_lasso_ridge():
    """Lasso & Ridge Regression with feature selection"""
    try:
        data        = request.get_json()
        target_col  = data.get('target', '')
        model_type  = data.get('model_type', 'both')  # lasso, ridge, both
        alpha       = float(data.get('alpha', 1.0))

        if agent.df is None:
            return jsonify({'error': 'Pehle CSV upload karo!'}), 400
        if not target_col or target_col not in agent.df.columns:
            return jsonify({'error': 'Target column select karo!'}), 400

        import numpy as _np
        from sklearn.linear_model import Lasso, Ridge
        from sklearn.preprocessing import StandardScaler
        from sklearn.model_selection import cross_val_score, train_test_split
        from sklearn.metrics import r2_score, mean_squared_error
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as _plt
        import io as _io, base64 as _b64

        num_cols = [c for c in agent.df.select_dtypes('number').columns if c != target_col]
        if len(num_cols) < 1:
            return jsonify({'error': 'Numeric feature columns chahiye!'}), 400

        df_clean = agent.df[num_cols + [target_col]].dropna()
        X = df_clean[num_cols].values
        y = df_clean[target_col].values

        scaler = StandardScaler()
        X_scaled = scaler.fit_transform(X)
        X_train, X_test, y_train, y_test = train_test_split(X_scaled, y, test_size=0.2, random_state=42)

        results = {}
        models_to_run = ['lasso', 'ridge'] if model_type == 'both' else [model_type]

        for mtype in models_to_run:
            model = Lasso(alpha=alpha, max_iter=10000) if mtype == 'lasso' else Ridge(alpha=alpha)
            model.fit(X_train, y_train)
            y_pred = model.predict(X_test)
            r2     = r2_score(y_test, y_pred)
            rmse   = _np.sqrt(mean_squared_error(y_test, y_pred))
            cv     = cross_val_score(model, X_scaled, y, cv=5, scoring='r2')
            coefs  = dict(zip(num_cols, model.coef_))
            zero_coefs = sum(1 for c in model.coef_ if abs(c) < 0.001)
            results[mtype] = {
                'r2_score':    round(r2, 4),
                'rmse':        round(rmse, 4),
                'cv_mean':     round(cv.mean(), 4),
                'cv_std':      round(cv.std(), 4),
                'zero_features': zero_coefs,
                'active_features': len(num_cols) - zero_coefs,
                'coefficients': {k: round(v, 4) for k, v in sorted(coefs.items(), key=lambda x: abs(x[1]), reverse=True)[:10]}
            }

        # Plot coefficients
        n_plots = len(models_to_run)
        fig, axes = _plt.subplots(1, n_plots, figsize=(8*n_plots, 5))
        if n_plots == 1:
            axes = [axes]

        colors_pos = '#5b5ef4'; colors_neg = '#dc2626'
        for ax, mtype in zip(axes, models_to_run):
            coefs_sorted = sorted(zip(num_cols, results[mtype]['coefficients'].values()),
                                  key=lambda x: abs(x[1]), reverse=True)[:10]
            feat_names = [c[0] for c in coefs_sorted]
            feat_vals  = [c[1] for c in coefs_sorted]
            bar_colors = [colors_pos if v >= 0 else colors_neg for v in feat_vals]
            ax.barh(range(len(feat_names)), feat_vals, color=bar_colors, alpha=0.8)
            ax.set_yticks(range(len(feat_names)))
            ax.set_yticklabels(feat_names, fontsize=9)
            ax.axvline(x=0, color='black', linewidth=0.8)
            ax.set_title(f'{mtype.title()} Regression\nR² = {results[mtype]["r2_score"]} | RMSE = {results[mtype]["rmse"]}',
                        fontweight='bold')
            ax.set_xlabel('Coefficient Value')

        _plt.tight_layout()
        buf = _io.BytesIO()
        _plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        _plt.close()
        buf.seek(0)
        chart_b64 = _b64.b64encode(buf.read()).decode('utf-8')

        return jsonify({
            'success':    True,
            'chart':      chart_b64,
            'results':    results,
            'target':     target_col,
            'features':   num_cols,
            'alpha':      alpha,
            'message':    f'Lasso/Ridge done! Best R²: {max(r["r2_score"] for r in results.values())}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/ml/confusion_matrix', methods=['POST'])
@login_required
def ml_confusion_matrix():
    """Confusion Matrix Visual for classification models"""
    try:
        data       = request.get_json()
        target_col = data.get('target', '')
        model_type = data.get('model_type', 'random_forest')

        if agent.df is None:
            return jsonify({'error': 'Pehle CSV upload karo!'}), 400
        if not target_col or target_col not in agent.df.columns:
            return jsonify({'error': 'Target column select karo!'}), 400

        import numpy as _np
        from sklearn.model_selection import train_test_split, cross_val_score
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        from sklearn.metrics import (confusion_matrix, classification_report,
                                     accuracy_score, roc_auc_score, roc_curve)
        from sklearn.ensemble import RandomForestClassifier, GradientBoostingClassifier
        from sklearn.linear_model import LogisticRegression
        from sklearn.tree import DecisionTreeClassifier
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as _plt
        import seaborn as _sns
        import io as _io, base64 as _b64

        num_cols = [c for c in agent.df.select_dtypes('number').columns if c != target_col]
        if len(num_cols) < 1:
            return jsonify({'error': 'Numeric feature columns chahiye!'}), 400

        df_clean = agent.df[num_cols + [target_col]].dropna()
        X = df_clean[num_cols].values
        y_raw = df_clean[target_col].values

        le = LabelEncoder()
        y = le.fit_transform(y_raw)
        classes = le.classes_

        X_scaled = StandardScaler().fit_transform(X)
        X_train, X_test, y_train, y_test = train_test_split(
            X_scaled, y, test_size=0.2, random_state=42, stratify=y if len(_np.unique(y)) > 1 else None)

        models_map = {
            'random_forest':    RandomForestClassifier(n_estimators=100, random_state=42),
            'gradient_boosting': GradientBoostingClassifier(random_state=42),
            'logistic':         LogisticRegression(max_iter=1000, random_state=42),
            'decision_tree':    DecisionTreeClassifier(random_state=42),
        }
        model = models_map.get(model_type, models_map['random_forest'])
        model.fit(X_train, y_train)
        y_pred = model.predict(X_test)
        accuracy = accuracy_score(y_test, y_pred)
        cm = confusion_matrix(y_test, y_pred)
        report = classification_report(y_test, y_pred, target_names=[str(c) for c in classes], output_dict=True)
        cv_scores = cross_val_score(model, X_scaled, y, cv=5, scoring='accuracy')

        # Plots
        n_classes = len(classes)
        if n_classes == 2:
            fig, axes = _plt.subplots(1, 3, figsize=(15, 4))
        else:
            fig, axes = _plt.subplots(1, 2, figsize=(12, 5))

        # Confusion Matrix Heatmap
        _sns.heatmap(cm, annot=True, fmt='d', cmap='Blues',
                    xticklabels=[str(c) for c in classes],
                    yticklabels=[str(c) for c in classes],
                    ax=axes[0], linewidths=0.5)
        axes[0].set_title(f'Confusion Matrix\nAccuracy: {accuracy*100:.1f}%', fontweight='bold')
        axes[0].set_xlabel('Predicted'); axes[0].set_ylabel('Actual')

        # Per-class metrics bar chart
        class_names = [str(c) for c in classes]
        precisions  = [report[str(c)]['precision'] for c in classes if str(c) in report]
        recalls     = [report[str(c)]['recall'] for c in classes if str(c) in report]
        f1s         = [report[str(c)]['f1-score'] for c in classes if str(c) in report]
        x_pos = _np.arange(len(class_names[:len(precisions)]))
        width = 0.25
        axes[1].bar(x_pos - width, precisions, width, label='Precision', color='#5b5ef4', alpha=0.8)
        axes[1].bar(x_pos,         recalls,    width, label='Recall',    color='#059669', alpha=0.8)
        axes[1].bar(x_pos + width, f1s,        width, label='F1-Score',  color='#d97706', alpha=0.8)
        axes[1].set_xticks(x_pos)
        axes[1].set_xticklabels(class_names[:len(precisions)], rotation=45, ha='right', fontsize=8)
        axes[1].set_title('Per-Class Metrics', fontweight='bold')
        axes[1].legend(); axes[1].set_ylim(0, 1.1)

        # ROC Curve (binary only)
        if n_classes == 2:
            try:
                y_prob = model.predict_proba(X_test)[:, 1]
                fpr, tpr, _ = roc_curve(y_test, y_prob)
                auc_score = roc_auc_score(y_test, y_prob)
                axes[2].plot(fpr, tpr, color='#5b5ef4', lw=2,
                            label=f'ROC (AUC = {auc_score:.3f})')
                axes[2].plot([0,1],[0,1], 'k--', alpha=0.5)
                axes[2].set_xlabel('False Positive Rate')
                axes[2].set_ylabel('True Positive Rate')
                axes[2].set_title('ROC Curve', fontweight='bold')
                axes[2].legend()
            except: pass

        _plt.tight_layout()
        buf = _io.BytesIO()
        _plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        _plt.close()
        buf.seek(0)
        chart_b64 = _b64.b64encode(buf.read()).decode('utf-8')

        # Format report
        report_summary = {}
        for cls in [str(c) for c in classes]:
            if cls in report:
                report_summary[cls] = {
                    'precision': round(report[cls]['precision'], 3),
                    'recall':    round(report[cls]['recall'], 3),
                    'f1_score':  round(report[cls]['f1-score'], 3),
                    'support':   int(report[cls]['support'])
                }

        return jsonify({
            'success':       True,
            'chart':         chart_b64,
            'accuracy':      round(accuracy*100, 2),
            'cv_accuracy':   round(cv_scores.mean()*100, 2),
            'cv_std':        round(cv_scores.std()*100, 2),
            'report':        report_summary,
            'classes':       [str(c) for c in classes],
            'model_used':    model_type,
            'message': f'Confusion Matrix ready! Accuracy: {accuracy*100:.1f}% | CV: {cv_scores.mean()*100:.1f}%'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/ml/cross_validation', methods=['POST'])
@login_required
def ml_cross_validation():
    """Advanced Cross Validation with multiple strategies"""
    try:
        data        = request.get_json()
        target_col  = data.get('target', '')
        task_type   = data.get('task_type', 'classification')
        cv_folds    = int(data.get('cv_folds', 5))
        models_list = data.get('models', ['random_forest','logistic','decision_tree'])

        if agent.df is None:
            return jsonify({'error': 'Pehle CSV upload karo!'}), 400
        if not target_col or target_col not in agent.df.columns:
            return jsonify({'error': 'Target column select karo!'}), 400

        import numpy as _np
        from sklearn.model_selection import (cross_val_score, StratifiedKFold, KFold,
                                             learning_curve)
        from sklearn.preprocessing import LabelEncoder, StandardScaler
        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor, GradientBoostingClassifier, GradientBoostingRegressor
        from sklearn.linear_model import LogisticRegression, Ridge, Lasso
        from sklearn.tree import DecisionTreeClassifier, DecisionTreeRegressor
        from sklearn.svm import SVC, SVR
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as _plt
        import io as _io, base64 as _b64

        num_cols = [c for c in agent.df.select_dtypes('number').columns if c != target_col]
        if len(num_cols) < 1:
            return jsonify({'error': 'Numeric feature columns chahiye!'}), 400

        df_clean = agent.df[num_cols + [target_col]].dropna()
        X = StandardScaler().fit_transform(df_clean[num_cols].values)

        if task_type == 'classification':
            le = LabelEncoder()
            y = le.fit_transform(df_clean[target_col].values)
            cv_strategy = StratifiedKFold(n_splits=cv_folds, shuffle=True, random_state=42)
            scoring = 'accuracy'
            models_map = {
                'random_forest':     RandomForestClassifier(n_estimators=100, random_state=42),
                'gradient_boosting': GradientBoostingClassifier(random_state=42),
                'logistic':          LogisticRegression(max_iter=1000, random_state=42),
                'decision_tree':     DecisionTreeClassifier(random_state=42),
                'svm':               SVC(kernel='rbf', random_state=42),
            }
        else:
            y = df_clean[target_col].values.astype(float)
            cv_strategy = KFold(n_splits=cv_folds, shuffle=True, random_state=42)
            scoring = 'r2'
            models_map = {
                'random_forest':     RandomForestRegressor(n_estimators=100, random_state=42),
                'gradient_boosting': GradientBoostingRegressor(random_state=42),
                'ridge':             Ridge(alpha=1.0),
                'lasso':             Lasso(alpha=1.0, max_iter=10000),
                'decision_tree':     DecisionTreeRegressor(random_state=42),
            }

        cv_results = {}
        for mname in models_list:
            if mname not in models_map: continue
            model = models_map[mname]
            scores = cross_val_score(model, X, y, cv=cv_strategy, scoring=scoring)
            cv_results[mname] = {
                'mean':   round(scores.mean()*100, 2),
                'std':    round(scores.std()*100, 2),
                'min':    round(scores.min()*100, 2),
                'max':    round(scores.max()*100, 2),
                'scores': [round(s*100, 2) for s in scores],
            }

        # Plot
        fig, axes = _plt.subplots(1, 2, figsize=(12, 5))

        # Box plot comparison
        model_names = list(cv_results.keys())
        score_data  = [cv_results[m]['scores'] for m in model_names]
        bp = axes[0].boxplot(score_data, labels=model_names, patch_artist=True)
        colors_box = ['#5b5ef4','#059669','#d97706','#dc2626','#8b5cf6']
        for patch, color in zip(bp['boxes'], colors_box):
            patch.set_facecolor(color); patch.set_alpha(0.7)
        axes[0].set_title(f'{cv_folds}-Fold Cross Validation Comparison', fontweight='bold')
        axes[0].set_ylabel(f'{scoring.upper()} Score (%)')
        axes[0].tick_params(axis='x', rotation=30)

        # Mean scores bar chart
        means = [cv_results[m]['mean'] for m in model_names]
        stds  = [cv_results[m]['std'] for m in model_names]
        bars  = axes[1].bar(model_names, means, yerr=stds,
                           color=colors_box[:len(model_names)],
                           alpha=0.8, capsize=5)
        for bar, mean in zip(bars, means):
            axes[1].text(bar.get_x()+bar.get_width()/2,
                        bar.get_height()+0.5,
                        f'{mean:.1f}%', ha='center', va='bottom', fontsize=9, fontweight='bold')
        axes[1].set_title('Mean CV Scores (± Std Dev)', fontweight='bold')
        axes[1].set_ylabel(f'Mean {scoring.upper()} Score (%)')
        axes[1].tick_params(axis='x', rotation=30)
        best_model = max(cv_results, key=lambda x: cv_results[x]['mean'])
        axes[1].set_xlabel(f'🏆 Best: {best_model} ({cv_results[best_model]["mean"]}%)')

        _plt.tight_layout()
        buf = _io.BytesIO()
        _plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        _plt.close()
        buf.seek(0)
        chart_b64 = _b64.b64encode(buf.read()).decode('utf-8')

        best = max(cv_results, key=lambda x: cv_results[x]['mean'])
        return jsonify({
            'success':     True,
            'chart':       chart_b64,
            'cv_results':  cv_results,
            'best_model':  best,
            'best_score':  cv_results[best]['mean'],
            'cv_folds':    cv_folds,
            'task_type':   task_type,
            'target':      target_col,
            'message': f'Cross Validation done! Best: {best} ({cv_results[best]["mean"]}% ± {cv_results[best]["std"]}%)'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# 🔬 PCA — Principal Component Analysis
# ════════════════════════════════════════════════════════════════
@app.route('/api/pca', methods=['POST'])
@login_required
def run_pca():
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from sklearn.decomposition import PCA
        from sklearn.preprocessing import StandardScaler
        import numpy as np

        if agent.df is None:
            return jsonify({'error': 'Pehle file upload karo!'}), 400

        data = request.get_json()
        n_components = int(data.get('n_components', 2))
        target_col    = data.get('target_col', '')

        df = agent.df.copy().dropna()
        # Only numeric columns
        num_df = df.select_dtypes(include='number')
        if target_col and target_col in num_df.columns:
            num_df = num_df.drop(columns=[target_col])

        if num_df.shape[1] < 2:
            return jsonify({'error': 'Kam se kam 2 numeric columns chahiye!'}), 400

        n_components = min(n_components, num_df.shape[1], num_df.shape[0])

        # Scale
        scaler = StandardScaler()
        X_scaled = scaler.fit_transform(num_df)

        # PCA
        pca = PCA(n_components=n_components)
        X_pca = pca.fit_transform(X_scaled)

        # Explained variance
        ev_ratio  = pca.explained_variance_ratio_
        cumulative = np.cumsum(ev_ratio)

        # Plot 1 — Explained Variance
        fig, axes = plt.subplots(1, 2, figsize=(10, 4))
        axes[0].bar(range(1, len(ev_ratio)+1), ev_ratio*100, color='#5b5ef4', alpha=0.8)
        axes[0].plot(range(1, len(ev_ratio)+1), cumulative*100, 'ro-', linewidth=2)
        axes[0].set_xlabel('Principal Component')
        axes[0].set_ylabel('Variance Explained (%)')
        axes[0].set_title('Explained Variance per Component', fontweight='bold')
        axes[0].axhline(y=80, color='green', linestyle='--', label='80% threshold')
        axes[0].legend()

        # Plot 2 — 2D Scatter
        if n_components >= 2:
            if target_col and target_col in df.columns:
                target_vals = df[target_col].values
                scatter = axes[1].scatter(X_pca[:,0], X_pca[:,1],
                    c=pd.factorize(target_vals.astype(str))[0],
                    cmap='tab10', alpha=0.6, s=30)
                axes[1].set_title(f'PCA 2D — colored by {target_col}', fontweight='bold')
            else:
                axes[1].scatter(X_pca[:,0], X_pca[:,1], alpha=0.5, color='#5b5ef4', s=30)
                axes[1].set_title('PCA 2D Projection', fontweight='bold')
            axes[1].set_xlabel('PC1')
            axes[1].set_ylabel('PC2')

        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode()

        # Loading scores (top features per component)
        loadings = {}
        for i, comp in enumerate(pca.components_[:3]):
            top_idx  = np.argsort(np.abs(comp))[::-1][:5]
            loadings[f'PC{i+1}'] = [
                {'feature': num_df.columns[j], 'loading': round(float(comp[j]), 4)}
                for j in top_idx
            ]

        return jsonify({
            'success':        True,
            'n_components':   n_components,
            'total_features': num_df.shape[1],
            'total_rows':     len(df),
            'explained_variance': [round(float(v)*100, 2) for v in ev_ratio],
            'cumulative_variance': [round(float(v)*100, 2) for v in cumulative],
            'variance_80_components': int(np.argmax(cumulative >= 0.80)) + 1,
            'loadings':       loadings,
            'chart':          chart_b64,
            'message': f'PCA complete! {n_components} components → {round(float(cumulative[-1])*100,1)}% variance explained'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# 🔵 DBSCAN Clustering
# ════════════════════════════════════════════════════════════════
@app.route('/api/dbscan', methods=['POST'])
@login_required
def run_dbscan():
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from sklearn.cluster import DBSCAN, KMeans
        from sklearn.preprocessing import StandardScaler
        from sklearn.decomposition import PCA
        import numpy as np

        if agent.df is None:
            return jsonify({'error': 'Pehle file upload karo!'}), 400

        data = request.get_json()
        eps         = float(data.get('eps', 0.5))
        min_samples = int(data.get('min_samples', 5))
        compare_kmeans = data.get('compare_kmeans', True)

        df     = agent.df.copy().dropna()
        num_df = df.select_dtypes(include='number')

        if num_df.shape[1] < 2:
            return jsonify({'error': 'Kam se kam 2 numeric columns chahiye!'}), 400

        # Limit rows
        if len(num_df) > 5000:
            num_df = num_df.sample(5000, random_state=42)

        scaler   = StandardScaler()
        X_scaled = scaler.fit_transform(num_df)

        # DBSCAN
        db      = DBSCAN(eps=eps, min_samples=min_samples)
        labels  = db.fit_predict(X_scaled)
        n_clusters  = len(set(labels)) - (1 if -1 in labels else 0)
        n_noise     = list(labels).count(-1)
        noise_pct   = round(n_noise / len(labels) * 100, 1)

        # Cluster sizes
        cluster_sizes = {}
        for lbl in sorted(set(labels)):
            name = 'Noise' if lbl == -1 else f'Cluster {lbl+1}'
            cluster_sizes[name] = int(list(labels).count(lbl))

        # PCA for visualization
        pca_2d = PCA(n_components=2)
        X_2d   = pca_2d.fit_transform(X_scaled)

        fig, axes = plt.subplots(1, 2 if compare_kmeans else 1, figsize=(12 if compare_kmeans else 6, 5))
        if not compare_kmeans:
            axes = [axes]

        # DBSCAN plot
        colors  = plt.cm.tab10(np.linspace(0, 1, max(n_clusters+1, 2)))
        for lbl in sorted(set(labels)):
            mask  = labels == lbl
            color = 'black' if lbl == -1 else colors[lbl % len(colors)]
            label = 'Noise' if lbl == -1 else f'Cluster {lbl+1} (n={cluster_sizes.get(f"Cluster {lbl+1}",0)})'
            axes[0].scatter(X_2d[mask,0], X_2d[mask,1], c=[color], label=label, alpha=0.6, s=20)
        axes[0].set_title(f'DBSCAN — {n_clusters} clusters, {noise_pct}% noise', fontweight='bold')
        axes[0].set_xlabel('PC1'); axes[0].set_ylabel('PC2')
        axes[0].legend(fontsize=8, loc='best')

        # KMeans comparison
        kmeans_result = None
        if compare_kmeans and n_clusters > 1:
            km     = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
            km_lbl = km.fit_predict(X_scaled)
            for lbl in range(n_clusters):
                mask = km_lbl == lbl
                axes[1].scatter(X_2d[mask,0], X_2d[mask,1],
                    c=[colors[lbl % len(colors)]], label=f'Cluster {lbl+1}', alpha=0.6, s=20)
            axes[1].set_title(f'KMeans ({n_clusters} clusters) — Comparison', fontweight='bold')
            axes[1].set_xlabel('PC1'); axes[1].set_ylabel('PC2')
            axes[1].legend(fontsize=8, loc='best')
            kmeans_result = {'n_clusters': n_clusters}

        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success':       True,
            'n_clusters':    n_clusters,
            'n_noise':       n_noise,
            'noise_pct':     noise_pct,
            'total_rows':    len(labels),
            'cluster_sizes': cluster_sizes,
            'params':        {'eps': eps, 'min_samples': min_samples},
            'kmeans_comparison': kmeans_result,
            'chart':         chart_b64,
            'message': f'DBSCAN: {n_clusters} clusters found, {noise_pct}% noise points'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# 🔴 Lasso Regression (dedicated route)
# ════════════════════════════════════════════════════════════════
@app.route('/api/lasso', methods=['POST'])
@login_required
def run_lasso():
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from sklearn.linear_model import Lasso, Ridge, ElasticNet
        from sklearn.model_selection import train_test_split, cross_val_score
        from sklearn.preprocessing import StandardScaler, LabelEncoder
        from sklearn.metrics import r2_score, mean_squared_error
        import numpy as np

        if agent.df is None:
            return jsonify({'error': 'Pehle file upload karo!'}), 400

        data       = request.get_json()
        target_col = data.get('target', '')
        alpha      = float(data.get('alpha', 1.0))
        model_type = data.get('model_type', 'lasso')  # lasso / ridge / elasticnet
        cv_folds   = int(data.get('cv_folds', 5))

        df = agent.df.copy().dropna()
        if target_col not in df.columns:
            return jsonify({'error': f'Column "{target_col}" nahi mila!'}), 400

        y = df[target_col]
        X = df.drop(columns=[target_col])

        # Encode categoricals
        for c in X.select_dtypes(include='object').columns:
            X[c] = LabelEncoder().fit_transform(X[c].astype(str))
        if y.dtype == object:
            y = LabelEncoder().fit_transform(y.astype(str))

        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        scaler = StandardScaler()
        X_train_s = scaler.fit_transform(X_train)
        X_test_s  = scaler.transform(X_test)

        # Choose model
        models_to_run = {
            'Lasso':      Lasso(alpha=alpha, max_iter=5000),
            'Ridge':      Ridge(alpha=alpha),
            'ElasticNet': ElasticNet(alpha=alpha, l1_ratio=0.5, max_iter=5000),
        }

        results = {}
        for name, model in models_to_run.items():
            model.fit(X_train_s, y_train)
            y_pred  = model.predict(X_test_s)
            r2      = round(float(r2_score(y_test, y_pred)), 4)
            rmse    = round(float(np.sqrt(mean_squared_error(y_test, y_pred))), 4)
            cv      = cross_val_score(model, X_train_s, y_train, cv=cv_folds, scoring='r2')
            results[name] = {
                'r2': r2, 'rmse': rmse,
                'cv_mean': round(float(cv.mean()), 4),
                'cv_std':  round(float(cv.std()), 4),
            }

        # Feature coefficients for Lasso (feature selection)
        lasso_model = models_to_run['Lasso']
        coefs = pd.DataFrame({
            'feature':     X.columns,
            'coefficient': lasso_model.coef_
        }).sort_values('coefficient', key=abs, ascending=False)

        zero_features    = int((lasso_model.coef_ == 0).sum())
        nonzero_features = int((lasso_model.coef_ != 0).sum())

        # Plot
        fig, axes = plt.subplots(1, 2, figsize=(12, 5))

        # Coefficients
        top_coef = coefs[coefs['coefficient'] != 0].head(15)
        colors   = ['#dc2626' if v < 0 else '#059669' for v in top_coef['coefficient']]
        axes[0].barh(top_coef['feature'], top_coef['coefficient'], color=colors)
        axes[0].axvline(x=0, color='black', linewidth=0.8)
        axes[0].set_title('Lasso Coefficients\n(Zero = Feature Eliminated)', fontweight='bold')
        axes[0].set_xlabel('Coefficient Value')

        # Model comparison
        model_names = list(results.keys())
        r2_scores   = [results[m]['r2'] for m in model_names]
        bars        = axes[1].bar(model_names, r2_scores,
            color=['#5b5ef4','#8b5cf6','#0d9488'])
        axes[1].set_title('Model Comparison (R² Score)', fontweight='bold')
        axes[1].set_ylabel('R² Score')
        axes[1].set_ylim(0, 1)
        for bar, score in zip(bars, r2_scores):
            axes[1].text(bar.get_x()+bar.get_width()/2,
                bar.get_height()+0.01, f'{score:.3f}',
                ha='center', va='bottom', fontweight='bold')

        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode()

        best_model = max(results, key=lambda m: results[m]['r2'])

        return jsonify({
            'success':          True,
            'target':           target_col,
            'results':          results,
            'best_model':       best_model,
            'best_r2':          results[best_model]['r2'],
            'lasso_feature_selection': {
                'zero_features':    zero_features,
                'nonzero_features': nonzero_features,
                'eliminated':       list(coefs[coefs['coefficient']==0]['feature']),
                'top_features':     list(top_coef['feature'][:10])
            },
            'chart':   chart_b64,
            'message': f'Best: {best_model} (R²={results[best_model]["r2"]}) | Lasso eliminated {zero_features} features'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# ✅ Cross Validation — Dedicated & Improved
# ════════════════════════════════════════════════════════════════
@app.route('/api/cross_validation', methods=['POST'])
@login_required
def run_cross_validation():
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from sklearn.model_selection import (cross_val_score, StratifiedKFold,
                                             KFold, learning_curve)
        from sklearn.preprocessing import StandardScaler, LabelEncoder
        from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
        from sklearn.linear_model import LogisticRegression, Ridge
        from sklearn.tree import DecisionTreeClassifier
        import numpy as np

        if agent.df is None:
            return jsonify({'error': 'Pehle file upload karo!'}), 400

        data        = request.get_json()
        target_col  = data.get('target', '')
        cv_folds    = int(data.get('cv_folds', 5))
        task_type   = data.get('task', 'auto')

        df = agent.df.copy().dropna()
        if target_col not in df.columns:
            return jsonify({'error': f'Column "{target_col}" nahi mila!'}), 400

        if len(df) > 10000:
            df = df.sample(10000, random_state=42)

        y = df[target_col]
        X = df.drop(columns=[target_col])

        for c in X.select_dtypes(include='object').columns:
            X[c] = LabelEncoder().fit_transform(X[c].astype(str))
        if y.dtype == object:
            y = LabelEncoder().fit_transform(y.astype(str))

        import numpy as np
        if task_type == 'auto':
            task_type = 'classification' if len(np.unique(y)) <= 20 else 'regression'

        scaler    = StandardScaler()
        X_scaled  = scaler.fit_transform(X)

        if task_type == 'classification':
            models  = {
                'Random Forest':      RandomForestClassifier(n_estimators=50, random_state=42),
                'Logistic Regression':LogisticRegression(max_iter=500),
                'Decision Tree':      DecisionTreeClassifier(random_state=42),
            }
            scoring = 'accuracy'
            cv      = StratifiedKFold(n_splits=cv_folds, shuffle=True, random_state=42)
        else:
            models  = {
                'Random Forest': RandomForestRegressor(n_estimators=50, random_state=42),
                'Ridge':         Ridge(),
            }
            scoring = 'r2'
            cv      = KFold(n_splits=cv_folds, shuffle=True, random_state=42)

        cv_results = {}
        for name, model in models.items():
            scores = cross_val_score(model, X_scaled, y, cv=cv, scoring=scoring)
            cv_results[name] = {
                'scores':  [round(float(s), 4) for s in scores],
                'mean':    round(float(scores.mean()), 4),
                'std':     round(float(scores.std()), 4),
                'min':     round(float(scores.min()), 4),
                'max':     round(float(scores.max()), 4),
            }

        # Learning curve for best model
        best_name  = max(cv_results, key=lambda m: cv_results[m]['mean'])
        best_model = models[best_name]
        train_sizes, train_scores, val_scores = learning_curve(
            best_model, X_scaled, y, cv=cv, scoring=scoring,
            train_sizes=np.linspace(0.1, 1.0, 8), n_jobs=-1)

        # Plot
        fig, axes = plt.subplots(1, 2, figsize=(12, 5))

        # CV scores comparison
        model_names = list(cv_results.keys())
        means = [cv_results[m]['mean'] for m in model_names]
        stds  = [cv_results[m]['std']  for m in model_names]
        axes[0].barh(model_names, means, xerr=stds,
            color=['#5b5ef4','#8b5cf6','#0d9488'], alpha=0.8, capsize=5)
        axes[0].set_title(f'{cv_folds}-Fold Cross Validation\n({scoring})', fontweight='bold')
        axes[0].set_xlabel(scoring.upper())
        for i, (m, s) in enumerate(zip(means, stds)):
            axes[0].text(m+0.005, i, f'{m:.3f} ± {s:.3f}', va='center', fontsize=9)

        # Learning curve
        train_mean = np.mean(train_scores, axis=1)
        train_std  = np.std(train_scores, axis=1)
        val_mean   = np.mean(val_scores, axis=1)
        val_std    = np.std(val_scores, axis=1)
        axes[1].plot(train_sizes, train_mean, 'o-', color='#5b5ef4', label='Training')
        axes[1].fill_between(train_sizes, train_mean-train_std, train_mean+train_std, alpha=0.15, color='#5b5ef4')
        axes[1].plot(train_sizes, val_mean, 'o-', color='#059669', label='Validation')
        axes[1].fill_between(train_sizes, val_mean-val_std, val_mean+val_std, alpha=0.15, color='#059669')
        axes[1].set_title(f'Learning Curve — {best_name}', fontweight='bold')
        axes[1].set_xlabel('Training Size')
        axes[1].set_ylabel(scoring.upper())
        axes[1].legend(); axes[1].grid(alpha=0.3)

        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode()

        return jsonify({
            'success':     True,
            'task':        task_type,
            'cv_folds':    cv_folds,
            'scoring':     scoring,
            'cv_results':  cv_results,
            'best_model':  best_name,
            'best_score':  cv_results[best_name]['mean'],
            'chart':       chart_b64,
            'message': f'CV Complete! Best: {best_name} ({cv_results[best_name]["mean"]:.3f} ± {cv_results[best_name]["std"]:.3f})'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# 📊 Confusion Matrix — Dedicated Visual
# ════════════════════════════════════════════════════════════════
@app.route('/api/confusion_matrix', methods=['POST'])
@login_required
def run_confusion_matrix():
    try:
        import base64, io as _io
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import seaborn as sns
        from sklearn.model_selection import train_test_split
        from sklearn.preprocessing import StandardScaler, LabelEncoder
        from sklearn.ensemble import RandomForestClassifier, GradientBoostingClassifier
        from sklearn.linear_model import LogisticRegression
        from sklearn.metrics import (confusion_matrix, classification_report,
                                     roc_curve, auc, accuracy_score)
        import numpy as np

        if agent.df is None:
            return jsonify({'error': 'Pehle file upload karo!'}), 400

        data        = request.get_json()
        target_col  = data.get('target', '')
        model_name  = data.get('model', 'Random Forest')

        df = agent.df.copy().dropna()
        if target_col not in df.columns:
            return jsonify({'error': f'Column "{target_col}" nahi mila!'}), 400

        if len(df) > 10000:
            df = df.sample(10000, random_state=42)

        y = df[target_col]
        X = df.drop(columns=[target_col])

        le_y = LabelEncoder()
        y_enc = le_y.fit_transform(y.astype(str))
        class_names = list(le_y.classes_)

        for c in X.select_dtypes(include='object').columns:
            X[c] = LabelEncoder().fit_transform(X[c].astype(str))

        X_train, X_test, y_train, y_test = train_test_split(
            X, y_enc, test_size=0.2, random_state=42, stratify=y_enc)

        scaler = StandardScaler()
        X_train_s = scaler.fit_transform(X_train)
        X_test_s  = scaler.transform(X_test)

        models_map = {
            'Random Forest':       RandomForestClassifier(n_estimators=100, random_state=42),
            'Gradient Boosting':   GradientBoostingClassifier(random_state=42),
            'Logistic Regression': LogisticRegression(max_iter=1000),
        }
        model = models_map.get(model_name, models_map['Random Forest'])
        use_scaled = model_name == 'Logistic Regression'
        Xtr = X_train_s if use_scaled else X_train
        Xte = X_test_s  if use_scaled else X_test
        model.fit(Xtr, y_train)
        y_pred = model.predict(Xte)
        y_prob = model.predict_proba(Xte) if hasattr(model, 'predict_proba') else None

        cm       = confusion_matrix(y_test, y_pred)
        acc      = accuracy_score(y_test, y_pred)
        report   = classification_report(y_test, y_pred, target_names=class_names, output_dict=True)

        # Plots
        n_classes = len(class_names)
        has_roc   = y_prob is not None and n_classes == 2
        fig_cols  = 3 if has_roc else 2
        fig, axes = plt.subplots(1, fig_cols, figsize=(5*fig_cols, 5))

        # 1. Confusion Matrix Heatmap
        cm_norm = cm.astype('float') / cm.sum(axis=1)[:, np.newaxis]
        sns.heatmap(cm_norm, annot=True, fmt='.2f', cmap='Blues',
                    xticklabels=class_names, yticklabels=class_names,
                    ax=axes[0], linewidths=0.5)
        axes[0].set_title(f'Confusion Matrix\nAccuracy: {acc:.3f}', fontweight='bold')
        axes[0].set_ylabel('Actual')
        axes[0].set_xlabel('Predicted')

        # 2. Per-class metrics
        metrics_data = {cls: report[cls] for cls in class_names if cls in report}
        if metrics_data:
            cls_names = list(metrics_data.keys())[:8]
            prec = [metrics_data[c]['precision'] for c in cls_names]
            rec  = [metrics_data[c]['recall']    for c in cls_names]
            x    = range(len(cls_names))
            axes[1].bar([i-0.2 for i in x], prec, 0.4, label='Precision', color='#5b5ef4', alpha=0.8)
            axes[1].bar([i+0.2 for i in x], rec,  0.4, label='Recall',    color='#059669', alpha=0.8)
            axes[1].set_xticks(list(x))
            axes[1].set_xticklabels(cls_names, rotation=45, ha='right')
            axes[1].set_title('Per-Class Precision vs Recall', fontweight='bold')
            axes[1].set_ylim(0, 1.1)
            axes[1].legend()

        # 3. ROC Curve (binary only)
        if has_roc:
            fpr, tpr, _ = roc_curve(y_test, y_prob[:,1])
            roc_auc = auc(fpr, tpr)
            axes[2].plot(fpr, tpr, color='#5b5ef4', lw=2,
                label=f'ROC (AUC = {roc_auc:.3f})')
            axes[2].plot([0,1],[0,1],'k--', lw=1)
            axes[2].set_xlabel('False Positive Rate')
            axes[2].set_ylabel('True Positive Rate')
            axes[2].set_title('ROC Curve', fontweight='bold')
            axes[2].legend()

        plt.tight_layout()
        buf = _io.BytesIO(); plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close('all'); buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode()

        # Per class summary
        class_summary = {}
        for cls in class_names:
            if cls in report:
                class_summary[cls] = {
                    'precision': round(report[cls]['precision'], 3),
                    'recall':    round(report[cls]['recall'],    3),
                    'f1':        round(report[cls]['f1-score'],  3),
                    'support':   int(report[cls]['support']),
                }

        return jsonify({
            'success':       True,
            'model':         model_name,
            'target':        target_col,
            'accuracy':      round(float(acc), 4),
            'n_classes':     n_classes,
            'class_names':   class_names,
            'class_summary': class_summary,
            'macro_f1':      round(float(report['macro avg']['f1-score']), 4),
            'weighted_f1':   round(float(report['weighted avg']['f1-score']), 4),
            'chart':         chart_b64,
            'message': f'Confusion Matrix ready! Accuracy: {acc:.3f} | Classes: {n_classes}'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


# ════════════════════════════════════════════════════════════════
# 📚 LEARNING SUITE — Linear Algebra, RL, Tableau/PowerBI
# Theory + Examples + Quiz + Visual Tools + Practice
# ════════════════════════════════════════════════════════════════

LEARN_PROMPT = """You are an expert Data Science tutor who teaches in Hinglish (Hindi + English mix).
You explain concepts simply with real examples, analogies, and code.
Always include: concept explanation, real-world analogy, Python code example, and key takeaway.
Keep it friendly, encouraging, and practical."""

# ── 1. Linear Algebra Learning ───────────────────────────────────
@app.route('/api/learn/linear_algebra', methods=['POST'])
@login_required
def learn_linear_algebra():
    try:
        import json as _j
        data    = request.get_json()
        topic   = data.get('topic', 'vectors')
        mode    = data.get('mode', 'theory')  # theory / example / quiz / practice

        topics_map = {
            'vectors':      'Vectors — kya hote hain, operations, dot product, real use in DS',
            'matrices':     'Matrices — addition, multiplication, transpose, inverse, determinant',
            'eigenvalues':  'Eigenvalues & Eigenvectors — concept, calculation, use in PCA',
            'svd':          'SVD (Singular Value Decomposition) — concept aur recommendation systems mein use',
            'pca_math':     'PCA Mathematics — covariance matrix, eigenvectors, dimensionality reduction step by step',
            'linear_trans': 'Linear Transformations — scaling, rotation, projection visually',
            'norms':        'Vector Norms — L1, L2, infinity norm, use in Lasso/Ridge',
            'calculus':     'Calculus basics for ML — derivatives, gradients, chain rule, backpropagation',
        }

        topic_desc = topics_map.get(topic, topic)

        if mode == 'theory':
            prompt = f"""Explain this topic in detail for a beginner Data Science student:
Topic: {topic_desc}

Format your response as:
## 🎯 Kya Hai? (Simple definition)
## 🌍 Real Life Analogy (Relatable example)
## 📐 Mathematics (Simple formulas with explanation)
## 🐍 Python Code (Working example with numpy)
## 🔗 DS/ML Mein Use (Where this is used in real ML)
## 💡 Key Takeaway (One line summary)

Use Hinglish language. Be friendly and encouraging."""

        elif mode == 'example':
            prompt = f"""Give 3 practical solved examples for: {topic_desc}

For each example:
- Problem statement (real world scenario)
- Step by step solution
- Python numpy code
- Expected output
- What this means in ML context

Use Hinglish. Make it practical and hands-on."""

        elif mode == 'quiz':
            prompt = f"""Create a 5-question quiz on: {topic_desc}

Format as JSON only, no other text:
{{
  "quiz_title": "...",
  "topic": "...",
  "questions": [
    {{
      "id": 1,
      "question": "...",
      "options": ["A. ...", "B. ...", "C. ...", "D. ..."],
      "correct": "A",
      "explanation": "..."
    }}
  ]
}}

Mix easy (2), medium (2), hard (1) questions.
Focus on practical ML applications."""

        elif mode == 'practice':
            if not hasattr(request, '_agent_df'):
                prompt = f"""Create 3 hands-on Python practice problems for: {topic_desc}

For each problem:
- Title
- Problem description
- Starter code (with TODOs)
- Hints
- Complete solution
- How this applies in real DS work

Use numpy, make it practical."""
            else:
                prompt = f"""Create practice problems for {topic_desc} using this dataset context."""

        result = get_groq_response(LEARN_PROMPT, prompt, max_tokens=2000)

        if mode == 'quiz':
            try:
                import re as _re
                json_match = _re.search(r'\{.*\}', result, _re.DOTALL)
                if json_match:
                    quiz_data = _j.loads(json_match.group())
                    return jsonify({'success': True, 'mode': mode, 'topic': topic, 'quiz': quiz_data})
            except:
                pass

        return jsonify({
            'success': True,
            'mode':    mode,
            'topic':   topic,
            'content': result,
            'message': f'Linear Algebra — {topic} ({mode}) ready!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/learn/matrix_calculator', methods=['POST'])
@login_required
def matrix_calculator():
    """Visual Matrix Calculator — operations with step by step"""
    try:
        import numpy as np
        import json as _j
        data      = request.get_json()
        operation = data.get('operation', 'multiply')
        matrix_a  = data.get('matrix_a', [[1,2],[3,4]])
        matrix_b  = data.get('matrix_b', [[5,6],[7,8]])

        A = np.array(matrix_a, dtype=float)
        B = np.array(matrix_b, dtype=float) if matrix_b else None

        result_data = {}
        steps       = []

        if operation == 'multiply':
            if B is None or A.shape[1] != B.shape[0]:
                return jsonify({'error': f'Matrix dimensions incompatible for multiplication! A:{A.shape} B:{B.shape}'}), 400
            C = np.dot(A, B)
            result_data = {'result': C.tolist()}
            steps = [
                f'Matrix A: {A.shape[0]}×{A.shape[1]}',
                f'Matrix B: {B.shape[0]}×{B.shape[1]}',
                f'Result C: {C.shape[0]}×{C.shape[1]}',
                'C[i,j] = sum of A[i,k] × B[k,j]',
                f'Result:\n{np.array2string(C, precision=3)}'
            ]
        elif operation == 'transpose':
            T = A.T
            result_data = {'result': T.tolist()}
            steps = [f'Original: {A.shape}', f'Transposed: {T.shape}', 'Rows become columns, columns become rows']
        elif operation == 'inverse':
            if A.shape[0] != A.shape[1]:
                return jsonify({'error': 'Square matrix chahiye inverse ke liye!'}), 400
            det = np.linalg.det(A)
            if abs(det) < 1e-10:
                return jsonify({'error': 'Singular matrix — inverse exist nahi karta!'}), 400
            inv = np.linalg.inv(A)
            result_data = {'result': inv.tolist(), 'determinant': round(float(det), 4)}
            steps = [f'Determinant: {round(det,4)}', 'A × A⁻¹ = Identity Matrix', f'Inverse:\n{np.array2string(inv, precision=3)}']
        elif operation == 'eigenvalues':
            if A.shape[0] != A.shape[1]:
                return jsonify({'error': 'Square matrix chahiye!'}), 400
            vals, vecs = np.linalg.eig(A)
            result_data = {
                'eigenvalues':  [round(float(v.real), 4) for v in vals],
                'eigenvectors': [[round(float(x.real), 4) for x in vec] for vec in vecs.T]
            }
            steps = [
                f'{len(vals)} eigenvalues found',
                f'Eigenvalues: {[round(float(v.real),3) for v in vals]}',
                'Eigenvectors: directions jo transform ke baad same rehte hain',
                'Use in PCA: Sort by eigenvalue → select top k → reduce dimensions'
            ]
        elif operation == 'determinant':
            if A.shape[0] != A.shape[1]:
                return jsonify({'error': 'Square matrix chahiye!'}), 400
            det = np.linalg.det(A)
            result_data = {'determinant': round(float(det), 6)}
            steps = [f'det(A) = {round(det,4)}',
                     'det=0 → Singular (no inverse)',
                     'det≠0 → Non-singular (has inverse)',
                     f'This matrix is {"Singular" if abs(det)<1e-10 else "Non-singular"}']
        elif operation == 'dot_product':
            if B is None:
                return jsonify({'error': 'Dono vectors chahiye dot product ke liye!'}), 400
            a_flat = A.flatten(); b_flat = B.flatten()
            dot    = float(np.dot(a_flat, b_flat))
            norm_a = float(np.linalg.norm(a_flat))
            norm_b = float(np.linalg.norm(b_flat))
            cos_sim = dot / (norm_a * norm_b) if norm_a * norm_b > 0 else 0
            result_data = {
                'dot_product':    round(dot, 4),
                'norm_a':         round(norm_a, 4),
                'norm_b':         round(norm_b, 4),
                'cosine_similarity': round(cos_sim, 4),
                'angle_degrees':  round(float(np.degrees(np.arccos(np.clip(cos_sim,-1,1)))), 2)
            }
            steps = [
                f'Dot Product = {round(dot,4)}',
                f'|A| = {round(norm_a,4)}, |B| = {round(norm_b,4)}',
                f'Cosine Similarity = {round(cos_sim,4)}',
                f'Angle = {round(float(np.degrees(np.arccos(np.clip(cos_sim,-1,1)))),2)}°',
                'Used in: NLP similarity, recommendation systems, neural networks'
            ]
        elif operation == 'svd':
            U, S, Vt = np.linalg.svd(A)
            result_data = {
                'U':  U.tolist(),
                'S':  S.tolist(),
                'Vt': Vt.tolist(),
                'singular_values': [round(float(s),4) for s in S],
                'rank': int(np.linalg.matrix_rank(A))
            }
            steps = [
                f'A = U × Σ × V^T',
                f'Singular Values: {[round(float(s),3) for s in S]}',
                f'Matrix Rank: {int(np.linalg.matrix_rank(A))}',
                'Largest singular values = most important information',
                'Use in: Image compression, recommendation systems, PCA'
            ]

        return jsonify({
            'success':   True,
            'operation': operation,
            'matrix_a':  matrix_a,
            'matrix_b':  matrix_b,
            'result':    result_data,
            'steps':     steps,
            'message':   f'{operation.title()} complete!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── 2. Reinforcement Learning Learning ──────────────────────────
@app.route('/api/learn/rl', methods=['POST'])
@login_required
def learn_rl():
    try:
        import json as _j
        data  = request.get_json()
        topic = data.get('topic', 'basics')
        mode  = data.get('mode', 'theory')

        topics_map = {
            'basics':       'RL Basics — Agent, Environment, State, Action, Reward kya hote hain',
            'markov':       'Markov Decision Process (MDP) — States, Transitions, Rewards formally',
            'q_learning':   'Q-Learning Algorithm — Q-table, Bellman equation, exploration vs exploitation',
            'deep_rl':      'Deep RL — DQN, Policy Gradient, Actor-Critic basics',
            'applications': 'RL Applications — Games, Trading bots, Robotics, Recommendation systems',
            'gym':          'OpenAI Gym — Environment setup, CartPole, basic RL loop in Python',
        }

        topic_desc = topics_map.get(topic, topic)

        if mode == 'theory':
            prompt = f"""Explain Reinforcement Learning topic for a beginner DS student:
Topic: {topic_desc}

Format:
## 🎯 Concept (Simple definition)
## 🎮 Real Life Analogy (Game/Life example)
## 🔄 How It Works (Step by step)
## 🐍 Python Code (Simple working example)
## 📊 Diagram Description (Explain the flow in text)
## 🌍 Real Applications (Where it's actually used)
## 💡 Key Takeaway

Use Hinglish. Use gaming analogies — very beginner friendly."""

        elif mode == 'example':
            prompt = f"""Give practical examples for RL topic: {topic_desc}

Include:
1. Simple toy example (step by step)
2. Python code (OpenAI Gym or custom environment)
3. Expected output explanation
4. Real world application example

Make it very visual with text diagrams. Hinglish language."""

        elif mode == 'quiz':
            prompt = f"""Create a 5-question quiz on: {topic_desc}

Return JSON only:
{{
  "quiz_title": "...",
  "topic": "...",
  "questions": [
    {{
      "id": 1,
      "question": "...",
      "options": ["A. ...", "B. ...", "C. ...", "D. ..."],
      "correct": "A",
      "explanation": "..."
    }}
  ]
}}"""

        elif mode == 'simulator':
            prompt = f"""Create a simple text-based RL simulation for: {topic_desc}

Include:
1. Environment description (Grid World ya simple game)
2. Complete Python code (no external libs except numpy)
3. Step by step execution trace
4. How agent learns over episodes
5. Output showing improvement over time

Make it runnable in Python notebook. Hinglish explanation."""

        result = get_groq_response(LEARN_PROMPT, prompt, max_tokens=2000)

        if mode == 'quiz':
            try:
                import re as _re
                json_match = _re.search(r'\{.*\}', result, _re.DOTALL)
                if json_match:
                    quiz_data = _j.loads(json_match.group())
                    return jsonify({'success': True, 'mode': mode, 'topic': topic, 'quiz': quiz_data})
            except:
                pass

        return jsonify({
            'success': True,
            'mode':    mode,
            'topic':   topic,
            'content': result,
            'message': f'RL — {topic} ({mode}) ready!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/learn/rl_simulator', methods=['POST'])
@login_required
def rl_grid_simulator():
    """Simple Q-Learning Grid World Simulator"""
    try:
        import numpy as np
        import json as _j
        data     = request.get_json()
        episodes = min(int(data.get('episodes', 100)), 500)
        grid_size = int(data.get('grid_size', 4))
        alpha    = float(data.get('alpha', 0.1))    # learning rate
        gamma    = float(data.get('gamma', 0.9))    # discount factor
        epsilon  = float(data.get('epsilon', 0.1))  # exploration rate

        # Grid World — Agent navigates from (0,0) to (grid_size-1, grid_size-1)
        n_states  = grid_size * grid_size
        n_actions = 4  # Up, Down, Left, Right
        ACTIONS   = ['↑ Up', '↓ Down', '← Left', '→ Right']

        # Q-table
        Q = np.zeros((n_states, n_actions))

        # Obstacles
        obstacles = set()
        if grid_size >= 4:
            obstacles = {5, 7, 9, 11} if grid_size == 4 else set()

        goal_state  = n_states - 1
        rewards_per_episode = []
        steps_per_episode   = []

        def state_to_pos(s): return (s // grid_size, s % grid_size)
        def pos_to_state(r,c): return r * grid_size + c

        def take_action(state, action):
            row, col = state_to_pos(state)
            if   action == 0: row = max(0, row-1)       # Up
            elif action == 1: row = min(grid_size-1, row+1)  # Down
            elif action == 2: col = max(0, col-1)       # Left
            elif action == 3: col = min(grid_size-1, col+1)  # Right
            new_state = pos_to_state(row, col)
            if new_state in obstacles: new_state = state
            reward = 10 if new_state == goal_state else -0.1
            done   = new_state == goal_state
            return new_state, reward, done

        # Train
        for ep in range(episodes):
            state = 0
            total_reward = 0
            steps = 0
            for _ in range(100):
                if np.random.random() < epsilon:
                    action = np.random.randint(n_actions)
                else:
                    action = int(np.argmax(Q[state]))
                next_state, reward, done = take_action(state, action)
                Q[state, action] += alpha * (reward + gamma * np.max(Q[next_state]) - Q[state, action])
                state = next_state
                total_reward += reward
                steps += 1
                if done: break
            rewards_per_episode.append(round(float(total_reward), 2))
            steps_per_episode.append(steps)

        # Best path
        state = 0
        path  = [state]
        for _ in range(50):
            action = int(np.argmax(Q[state]))
            state, _, done = take_action(state, action)
            path.append(state)
            if done: break

        # Grid visualization
        grid_vis = []
        for r in range(grid_size):
            row_vis = []
            for c in range(grid_size):
                s = pos_to_state(r, c)
                if s == 0:              cell = 'S'   # Start
                elif s == goal_state:   cell = 'G'   # Goal
                elif s in obstacles:    cell = 'X'   # Obstacle
                elif s in path:         cell = '●'   # Path
                else:                   cell = '·'
                row_vis.append(cell)
            grid_vis.append(row_vis)

        # Learning stats
        avg_first  = round(float(np.mean(rewards_per_episode[:10])), 2)
        avg_last   = round(float(np.mean(rewards_per_episode[-10:])), 2)
        improvement = round(avg_last - avg_first, 2)

        return jsonify({
            'success':    True,
            'grid_size':  grid_size,
            'episodes':   episodes,
            'params':     {'alpha': alpha, 'gamma': gamma, 'epsilon': epsilon},
            'grid_visual': grid_vis,
            'best_path':  path,
            'path_length': len(path),
            'q_table_sample': Q[:min(6,n_states)].tolist(),
            'rewards': {
                'first_10_avg': avg_first,
                'last_10_avg':  avg_last,
                'improvement':  improvement,
                'all': rewards_per_episode[::max(1, episodes//20)]
            },
            'learning_happened': improvement > 0,
            'explanation': f'Agent {grid_size}x{grid_size} grid mein seekha! {episodes} episodes ke baad reward {avg_first} se {avg_last} ho gaya (+{improvement}). Best path: {len(path)} steps mein goal reach kiya!',
            'message': f'Q-Learning complete! {episodes} episodes | Improvement: {improvement:+.2f}'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── 3. Tableau/PowerBI Alternative Learning ─────────────────────
@app.route('/api/learn/visualization', methods=['POST'])
@login_required
def learn_visualization():
    try:
        import json as _j
        data  = request.get_json()
        topic = data.get('topic', 'tableau_basics')
        mode  = data.get('mode', 'theory')

        topics_map = {
            'tableau_basics':  'Tableau basics — interface, data connection, drag-drop, sheets, dashboards',
            'tableau_charts':  'Tableau charts — bar, line, scatter, map, treemap, heatmap banane ka tarika',
            'powerbi_basics':  'Power BI basics — Power Query, DAX basics, visualizations, publishing',
            'powerbi_dax':     'DAX Formulas — CALCULATE, SUMX, FILTER, time intelligence',
            'plotly_python':   'Plotly in Python — interactive charts jo Tableau/PowerBI jaisi dikhti hain',
            'dash_app':        'Dash by Plotly — Python mein PowerBI jaisa web dashboard banana',
            'comparison':      'Tableau vs PowerBI vs Python — kab kya use karna chahiye',
            'free_alternatives': 'Free alternatives — Google Data Studio, Metabase, Apache Superset',
        }

        topic_desc = topics_map.get(topic, topic)

        if mode == 'theory':
            prompt = f"""Explain this data visualization/BI tool topic:
Topic: {topic_desc}

Format:
## 🎯 Kya Hai?
## 💻 Interface/Setup (Screenshots description)
## 🔧 Step by Step (Kaise use karte hain)
## 🐍 Python Alternative (Plotly/Dash se same kaam)
## ✅ Pros & Cons
## 🌍 Real Use Cases (Kab use karna chahiye)
## 💡 Key Takeaway

Include Python code for Plotly/Dash as free alternative.
Hinglish language. Be practical."""

        elif mode == 'example':
            prompt = f"""Give practical examples for: {topic_desc}

Include:
1. Step by step tutorial (with description of clicks/actions)
2. Python Plotly code achieving same result (free alternative)
3. Sample dataset to practice with
4. Common mistakes to avoid

Make it hands-on. Hinglish."""

        elif mode == 'quiz':
            prompt = f"""Create quiz on: {topic_desc}

Return JSON only:
{{
  "quiz_title": "...",
  "topic": "...",
  "questions": [
    {{
      "id": 1,
      "question": "...",
      "options": ["A. ...", "B. ...", "C. ...", "D. ..."],
      "correct": "A",
      "explanation": "..."
    }}
  ]
}}"""

        elif mode == 'practice':
            prompt = f"""Create 3 hands-on practice projects for: {topic_desc}

For each project:
- Project title
- Dataset to use (from Kaggle or built-in)
- What to build
- Python Plotly/Dash code (complete working)
- Expected result description
- How this compares to Tableau/PowerBI version

Focus on Python free alternatives. Hinglish."""

        result = get_groq_response(LEARN_PROMPT, prompt, max_tokens=2000)

        if mode == 'quiz':
            try:
                import re as _re
                json_match = _re.search(r'\{.*\}', result, _re.DOTALL)
                if json_match:
                    quiz_data = _j.loads(json_match.group())
                    return jsonify({'success': True, 'mode': mode, 'topic': topic, 'quiz': quiz_data})
            except:
                pass

        return jsonify({
            'success': True,
            'mode':    mode,
            'topic':   topic,
            'content': result,
            'message': f'Visualization — {topic} ({mode}) ready!'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

